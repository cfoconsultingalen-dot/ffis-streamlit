
from __future__ import annotations
import streamlit as st
from typing import Iterable, List, Optional, Set,Dict
from datetime import date, datetime, timedelta,timezone
import pandas as pd
from supabase import create_client
import altair as alt
import calendar
import os
import time
import math
import numpy as np
import io
from postgrest.exceptions import APIError
import matplotlib.pyplot as plt
from pandas.tseries.offsets import MonthEnd

#import Fundtion for Board Reportinf
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ---------- Supabase config ----------

SUPABASE_URL = st.secrets["SUPABASE_URL"]
SUPABASE_ANON_KEY = st.secrets["SUPABASE_ANON_KEY"]

@st.cache_resource
def get_supabase_client():
    return create_client(SUPABASE_URL, SUPABASE_ANON_KEY)


supabase = get_supabase_client()

st.set_page_config(page_title="Matfina – Founder Financial Intelligence", layout="wide")

#......Month Seection.......$


def _parse_any_date(raw) -> Optional[date]:
    """
    Try to parse a date from various formats:
    - Already a date object -> return as-is
    - '2025-08-21' (ISO)
    - '21-08-2025' or '21/08/2025' (dd-mm-YYYY / dd/mm/YYYY)

    Returns a `date` or None if parsing fails.
    """
    if raw is None:
        return None

    if isinstance(raw, date):
        return raw

    if not isinstance(raw, str):
        return None

    raw = raw.strip()
    if not raw:
        return None

    # 1) ISO first (YYYY-MM-DD)
    try:
        return date.fromisoformat(raw)
    except ValueError:
        pass

    # 2) dd-mm-YYYY or dd/mm/YYYY
    for fmt in ("%d-%m-%Y", "%d/%m/%Y"):
        try:
            return datetime.strptime(raw, fmt).date()
        except ValueError:
            continue

    return None

def _collect_months_from_table(
    sb,
    table_name: str,
    client_id: str,
    date_columns: Iterable[str],
) -> Set[date]:
    """
    Fetch all rows for a client from a table and extract month-level dates
    from the given date_columns.

    Returns a set of `date` objects (normalised to first-of-month).
    """
    months: Set[date] = set()

    columns = ["client_id", *date_columns]
    try:
        resp = (
            sb.table(table_name)
              .select(",".join(columns))
              .eq("client_id", str(client_id))
              .execute()
        )
    except Exception as e:
        print(f"[ERROR] _collect_months_from_table({table_name}) -> {e}")
        return months

    data = getattr(resp, "data", None) if hasattr(resp, "data") else resp.get("data", [])
    if not data:
        return months

    for row in data:
        for col in date_columns:
            raw = row.get(col)
            d = _parse_any_date(raw)
            if not d:
                continue

            # Normalise to first of month
            d = d.replace(day=1)
            months.add(d)

    return months



def get_focus_month_window(focus_month: date) -> tuple[date, date]:
    """
    Given a focus month (any date within the month),
    return (month_start, month_end) as pure date objects.
    """
    if focus_month is None:
        return None, None

    # Normalise to 1st of month
    month_start_ts = pd.to_datetime(focus_month).to_period("M").to_timestamp()
    month_start = month_start_ts.date()

    # Last day of that month
    month_end_ts = (month_start_ts + pd.offsets.MonthEnd(0))
    month_end = month_end_ts.date()

    return month_start, month_end


def get_client_focus_month_options(client_id: str) -> List[str]:
    """
    Build list of Focus Month labels for a client using ONLY data in Supabase.

    Sources (per client_id), example:
      - ar_ap_tracker: due_date, expected_payment_date, issued_date
      - revenue_pipeline: start_month, end_month
      - operating_opex: month_date
      - operating_other_income: month_date
      - investing_flows: month_date
      - financing_flows: month_date
      - payroll_positions: start_date, end_date

    Rules:
      - Only months where data actually exists.
      - NO artificial caps: if data exists for 2026, 2030, etc., show it.
      - No synthetic months; if no data, return [].

    Returns:
      List of strings like "Jul 2025", sorted ascending.
    """
    if not client_id:
        return []

    sb = get_supabase_client()  # your existing helper that returns a Supabase client

    all_months: Set[date] = set()

    # 1) AR/AP tracker
    all_months |= _collect_months_from_table(
        sb,
        table_name="ar_ap_tracker",
        client_id=client_id,
        date_columns=["due_date", "expected_payment_date", "issued_date"],
    )

    # 2) Revenue pipeline
    all_months |= _collect_months_from_table(
        sb,
        table_name="revenue_pipeline",
        client_id=client_id,
        date_columns=["start_month", "end_month"],
    )

    # 3) Operating OPEX
    all_months |= _collect_months_from_table(
        sb,
        table_name="operating_opex",
        client_id=client_id,
        date_columns=["month_date"],
    )

    # 4) Operating other income
    all_months |= _collect_months_from_table(
        sb,
        table_name="operating_other_income",
        client_id=client_id,
        date_columns=["month_date"],
    )

    # 5) Investing flows
    all_months |= _collect_months_from_table(
        sb,
        table_name="investing_flows",
        client_id=client_id,
        date_columns=["month_date"],
    )

    # 6) Financing flows
    all_months |= _collect_months_from_table(
        sb,
        table_name="financing_flows",
        client_id=client_id,
        date_columns=["month_date"],
    )

    # 7) Payroll positions (start/end)
    all_months |= _collect_months_from_table(
        sb,
        table_name="payroll_positions",
        client_id=client_id,
        date_columns=["start_date", "end_date"],
    )

    if not all_months:
        # No dated data at all -> no focus months
        return []

    sorted_months = sorted(all_months)
    labels = [d.strftime("%b %Y") for d in sorted_months]  # e.g. "Jul 2025"
    return labels





def fetch_scenarios_for_client(client_id: str):
    """Return list of saved scenarios for this client, newest first."""
    try:
        resp = (
            supabase.table("scenarios")
            .select("*")
            .eq("client_id", str(client_id))
            .order("created_at", desc=True)
            .execute()
        )
        return resp.data or []
    except Exception as e:
        print("Error fetching scenarios:", e)
        return []


def save_scenario_for_client(
    client_id: str,
    name: str,
    controls: dict,
    base_month: date | None = None,
    description: str | None = None,
):
    """Insert a new named scenario for this client."""
    try:
        payload = {
            "client_id": str(client_id),
            "name": name.strip(),
            "controls": controls,
            "base_month": base_month.isoformat() if base_month else None,
            "description": (description or "").strip() or None,
        }
        supabase.table("scenarios").insert(payload).execute()
        return True
    except Exception as e:
        print("Error saving scenario:", e)
        return False


def _download_csv_template(filename: str, df_template: pd.DataFrame):
    csv = df_template.to_csv(index=False)
    st.download_button(
        label=f"⬇️ Download {filename}",
        data=csv,
        file_name=filename,
        mime="text/csv",
        use_container_width=True,
    )

def _read_uploaded_csv(uploaded_file) -> pd.DataFrame:
    # robust csv reader (handles Excel-export csv)
    raw = uploaded_file.read()
    return pd.read_csv(io.BytesIO(raw))

def _normalize_dates(df: pd.DataFrame, date_cols: list[str]):
    for c in date_cols:
        if c in df.columns:
            df[c] = pd.to_datetime(df[c], errors="coerce").dt.date
    return df

def _normalize_month_start(df: pd.DataFrame, col: str):
    # converts to first of month date
    if col in df.columns:
        dt = pd.to_datetime(df[col], errors="coerce")
        df[col] = dt.dt.to_period("M").dt.to_timestamp().dt.date
    return df

def _require_columns(df: pd.DataFrame, required: list[str]):
    missing = [c for c in required if c not in df.columns]
    if missing:
        raise ValueError(f"Missing required columns: {missing}")

def _keep_columns(df: pd.DataFrame, allowed: list[str]):
    # keep only known columns; ignore extras safely
    keep = [c for c in allowed if c in df.columns]
    return df[keep].copy()

def _clean_empty_rows(df: pd.DataFrame):
    # drop completely empty rows
    df = df.dropna(how="all")
    # strip string columns
    for c in df.columns:
        if pd.api.types.is_object_dtype(df[c]):
            df[c] = df[c].astype(str).str.strip()
            df.loc[df[c].isin(["", "nan", "None"]), c] = None
    return df

def _supabase_upsert_rows(table: str, rows: list[dict]):
    # assumes you have global `supabase`
    if not rows:
        return
    supabase.table(table).upsert(rows).execute()


def fetch_scenario_by_id(scenario_id: int):
    """Get a single scenario by numeric ID."""
    try:
        resp = (
            supabase.table("scenarios")
            .select("*")
            .eq("id", scenario_id)
            .limit(1)
            .execute()
        )
        rows = resp.data or []
        return rows[0] if rows else None
    except Exception as e:
        print("Error fetching scenario by id:", e)
        return None


def delete_scenario_by_id(scenario_id: int) -> bool:
    """Optional: allow deleting a saved scenario from the UI."""
    try:
        supabase.table("scenarios").delete().eq("id", scenario_id).execute()
        return True
    except Exception as e:
        print("Error deleting scenario:", e)
        return False


@st.cache_data(ttl=60)
def fetch_operating_other_income_for_client(client_id) -> pd.DataFrame:
    """
    Load operating_other_income rows for a client.

    Supabase table columns:
      - id
      - client_id
      - month_date      (date: cash month)
      - description
      - cash_in         (numeric, positive = cash in)
      - notes
      - created_at
      - updated_at

    Returns a DataFrame with:
      - month_date (Timestamp)
      - cash_in (float)
      - month_bucket (Timestamp, first day of month)
    """
    if not client_id:
        return pd.DataFrame(columns=["month_bucket", "cash_in"])

    try:
        supabase = get_supabase_client()
        resp = (
            supabase.table("operating_other_income")
            .select("*")
            .eq("client_id", str(client_id))
            .execute()
        )
        data = getattr(resp, "data", None) if hasattr(resp, "data") else resp.get("data", [])
    except Exception as e:
        print("Error fetching operating_other_income:", e)
        return pd.DataFrame(columns=["month_bucket", "cash_in"])

    if not data:
        print(f"[DEBUG] fetch_operating_other_income_for_client -> 0 rows for {client_id}")
        return pd.DataFrame(columns=["month_bucket", "cash_in"])

    df = pd.DataFrame(data)

    # Parse month_date
    df["month_date"] = pd.to_datetime(df.get("month_date"), errors="coerce")
    df = df[df["month_date"].notna()]

    # Ensure cash_in is numeric
    df["cash_in"] = pd.to_numeric(df.get("cash_in"), errors="coerce").fillna(0.0)

    # Month bucket = first of month (for engine)
    df["month_bucket"] = df["month_date"].dt.to_period("M").dt.to_timestamp()

    print(f"[DEBUG] fetch_operating_other_income_for_client -> {len(df)} rows for {client_id}")
    return df


from datetime import date
import pandas as pd

def _month_start(d: date) -> date:
    return pd.to_datetime(d).to_period("M").to_timestamp().date()

def fetch_client_monthly_targets(client_id: str) -> pd.DataFrame:
    if not client_id:
        return pd.DataFrame()

    resp = (
        supabase.table("client_monthly_targets")
        .select("*")
        .eq("client_id", str(client_id))
        .execute()
    )

    rows = resp.data or []
    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)
    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    return df



PAGE_KPI_MAP = {
    "business_overview": {
        "label": "Business overview",
        "kpis": {
            "revenue": "Revenue (this month)",
            "burn": "Burn / cash out (this month)",
            "cash_balance": "Cash balance",
            "runway_months": "Runway (months)",
        },
    },
    "sales_deals": {
        "label": "Sales & deals",
        "kpis": {
            "recognised_revenue": "Recognised revenue",
            "pipeline_total": "Total pipeline",
            "pipeline_weighted": "Weighted pipeline",
            "open_deals": "Open deals count",
        },
    },
    "team_spending": {
        "label": "Team & spending",
        "kpis": {
            "payroll_total": "Payroll (this month)",
            "headcount_fte": "Headcount (FTE)",
            "opex_total": "Operating expenses",
        },
    },
    "cash_bills": {
        "label": "Cash & bills",
        "kpis": {
            "ar_in": "AR cash-in (expected)",
            "ap_out": "AP cash-out (expected)",
            "overdue_ar": "Overdue AR",
        },
    },
    "funding_strategy": {
        "label": "Funding strategy",
        "kpis": {
            "cash_cliff": "Cash cliff date",
            "effective_burn": "Effective burn",
            "runway_months": "Runway (months)",
        },
    },
    
}




# ---------- Helpers (data + utilities) ----------

@st.cache_data(ttl=60)
def load_clients():
    """
    Load clients from Supabase as a list of dicts: [{'id': ..., 'name': ...}, ...]
    Falls back to a single demo client if table is empty or fails.
    """
    try:
        response = supabase.table("clients").select("id, name").order("created_at").execute()
        rows = response.data or []
        if rows:
            return rows
    except Exception:
        pass

    # Fallback demo client (no id)
    return [{"id": None, "name": "Demo Startup"}]



# ---------- Payroll helpers ----------

@st.cache_data(ttl=60)
def fetch_payroll_positions_for_client(client_id):
    """
    Return payroll positions for a client as a DataFrame.

    Matches Supabase schema:

        base_salary_annual  -- annual base at 1.0 FTE
        super_rate_pct      -- e.g. 11.0
        payroll_tax_pct     -- e.g. 3.0
        fte                 -- FTE fraction
        start_date / end_date

    Handles missing columns defensively.
    """
    if client_id is None:
        return pd.DataFrame()

    try:
        res = (
            supabase
            .table("payroll_positions")
            .select("*")
            .eq("client_id", str(client_id))
            .execute()
        )
        rows = res.data or []
    except Exception as e:
        print("Error fetching payroll positions:", e)
        rows = []

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)

    # Numeric fields
    for col in ["base_salary_annual", "super_rate_pct", "payroll_tax_pct", "fte"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")
        else:
            # Fallback defaults
            if col == "base_salary_annual":
                df[col] = 0.0
            elif col == "super_rate_pct":
                df[col] = 0.0
            elif col == "payroll_tax_pct":
                df[col] = 0.0
            elif col == "fte":
                df[col] = 1.0

    # Dates
    for col in ["start_date", "end_date"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce").dt.date
        else:
            df[col] = None

    # Optional contractor flag – not in your schema, so default False
    if "is_contractor" not in df.columns:
        df["is_contractor"] = False
    else:
        df["is_contractor"] = df["is_contractor"].fillna(False).astype(bool)

    return df


def compute_annual_income_tax_au(taxable_income: float) -> float:
    """
    Australian resident income tax for a full year (excluding Medicare levy),
    based on the slab you provided.

    Taxable income:
      0 – 18,200       -> 0
      18,201 – 45,000  -> 16c for each $1 over 18,200
      45,001 – 135,000 -> 4,288 + 30c for each $1 over 45,000
      135,001 – 190,000-> 31,288 + 37c for each $1 over 135,000
      190,001+         -> 51,638 + 45c for each $1 over 190,000
    """
    ti = float(max(taxable_income, 0.0))

    if ti <= 18_200:
        base_tax = 0.0
    elif ti <= 45_000:
        base_tax = (ti - 18_200) * 0.16
    elif ti <= 135_000:
        base_tax = 4_288 + (ti - 45_000) * 0.30
    elif ti <= 190_000:
        base_tax = 31_288 + (ti - 135_000) * 0.37
    else:
        base_tax = 51_638 + (ti - 190_000) * 0.45

    # Medicare levy: 2% of taxable income (simplified)
    medicare = ti * 0.02

    return base_tax + medicare


from typing import List, Optional

# -----------------------------------------
# Central task fetcher for Collaboration Hub
# -----------------------------------------
def fetch_tasks_for_collab_hub(
    client_id: str,
    status_filter: Optional[List[str]] = None,
    page_filter: Optional[List[str]] = None,
    owner_filter: Optional[List[str]] = None,
    search_text: str = "",
) -> pd.DataFrame:
    """
    Fetch tasks for the central Collaboration Hub board.
    Uses existing `tasks` table:
      - status: open / in_progress / done / cancelled
      - owner_name: assignee
      - page_name: originating page (business_overview, sales_deals, etc)
    """
    query = (
        supabase.table("tasks")
        .select("*")
        .eq("client_id", client_id)
    )

    if status_filter:
        query = query.in_("status", status_filter)

    if page_filter:
        query = query.in_("page_name", page_filter)

    if owner_filter:
        query = query.in_("owner_name", owner_filter)

    search_text = (search_text or "").strip()
    if search_text:
        like_expr = f"%{search_text}%"
        # OR across title/description
        query = query.or_(
            f"title.ilike.{like_expr},description.ilike.{like_expr}"
        )

    resp = (
        query
        .order("due_date", desc=False)
        .order("priority", desc=False)
        .order("created_at", desc=False)
        .execute()
    )

    data = resp.data or []
    return pd.DataFrame(data)


# -----------------------------------------
# Distinct owners for filter dropdown
# -----------------------------------------
def get_task_owner_options(client_id: str) -> List[str]:
    """
    Distinct owner_name values, for filter dropdowns.
    We deduplicate in Python instead of using `distinct=` in select().
    """
    resp = (
        supabase.table("tasks")
        .select("owner_name")
        .eq("client_id", client_id)
        .execute()
    )
    rows = resp.data or []
    owners = [r["owner_name"] for r in rows if r.get("owner_name")]
    return sorted(set(owners))


# -----------------------------------------
# Distinct pages for filter dropdown
# -----------------------------------------
def get_task_page_options(client_id: str) -> List[str]:
    """
    Known page slugs + any additional page_name values found in tasks.
    """
    base_pages = [
        "business_overview",
        "sales_deals",
        "team_spending",
        "cash_bills",
        "alerts_todos",
        "collaboration_hub",
    ]

    resp = (
        supabase.table("tasks")
        .select("page_name")
        .eq("client_id", client_id)
        .execute()
    )
    rows = resp.data or []
    from_db = [r["page_name"] for r in rows if r.get("page_name")]

    all_pages = sorted(set(base_pages + from_db))
    return all_pages

def upsert_payroll_position(
    client_id,
    position_id,
    role_name,
    employee_name,
    fte,
    base_salary_annual,
    super_rate_pct,
    start_date,
    end_date,
    notes,
) -> bool:
    """
    Insert or update a payroll position row.
    If position_id is None -> insert; otherwise update that row.
    """
    if client_id is None:
        return False

    payload = {
        "client_id": str(client_id),
        "role_name": role_name or None,
        "employee_name": employee_name or None,
        "fte": float(fte) if fte is not None else 1.0,
        "base_salary_annual": float(base_salary_annual or 0),
        "super_rate_pct": float(super_rate_pct or 0),
        # payroll_tax_pct REMOVED (now in client_settings)
        "start_date": start_date.isoformat() if start_date else None,
        "end_date": end_date.isoformat() if end_date else None,
        "notes": notes or None,
    }

    try:
        if position_id is None:
            supabase.table("payroll_positions").insert(payload).execute()
        else:
            supabase.table("payroll_positions").update(payload).eq("id", position_id).execute()
        return True
    except Exception as e:
        print("Error upserting payroll position:", e)
        return False



def delete_payroll_position(position_id) -> bool:
    if not position_id:
        return False
    try:
        supabase.table("payroll_positions").delete().eq("id", position_id).execute()
        return True
    except Exception as e:
        print("Error deleting payroll position:", e)
        return False


def get_month_options(n_months: int = 18):
    """Return a list of month labels (most recent first)."""
    today = date.today().replace(day=1)
    months = []
    for i in range(n_months):
        m = today - pd.DateOffset(months=i)
        months.append(m.strftime("%b %Y"))  # e.g. "Dec 2025"
    return months


def parse_month_label_to_date(label: str) -> date:
    """
    Convert 'Dec 2025' → date(2025, 12, 1)
    """
    return pd.to_datetime(label, format="%b %Y").date().replace(day=1)


def fetch_hiring_monthly_for_client(client_id):
    """
    Fetch hiring + headcount monthly data for this client from hiring_monthly.
    Expected columns (flexible):
      - month_date or month
      - department
      - roles_count
      - new_payroll
      - cumulative_roles
      - cumulative_payroll
    """
    if client_id is None:
        return pd.DataFrame()

    try:
        res = (
            supabase
            .table("hiring_monthly")
            .select("*")
            .eq("client_id", str(client_id))
            .execute()
        )
        rows = res.data or []
    except Exception:
        rows = []

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)

    # Parse month column
    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    elif "month" in df.columns:
        df["month_date"] = pd.to_datetime(df["month"], errors="coerce")

    return df



def build_payroll_monthly_for_client(client_id, month_index: pd.DatetimeIndex) -> pd.DataFrame:
    """
    Returns one row per month with the 'economic' payroll split into:
      - net_salary      (cash paid to employees for that month)
      - tax_payable     (PAYG generated for that month)
      - super_payable   (super accrued for that month)

    Schema:
      month_date   (Timestamp, first of month)
      net_salary   (float)
      tax_payable  (float)
      super_payable(float)
    """
    # TODO: replace this body with your real existing payroll engine logic.
    # The key is: this function should reproduce your Excel schedule
    # per month, including pro-rated salaries for start/end dates.
    df_positions = fetch_payroll_positions_for_client(client_id)

    # ... build a monthly schedule here ...

    # For now, imagine we end with a DataFrame like:
    # month_date | net_salary | tax_payable | super_payable

    payroll_df = ...  # <- your existing logic

    # Normalise month_date to first-of-month Timestamp
    payroll_df["month_date"] = (
        pd.to_datetime(payroll_df["month_date"], errors="coerce")
        .dt.to_period("M")
        .dt.to_timestamp()
    )

    # Only keep the months we care about
    wanted = set(pd.to_datetime(m).to_period("M").to_timestamp() for m in month_index)
    payroll_df = payroll_df[payroll_df["month_date"].isin(wanted)]

    return payroll_df


def compute_payroll_by_month(
    client_id,
    month_index,
) -> dict[pd.Timestamp, float]:
    """
    Compute total payroll cash-out for each month in month_index, with:

      - Pro-rated gross salary by days active in each calendar month.
      - Income tax + Medicare from AU slabs (annual -> monthly -> pro-rated).
      - Superannuation as % of salary (annual -> monthly -> pro-rated).
      - Timing lags:
          * Net salary paid in the same month the work happens.
          * Tax+Medicare cash paid 'payroll_tax_lag_months' later.
          * Super cash paid 'super_lag_months' later.

    Returns:
      dict { month_start (Timestamp, 1st of month) -> total cash out in that month }
    """

    # 1) Normalise month_index -> unique first day of each month (Timestamp)
    month_list: list[pd.Timestamp] = [
        pd.to_datetime(m).to_period("M").to_timestamp()
        for m in month_index
    ]
    if not month_list:
        return {}

    month_list = sorted(set(month_list))
    month_set = set(month_list)

    # 2) Client settings: lags for super + tax remittance
    settings = get_client_settings(client_id)
    super_lag = int(settings.get("super_lag_months", 1))          # e.g. 1 = pay next month
    payroll_tax_lag = int(settings.get("payroll_tax_lag_months", 1))

    # We may need to look back up to the max lag so that
    # the FIRST month in our horizon gets prior-month tax/super.
    lookback = max(super_lag, payroll_tax_lag, 0)

    # 3) Load positions
    df_positions = fetch_payroll_positions_for_client(client_id)
    if df_positions is None or df_positions.empty:
        return {m: 0.0 for m in month_list}

    pos = df_positions.copy()

    # --- DATES: parse in AU format (day-first) ---
    pos["start_date"] = pd.to_datetime(
        pos.get("start_date"),
        errors="coerce",
        dayfirst=True,
    )
    pos["end_date"] = pd.to_datetime(
        pos.get("end_date"),
        errors="coerce",
        dayfirst=True,
    )

    # Drop roles with no start_date
    pos = pos[pos["start_date"].notna()]
    if pos.empty:
        return {m: 0.0 for m in month_list}

    # Numerics
    for col in ["base_salary_annual", "fte", "super_rate_pct"]:
        if col not in pos.columns:
            pos[col] = 0.0
        pos[col] = pd.to_numeric(pos[col], errors="coerce").fillna(0.0)

    # 4) Per-role annual amounts (full year, if fully active)
    #    base salary pro-rated by FTE
    pos["salary_annual_fte"] = pos["base_salary_annual"] * pos["fte"]

    #    annual income tax + Medicare using your slab function
    pos["tax_and_medicare_annual"] = pos["salary_annual_fte"].apply(
        compute_annual_income_tax_au
    )

    #    annual superannuation (e.g. 11%)
    pos["super_annual"] = pos["salary_annual_fte"] * (pos["super_rate_pct"] / 100.0)

    # Convert to "full-month" values (if role active entire month)
    # These correspond to the "Monthly salary full", "Monthly tax payable" and
    # "Monthly Super" columns in your Excel.
    pos["gross_month_full"] = pos["salary_annual_fte"] / 12.0
    pos["tax_month_full"] = pos["tax_and_medicare_annual"] / 12.0
    pos["super_month_full"] = pos["super_annual"] / 12.0

    # 5) Helper: fraction of a given month that the role is active
    def _month_fraction(row, month_start: pd.Timestamp, month_end: pd.Timestamp) -> float:
        role_start = row["start_date"]
        role_end = row["end_date"]

        # Treat no end_date as open-ended
        if pd.isna(role_end):
            role_end = month_end

        # Effective overlap with this month
        start_eff = max(role_start, month_start)
        end_eff = min(role_end, month_end)

        if pd.isna(start_eff) or pd.isna(end_eff):
            return 0.0
        if end_eff < month_start or start_eff > month_end:
            return 0.0

        # Inclusive day count, same as your sheet (e.g. 8 Aug – 19 Aug = 12 days)
        days_in_month = (month_end - month_start).days + 1
        days_active = (end_eff - start_eff).days + 1
        if days_in_month <= 0 or days_active <= 0:
            return 0.0

        return float(days_active / days_in_month)

    # 6) Build an extended month range to capture prior-month work for lags
    if lookback > 0:
        first_month = min(month_list)
        extended_months = [
            (first_month - pd.DateOffset(months=k)) for k in range(lookback, 0, -1)
        ] + month_list
    else:
        extended_months = month_list

    # Per-month "work" values across the extended window
    # These match your Excel's *work month* logic:
    #   gross_work[m] ~ Σ Monthly Salary pro rata (for month m)
    #   tax_work[m]   ~ Σ Monthly tax payable pro rata (for month m)
    #   super_work[m] ~ Σ Monthly super pro rata (for month m)
    gross_work: dict[pd.Timestamp, float] = {m: 0.0 for m in extended_months}
    tax_work: dict[pd.Timestamp, float] = {m: 0.0 for m in extended_months}
    super_work: dict[pd.Timestamp, float] = {m: 0.0 for m in extended_months}

    for m in extended_months:
        month_start = m
        month_end = m + pd.offsets.MonthEnd(0)  # last day of month

        # Fraction of this month each role is active
        fractions = pos.apply(
            lambda row: _month_fraction(row, month_start, month_end),
            axis=1,
        )

        # Work-based amounts for this month
        # (no rounding here, to match your sheet's use of full precision)
        gross_work[m] = float((pos["gross_month_full"] * fractions).sum())
        tax_work[m] = float((pos["tax_month_full"] * fractions).sum())
        super_work[m] = float((pos["super_month_full"] * fractions).sum())

    # 7) Now compute actual cash-out for the ORIGINAL horizon months only
    totals: dict[pd.Timestamp, float] = {}

    for m in month_list:
        # Net salary for work done IN this month (cash to staff now)
        gross_this = gross_work.get(m, 0.0)
        tax_this = tax_work.get(m, 0.0)
        net_salary = gross_this - tax_this

        # Tax cash: for work done 'payroll_tax_lag' months BEFORE this month
        if payroll_tax_lag > 0:
            prior_tax_month = m - pd.DateOffset(months=payroll_tax_lag)
            tax_cash = tax_work.get(prior_tax_month, 0.0)
        else:
            tax_cash = tax_this  # no lag => pay in same month

        # Super cash: for work done 'super_lag' months BEFORE this month
        if super_lag > 0:
            prior_super_month = m - pd.DateOffset(months=super_lag)
            super_cash = super_work.get(prior_super_month, 0.0)
        else:
            super_cash = super_work.get(m, 0.0)

        # Total payroll cash-out for this month:
        #   = net salary (this month's work)
        #   + tax on prior month(s) work (depending on lag)
        #   + super on prior month(s) work (depending on lag)
        totals[m] = net_salary + tax_cash + super_cash

    return totals


def debug_payroll_breakdown_for_months(client_id, months):
    """
    Print per-role gross, tax, net and fractions for the given months
    so you can line them up against your Excel.
    """
    from pprint import pprint

    # Normalise months to first-day-of-month Timestamps
    month_list = sorted(
        {pd.to_datetime(m).to_period("M").to_timestamp() for m in months}
    )

    # Load client settings
    settings = get_client_settings(client_id)
    super_lag = int(settings.get("super_lag_months", 1))
    payroll_tax_lag = int(settings.get("payroll_tax_lag_months", 1))

    # Load positions
    df = fetch_payroll_positions_for_client(client_id)
    if df is None or df.empty:
        print("[DEBUG] No payroll_positions for client:", client_id)
        return

    pos = df.copy()
    pos["start_date"] = pd.to_datetime(pos.get("start_date"), errors="coerce", dayfirst=True)
    pos["end_date"]   = pd.to_datetime(pos.get("end_date"),   errors="coerce", dayfirst=True)

    pos = pos[pos["start_date"].notna()]
    for col in ["base_salary_annual", "fte", "super_rate_pct"]:
        if col not in pos.columns:
            pos[col] = 0.0
        pos[col] = pd.to_numeric(pos[col], errors="coerce").fillna(0.0)

    pos["salary_annual_fte"] = pos["base_salary_annual"] * pos["fte"]
    pos["tax_annual"]        = pos["salary_annual_fte"].apply(compute_annual_income_tax_au)
    pos["super_annual"]      = pos["salary_annual_fte"] * (pos["super_rate_pct"] / 100.0)

    pos["gross_month_full"] = pos["salary_annual_fte"] / 12.0
    pos["tax_month_full"]   = pos["tax_annual"]        / 12.0
    pos["super_month_full"] = pos["super_annual"]      / 12.0

    def month_fraction(row, month_start, month_end):
        role_start = row["start_date"]
        role_end   = row["end_date"]
        if pd.isna(role_end):
            role_end = month_end
        start_eff = max(role_start, month_start)
        end_eff   = min(role_end,   month_end)
        if pd.isna(start_eff) or pd.isna(end_eff):
            return 0.0
        if end_eff < month_start or start_eff > month_end:
            return 0.0
        days_in_month = (month_end - month_start).days + 1
        days_active   = (end_eff   - start_eff).days + 1
        if days_in_month <= 0 or days_active <= 0:
            return 0.0
        return float(days_active / days_in_month)

    for m in month_list:
        month_start = m
        month_end   = m + pd.offsets.MonthEnd(0)
        print("\n========== Payroll breakdown for", month_start.date(), "==========")

        rows = []
        for _, row in pos.iterrows():
            frac = month_fraction(row, month_start, month_end)
            gross = row["gross_month_full"] * frac
            tax   = row["tax_month_full"]   * frac
            net   = gross - tax
            super_amt = row["super_month_full"] * frac

            rows.append({
                "role_name":    row.get("role_name"),
                "employee":     row.get("employee_name"),
                "start_date":   row["start_date"].date() if pd.notna(row["start_date"]) else None,
                "end_date":     row["end_date"].date()   if pd.notna(row["end_date"])   else None,
                "fraction":     frac,
                "gross_work":   gross,
                "tax_work":     tax,
                "net_work":     net,
                "super_work":   super_amt,
            })

        # Pretty print rows
        for r in rows:
            print(
                f"{r['role_name'] or ''} | {r['employee'] or ''} | "
                f"{r['start_date']} -> {r['end_date']} | frac={r['fraction']:.6f} | "
                f"gross={r['gross_work']:.2f} | tax={r['tax_work']:.2f} | "
                f"net={r['net_work']:.2f} | super={r['super_work']:.2f}"
            )

        total_gross = sum(r["gross_work"] for r in rows)
        total_tax   = sum(r["tax_work"]   for r in rows)
        total_net   = sum(r["net_work"]   for r in rows)
        total_super = sum(r["super_work"] for r in rows)

        print(f"---- Totals for {month_start.date()} ----")
        print(f"  gross_work = {total_gross:,.2f}")
        print(f"  tax_work   = {total_tax:,.2f}")
        print(f"  net_work   = {total_net:,.2f}")
        print(f"  super_work = {total_super:,.2f}")


@st.cache_data(ttl=60)
def fetch_baseline_monthly_for_client(client_id):
    """
    Fetch baseline monthly numbers for this client from baseline_monthly.

    Expected (flexible) columns:
      - month_date or month or period
      - closing_cash  (used for the cash curve + CF engine)
      - optional: operating_cf, investing_cf, financing_cf, free_cash_flow
    """
    if client_id is None:
        return pd.DataFrame()

    try:
        res = (
            supabase
            .table("baseline_monthly")
            .select("*")
            .eq("client_id", str(client_id))
            .order("month_date", desc=False)
            .execute()
        )
        rows = res.data or []
    except Exception:
        rows = []

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)

    # Normalise month_date
    month_col = None
    for cand in ["month_date", "month", "period"]:
        if cand in df.columns:
            month_col = cand
            break

    if month_col is None:
        return pd.DataFrame()

    df["month_date"] = pd.to_datetime(df[month_col], errors="coerce")
    df = df[df["month_date"].notna()]
    if df.empty:
        return pd.DataFrame()

    return df


def upload_csv_to_table_for_client(client_id, csv_file, table_name: str):
    """
    Simple helper to upload a CSV into a Supabase table for the current client.
    - Expects the CSV columns to match the table columns (except client_id).
    - Adds/overwrites client_id on each row.
    """
    if client_id is None:
        return False, "No client selected."

    try:
        df = pd.read_csv(csv_file)

        # Force client_id column
        df["client_id"] = str(client_id)

        records = df.to_dict(orient="records")
        if not records:
            return False, "CSV has no rows."

        supabase.table(table_name).insert(records).execute()
        return True, f"Inserted {len(records)} rows into {table_name}."
    except Exception as e:
        return False, f"Upload failed: {e}"

def get_engine_window_for_scenarios(
    client_id: str,
    base_month: date,
    n_months: int = 12,
) -> pd.DataFrame | None:
    """
    Pull a clean 12-month baseline window from cashflow_summary
    for scenarios. Does NOT write anything back.

    Ensures numeric columns and reconstructs opening cash.
    """
    engine_df = fetch_cashflow_summary_for_client_as_of(client_id, base_month)
    if engine_df is None or engine_df.empty:
        return None

    df = engine_df.copy()
    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    df = df.dropna(subset=["month_date"])

    start = pd.to_datetime(base_month).replace(day=1)
    end = start + pd.DateOffset(months=n_months)

    df = df[(df["month_date"] >= start) & (df["month_date"] < end)].copy()
    if df.empty:
        return None

    df = df.sort_values("month_date").reset_index(drop=True)

    # Ensure required numeric fields
    for col in ["operating_cf", "investing_cf", "financing_cf", "free_cash_flow", "closing_cash"]:
        if col not in df.columns:
            df[col] = 0.0
        df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0.0)

    # Reconstruct opening cash series
    closing_0 = float(df.loc[0, "closing_cash"] or 0.0)
    free_0 = float(df.loc[0, "free_cash_flow"] or 0.0)
    opening_0 = closing_0 - free_0

    openings = [opening_0]
    for i in range(1, len(df)):
        openings.append(float(df.loc[i - 1, "closing_cash"] or 0.0))

    df["opening_cash"] = openings

    return df

def assert_single_plan_version(df: pd.DataFrame, context=""):
    if df is None or df.empty:
        return
    if "as_of_month" in df.columns:
        n = df["as_of_month"].nunique(dropna=True)
        if n > 1:
            print(f"[WARN] Mixed plan versions detected ({n}) in {context}. You probably need as_of filtering.")


#try Supabase servcer again if error arise
def sb_execute_with_retry(fn, label: str, tries: int = 3, base_wait: float = 0.6):
    """
    Runs a Supabase query with retry.
    Returns:
      - result (whatever fn returns) on success
      - None on failure after retries
    """
    last_err = None
    for i in range(tries):
        try:
            return fn()
        except Exception as e:
            last_err = e
            print(f"[WARN] {label} attempt {i+1}/{tries} failed: {e}")
            time.sleep(base_wait * (i + 1))

    print(f"[ERROR] {label} failed after {tries} tries: {last_err}")
    return None


def fetch_cashflow_summary_for_client(client_id: str) -> pd.DataFrame | None:
    """
    Fetch cashflow_summary rows for a client.

    Returns:
        - None  -> Supabase connection problem (sb_execute_with_retry returned None)
        - empty DataFrame -> call succeeded but no rows
        - non-empty DataFrame -> valid engine data
    """
    if client_id is None:
        return pd.DataFrame()

    def _call():
        return (
            supabase
            .table("cashflow_summary")
            .select("*")
            .eq("client_id", str(client_id))
            .order("month_date")
            .execute()
        )

    res = sb_execute_with_retry(_call, label="fetch cashflow_summary")
    if res is None:
        print("[ERROR] fetch_cashflow_summary_for_client -> connection problem (None)")
        return None

    rows = res.data or []
    if not rows:
        print("[DEBUG] fetch_cashflow_summary_for_client -> 0 rows for", client_id)
        return pd.DataFrame()

    df = pd.DataFrame(rows)

    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
        df = df[df["month_date"].notna()].sort_values("month_date")

    print(f"[DEBUG] fetch_cashflow_summary_for_client -> {len(df)} rows for {client_id}")
    return df


def fetch_kpis_for_client_month(
    client_id: str,
    month_date: date | datetime,
) -> dict | None:
    """
    Get a single KPI row for (client, month_date) from kpi_monthly.

    If nothing exists, derive a KPI row from:
      - AR (cash in)
      - AP + payroll (+ optional opex) (cash out)
      - cashflow_summary (closing cash + runway)

    and just RETURN it (no upsert).

    IMPORTANT:
      - We treat 'month_date' purely as a focus month
      - Money in/out are for THAT month only
      - Overdue / ageing is computed as of month END
    """
    if client_id is None or month_date is None:
        return None

    # --- Normalise month + window ---
    # month_dt is 1st of the focus month
    month_dt = pd.to_datetime(month_date).date().replace(day=1)
    month_iso = month_dt.isoformat()

    # start = 1st, end = last day of that month
    month_start, month_end = get_focus_month_window(month_dt)
    base_month_ts = pd.to_datetime(month_start)

    # ---------- 1) Try to read from kpi_monthly directly ----------
    try:
        resp = (
            supabase.table("kpi_monthly")
            .select("*")
            .eq("client_id", client_id)
            .eq("month_date", month_iso)
            .execute()
        )

        rows = []
        if hasattr(resp, "data"):
            rows = resp.data or []
        elif isinstance(resp, dict):
            rows = resp.get("data", []) or []

        if rows:
            return rows[0]

    except Exception as e:
        print("Error fetching from kpi_monthly:", e)
        # Don't return yet — we can still derive

    # ---------- 2) Derive KPIs from AR/AP + payroll + engine ----------
    try:
        # Settings + AR/AP
        settings = get_client_settings(client_id) or {}
        ar_days = int(settings.get("ar_default_days", 30))
        ap_days = int(settings.get("ap_default_days", 30))

        df_ar, df_ap = fetch_ar_ap_for_client(client_id)

        # ----- AR: expected cash-in THIS MONTH (Money In) -----
        money_in = 0.0
        if df_ar is not None and not df_ar.empty:
            df_ar = df_ar.copy()
            # Age as at month END (not start, not today)
            df_ar = add_ar_aging(df_ar, as_of=month_end, ar_default_days=ar_days)

            df_ar["expected_date"] = pd.to_datetime(
                df_ar.get("expected_date"), errors="coerce"
            )
            df_ar["amount"] = pd.to_numeric(
                df_ar.get("amount"), errors="coerce"
            ).fillna(0.0)

            # Ignore fully paid/closed AR
            if "status" in df_ar.columns:
                df_ar = df_ar[
                    ~df_ar["status"]
                    .astype(str)
                    .str.lower()
                    .isin(["paid", "closed", "settled"])
                ]

            # Only rows whose expected_date lies in THIS month
            mask_ar = df_ar["expected_date"].dt.to_period("M") == base_month_ts.to_period("M")
            money_in = float(df_ar.loc[mask_ar, "amount"].sum())

        # ----- AP: expected cash-out THIS MONTH (Money Out) -----
        ap_cash = 0.0
        if df_ap is not None and not df_ap.empty:
            df_ap = df_ap.copy()
            # Age as at month END
            df_ap = add_ap_aging(df_ap, as_of=month_end, ap_default_days=ap_days)

            # Decide which date column we use for "cash out" timing
            if "pay_expected_date" in df_ap.columns:
                pay_col = "pay_expected_date"
            elif "expected_payment_date" in df_ap.columns:
                pay_col = "expected_payment_date"
            else:
                # IMPORTANT: we now use due_date to align with the engine
                pay_col = "due_date"

            df_ap[pay_col] = pd.to_datetime(df_ap.get(pay_col), errors="coerce")
            df_ap["amount"] = pd.to_numeric(
                df_ap.get("amount"), errors="coerce"
            ).fillna(0.0)

            # Ignore already paid AP
            if "status" in df_ap.columns:
                df_ap = df_ap[
                    ~df_ap["status"]
                    .astype(str)
                    .str.lower()
                    .isin(["paid", "closed", "settled"])
                ]

            # Only rows whose pay date lies in THIS month
            mask_ap = df_ap[pay_col].dt.to_period("M") == base_month_ts.to_period("M")
            ap_cash = float(df_ap.loc[mask_ap, "amount"].sum())

        # ----- Payroll this month -----
        payroll_cash = 0.0
        try:
            payroll_by_month = compute_payroll_by_month(client_id, [base_month_ts])
            payroll_cash = float(payroll_by_month.get(base_month_ts, 0.0))
        except Exception as e:
            print("[WARN] compute_payroll_by_month failed:", e)
            payroll_cash = 0.0

        # ----- Optional Opex this month -----
        opex_cash = 0.0
        try:
            df_opex = fetch_opex_for_client(client_id)
            if df_opex is not None and not df_opex.empty:
                df_opex = df_opex.copy()
                df_opex["month_bucket"] = (
                    pd.to_datetime(df_opex.get("month_date"), errors="coerce")
                    .dt.to_period("M")
                    .dt.to_timestamp()
                )
                mask_ox = df_opex["month_bucket"] == base_month_ts
                opex_cash = float(
                    pd.to_numeric(
                        df_opex.loc[mask_ox, "amount"], errors="coerce"
                    )
                    .fillna(0.0)
                    .sum()
                )
        except Exception:
            opex_cash = 0.0

        money_out = ap_cash + payroll_cash + opex_cash

        # ----- Engine: closing cash + runway for THIS month -----
        closing_cash = None
        runway_months = None
        effective_burn = None

        engine_df = fetch_cashflow_summary_for_client_as_of(client_id, month_dt)

        if engine_df is not None and not engine_df.empty:
            df_eng = engine_df.copy()
            df_eng["month_date"] = pd.to_datetime(
                df_eng["month_date"], errors="coerce"
            ).dt.date
            row_eng = df_eng[df_eng["month_date"] == month_dt]

            if not row_eng.empty and "closing_cash" in row_eng.columns:
                closing_cash = float(row_eng["closing_cash"].iloc[0] or 0.0)

            # IMPORTANT: runway based on THIS month, not global sidebar
            runway_months, effective_burn = compute_runway_and_effective_burn_from_df(
                engine_df,
                month_dt,   # <= key change here
            )

        kpi_row = {
            "client_id": client_id,
            "month_date": month_iso,
            # Names used on Business overview:
            # Money in  -> revenue
            # Money out -> burn
            "revenue": float(money_in),
            "burn": float(money_out),
            "cash_balance": closing_cash,      # None if engine unavailable
            "runway_months": runway_months,    # None if engine unavailable
        }

        print(
            f"Derived KPIs from cashflow_summary/AR/AP for client={client_id}, "
            f"month_date={month_iso}: {kpi_row}"
        )

        return kpi_row

    except Exception as e:
        print("Error deriving KPIs from cashflow_summary:", e)
        return None


def fetch_pipeline_for_client(client_id):
    """
    Fetch all revenue pipeline deals for the given client.
    Adds a normalized numeric column: contract_months_norm
    """
    if client_id is None:
        return pd.DataFrame()

    try:
        res = (
            supabase
            .table("revenue_pipeline")
            .select("*")
            .eq("client_id", str(client_id))
            .order("created_at", desc=True)
            .execute()
        )
        rows = res.data or []
    except Exception:
        rows = []

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)

    # Parse dates
    for col in ["start_month", "end_month", "created_at", "updated_at"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce")

    # ---------- Normalize contract months ----------
    # Accept multiple possible column names (DB drift protection)
    candidates = [
        "contract_months",
        "contract_length_months",
        "contract_length",
        "contract_term_months",
        "term_months",
    ]

    def _pick_contract_months(row):
        for c in candidates:
            if c in row and pd.notna(row.get(c)):
                v = row.get(c)
                # handle strings like "24", "24 months"
                try:
                    v = str(v).strip().lower().replace("months", "").replace("month", "").strip()
                except Exception:
                    pass
                try:
                    n = int(float(v))
                    if n > 0:
                        return n
                except Exception:
                    continue
        return None

    df["contract_months_norm"] = df.apply(_pick_contract_months, axis=1)

    return df



#POC Revenue Recongition
@st.cache_data(ttl=60)
def fetch_poc_progress_for_client(client_id) -> pd.DataFrame:
    """
    Fetch all percentage-of-completion rows for this client.
    Each row is cumulative pct_complete for a given deal + month_date.
    """
    if client_id is None:
        return pd.DataFrame()

    try:
        res = (
            supabase
            .table("deal_poc_progress")
            .select("*")
            .eq("client_id", str(client_id))
            .order("deal_id")
            .order("month_date")
            .execute()
        )
        rows = res.data or []
    except Exception as e:
        print("Error fetching POC progress:", e)
        rows = []

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)
    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    if "pct_complete" in df.columns:
        df["pct_complete"] = pd.to_numeric(df["pct_complete"], errors="coerce").fillna(0.0)
    return df

def upsert_poc_progress_row(
    client_id,
    deal_id: int,
    month_date: date,
    pct_complete: float,
    notes: str | None = None,
) -> bool:
    """
    Upsert a single POC row for (deal_id, month_date).
    pct_complete is cumulative 0–100.
    """
    if client_id is None or deal_id is None or month_date is None:
        return False

    try:
        payload = {
            "client_id": str(client_id),
            "deal_id": int(deal_id),
            "month_date": month_date.isoformat(),
            "pct_complete": float(pct_complete),
            "notes": notes or None,
        }

        (
            supabase
            .table("deal_poc_progress")
            .upsert(payload, on_conflict="deal_id,month_date")
            .execute()
        )

        # Clear cache so engine picks up new % immediately
        st.cache_data.clear()
        return True
    except Exception as e:
        print("Error upserting POC progress:", e)
        return False



#Milstone Revenue Recongition
@st.cache_data(ttl=60)
def fetch_milestones_for_deal(client_id, deal_id) -> pd.DataFrame:
    """
    Load all milestone rows for a single deal.
    """
    if client_id is None or deal_id is None:
        return pd.DataFrame()

    try:
        res = (
            supabase
            .table("deal_milestones")
            .select("*")
            .eq("client_id", str(client_id))
            .eq("deal_id", int(deal_id))
            .order("month_date")
            .execute()
        )
        rows = res.data or []
    except Exception:
        rows = []

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)
    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce").dt.date
    df["percent_of_contract"] = pd.to_numeric(
        df.get("percent_of_contract"), errors="coerce"
    ).fillna(0.0)

    return df


def save_milestones_for_deal(client_id, deal_id, milestones_df: pd.DataFrame) -> bool:
    """
    Replace milestones for this deal with the rows from milestones_df.
    Expected columns in milestones_df:
      - month_date (date or string)
      - percent_of_contract (float, 0–100)
      - milestone_name (optional)
    """
    if client_id is None or deal_id is None:
        return False

    if milestones_df is None or milestones_df.empty:
        # nothing to save
        return False

    df = milestones_df.copy()

    # Normalise and clean
    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce").dt.date
    df["percent_of_contract"] = pd.to_numeric(
        df.get("percent_of_contract"), errors="coerce"
    ).fillna(0.0)

    df = df[df["month_date"].notna() & (df["percent_of_contract"] > 0)]
    if df.empty:
        return False

    # Build payload
    payload = []
    for _, row in df.iterrows():
        payload.append(
            {
                "client_id": str(client_id),
                "deal_id": int(deal_id),
                "month_date": row["month_date"].isoformat(),
                "percent_of_contract": float(row["percent_of_contract"]),
                "milestone_name": row.get("milestone_name") or None,
            }
        )

    try:
        # Delete existing milestones for this deal
        supabase.table("deal_milestones").delete().eq(
            "client_id", str(client_id)
        ).eq("deal_id", int(deal_id)).execute()

        # Insert new ones
        supabase.table("deal_milestones").insert(payload).execute()

        # Clear cache so fetch_milestones_for_deal + revenue schedule recompute
        st.cache_data.clear()

        return True
    except Exception as e:
        print("Error saving milestones:", e)
        return False


def build_expected_revenue_curve(
    df_pipeline: pd.DataFrame,
    base_month: date,
    n_months: int = 12,
):
    """
    Build a simple expected revenue curve:
    For each deal, we place value_total * (probability_pct/100)
    into its start_month (or created_at month if start_month missing).
    """
    if df_pipeline is None or df_pipeline.empty:
        return pd.DataFrame()

    # Generate month buckets starting from base_month
    month_index = pd.date_range(
        start=pd.to_datetime(base_month),
        periods=n_months,
        freq="MS",  # month start
    )
    months = {m: 0.0 for m in month_index}

    for _, row in df_pipeline.iterrows():
        # Decide which month to use
        start_dt = row.get("start_month")
        if pd.isna(start_dt):
            start_dt = row.get("created_at")

        if pd.isna(start_dt):
            continue

        # Snap to first day of month
        month_key = pd.to_datetime(start_dt).to_period("M").to_timestamp()

        if month_key not in months:
            # Ignore deals outside our window for now
            continue

        value = row.get("value_total") or 0.0
        prob = row.get("probability_pct") or 0.0
        expected_value = float(value) * float(prob) / 100.0

        months[month_key] += expected_value

    # Build DataFrame for chart
    out_df = pd.DataFrame(
        {
            "month_date": month_index,
            "month_label": [m.strftime("%b %Y") for m in month_index],
            "expected_revenue": [months[m] for m in month_index],
        }
    )
    return out_df

def get_recognised_revenue_for_month(
        
    client_id: str,
    focus_month: date,
) -> float:
    """
    Return total recognised revenue for the given focus_month
    using the same schedule logic as page_sales_deals.

    Uses build_revenue_schedule_for_client(...).
    """
    if not client_id or focus_month is None:
        return 0.0

    snap_df = build_revenue_schedule_for_client(
        client_id,
        base_month=focus_month,
        n_months=12,  # same as Sales & Deals
    )

    if snap_df is None or snap_df.empty:
        return 0.0

    df = snap_df.copy()
    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    df = df[df["month_date"].notna()]

    # Decide which column holds the recognised revenue
    value_col = None
    for cand in ["recognised_revenue", "revenue_amount", "amount"]:
        if cand in df.columns:
            value_col = cand
            break

    if value_col is None:
        return 0.0

    # Bucket to the first of the focus month
    focus_start = pd.to_datetime(focus_month).replace(day=1)

    monthly = (
        df.groupby("month_date", as_index=False)[value_col]
        .sum()
        .sort_values("month_date")
    )

    this_row = monthly[monthly["month_date"] == focus_start]
    this_rev = float(this_row[value_col].iloc[0]) if not this_row.empty else 0.0
    return this_rev



# ---------- Revenue recognition engine (6 modes) ----------

def _determine_method_key(method_raw: str) -> str:
    """
    Take whatever is stored in revenue_pipeline.method and normalise to a key:
      'saas', 'milestone', 'poc', 'straight', 'usage', 'point'
    """
    m = (method_raw or "").lower()

    if "usage" in m:
        return "usage"
    if "percent" in m or "poc" in m:
        return "poc"
    if "straight" in m or "service" in m:
        return "straight"
    if "point" in m or "goods" in m:
        return "point"
    if "mile" in m or "project" in m:
        return "milestone"
    # default: treat as SaaS
    return "saas"
def _append_revenue_row(
    rows: list,
    window_start: pd.Timestamp,
    window_end: pd.Timestamp,
    deal_row: pd.Series,
    method_label: str,
    month_value,
    amount: float,
    contract_months_used: int | None = None,
):
    if amount is None:
        return
    try:
        amount = float(amount)
    except Exception:
        return
    if abs(amount) < 1e-6:
        return

    if pd.isna(month_value):
        return

    month_ts = pd.to_datetime(month_value, errors="coerce")
    if pd.isna(month_ts):
        return

    month_ts = month_ts.to_period("M").to_timestamp()

    if not (window_start <= month_ts < window_end):
        return

    rows.append(
        {
            "month_date": month_ts,
            "revenue_amount": amount,
            "revenue_type": method_label,
            "deal_id": deal_row.get("id"),
            "deal_name": deal_row.get("deal_name"),
            "customer_name": deal_row.get("customer_name"),
            "contract_months_used": contract_months_used,
        }
    )

def _get_contract_months(deal: pd.Series) -> int | None:
    candidates = [
        "contract_months",
        "contract_length_months",
        "contract_length",
        "contract_term_months",
        "term_months",
        "contract_length_in_months",
    ]
    for c in candidates:
        if c in deal.index:
            v = deal.get(c)
            if v is None or (isinstance(v, float) and pd.isna(v)):
                continue
            try:
                s = str(v).strip().lower().replace("months", "").replace("month", "").strip()
                n = int(float(s))
                if n > 0:
                    return n
            except Exception:
                continue
    return None



def _month_span_inclusive(start_ts: pd.Timestamp, end_ts: pd.Timestamp) -> int:
    """
    Number of months inclusive between start and end, month-bucketed.
    Example: Jan to Mar = 3.
    """
    s = pd.to_datetime(start_ts).to_period("M")
    e = pd.to_datetime(end_ts).to_period("M")
    return int((e - s).n) + 1


def _schedule_saas(
    deal: pd.Series,
    window_start: pd.Timestamp,
    window_end: pd.Timestamp,
    rows: list,
    prob_factor: float,
    default_months: int | None = None,
    require_explicit_months: bool = False,
):
    """
    SaaS subscription revenue:
      - value_total spread evenly over contract months (NO forced 12)
      - optional annual uplift (kept, but uplift affects totals unless you model price increases intentionally)
    """
    total = float(deal.get("value_total") or 0.0) * prob_factor
    if total == 0:
        return

    start_raw = deal.get("start_month") or deal.get("created_at")
    start_ts = pd.to_datetime(start_raw, errors="coerce")
    if pd.isna(start_ts):
        return
    start_ts = start_ts.to_period("M").to_timestamp()

    contract_len = _get_contract_months(deal)

    # ✅ Only fallback in forecast mode (if allowed)
    if (contract_len is None or contract_len <= 0):
        if require_explicit_months:
            return
        contract_len = int(default_months or 12)

    uplift = 0.0
    try:
        uplift = float(deal.get("annual_uplift_pct") or 0.0)
    except Exception:
        uplift = 0.0

    base_mrr = total / contract_len if contract_len else 0.0

    for i in range(contract_len):
        month_ts = start_ts + pd.DateOffset(months=i)

        year_idx = i // 12
        factor = (1 + uplift / 100.0) ** year_idx if uplift > 0 else 1.0
        amount = base_mrr * factor

        _append_revenue_row(
            rows, window_start, window_end, deal,
            "SaaS subscription",
            month_ts, amount,
            contract_months_used=contract_len,
        )


def _schedule_straight_line(
    deal: pd.Series,
    window_start: pd.Timestamp,
    window_end: pd.Timestamp,
    rows: list,
    prob_factor: float,
    default_months: int | None = None,
    require_explicit: bool = False,
):
    """
    Straight-line service revenue:
      - If end month exists: spread from start->end (inclusive)
      - Else spread over contract months
      - Else optional fallback to default months (forecast only)
    """
    total = float(deal.get("value_total") or 0.0) * prob_factor
    if total == 0:
        return

    start_raw = deal.get("start_month") or deal.get("created_at")
    start_ts = pd.to_datetime(start_raw, errors="coerce")
    if pd.isna(start_ts):
        return
    start_ts = start_ts.to_period("M").to_timestamp()

    end_raw = deal.get("service_end_month") or deal.get("end_month")
    end_ts = pd.to_datetime(end_raw, errors="coerce")

    # Prefer explicit end-date span
    if not pd.isna(end_ts):
        end_ts = end_ts.to_period("M").to_timestamp()
        if end_ts < start_ts:
            return
        months_to_spread = _month_span_inclusive(start_ts, end_ts)
    else:
        # Else use explicit duration
        months_to_spread = _get_contract_months(deal)

        # Forecast-only fallback
        if (months_to_spread is None or months_to_spread <= 0):
            if require_explicit:
                return
            months_to_spread = int(default_months or 6)

    per_month = total / months_to_spread if months_to_spread else 0.0

    for i in range(months_to_spread):
        month_ts = start_ts + pd.DateOffset(months=i)
        _append_revenue_row(
            rows, window_start, window_end, deal,
            "Straight-line service",
            month_ts, per_month,
            contract_months_used=months_to_spread,
        )

@st.cache_data(ttl=60)
def build_revenue_schedule_for_client(
    client_id,
    base_month: date,
    n_months: int = 12,
    mode: str = "forecast",  # "forecast" or "secured"
) -> pd.DataFrame:

    if client_id is None or base_month is None:
        return pd.DataFrame()

    settings = get_client_settings(client_id) or {}
    min_prob_pct = float(settings.get("min_revenue_prob_pct", 20) or 20)

    window_start = pd.to_datetime(base_month).to_period("M").to_timestamp()
    window_end = window_start + pd.DateOffset(months=n_months)

    df = fetch_pipeline_for_client(client_id)
    if df is None or df.empty:
        return pd.DataFrame()

    df = df.copy()
    df["stage"] = df.get("stage", "").astype(str).str.lower().fillna("")
    df["value_total"] = pd.to_numeric(df.get("value_total"), errors="coerce").fillna(0.0)
    df["probability_pct"] = pd.to_numeric(df.get("probability_pct"), errors="coerce").fillna(0.0)

    for col in ["start_month", "end_month", "service_end_month", "delivery_date", "created_at"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce")

    # MODE FILTERING / PROBABILITY
    if mode == "secured":
        df = df[df["stage"] == "closed_won"].copy()
        if df.empty:
            return pd.DataFrame()
        df["prob_factor"] = 1.0
        require_explicit = True
    else:
        def _prob_factor(row):
            stg = str(row.get("stage", "")).lower()
            if stg == "closed_lost":
                return 0.0
            if stg == "closed_won":
                return 1.0
            try:
                return max(0.0, min(float(row.get("probability_pct", 0.0)) / 100.0, 1.0))
            except Exception:
                return 0.0

        df["prob_factor"] = df.apply(_prob_factor, axis=1)
        df["prob_pct_eff"] = df["prob_factor"] * 100.0
        df = df[df["prob_pct_eff"] >= min_prob_pct].copy()
        if df.empty:
            return pd.DataFrame()
        require_explicit = False

    rows: list[dict] = []

    for _, deal in df.iterrows():
        prob_factor = float(deal.get("prob_factor") or 0.0)
        if prob_factor <= 0:
            continue

        method_raw = deal.get("method") or settings.get("revenue_recognition_method") or "saas"
        method_key = _determine_method_key(method_raw)

        if method_key == "saas":
            _schedule_saas(
                deal, window_start, window_end, rows,
                prob_factor=prob_factor,
                default_months=int(settings.get("saas_default_months", 12) or 12),
                require_explicit_months=require_explicit,
            )

        elif method_key == "straight_line":
            _schedule_straight_line(
                deal, window_start, window_end, rows,
                prob_factor=prob_factor,
                default_months=int(settings.get("project_default_months", 6) or 6),
                require_explicit=require_explicit,
            )

        elif method_key == "milestone":
            _schedule_milestone(
                client_id, deal, window_start, window_end, rows,
                prob_factor=prob_factor,
                require_explicit=require_explicit,
            )

        elif method_key == "poc":
            _schedule_poc(
                client_id, deal, window_start, window_end, rows,
                prob_factor=prob_factor,
                require_explicit=require_explicit,
            )

        elif method_key == "usage":
            _schedule_usage(deal, window_start, window_end, rows, prob_factor)

        elif method_key == "point_in_time":
            _schedule_point_in_time(deal, window_start, window_end, rows, prob_factor)

        else:
            # safest fallback
            _schedule_straight_line(
                deal, window_start, window_end, rows,
                prob_factor=prob_factor,
                default_months=None,
                require_explicit=True,
            )

    if not rows:
        return pd.DataFrame()

    out_df = pd.DataFrame(rows)
    out_df["month_date"] = pd.to_datetime(out_df["month_date"], errors="coerce")
    out_df = out_df[out_df["month_date"].notna()].sort_values("month_date")
    out_df["month_label"] = out_df["month_date"].dt.strftime("%b %Y")
    return out_df



def _schedule_poc(
    client_id,
    deal: pd.Series,
    window_start,
    window_end,
    rows: list,
    prob_factor: float,
    require_explicit: bool = False,
):
    """
    POC revenue using deal_poc_progress table (cumulative %):
      revenue_this_month = (pct_complete_this_month - pct_complete_prev_month) * total
    """
    total = float(deal.get("value_total") or 0.0) * prob_factor
    if total == 0:
        return

    deal_id = deal.get("id")
    poc_df = fetch_poc_progress_for_client(client_id)
    if poc_df is None or poc_df.empty:
        if require_explicit:
            return
        _schedule_straight_line(deal, window_start, window_end, rows, prob_factor)
        return

    poc_df = poc_df.copy()
    poc_df = poc_df[poc_df["deal_id"] == deal_id].copy()
    if poc_df.empty:
        if require_explicit:
            return
        _schedule_straight_line(deal, window_start, window_end, rows, prob_factor)
        return

    poc_df["month_date"] = pd.to_datetime(poc_df["month_date"], errors="coerce")
    poc_df["pct_complete"] = pd.to_numeric(poc_df["pct_complete"], errors="coerce").fillna(0.0).clip(0, 100)
    poc_df = poc_df[poc_df["month_date"].notna()].sort_values("month_date")

    prev_pct = 0.0
    for _, r in poc_df.iterrows():
        pct = float(r["pct_complete"])
        incr = pct - prev_pct
        prev_pct = pct
        if incr <= 0:
            continue

        amt = total * (incr / 100.0)
        _append_revenue_row(
            rows, window_start, window_end, deal,
            "Percentage-of-completion",
            r["month_date"],
            amt
        )



def _schedule_usage(deal: pd.Series, window_start, window_end, rows: list, prob_factor: float):
    """
    Usage-based revenue:
      usage_data JSONB = list of {month: '2025-12-01', units: 320, rate: 0.1}
    """
    usage_data = deal.get("usage_data")
    if not usage_data:
        return

    for item in usage_data:
        try:
            units = float(item.get("units", 0.0))
            rate = float(item.get("rate", 0.0))
        except Exception:
            continue
        amount = units * rate * prob_factor
        month_val = item.get("month") or item.get("month_date")
        _append_revenue_row(
            rows,
            window_start,
            window_end,
            deal,
            "Usage-based",
            month_val,
            amount,
        )


def _schedule_point_in_time(deal: pd.Series, window_start, window_end, rows: list, prob_factor: float):
    """
    Point-in-time goods delivery:
      - if delivery_date present, use that
      - else use start_month
    """
    total = float(deal.get("value_total") or 0.0) * prob_factor
    if total == 0:
        return

    delivery_raw = deal.get("delivery_date") or deal.get("start_month")
    if pd.isna(delivery_raw):
        return

    _append_revenue_row(
        rows,
        window_start,
        window_end,
        deal,
        "Point-in-time goods",
        delivery_raw,
        total,
    )



def get_effective_probability(row: dict, min_prob: float) -> float:
    """
    Compute effective probability for revenue recognition.

    - closed_lost -> 0
    - closed_won  -> 100
    - other stages -> probability_pct if >= min_prob, else 0 (paused)
    """
    stage = str(row.get("stage", "")).lower()
    prob_raw = row.get("probability_pct", 0) or 0
    try:
        prob = float(prob_raw)
    except Exception:
        prob = 0.0

    if stage == "closed_lost":
        return 0.0
    if stage == "closed_won":
        return 100.0
    if prob < min_prob:
        return 0.0
    return prob

# --- NEW/UPDATED: Board narrative helpers (no impact on existing features) ---

def _fmt_month(ts) -> str:
    if ts is None:
        return "—"
    try:
        return pd.to_datetime(ts).strftime("%b %Y")
    except Exception:
        return "—"

def _months_between(a: pd.Timestamp | None, b: pd.Timestamp | None) -> int | None:
    if a is None or b is None:
        return None
    try:
        pa = pd.to_datetime(a).to_period("M")
        pb = pd.to_datetime(b).to_period("M")
        return int(pb.ordinal - pa.ordinal)
    except Exception:
        return None

def _detect_debt_used(fin_df: pd.DataFrame, focus_ts: pd.Timestamp) -> bool:
    if fin_df is None or fin_df.empty:
        return False
    f = fin_df.copy()
    if "month_date" not in f.columns:
        return False
    f["month_date"] = pd.to_datetime(f["month_date"], errors="coerce")
    f = f[f["month_date"].notna()].copy()
    f = f[(f["month_date"] >= focus_ts) & (f["month_date"] < (focus_ts + pd.DateOffset(months=12)))]
    if f.empty:
        return False

    cat = f["category"].astype(str).str.lower() if "category" in f.columns else pd.Series([""] * len(f))
    notes = f["notes"].astype(str).str.lower() if "notes" in f.columns else pd.Series([""] * len(f))
    text = (cat.fillna("") + " " + notes.fillna("")).str.lower()
    return text.str.contains("loan|interest|repay|covenant|debt").any()

def _build_scenario_impact_bullet(
    *,
    scenario_summaries: list[dict] | None,
    base_runway: float | None,
    base_cliff: pd.Timestamp | None,
) -> str | None:
    """
    scenario_summaries expected shape (per scenario):
      {
        "name": "Scenario A",
        "runway_months": 8.5,             # months to buffer breach
        "stress_date_exact": <datetime>   # exact or estimated
      }
    Returns a single CFO-grade bullet (brief) comparing up to 2 scenarios vs base.
    """
    if not scenario_summaries:
        return None

    # Keep only valid items with at least one meaningful KPI
    valid = []
    for s in scenario_summaries:
        if not isinstance(s, dict):
            continue
        nm = str(s.get("name") or "").strip()
        rw = s.get("runway_months")
        cd = s.get("stress_date_exact")
        if not nm:
            continue
        if rw is None and cd is None:
            continue
        valid.append({"name": nm, "runway_months": rw, "stress_date_exact": cd})

    if not valid:
        return None

    # Compare only top 1–2 scenarios (A/B). Prefer higher runway if available.
    def _sort_key(x):
        rw = x["runway_months"]
        return -(float(rw) if rw is not None else -1e9)

    valid = sorted(valid, key=_sort_key)[:2]

    parts = []
    for s in valid:
        srw = s["runway_months"]
        scliff = pd.to_datetime(s["stress_date_exact"]) if s["stress_date_exact"] is not None else None

        # runway delta
        runway_phrase = None
        if srw is not None and base_runway is not None:
            delta = float(srw) - float(base_runway)
            if abs(delta) < 0.05:
                runway_phrase = f"runway is ~unchanged at **{float(srw):.1f} months**"
            elif delta > 0:
                runway_phrase = f"extends runway to **{float(srw):.1f} months** (+{delta:.1f}m)"
            else:
                runway_phrase = f"reduces runway to **{float(srw):.1f} months** ({delta:.1f}m)"
        elif srw is not None:
            runway_phrase = f"runway is **{float(srw):.1f} months**"

        # cliff shift
        cliff_phrase = None
        if scliff is not None:
            if base_cliff is not None:
                shift_m = _months_between(base_cliff, scliff)
                if shift_m is None:
                    cliff_phrase = f"moves the cash stress point to **{_fmt_month(scliff)}**"
                elif shift_m > 0:
                    cliff_phrase = f"pushes cash stress to **{_fmt_month(scliff)}** (+{shift_m}m)"
                elif shift_m < 0:
                    cliff_phrase = f"brings cash stress forward to **{_fmt_month(scliff)}** ({shift_m}m)"
                else:
                    cliff_phrase = f"keeps cash stress around **{_fmt_month(scliff)}**"
            else:
                cliff_phrase = f"sets cash stress around **{_fmt_month(scliff)}**"

        # Build per-scenario micro-summary
        bits = [b for b in [runway_phrase, cliff_phrase] if b]
        if bits:
            parts.append(f"{s['name']}: " + "; ".join(bits))

    if not parts:
        return None

    # Single bullet, board-grade, no fluff
    return "Scenario impact: " + " | ".join(parts)

def _build_what_this_means_bullets(
    *,
    runway_months_to_buffer: float | None,
    cliff_ts_exact: pd.Timestamp | None,
    cash_buffer: float,
    debt_used: bool,
    scenario_summaries: list[dict] | None = None,   # ✅ NEW
) -> list[str]:
    bullets: list[str] = []

    # 1) Survival risk
    if cash_buffer > 0:
        if cliff_ts_exact is not None:
            bullets.append(
                f"Survival risk: Under the base case, cash stress (below buffer) occurs around **{_fmt_month(cliff_ts_exact)}**."
            )
        else:
            bullets.append(
                "Survival risk: Under the base case, cash stays above the safety buffer across the current 12-month window."
            )
    else:
        bullets.append(
            "Survival risk: No safety buffer is set, so an ‘operationally unsafe’ cliff cannot be assessed defensibly."
        )

    # 2) Timing pressure
    if cash_buffer > 0 and cliff_ts_exact is not None:
        start_raise_ts = (pd.to_datetime(cliff_ts_exact) - pd.DateOffset(months=3)).to_period("M").to_timestamp()
        bullets.append(
            f"Timing pressure: If a raise is required, fundraising should begin no later than **{_fmt_month(start_raise_ts)}** to stay ahead of the buffer breach."
        )
    else:
        bullets.append(
            "Timing pressure: Even when a cliff is not visible in-window, funding cycles and documentation lead times still create timing risk."
        )

    # 3) Funding sufficiency
    if cash_buffer > 0:
        if runway_months_to_buffer is not None:
            bullets.append(
                f"Funding sufficiency: Runway to buffer breach is **{runway_months_to_buffer:.1f} months** under the base case."
            )
        else:
            bullets.append(
                "Funding sufficiency: Runway to buffer breach is not available yet — validate inputs and rebuild the 14-week cash table."
            )
    else:
        bullets.append(
            "Funding sufficiency: Set a cash safety buffer to measure runway in a board-defensible way (buffer breach, not zero cash)."
        )

    # 4) Trade-off clarity OR Scenario impact (replacement rule)
    scenario_bullet = _build_scenario_impact_bullet(
        scenario_summaries=scenario_summaries,
        base_runway=runway_months_to_buffer,
        base_cliff=cliff_ts_exact,
    )
    if scenario_bullet:
        bullets.append(scenario_bullet)
    else:
        if debt_used:
            bullets.append(
                "Trade-off clarity: Debt improves near-term liquidity but introduces repayment/covenant risk that must be stress-tested before committing."
            )
        else:
            bullets.append(
                "Trade-off clarity: Runway is primarily a function of burn control and the timing/size of funding — small timing shifts can materially change outcomes."
            )

    # 5) Decision priority
    if cash_buffer > 0 and runway_months_to_buffer is not None and runway_months_to_buffer < 6:
        bullets.append(
            "Decision priority: The most defensible path is to **start fundraising immediately** while tightening burn levers to protect runway."
        )
    else:
        bullets.append(
            "Decision priority: Lock a funding plan (amount + timing + instrument) and align spend/hiring decisions to runway protection."
        )

    return bullets[:5]

def _get_saved_scenario_summaries_safe(selected_client_id, selected_month_start) -> list[dict]:
    """
    Non-breaking:
    - If your scenario module puts summaries into st.session_state, we use it.
    - Otherwise returns [] and the page falls back to Trade-off clarity.
    Expected list item shape:
      {"name": "...", "runway_months": 8.5, "stress_date_exact": <dt>}
    """
    try:
        ss = st.session_state.get("scenario_summaries")
        if isinstance(ss, list):
            return ss
    except Exception:
        pass

    # If you later add a real DB fetch, plug it here:
    # try:
    #     return fetch_saved_scenarios_summary(selected_client_id, selected_month_start) or []
    # except Exception:
    #     return []

    return []


def _build_action_plan_30_90(
    *,
    runway_months_to_buffer: float | None,
    cliff_ts_exact: pd.Timestamp | None,
    cash_buffer: float,
    debt_used: bool,
) -> list[str]:
    """
    4–6 time-bound, board-defensible actions.
    """
    actions: list[str] = []

    # If buffer not set, force it first (otherwise cliff/runway is not defensible)
    if cash_buffer <= 0:
        actions.append("Within 7 days: Set a formal cash safety buffer in Client Settings and re-run runway/cliff analysis.")
        actions.append("Within 14 days: Validate the cashflow engine inputs driving burn (payroll, opex, AR assumptions).")

    # Cash cliff < 6 months
    if runway_months_to_buffer is not None and runway_months_to_buffer < 6:
        actions.append("Within 7 days: Begin fundraising immediately (define target amount, instrument, and investor list).")
        actions.append("Within 14 days: Prepare investor pack (deck, model, use-of-funds, milestones) and schedule first meetings.")
        actions.append("Within 30 days: Tighten burn controls (freeze non-critical spend; require approval gates for new commitments).")
        actions.append("Within 30 days: Secure bridge options (SAFE/convertible/credit line) to remove timing failure risk.")

    # Funding timing risk exists (conservative: cliff exists OR runway < 9 months)
    elif (cliff_ts_exact is not None) or (runway_months_to_buffer is not None and runway_months_to_buffer < 9):
        actions.append("Within 14 days: Bring forward fundraising timeline and set a board-approved raise start date.")
        actions.append("Within 30 days: Build a bridge contingency plan (terms, counterparties, documentation readiness).")
        actions.append("Within 30 days: Accelerate collections (tighten payment terms, follow-ups, and invoicing cadence).")

    else:
        actions.append("Within 30 days: Define funding strategy (raise amount + timing) even if not urgent, and align it to milestones.")
        actions.append("Within 60–90 days: Run at least 3 scenarios (base, downside, raise-now) and approve guardrails for spend/hiring.")

    # If debt used
    if debt_used:
        actions.append("Within 14 days: Stress-test repayments and confirm covenant compliance under downside scenarios.")

    # Enforce 4–6 actions
    actions = actions[:6]
    if len(actions) < 4:
        # pad safely without vagueness
        actions.append("Within 30 days: Confirm the earliest ‘decision deadline’ for funding (lead time + legal + investor cycle).")
    return actions[:6]


def page_investing_financing():
    top_header("Funding Strategy")

    # -------------------------------
    # Guards
    # -------------------------------
    if not selected_client_id:
        st.info("Select a business from the sidebar to see funding strategy.")
        return

    if not selected_month_start:
        st.info("Pick a focus month in the navbar to analyse funding strategy.")
        return

    currency_code, currency_symbol = get_client_currency(selected_client_id)
    settings = get_client_settings(selected_client_id) or {}

    # Cash safety buffer (USED for runway + cliff)
    cash_buffer = float(settings.get("min_cash_buffer") or 0.0)

    # -------------------------------
    # Load engine (AS-OF focus month) + flows
    # -------------------------------
    engine_df = fetch_cashflow_summary_for_client_as_of(selected_client_id, selected_month_start)
    if engine_df is None:
        st.error("Could not load cashflow engine (Supabase disconnected). Please refresh.")
        return
    if engine_df.empty:
        st.info("No cashflow engine rows yet. Rebuild cashflow from Business overview.")
        return

    df_inv = fetch_investing_flows_for_client(selected_client_id)
    if df_inv is None:
        df_inv = pd.DataFrame()

    df_fin = fetch_financing_flows_for_client(selected_client_id)
    if df_fin is None:
        df_fin = pd.DataFrame()
    # -------------------------------
    # Helpers (KEEP only these page-local cash helpers)
    # -------------------------------
    def _month_start(x):
        return pd.to_datetime(x).to_period("M").to_timestamp()

    def _prep_engine_monthly(df: pd.DataFrame) -> pd.DataFrame:
        ef = df.copy()
        ef["month_date"] = pd.to_datetime(ef.get("month_date"), errors="coerce")
        ef = ef[ef["month_date"].notna()].sort_values("month_date")
        ef["month_date"] = ef["month_date"].dt.to_period("M").dt.to_timestamp()

        for c in ["operating_cf", "investing_cf", "financing_cf", "net_cash", "closing_cash"]:
            if c in ef.columns:
                ef[c] = pd.to_numeric(ef[c], errors="coerce").fillna(0.0)

        if "net_cash" not in ef.columns or ef["net_cash"].isna().all():
            if all(c in ef.columns for c in ["operating_cf", "investing_cf", "financing_cf"]):
                ef["net_cash"] = ef["operating_cf"] + ef["investing_cf"] + ef["financing_cf"]
            else:
                ef["net_cash"] = 0.0

        if "closing_cash" not in ef.columns:
            ef["closing_cash"] = np.nan

        return ef

    def _recompute_closing_cash_series(ef: pd.DataFrame, opening_cash: float) -> pd.DataFrame:
        ef2 = ef.copy().sort_values("month_date")
        closes, opens = [], []
        o = float(opening_cash)

        for _, r in ef2.iterrows():
            opens.append(o)
            net = float(r.get("net_cash", 0.0) or 0.0)
            c = o + net
            closes.append(c)
            o = c

        ef2["opening_cash_sim"] = opens
        ef2["closing_cash_sim"] = closes
        return ef2

    def _compute_cash_cliff_and_lowest(base_series: pd.DataFrame, buffer_amt: float):
        if base_series is None or base_series.empty:
            return None, None, None

        s = base_series.copy()
        if "closing_cash_sim" not in s.columns:
            return None, None, None

        s["closing_cash_sim"] = pd.to_numeric(s["closing_cash_sim"], errors="coerce")
        s = s[s["closing_cash_sim"].notna()].copy()
        if s.empty:
            return None, None, None

        idx_min = s["closing_cash_sim"].idxmin()
        lowest_ts = s.loc[idx_min, "month_date"]
        min_cash = float(s.loc[idx_min, "closing_cash_sim"])

        cliff_ts = None
        if buffer_amt is not None:
            breach = s[s["closing_cash_sim"] <= float(buffer_amt)]
            if not breach.empty:
                cliff_ts = breach["month_date"].iloc[0]

        return cliff_ts, lowest_ts, min_cash

    # -------------------------------
    # Prepare engine window (12 months)
    # -------------------------------
    ef = _prep_engine_monthly(engine_df)

    focus_ts = _month_start(selected_month_start)
    ef_12 = ef[(ef["month_date"] >= focus_ts) & (ef["month_date"] < (focus_ts + pd.DateOffset(months=12)))].copy()
    if ef_12.empty:
        st.info("No cashflow engine data in the next 12 months for this focus month.")
        return

    # -------------------------------
    # KPI 1 — ACTUAL cash today (from KPIs table)
    # -------------------------------
    kpis = fetch_kpis_for_client_month(selected_client_id, selected_month_start)
    actual_cash = float(kpis.get("cash_balance") or 0.0) if kpis else 0.0

    base_series = _recompute_closing_cash_series(ef_12, opening_cash=actual_cash)

    # -------------------------------
    # KPI 2 — Burn (base case)
    # -------------------------------
    runway_from_engine, effective_burn = compute_runway_and_effective_burn_from_df(engine_df, selected_month_start)
    if runway_from_engine is None or effective_burn is None:
        st.info("Not enough cashflow engine data yet. Rebuild cashflow from Business overview.")
        return
    current_burn = float(effective_burn)

    # -------------------------------
    # KPI 3/4 — Runway + Cliff (use compute_cash_stress_kpis)
    # -------------------------------
    week_df = pd.DataFrame()
    try:
        week_df = build_14_week_cash_table(selected_client_id, selected_month_start) or pd.DataFrame()
    except Exception:
        week_df = pd.DataFrame()

    # Anchor "today" to weekly table horizon (prevents UTC-now skew)
    today_anchor = None
    try:
        if week_df is not None and not week_df.empty and "week_start" in week_df.columns:
            today_anchor = pd.to_datetime(week_df["week_start"], errors="coerce").min()
    except Exception:
        today_anchor = None
    if today_anchor is None or pd.isna(today_anchor):
        today_anchor = pd.to_datetime(selected_month_start).to_period("M").to_timestamp()

    stress_kpis = compute_cash_stress_kpis(
        week_df=week_df,
        cash_buffer=cash_buffer,
        today=today_anchor,
        client_id=selected_client_id,
        as_of_month=selected_month_start,
        fetch_cashflow_summary_for_client_as_of=fetch_cashflow_summary_for_client_as_of,
    )

    runway_display = stress_kpis.get("runway_months")
    cliff_date_exact = stress_kpis.get("stress_date_exact")

    cliff_label = pd.to_datetime(cliff_date_exact) if cliff_date_exact is not None else None
    cliff_month_ts = cliff_label.to_period("M").to_timestamp() if cliff_label is not None else None

    _, lowest_ts, _ = _compute_cash_cliff_and_lowest(
        base_series,
        buffer_amt=cash_buffer if cash_buffer > 0 else 0.0
    )

    # -------------------------------
    # Layout (ONLY LEFT)
    # -------------------------------
    st.subheader("🚀 Funding decision panel")

    cA, cB, cC, cD = st.columns(4)
    with cA:
        st.metric("Current cash balance", f"{currency_symbol}{actual_cash:,.0f}")
        st.caption("Actual cash (KPIs), not forecast.")
    with cB:
        st.metric("Net monthly burn (base)", f"{currency_symbol}{current_burn:,.0f} / month")
    with cC:
        if cash_buffer and cash_buffer > 0:
            st.metric(
                "Runway (to buffer breach)",
                f"{runway_display:.1f} months" if runway_display is not None else "≥ 12m"
            )
        else:
            st.metric(
                "Runway (forward)",
                f"{float(runway_from_engine):.1f} months" if runway_from_engine is not None else "—"
            )
    with cD:
        st.metric(
            "Cash cliff date",
            cliff_label.strftime("%b %Y") if cliff_label is not None else "No cliff in forecast",
        )

    if cash_buffer and cash_buffer > 0:
        if cliff_label is not None:
            st.warning(f"Cash becomes operationally unsafe (below buffer) around **{cliff_label.strftime('%b %Y')}**.")
        else:
            st.success("Cash stays above your safety buffer in the current forecast window.")
    else:
        st.info("Cash safety buffer is not set (or is 0). Set it in Client Settings to get a true ‘operationally unsafe’ cliff.")

    st.markdown("---")

    # -------------------------------
    # Cash cliff chart (base case)
    # -------------------------------
    st.subheader("📉 Cash cliff chart (base case)")

    plot_base = base_series[["month_date", "closing_cash_sim"]].copy()
    plot_base = plot_base.rename(columns={"closing_cash_sim": "Closing cash"})

    base_line = (
        alt.Chart(plot_base)
        .mark_line()
        .encode(
            x=alt.X("month_date:T", title="Month", axis=alt.Axis(format="%b %Y")),
            y=alt.Y(
                "Closing cash:Q",
                title=f"Projected closing cash ({currency_code})",
                axis=alt.Axis(tickMinStep=500_000, format=",.0f"),
            ),
            tooltip=[
                alt.Tooltip("month_date:T", title="Month", format="%b %Y"),
                alt.Tooltip("Closing cash:Q", title="Closing cash", format=",.0f"),
            ],
        )
    )

    layers = [base_line]

    if cash_buffer and cash_buffer > 0:
        buffer_rule = (
            alt.Chart(pd.DataFrame({"y": [cash_buffer]}))
            .mark_rule(strokeDash=[6, 6])
            .encode(y="y:Q")
        )
        layers.append(buffer_rule)

    if lowest_ts is not None:
        low_df = plot_base[plot_base["month_date"] == lowest_ts]
        if not low_df.empty:
            layers.append(alt.Chart(low_df).mark_point(size=120).encode(x="month_date:T", y="Closing cash:Q"))

    if cliff_month_ts is not None:
        cliff_df = plot_base[plot_base["month_date"] == cliff_month_ts]
        if not cliff_df.empty:
            layers.append(alt.Chart(cliff_df).mark_point(size=120).encode(x="month_date:T", y="Closing cash:Q"))

    st.altair_chart(alt.layer(*layers).interactive(), width="stretch")
    st.caption("Dashed line = safety buffer. Markers = lowest cash and first cliff month.")

    # -------------------------------
    # Scenarios module (kept)
    # -------------------------------
    scenario_summaries = render_scenarios_what_if(
        selected_client_id=selected_client_id,
        selected_month_start=selected_month_start,
    ) or []


    # ✅ NEW: pull saved scenario summaries (safe; falls back gracefully)
    scenario_summaries = _get_saved_scenario_summaries_safe(selected_client_id, selected_month_start)

    # -------------------------------
    # Board-level narrative + Action plan
    # -------------------------------
    st.markdown("---")
    st.subheader("🧭 Board-level summary")

    debt_used = _detect_debt_used(df_fin, focus_ts)

    what_this_means = _build_what_this_means_bullets(
        runway_months_to_buffer=runway_display,
        cliff_ts_exact=cliff_label,
        cash_buffer=cash_buffer,
        debt_used=debt_used,
        scenario_summaries=scenario_summaries,  # ✅ this is the missing link
    )


    action_plan = _build_action_plan_30_90(
        runway_months_to_buffer=runway_display,
        cliff_ts_exact=cliff_label,
        cash_buffer=cash_buffer,
        debt_used=debt_used,
    )

    colW, colP = st.columns([1, 1])
    with colW:
        st.markdown("### ✅ WHAT THIS MEANS")
        for b in what_this_means[:5]:
            st.markdown(f"- {b}")

    with colP:
        st.markdown("### 📌 ACTION PLAN — NEXT 30–90 DAYS")
        for a in action_plan[:6]:
            st.markdown(f"- {a}")

    # -------------------------------
    # Admin expander (kept)
    # -------------------------------
    st.markdown("---")
    with st.expander("📚 Detailed investing & financing entries (admin)", expanded=False):
        col_left_admin, col_right_admin = st.columns([2, 1])

        with col_left_admin:
            st.markdown("### Investing flows")
            if df_inv is None or df_inv.empty:
                st.caption("No investing flows yet.")
            else:
                inv_view = df_inv.copy()
                inv_view["month_date"] = pd.to_datetime(inv_view["month_date"], errors="coerce")
                inv_view = inv_view.sort_values("month_date")
                inv_view["Month"] = inv_view["month_date"].dt.strftime("%b %Y")
                cols = ["id", "Month"] + [c for c in ["amount", "category", "notes"] if c in inv_view.columns]
                st.dataframe(inv_view[cols], width="stretch")

            st.markdown("### Financing flows")
            if df_fin is None or df_fin.empty:
                st.caption("No financing flows yet.")
            else:
                fin_view = df_fin.copy()
                fin_view["month_date"] = pd.to_datetime(fin_view["month_date"], errors="coerce")
                fin_view = fin_view.sort_values("month_date")
                fin_view["Month"] = fin_view["month_date"].dt.strftime("%b %Y")
                cols = ["id", "Month"] + [c for c in ["amount", "category", "notes"] if c in fin_view.columns]
                st.dataframe(fin_view[cols], width="stretch")

        with col_right_admin:
            st.markdown("### ➕ Add investing flow")
            with st.form(key="new_investing_flow_form"):
                inv_date = st.date_input("Month", value=selected_month_start, key="inv_date")
                inv_amount = st.number_input("Amount", value=0.0, step=1000.0, key="inv_amount")
                inv_category = st.selectbox(
                    "Category",
                    ["Capex", "R&D / product build", "Equipment & tools", "Fitout / leasehold", "Software / intangibles", "Other"],
                    key="inv_category",
                )
                inv_notes = st.text_area("Notes", value="", key="inv_notes")
                if st.form_submit_button("Save investing flow"):
                    ok = create_investing_flow(
                        selected_client_id, month_date=inv_date, amount=inv_amount, category=inv_category, notes=inv_notes
                    )
                    if ok:
                        st.success("Saved. Rebuild cashflow to refresh projections.")
                        st.rerun()
                    else:
                        st.error("Could not save investing flow.")

            st.markdown("### ➕ Add financing flow")
            with st.form(key="new_financing_flow_form"):
                fin_date = st.date_input("Month", value=selected_month_start, key="fin_date")
                fin_amount = st.number_input("Amount", value=0.0, step=1000.0, key="fin_amount")
                fin_category = st.selectbox(
                    "Category",
                    ["Equity raise", "SAFE / convertible", "Loan drawdown", "Loan repayment", "Interest payment", "Dividend", "Other"],
                    key="fin_category",
                )
                fin_notes = st.text_area("Notes", value="", key="fin_notes")
                if st.form_submit_button("Save financing flow"):
                    ok = create_financing_flow(
                        selected_client_id, month_date=fin_date, amount=fin_amount, category=fin_category, notes=fin_notes
                    )
                    if ok:
                        st.success("Saved. Rebuild cashflow to refresh projections.")
                        st.rerun()
                    else:
                        st.error("Could not save financing flow.")




def month_bucket_utc_ts(x) -> datetime | None:
    """
    Returns timezone-aware datetime at month-start UTC (YYYY-MM-01 00:00:00+00:00).
    """
    if x is None or x == "":
        return None
    ts = pd.to_datetime(x, errors="coerce")
    if pd.isna(ts):
        return None
    y, m = int(ts.year), int(ts.month)
    return datetime(y, m, 1, 0, 0, 0, tzinfo=timezone.utc)

def month_bucket_utc_iso(x) -> str | None:
    dt = month_bucket_utc_ts(x)
    return dt.isoformat() if dt else None
from datetime import datetime, timezone

def _utc_now_iso() -> str:
    """Cross-version safe UTC timestamp (Python 3.8+)."""
    return datetime.now(timezone.utc).isoformat()

## Pipeline health panel (UPDATED: removed weighted closing in ~90 days)
def build_pipeline_health_panel(df_pipeline, focus_month, currency_symbol):
    if df_pipeline is None or df_pipeline.empty:
        return ("warning", ["No deals yet. Add your top 5–10 active deals to see pipeline health."], [])

    df = df_pipeline.copy()
    df["value_total"] = pd.to_numeric(df.get("value_total"), errors="coerce").fillna(0.0)
    df["probability_pct"] = pd.to_numeric(df.get("probability_pct"), errors="coerce").fillna(0.0)
    df["stage"] = df.get("stage", "").astype(str).str.lower()

    # closing month logic (kept for missing-date checks)
    closing = pd.to_datetime(df.get("end_month"), errors="coerce")
    startm = pd.to_datetime(df.get("start_month"), errors="coerce")
    df["closing_month"] = closing.fillna(startm)

    open_mask = df["stage"].isin(["idea", "proposal", "demo", "contract"])
    df_open = df[open_mask].copy()
    df_open["weighted"] = df_open["value_total"] * df_open["probability_pct"] / 100.0

    # KPI
    weighted_total = float(df_open["weighted"].sum())

    # Warnings
    decisions = []
    metrics = []
    severity = "success"

    missing_dates = int(df_open["closing_month"].isna().sum())
    if missing_dates > 0:
        severity = "warning"
        decisions.append(
            f"{missing_dates} open deals have no expected close month. Add dates to make forecasting meaningful."
        )

    early_stage_weighted = float(df_open.loc[df_open["stage"].isin(["idea", "proposal"]), "weighted"].sum())
    if weighted_total > 0 and (early_stage_weighted / weighted_total) > 0.6:
        severity = "warning"
        decisions.append(
            "Pipeline is mostly early-stage (Idea/Proposal). Push deals to Demo/Contract to de-risk the forecast."
        )

    # Concentration
    top = df_open.sort_values("weighted", ascending=False).head(2)["weighted"].sum()
    if weighted_total > 0 and (top / weighted_total) > 0.5:
        severity = "warning"
        decisions.append("Pipeline is highly concentrated in 1–2 deals. Create backup deals to reduce risk.")

    # Metrics (removed 90-day metric)
    metrics.append(("Weighted pipeline (total)", f"{currency_symbol}{weighted_total:,.0f}"))

    if not decisions:
        decisions.append(
            "Pipeline health looks reasonable. Focus on closing the next best 2–3 deals and keeping dates/probabilities current."
        )

    return severity, decisions, metrics



def build_sales_next_actions(df_pipeline):
    if df_pipeline is None or df_pipeline.empty:
        return ["Add your top 5–10 active deals (name, value, stage, win chance, close month)."]

    df = df_pipeline.copy()
    df["stage"] = df.get("stage", "").astype(str).str.lower()

    actions = []
    if (df["stage"] == "idea").any():
        actions.append("Move **Idea → Proposal**: qualify ICP, confirm pain, book next meeting, set a close month.")
    if (df["stage"] == "proposal").any():
        actions.append("Move **Proposal → Demo**: confirm decision criteria, timeline, stakeholders; schedule demo on calendar.")
    if (df["stage"] == "demo").any():
        actions.append("Move **Demo → Contract**: identify champion, confirm procurement/legal, send pricing + implementation plan.")
    if (df["stage"] == "contract").any():
        actions.append("Close **Contract → Closed won**: push signature, confirm start date, invoice/PO, remove blockers weekly.")

    if not actions:
        actions.append("No active deals in selling stages. Add opportunities or restart outbound to build pipeline.")

    return actions[:4]

def build_pipeline_data_checks(df_pipeline):
    if df_pipeline is None or df_pipeline.empty:
        return []

    df = df_pipeline.copy()
    df["value_total"] = pd.to_numeric(df.get("value_total"), errors="coerce")
    df["probability_pct"] = pd.to_numeric(df.get("probability_pct"), errors="coerce")
    df["stage"] = df.get("stage", "").astype(str).str.lower()

    startm = pd.to_datetime(df.get("start_month"), errors="coerce")
    endm = pd.to_datetime(df.get("end_month"), errors="coerce")
    closing = endm.fillna(startm)

    issues = []
    if df["value_total"].isna().sum() > 0 or (df["value_total"].fillna(0) <= 0).sum() > 0:
        issues.append("Some deals have **$0 / blank value**. Add an estimate to make pipeline meaningful.")
    if df["probability_pct"].isna().sum() > 0:
        issues.append("Some deals have **blank win chance**. Set 10–90% for open deals.")
    if closing.isna().sum() > 0:
        issues.append("Some deals have **no expected close month**. Add close month to improve forecasting.")
    if ((df["stage"] == "closed_won") & (df["probability_pct"].fillna(0) < 90)).any():
        issues.append("Some **Closed won** deals have low probability. Set Closed won deals to 100% for consistency.")
    if ((df["stage"] == "closed_lost") & (df["probability_pct"].fillna(0) > 0)).any():
        issues.append("Some **Closed lost** deals still have probability > 0. Set to 0% to avoid confusion.")

    return issues

def build_recognised_revenue_debug_tables(
    client_id,
    focus_month: date,
    n_months_schedule: int = 12,
):
    """
    Debug tables for recognised revenue KPIs.

    Matches KPI rules:
      - Closed_Won only (enforced by build_revenue_schedule_for_client(mode="secured"))
      - No probability
      - Exclude if no explicit schedule rows (secured mode)
    Shows reconciliation columns:
      - Total value
      - Contract months (detected)
      - Expected MRR = Total value / Contract months (for SaaS)
      - Recognised revenue (from schedule) for focus month
      - Delta = Recognised - Expected
    """

    empty_pack = {
        "this_month_deals": pd.DataFrame(),
        "next3_deals": pd.DataFrame(),
        "next3_by_month": pd.DataFrame(),
        "this_month_rows": pd.DataFrame(),   # optional: raw schedule rows for focus month
    }

    # --- Build schedule (secured rules) ---
    sched = build_revenue_schedule_for_client(
        client_id,
        base_month=focus_month,
        n_months=n_months_schedule,
        mode="secured",
    )
    if sched is None or sched.empty:
        return empty_pack

    df = sched.copy()
    df["month_date"] = pd.to_datetime(df.get("month_date"), errors="coerce")
    df = df[df["month_date"].notna()]

    # Amount column
    amt_col = None
    for cand in ["revenue_amount", "recognised_revenue", "amount"]:
        if cand in df.columns:
            amt_col = cand
            break
    if amt_col is None or "deal_id" not in df.columns:
        return empty_pack

    # --- Pull pipeline (for labels + total value + contract months) ---
    pipe = fetch_pipeline_for_client(client_id)
    pipe_ok = None

    if pipe is not None and not pipe.empty:
        pipe = pipe.copy()
        pipe["deal_id"] = pipe.get("id")

        # detect deal/customer columns
        deal_name_candidates = ["deal_name", "name", "title", "deal", "opportunity_name"]
        cust_candidates = ["customer_name", "customer", "account_name", "client_name", "company_name"]

        deal_col = next((c for c in deal_name_candidates if c in pipe.columns), None)
        cust_col = next((c for c in cust_candidates if c in pipe.columns), None)

        # detect contract months columns (your schema variants)
        contract_cols = [
            "contract_months",
            "contract_length_months",
            "contract_length",
            "contract_term_months",
            "term_months",
        ]

        def _parse_int_months(v):
            if v is None or (isinstance(v, float) and pd.isna(v)):
                return None
            try:
                s = str(v).strip().lower()
                s = s.replace("months", "").replace("month", "").strip()
                n = int(float(s))
                return n if n > 0 else None
            except Exception:
                return None

        # create a single "contract_months_dbg" column by coalescing
        pipe["contract_months_dbg"] = None
        for c in contract_cols:
            if c in pipe.columns:
                parsed = pipe[c].apply(_parse_int_months)
                pipe["contract_months_dbg"] = pipe["contract_months_dbg"].fillna(parsed)

        # total value
        if "value_total" in pipe.columns:
            pipe["value_total_dbg"] = pd.to_numeric(pipe["value_total"], errors="coerce").fillna(0.0)
        else:
            pipe["value_total_dbg"] = 0.0

        labels = pipe[["deal_id"]].copy()
        labels["Deal"] = pipe[deal_col].astype(str) if deal_col else labels["deal_id"].astype(str)
        labels["Customer"] = pipe[cust_col].astype(str) if cust_col else ""
        labels["Total value"] = pipe["value_total_dbg"]
        labels["Contract months"] = pipe["contract_months_dbg"]

        labels = labels.drop_duplicates(subset=["deal_id"])
        pipe_ok = labels

        df = df.merge(pipe_ok, on="deal_id", how="left")
        df["Deal"] = df["Deal"].fillna(df["deal_id"].astype(str))
        df["Customer"] = df["Customer"].fillna("")
        df["Total value"] = pd.to_numeric(df.get("Total value"), errors="coerce").fillna(0.0)
    else:
        # fallback labels if pipeline not available
        df["Deal"] = df["deal_id"].astype(str)
        df["Customer"] = ""
        df["Total value"] = 0.0
        df["Contract months"] = None

    # --- Focus window (focus + next 3 months) ---
    focus_start = pd.to_datetime(focus_month).replace(day=1)
    next3_end = focus_start + pd.DateOffset(months=4)

    window = df[(df["month_date"] >= focus_start) & (df["month_date"] < next3_end)].copy()
    if window.empty:
        return empty_pack

    window["Month"] = window["month_date"].dt.strftime("%b %Y")

    # Raw schedule rows for THIS month (optional, super useful for recon)
    this_rows = window[window["month_date"] == focus_start].copy()
    # keep only the important columns if present
    keep_cols = [c for c in [
        "deal_id", "Deal", "Customer", "revenue_type", "stage", "month_date", "Month", amt_col,
        "Total value", "Contract months"
    ] if c in this_rows.columns]
    this_month_rows = this_rows[keep_cols].copy()

    # --- This month (per-deal contribution) ---
    this_m = window[window["month_date"] == focus_start].copy()

    this_month_deals = (
        this_m.groupby(["deal_id", "Deal", "Customer"], as_index=False)[amt_col]
        .sum()
        .rename(columns={amt_col: "Recognised revenue"})
    )

    # enrich per-deal totals/contract months (take first non-null per deal)
    meta_cols = ["deal_id"]
    for c in ["Total value", "Contract months"]:
        if c in window.columns:
            meta_cols.append(c)

    if len(meta_cols) > 1:
        meta = (
            window[meta_cols]
            .drop_duplicates(subset=["deal_id"])
            .copy()
        )
        this_month_deals = this_month_deals.merge(meta, on="deal_id", how="left")

    # expected MRR and delta (only meaningful for SaaS, but helpful anyway)
    this_month_deals["Contract months"] = pd.to_numeric(
        this_month_deals.get("Contract months"), errors="coerce"
    )

    this_month_deals["Expected (Total/Contract)"] = 0.0
    mask = this_month_deals["Contract months"].fillna(0) > 0
    this_month_deals.loc[mask, "Expected (Total/Contract)"] = (
        this_month_deals.loc[mask, "Total value"] / this_month_deals.loc[mask, "Contract months"]
    )

    this_month_deals["Delta (Recognised - Expected)"] = (
        this_month_deals["Recognised revenue"] - this_month_deals["Expected (Total/Contract)"]
    )

    # sort
    this_month_deals = this_month_deals.sort_values("Recognised revenue", ascending=False)

    # --- Next 3 months (per deal across next 3 months) ---
    next3 = window[window["month_date"] > focus_start].copy()
    next3_deals = (
        next3.groupby(["deal_id", "Deal", "Customer"], as_index=False)[amt_col]
        .sum()
        .rename(columns={amt_col: "Recognised revenue (next 3 months)"})
        .sort_values("Recognised revenue (next 3 months)", ascending=False)
    )

    # Month totals
    next3_by_month = (
        window.groupby(["month_date", "Month"], as_index=False)[amt_col]
        .sum()
        .rename(columns={amt_col: "Total recognised revenue"})
        .sort_values("month_date")
    )

    # Round for display
    def _round_money_cols(dfx: pd.DataFrame) -> pd.DataFrame:
        out = dfx.copy()
        for c in out.columns:
            if any(k in c.lower() for k in ["revenue", "total", "expected", "delta", "value"]):
                out[c] = pd.to_numeric(out[c], errors="coerce").fillna(0.0).round(2)
        if "Contract months" in out.columns:
            out["Contract months"] = pd.to_numeric(out["Contract months"], errors="coerce")
        return out

    return {
        "this_month_deals": _round_money_cols(this_month_deals),
        "next3_deals": _round_money_cols(next3_deals),
        "next3_by_month": _round_money_cols(next3_by_month),
        "this_month_rows": _round_money_cols(this_month_rows),
    }



# =========================
# MILESTONE SCHEDULER (MUST RECEIVE client_id)
# =========================

def _schedule_milestone(
    client_id,
    deal: pd.Series,
    window_start,
    window_end,
    rows: list,
    prob_factor: float,
    require_explicit: bool = False,
):
    """
    Milestone revenue using deal_milestones table:
      - month_date + percent_of_contract
      - allocates value_total * percent/100 into each milestone month
      - if missing milestones:
          - secured => exclude (require_explicit=True)
          - forecast => fallback: recognise full amount in start_month
    """
    total = float(deal.get("value_total") or 0.0) * prob_factor
    if total == 0:
        return

    deal_id = deal.get("id")
    if client_id is None or deal_id is None:
        return

    ms_df = fetch_milestones_for_deal(client_id, int(deal_id))

    if ms_df is None or ms_df.empty:
        if require_explicit:
            return
        start_raw = deal.get("start_month") or deal.get("created_at")
        _append_revenue_row(
            rows, window_start, window_end, deal,
            "Milestone project",
            start_raw, total
        )
        return

    ms_df = ms_df.copy()
    ms_df["month_date"] = pd.to_datetime(ms_df["month_date"], errors="coerce")
    ms_df["percent_of_contract"] = pd.to_numeric(ms_df.get("percent_of_contract"), errors="coerce").fillna(0.0)

    ms_df = ms_df[ms_df["month_date"].notna() & (ms_df["percent_of_contract"] > 0)]
    if ms_df.empty:
        if require_explicit:
            return
        start_raw = deal.get("start_month") or deal.get("created_at")
        _append_revenue_row(rows, window_start, window_end, deal, "Milestone project", start_raw, total)
        return

    pct_sum = float(ms_df["percent_of_contract"].sum())
    if pct_sum <= 0:
        if require_explicit:
            return
        start_raw = deal.get("start_month") or deal.get("created_at")
        _append_revenue_row(rows, window_start, window_end, deal, "Milestone project", start_raw, total)
        return

    # Normalise to 100 if user doesn't enter exactly 100 total
    scale = 100.0 / pct_sum

    for _, r in ms_df.iterrows():
        pct = float(r.get("percent_of_contract") or 0.0) * scale
        if pct <= 0:
            continue

        amt = total * (pct / 100.0)

        _append_revenue_row(
            rows, window_start, window_end, deal,
            "Milestone project",
            r["month_date"], amt
        )


# =========================
# UPDATED page_sales_deals()
# - UI editors based on SELECTED DEAL method (not business default)
# - Option A: checkbox to apply business default to ALL existing deals
# - cache clears handled via save_revenue_method + bulk update helper + milestone/poc saves
# =========================

def page_configuration():
    """
    Configuration page (NEW)

    Contains:
    1) Revenue recognition business default + Save (+ optional apply to all deals)
    2) Create / update deal editor (moved from Sales & Deals)

    Does NOT change any external function signatures.
    Reuses your existing save/bulk-update helpers.
    """

    top_header("Configuration")

    if not selected_client_id:
        st.info("Select a client from the sidebar to view this page.")
        return

    settings = get_client_settings(selected_client_id) or {}
    currency_code, currency_symbol = get_client_currency(selected_client_id)

    # ============================================================
    # 1) Revenue recognition method (business default)
    # ============================================================
    st.subheader("🧾 Revenue recognition settings")

    current_method = (settings.get("revenue_recognition_method") or "saas").lower()

    method_labels = {
        "saas": "SaaS subscription (spread over contract term)",
        "milestone": "Milestone / project (milestones by month)",
        "poc": "Percentage-of-completion (true-up by month)",
        "straight_line": "Straight-line service (start–end months)",
        "usage": "Usage-based (placeholder spreading)",
        "point_in_time": "Point-in-time goods delivery (lump on delivery month)",
    }

    method_display_map = {
        "saas": "SaaS subscription",
        "milestone": "Milestone project",
        "poc": "Percentage-of-completion",
        "straight_line": "Straight-line service",
        "usage": "Usage-based",
        "point_in_time": "Point-in-time goods",
    }

    method_options = list(method_labels.keys())
    default_index = method_options.index(current_method) if current_method in method_options else 0

    selected_method_label = st.selectbox(
        "Business default method",
        options=[method_labels[m] for m in method_options],
        index=default_index,
        help=(
            "Default method for this business. Deals can override this individually. "
            "The engine uses deal.method first, then falls back to this business default."
        ),
        key="cfg_rev_method_select",
    )

    reverse_map = {v: k for k, v in method_labels.items()}
    selected_method_code = reverse_map[selected_method_label]

    apply_to_all = st.checkbox(
        "Apply this method to ALL existing deals (overwrite per-deal overrides)",
        value=False,
        help="Use this if you want business default to drive the numbers across all deals.",
        key="cfg_apply_to_all",
    )

    cA, cB = st.columns([1, 3])
    with cA:
        if st.button("Save method", key="cfg_save_method"):
            ok = save_revenue_method(selected_client_id, selected_method_code)

            if ok and apply_to_all:
                ok2 = bulk_update_deal_methods_for_client(
                    selected_client_id,
                    method_display=method_display_map.get(selected_method_code, "Straight-line service"),
                )
                ok = ok and ok2

            if ok:
                try:
                    st.cache_data.clear()
                except Exception:
                    pass
                st.success("Revenue recognition method updated.")
                st.rerun()
            else:
                st.error("Could not update method. Please try again.")

    st.markdown("---")

    # ============================================================
    # 2) Create / update deal editor (moved here)
    # ============================================================
    st.subheader("✍️ Create / update a deal")

    # Load pipeline for editing list
    try:
        df_pipeline = fetch_pipeline_for_client(selected_client_id)
    except Exception as e:
        st.error("Could not load pipeline from Supabase (revenue_pipeline table).")
        st.caption(f"{type(e).__name__}: {e}")
        df_pipeline = pd.DataFrame()

    if df_pipeline is None:
        df_pipeline = pd.DataFrame()

    # ---- Helpers (same ones you had) ----
    def _parse_month_input(maybe_date) -> pd.Timestamp | None:
        if maybe_date is None or str(maybe_date).strip() == "":
            return None
        ts = pd.to_datetime(maybe_date, errors="coerce")
        if pd.isna(ts):
            return None
        return ts.to_period("M").to_timestamp()

    def _supabase_upsert_revenue_pipeline(row: dict) -> bool:
        try:
            sb = None
            if "get_supabase_authed" in globals():
                sb = get_supabase_authed()
            if sb is None and "get_supabase_public" in globals():
                sb = get_supabase_public()
            if sb is None:
                raise RuntimeError("No Supabase client available (get_supabase_authed/public missing).")
            sb.table("revenue_pipeline").upsert(row).execute()
            return True
        except Exception as e:
            st.error("Save failed (Supabase upsert).")
            st.caption(f"{type(e).__name__}: {e}")
            return False

    def _save_deal(row: dict) -> bool:
        for fn_name in [
            "save_pipeline_deal",
            "upsert_pipeline_deal",
            "save_deal_to_pipeline",
            "upsert_revenue_pipeline_row",
        ]:
            fn = globals().get(fn_name)
            if callable(fn):
                try:
                    out = fn(row)
                    return bool(out) if out is not None else True
                except Exception as e:
                    st.error(f"Save failed ({fn_name}).")
                    st.caption(f"{type(e).__name__}: {e}")
                    return False
        return _supabase_upsert_revenue_pipeline(row)


    def _delete_deal(deal_id: str) -> bool:
        if not deal_id:
            return False

        for fn_name in [
            "delete_pipeline_deal",
            "delete_deal_from_pipeline",
            "delete_revenue_pipeline_row",
        ]:
            fn = globals().get(fn_name)
            if callable(fn):
                try:
                    out = fn(deal_id, selected_client_id) if fn.__code__.co_argcount >= 2 else fn(deal_id)
                    return bool(out) if out is not None else True
                except Exception as e:
                    st.error(f"Delete failed ({fn_name}).")
                    st.caption(f"{type(e).__name__}: {e}")
                    return False

        try:
            sb = get_supabase_authed() if "get_supabase_authed" in globals() else None
            if sb is None and "get_supabase_public" in globals():
                sb = get_supabase_public()
            if sb is None:
                raise RuntimeError("No Supabase client available (get_supabase_authed/public missing).")
            sb.table("revenue_pipeline").delete().eq("id", deal_id).eq("client_id", selected_client_id).execute()
            return True
        except Exception as e:
            st.error("Delete failed (Supabase delete).")
            st.caption(f"{type(e).__name__}: {e}")
            return False

    # ---- Deal selector ----
    deal_id_col = "id" if ("id" in df_pipeline.columns) else None
    deal_choices = ["➕ New deal"]
    deal_map = {}

    if not df_pipeline.empty and deal_id_col:
        tmp = df_pipeline.copy()
        tmp["deal_name"] = tmp.get("deal_name", "").astype(str).fillna("")
        tmp["customer_name"] = tmp.get("customer_name", "").astype(str).fillna("")
        tmp["label"] = tmp.apply(
            lambda r: f"{r.get('deal_name','').strip() or 'Unnamed'}  —  {r.get('customer_name','').strip() or 'Unknown'}",
            axis=1,
        )
        tmp = tmp.sort_values(["label"])
        deal_choices = ["➕ New deal"] + tmp["label"].tolist()
        deal_map = {row["label"]: row for _, row in tmp.iterrows()}

    selected_label = st.selectbox("Select deal", deal_choices, index=0, key="cfg_deal_select")
    editing_existing = selected_label != "➕ New deal"
    row0 = deal_map.get(selected_label) if editing_existing else {}

    def _d(key, fallback=""):
        return row0.get(key, fallback) if isinstance(row0, dict) else fallback

    stage_options = ["idea", "demo", "proposal", "contract", "closed_won", "closed_lost"]
    default_stage = str(_d("stage", "idea") or "idea").lower()
    if default_stage not in stage_options:
        default_stage = "idea"

    # Use same display map as Sales page
    current_method_display = method_display_map.get(selected_method_code, "SaaS subscription")
    method_options_display = list(method_display_map.values())
    default_method_display = str(_d("method", "") or "").strip()
    if default_method_display not in method_options_display:
        default_method_display = current_method_display

    with st.form("cfg_deal_form", clear_on_submit=False):
        deal_name = st.text_input("Deal name", value=str(_d("deal_name", "") or ""))
        customer_name = st.text_input("Customer", value=str(_d("customer_name", "") or ""))

        c1, c2 = st.columns(2)
        with c1:
            stage = st.selectbox("Stage", stage_options, index=stage_options.index(default_stage))
        with c2:
            probability_pct = st.number_input(
                "Win chance (%)",
                min_value=0.0,
                max_value=100.0,
                value=float(pd.to_numeric(_d("probability_pct", 0.0), errors="coerce") or 0.0),
                step=5.0,
            )

        value_total = st.number_input(
            f"Contract value (total) ({currency_code})",
            min_value=0.0,
            value=float(pd.to_numeric(_d("value_total", 0.0), errors="coerce") or 0.0),
            step=1000.0,
        )

        method_display = st.selectbox(
            "Revenue recognition (deal override)",
            options=method_options_display,
            index=method_options_display.index(default_method_display),
            help="Overrides the business default above for this deal.",
        )

        c3, c4 = st.columns(2)
        with c3:
            start_month_in = st.text_input(
                "Expected close / start month (e.g., 2025-12-01 or Dec 2025)",
                value=str(_d("start_month", "") or ""),
            )
        with c4:
            end_month_in = st.text_input(
                "End month (optional; for service term)",
                value=str(_d("end_month", "") or ""),
            )

        contract_months_default = pd.to_numeric(_d("contract_months", 0), errors="coerce")
        contract_months_default = int(contract_months_default) if pd.notna(contract_months_default) else 0
        contract_months = st.number_input(
            "Contract months (optional)",
            min_value=0,
            value=int(contract_months_default),
            step=1,
        )

        stage_entry_default = _d("stage_entry_date", None)
        stage_entry_default = pd.to_datetime(stage_entry_default, errors="coerce")
        stage_entry_date = st.date_input(
            "Stage entry date",
            value=stage_entry_default.date() if pd.notna(stage_entry_default) else datetime.now(timezone.utc).date(),
            help="Used for ageing rules. Keep updated when stage changes.",
        )

        notes = st.text_area("Notes (optional)", value=str(_d("notes", "") or ""), height=120)

        submitted = st.form_submit_button("Save deal")

    if submitted:
        if not str(deal_name).strip():
            st.error("Deal name is required.")
        else:
            start_ts = _parse_month_input(start_month_in)
            end_ts = _parse_month_input(end_month_in)

            payload = {
                "client_id": selected_client_id,
                "deal_name": str(deal_name).strip(),
                "customer_name": str(customer_name).strip() if customer_name is not None else "",
                "stage": str(stage).lower().strip(),
                "probability_pct": float(probability_pct or 0.0),
                "value_total": float(value_total or 0.0),
                "method": str(method_display).strip(),
                "contract_months": int(contract_months or 0),
                "stage_entry_date": pd.to_datetime(stage_entry_date).strftime("%Y-%m-%d") if stage_entry_date else None,
                "notes": str(notes).strip() if notes is not None else "",
                "start_month": start_ts.strftime("%Y-%m-%d") if start_ts is not None else None,
                "end_month": end_ts.strftime("%Y-%m-%d") if end_ts is not None else None,
            }

            if editing_existing and deal_id_col and isinstance(row0, dict) and row0.get(deal_id_col):
                payload["id"] = row0.get(deal_id_col)

            ok = _save_deal(payload)
            if ok:
                try:
                    st.cache_data.clear()
                except Exception:
                    pass
                st.success("Deal saved.")
                st.rerun()

    if editing_existing and deal_id_col and isinstance(row0, dict) and row0.get(deal_id_col):
        st.markdown("---")
        if st.button("🗑️ Delete deal", type="secondary", key="cfg_delete_deal"):
            did = str(row0.get(deal_id_col))
            if _delete_deal(did):
                try:
                    st.cache_data.clear()
                except Exception:
                    pass
                st.success("Deal deleted.")
                st.rerun()

        # ============================================================
    # 1B) AR/AP: Record payments (moved from Cash & Bills)
    # ============================================================
    st.markdown("---")
    st.subheader("💳 AR/AP — record payments")

    # Load AR/AP (normalised) so the payment widgets have clean data
    try:
        df_ar_cfg, df_ap_cfg = fetch_ar_ap_for_client(selected_client_id)
    except Exception as e:
        st.error(f"Failed to load AR/AP: {e}")
        df_ar_cfg, df_ap_cfg = pd.DataFrame(), pd.DataFrame()

    # Normalise (safe)
    try:
        df_ar_cfg = normalise_ar_ap_for_cash(df_ar_cfg)
    except Exception:
        pass
    try:
        df_ap_cfg = normalise_ar_ap_for_cash(df_ap_cfg)
    except Exception:
        pass

    # Use Focus Month month-end as default payment date
    focus_start_cfg = selected_month_start
    focus_end_cfg = (pd.to_datetime(focus_start_cfg) + MonthEnd(0)).date()

    currency_code_cfg, currency_symbol_cfg = get_client_currency(selected_client_id)

    # ----------------------------
    # AR Payment
    # ----------------------------
    with st.expander("💳 Record a payment on an AR invoice", expanded=False):
        ar_open = (
            df_ar_cfg[df_ar_cfg.get("is_cash_relevant", True) & (df_ar_cfg.get("effective_amount", 0) > 0)].copy()
            if df_ar_cfg is not None and not df_ar_cfg.empty
            else pd.DataFrame()
        )

        if ar_open.empty:
            st.caption("No open customer invoices to record payments against.")
        else:
            AR_TABLE_NAME = "ar_ap_tracker"  # 🔁 adjust if needed

            label_map = {}
            for _, r in ar_open.iterrows():
                inv_id = r.get("id")
                inv_no = r.get("invoice_number", inv_id)
                cust = r.get("counterparty", "")
                due = r.get("due_date", "")
                bal = float(r.get("effective_amount") or 0.0)

                try:
                    due_txt = pd.to_datetime(due).strftime("%d %b %Y")
                except Exception:
                    due_txt = str(due)

                label_map[inv_id] = f"{inv_no} – {cust} – due {due_txt} – bal {currency_symbol_cfg}{bal:,.0f}"

            selected_id = st.selectbox(
                "Select invoice to update",
                options=list(label_map.keys()),
                format_func=lambda x: label_map.get(x, str(x)),
                key="cfg_ar_payment_invoice",
            )

            payment_type = st.radio(
                "Payment type",
                ["Full payment", "Partial payment"],
                horizontal=True,
                key="cfg_ar_payment_type",
            )

            current_row = ar_open[ar_open["id"] == selected_id].iloc[0]
            current_balance = float(current_row.get("effective_amount") or 0.0)

            default_amount = current_balance if payment_type == "Full payment" else 0.0

            payment_amount = st.number_input(
                "Amount paid",
                min_value=0.0,
                value=float(default_amount),
                step=100.0,
                key="cfg_ar_payment_amount",
            )

            payment_date = st.date_input(
                "Payment date",
                value=focus_end_cfg,
                key="cfg_ar_payment_date",
            )

            if st.button("Record AR payment", key="cfg_ar_payment_button"):
                try:
                    # Prefer your existing authed client if available
                    sb = None
                    if "get_supabase_authed" in globals():
                        sb = get_supabase_authed()
                    if sb is None and "get_supabase_client" in globals():
                        sb = get_supabase_client()
                    if sb is None:
                        raise RuntimeError("No Supabase client available.")

                    prev_paid = float(current_row.get("amount_paid") or 0.0)
                    new_paid_total = prev_paid + float(payment_amount)

                    total_amt = float(current_row.get("amount") or 0.0)
                    if new_paid_total >= total_amt - 0.01:
                        new_status = "paid"
                        partially_paid_flag = False
                    else:
                        new_status = "open"
                        partially_paid_flag = True

                    update_data = {
                        "amount_paid": new_paid_total,
                        "partially_paid": partially_paid_flag,
                        "status": new_status,
                        "payment_date": payment_date.isoformat(),
                        "updated_at": pd.Timestamp.utcnow().isoformat(),
                    }

                    sb.table(AR_TABLE_NAME).update(update_data).eq("id", selected_id).execute()

                    try:
                        st.cache_data.clear()
                    except Exception:
                        pass

                    st.success("AR payment recorded.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Failed to record payment: {e}")

    # ----------------------------
    # AP Payment
    # ----------------------------
    with st.expander("💳 Record a payment on an AP bill", expanded=False):
        ap_open = (
            df_ap_cfg[df_ap_cfg.get("is_cash_relevant", True) & (df_ap_cfg.get("effective_amount", 0) > 0)].copy()
            if df_ap_cfg is not None and not df_ap_cfg.empty
            else pd.DataFrame()
        )

        if ap_open.empty:
            st.caption("No open supplier bills to record payments against.")
        else:
            AP_TABLE_NAME = "ar_ap_tracker"  # 🔁 adjust if AP lives in a different table

            label_map_ap = {}
            for _, r in ap_open.iterrows():
                bill_id = r.get("id")
                bill_no = r.get("invoice_number", bill_id)
                supp = r.get("counterparty", "")
                due = r.get("due_date", "")
                bal = float(r.get("effective_amount") or 0.0)

                try:
                    due_txt = pd.to_datetime(due).strftime("%d %b %Y")
                except Exception:
                    due_txt = str(due)

                label_map_ap[bill_id] = f"{bill_no} – {supp} – due {due_txt} – bal {currency_symbol_cfg}{bal:,.0f}"

            selected_bill_id = st.selectbox(
                "Select bill to update",
                options=list(label_map_ap.keys()),
                format_func=lambda x: label_map_ap.get(x, str(x)),
                key="cfg_ap_payment_bill",
            )

            payment_type_ap = st.radio(
                "Payment type",
                ["Full payment", "Partial payment"],
                horizontal=True,
                key="cfg_ap_payment_type",
            )

            current_ap_row = ap_open[ap_open["id"] == selected_bill_id].iloc[0]
            current_balance_ap = float(current_ap_row.get("effective_amount") or 0.0)

            default_amount_ap = current_balance_ap if payment_type_ap == "Full payment" else 0.0

            payment_amount_ap = st.number_input(
                "Amount paid",
                min_value=0.0,
                value=float(default_amount_ap),
                step=100.0,
                key="cfg_ap_payment_amount",
            )

            payment_date_ap = st.date_input(
                "Payment date",
                value=focus_end_cfg,
                key="cfg_ap_payment_date",
            )

            if st.button("Record AP payment", key="cfg_ap_payment_button"):
                try:
                    sb = None
                    if "get_supabase_authed" in globals():
                        sb = get_supabase_authed()
                    if sb is None and "get_supabase_client" in globals():
                        sb = get_supabase_client()
                    if sb is None:
                        raise RuntimeError("No Supabase client available.")

                    prev_paid_ap = float(current_ap_row.get("amount_paid") or 0.0)
                    new_paid_total_ap = prev_paid_ap + float(payment_amount_ap)

                    total_amt_ap = float(current_ap_row.get("amount") or 0.0)
                    if new_paid_total_ap >= total_amt_ap - 0.01:
                        new_status_ap = "paid"
                        partially_paid_flag_ap = False
                    else:
                        new_status_ap = "open"
                        partially_paid_flag_ap = True

                    update_data_ap = {
                        "amount_paid": new_paid_total_ap,
                        "partially_paid": partially_paid_flag_ap,
                        "status": new_status_ap,
                        "payment_date": payment_date_ap.isoformat(),
                        "updated_at": pd.Timestamp.utcnow().isoformat(),
                    }

                    sb.table(AP_TABLE_NAME).update(update_data_ap).eq("id", selected_bill_id).execute()

                    try:
                        st.cache_data.clear()
                    except Exception:
                        pass

                    st.success("AP payment recorded.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Failed to record payment: {e}")




def page_sales_deals():
    """
    Sales & Deals — UPDATED (Create/update deal REMOVED)

    ✅ Pipeline overview KPIs
    ✅ Pipeline health panel
    ✅ Data checks warning
    ✅ “View full pipeline list” expander (kept)
    ✅ Waterfall chart (TABLE REMOVED)
    ✅ Revenue / Risk / What this means / Action plan (NO tables)

    REMOVED from this page:
    ❌ Revenue recognition method selector + save
    ❌ Create / update deal editor (moved to Configuration page)

    GUARANTEE:
    - Does NOT change any existing external function signatures.
    """

    top_header("Sales & Deals")

    if not selected_client_id:
        st.info("Select a client from the sidebar to view this page.")
        return

    # ---------------------------------------------------------------------
    # Settings + currency
    # ---------------------------------------------------------------------
    settings = get_client_settings(selected_client_id) or {}
    currency_code, currency_symbol = get_client_currency(selected_client_id)

    st.markdown("---")

    # ---------------------------------------------------------------------
    # Load pipeline
    # ---------------------------------------------------------------------
    with st.spinner("Loading your pipeline..."):
        try:
            df_pipeline = fetch_pipeline_for_client(selected_client_id)
        except Exception as e:
            st.error("Could not load pipeline from Supabase (revenue_pipeline table).")
            st.caption(f"{type(e).__name__}: {e}")
            df_pipeline = pd.DataFrame()

    if df_pipeline is None:
        df_pipeline = pd.DataFrame()

    # ---------------------------------------------------------------------
    # Normalise pipeline
    # ---------------------------------------------------------------------
    total_pipeline = 0.0
    open_deals_count = 0

    if not df_pipeline.empty:
        df_pipeline = df_pipeline.copy()
        df_pipeline["value_total"] = pd.to_numeric(df_pipeline.get("value_total"), errors="coerce").fillna(0.0)
        df_pipeline["probability_pct"] = pd.to_numeric(df_pipeline.get("probability_pct"), errors="coerce").fillna(0.0)
        df_pipeline["stage"] = df_pipeline.get("stage", "").astype(str).str.lower().fillna("")

        for col in ["start_month", "end_month", "stage_entry_date", "created_at"]:
            if col in df_pipeline.columns:
                df_pipeline[col] = pd.to_datetime(df_pipeline[col], errors="coerce")

        total_pipeline = float(df_pipeline["value_total"].sum())
        open_stages = ["idea", "proposal", "demo", "contract"]
        open_deals_count = int(df_pipeline["stage"].isin(open_stages).sum())

    # ---------------------------------------------------------------------
    # Helpers (local-only; no external APIs modified)
    # ---------------------------------------------------------------------
    def _safe_month_start(d) -> pd.Timestamp:
        return pd.to_datetime(d).to_period("M").to_timestamp()

    def _get_plan_for_month(client_id, focus_month) -> float:
        try:
            tdf = fetch_client_monthly_targets(client_id)
        except Exception:
            tdf = None
        if tdf is None or tdf.empty:
            return 0.0
        tdf = tdf.copy()
        tdf["month_date"] = pd.to_datetime(tdf["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
        m = _safe_month_start(focus_month)
        row = tdf[tdf["month_date"] == m]
        if row.empty:
            return 0.0
        return float(pd.to_numeric(row.get("revenue_target"), errors="coerce").fillna(0.0).iloc[0])

    def _compute_weighted_pipeline(df_pipe: pd.DataFrame) -> float:
        if df_pipe is None or df_pipe.empty:
            return 0.0
        df = df_pipe.copy()
        df["stage"] = df.get("stage", "").astype(str).str.lower().fillna("")
        df["value_total"] = pd.to_numeric(df.get("value_total"), errors="coerce").fillna(0.0)
        df["probability_pct"] = pd.to_numeric(df.get("probability_pct"), errors="coerce").fillna(0.0)

        df = df[df["stage"] != "closed_lost"].copy()
        if df.empty:
            return 0.0

        def _pf(row) -> float:
            stg = str(row.get("stage", "")).lower()
            if stg == "closed_won":
                return 1.0
            p = float(row.get("probability_pct", 0.0) or 0.0)
            p = max(0.0, min(100.0, p))
            return p / 100.0

        df["weighted_value"] = df["value_total"] * df.apply(_pf, axis=1)
        return float(df["weighted_value"].sum())

    def _weighted_pipeline_last_month(df_pipe: pd.DataFrame, focus_month) -> float:
        if df_pipe is None or df_pipe.empty:
            return 0.0
        df = df_pipe.copy()
        if "created_at" not in df.columns:
            return _compute_weighted_pipeline(df)

        df["created_at"] = pd.to_datetime(df["created_at"], errors="coerce", utc=True)
        cutoff = (
            pd.Timestamp(focus_month)
            .to_period("M")
            .to_timestamp()
            .tz_localize("UTC")
            + pd.DateOffset(months=1)
        )
        df = df[df["created_at"].notna() & (df["created_at"] < cutoff)].copy()
        return _compute_weighted_pipeline(df)

    def _get_actual_recognised_for_month_from_schedule_df(schedule_df: pd.DataFrame, focus_month) -> float:
        if schedule_df is None or schedule_df.empty:
            return 0.0
        df = schedule_df.copy()
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
        df = df[df["month_date"].notna()].copy()

        value_col = None
        for cand in ["recognised_revenue", "revenue_amount", "amount"]:
            if cand in df.columns:
                value_col = cand
                break
        if value_col is None:
            return 0.0

        m = _safe_month_start(focus_month)
        return float(pd.to_numeric(df.loc[df["month_date"] == m, value_col], errors="coerce").fillna(0.0).sum())

    def _get_next3_total_from_schedule_df(schedule_df: pd.DataFrame, focus_month) -> float:
        if schedule_df is None or schedule_df.empty:
            return 0.0
        df = schedule_df.copy()
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
        df = df[df["month_date"].notna()].copy()

        value_col = None
        for cand in ["recognised_revenue", "revenue_amount", "amount"]:
            if cand in df.columns:
                value_col = cand
                break
        if value_col is None:
            return 0.0

        focus_start = _safe_month_start(focus_month)
        end_4 = focus_start + pd.DateOffset(months=4)
        w = df[(df["month_date"] > focus_start) & (df["month_date"] < end_4)].copy()
        if w.empty:
            return 0.0
        return float(pd.to_numeric(w[value_col], errors="coerce").fillna(0.0).sum())

    def _avg_deal_size_active(df_pipe: pd.DataFrame) -> float:
        if df_pipe is None or df_pipe.empty:
            return 0.0
        active = df_pipe[df_pipe["stage"].isin(["proposal", "demo", "contract"])].copy()
        if active.empty:
            return 0.0
        active["value_total"] = pd.to_numeric(active.get("value_total"), errors="coerce").fillna(0.0)
        return float(active["value_total"].sum() / max(1, len(active)))

    def _cash_received_this_prev_month(client_id, focus_month) -> tuple[float, float]:
        df_ar, _ = fetch_ar_ap_for_client(client_id)
        if df_ar is None or df_ar.empty:
            return (0.0, 0.0)

        ar = df_ar.copy()
        if "payment_date" not in ar.columns or "amount" not in ar.columns:
            return (0.0, 0.0)

        ar["payment_date"] = pd.to_datetime(ar["payment_date"], errors="coerce")
        ar["amount"] = pd.to_numeric(ar["amount"], errors="coerce").fillna(0.0)
        ar["status"] = ar.get("status", "").astype(str).str.lower().fillna("")

        paid = ar[(ar["status"] == "paid") & ar["payment_date"].notna()].copy()
        if paid.empty:
            return (0.0, 0.0)

        focus_start = _safe_month_start(focus_month)
        prev_start = focus_start - pd.DateOffset(months=1)

        paid["pay_month"] = paid["payment_date"].dt.to_period("M").dt.to_timestamp()
        cash_this = float(paid.loc[paid["pay_month"] == focus_start, "amount"].sum())
        cash_prev = float(paid.loc[paid["pay_month"] == prev_start, "amount"].sum())
        return (cash_this, cash_prev)

    def _concentration_top2_weighted(df_pipe: pd.DataFrame) -> tuple[float, str]:
        if df_pipe is None or df_pipe.empty:
            return (0.0, "")
        df = df_pipe.copy()
        df["stage"] = df.get("stage", "").astype(str).str.lower().fillna("")
        df["value_total"] = pd.to_numeric(df.get("value_total"), errors="coerce").fillna(0.0)
        df["probability_pct"] = pd.to_numeric(df.get("probability_pct"), errors="coerce").fillna(0.0)

        df = df[df["stage"] != "closed_lost"].copy()
        if df.empty:
            return (0.0, "")

        def _pf(row) -> float:
            stg = str(row.get("stage", "")).lower()
            if stg == "closed_won":
                return 1.0
            p = float(row.get("probability_pct", 0.0) or 0.0)
            p = max(0.0, min(100.0, p))
            return p / 100.0

        df["weighted_value"] = df["value_total"] * df.apply(_pf, axis=1)
        total = float(df["weighted_value"].sum() or 0.0)
        if total <= 0:
            return (0.0, "")

        df = df.sort_values("weighted_value", ascending=False).copy()
        top2 = df.head(2)
        share = float(top2["weighted_value"].sum()) / total
        names = [str(r.get("deal_name") or "Unnamed") for _, r in top2.iterrows()]
        return (float(share), ", ".join(names[:2]))

    def _data_issues(df_pipe: pd.DataFrame) -> list[str]:
        issues = []
        if df_pipe is None or df_pipe.empty:
            return issues

        df = df_pipe.copy()
        df["stage"] = df.get("stage", "").astype(str).str.lower().fillna("")

        if "start_month" in df.columns:
            missing_close = pd.to_datetime(df["start_month"], errors="coerce").isna().sum()
            if missing_close > 0:
                issues.append(f"{missing_close} deal(s) missing expected close month")

        if "method" in df.columns:
            missing_method = df["method"].astype(str).str.strip().eq("").sum()
            if missing_method > 0:
                issues.append(f"{missing_method} deal(s) missing revenue recognition method")

        if "value_total" in df.columns:
            won = df[df["stage"] == "closed_won"].copy()
            if not won.empty:
                won_val = pd.to_numeric(won["value_total"], errors="coerce")
                bad = won_val.isna().sum() + int((won_val.fillna(0.0) <= 0.0).sum())
                if bad > 0:
                    issues.append(f"{bad} Closed_Won deal(s) missing/zero contract value")

        if (df["stage"] == "closed_lost").any():
            issues.append("Closed_Lost deals exist — ensure excluded from pipeline totals")

        return issues

    def _pipeline_ageing_risk_exists(df_pipe: pd.DataFrame, settings: dict) -> bool:
        if df_pipe is None or df_pipe.empty:
            return False

        max_proposal = int(settings.get("stage_max_days_proposal", 30) or 30)
        max_demo = int(settings.get("stage_max_days_demo", 45) or 45)
        max_contract = int(settings.get("stage_max_days_contract", 60) or 60)
        max_map = {"proposal": max_proposal, "demo": max_demo, "contract": max_contract}

        df = df_pipe.copy()
        df["stage"] = df.get("stage", "").astype(str).str.lower().fillna("")
        df = df[df["stage"].isin(["proposal", "demo", "contract"])].copy()
        if df.empty or "stage_entry_date" not in df.columns:
            return False

        df["stage_entry_date"] = pd.to_datetime(df["stage_entry_date"], errors="coerce")
        df = df[df["stage_entry_date"].notna()].copy()
        if df.empty:
            return False

        today = pd.Timestamp(datetime.now(timezone.utc).date())
        df["days_in_stage"] = (today - df["stage_entry_date"].dt.normalize()).dt.days.astype(int)
        df["max_days"] = df["stage"].map(max_map).fillna(0).astype(int)
        return bool((df["days_in_stage"] > df["max_days"]).any())

    def _pipeline_ageing_one_liner(df_pipe: pd.DataFrame, settings: dict, currency_symbol: str) -> str | None:
        if df_pipe is None or df_pipe.empty:
            return None

        max_proposal = int(settings.get("stage_max_days_proposal", 30) or 30)
        max_demo = int(settings.get("stage_max_days_demo", 45) or 45)
        max_contract = int(settings.get("stage_max_days_contract", 60) or 60)
        max_map = {"proposal": max_proposal, "demo": max_demo, "contract": max_contract}
        high_value_threshold = float(settings.get("pipeline_high_value_threshold", 100_000) or 100_000)

        active = df_pipe[df_pipe["stage"].isin(["proposal", "demo", "contract"])].copy()
        if active.empty:
            return None
        if "stage_entry_date" not in active.columns:
            return "⚠️ Pipeline ageing cannot be assessed (missing stage entry dates)."

        active["stage_entry_date"] = pd.to_datetime(active["stage_entry_date"], errors="coerce")
        if active["stage_entry_date"].isna().any():
            return "⚠️ Pipeline ageing data issue: missing stage entry date on one or more active deals."

        today = pd.Timestamp(datetime.now(timezone.utc).date())
        active["days_in_stage"] = (today - active["stage_entry_date"].dt.normalize()).dt.days.astype(int)
        active["max_days"] = active["stage"].map(max_map).fillna(0).astype(int)

        ageing = active[active["days_in_stage"] > active["max_days"]].copy()
        if ageing.empty:
            return None

        ageing["value_total"] = pd.to_numeric(ageing.get("value_total"), errors="coerce").fillna(0.0)
        late_stage_count = int((ageing["stage"].isin(["demo", "contract"])).sum())
        any_high_value = bool((ageing["value_total"] >= high_value_threshold).any())
        stalled_value = float(ageing["value_total"].sum())
        ageing_count = int(len(ageing))

        if late_stage_count >= 2:
            return f"🔴 {late_stage_count} late-stage deals ageing beyond expected stage duration"
        if any_high_value:
            return f"⚠️ High-value deal ageing beyond expected stage duration ({currency_symbol}{stalled_value:,.0f} impacted)"
        return f"ℹ️ {ageing_count} deal(s) ageing beyond expected stage duration ({currency_symbol}{stalled_value:,.0f})"

    def _rev_to_cash_one_liner(client_id, focus_month, settings: dict, secured_schedule_df: pd.DataFrame) -> str:
        focus_start = _safe_month_start(focus_month)
        prev_start = focus_start - pd.DateOffset(months=1)

        rev_this = 0.0
        rev_prev = 0.0
        if secured_schedule_df is not None and not secured_schedule_df.empty:
            sdf = secured_schedule_df.copy()
            sdf["month_date"] = pd.to_datetime(sdf["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
            sdf = sdf[sdf["month_date"].notna()].copy()

            val_col = None
            for cand in ["recognised_revenue", "revenue_amount", "amount"]:
                if cand in sdf.columns:
                    val_col = cand
                    break

            if val_col is not None:
                rev_this = float(pd.to_numeric(sdf.loc[sdf["month_date"] == focus_start, val_col], errors="coerce").fillna(0.0).sum())
                rev_prev = float(pd.to_numeric(sdf.loc[sdf["month_date"] == prev_start, val_col], errors="coerce").fillna(0.0).sum())

        df_ar, _df_ap = fetch_ar_ap_for_client(client_id)
        if df_ar is None or df_ar.empty:
            ar_days = int(settings.get("ar_default_days", 30) or 30)
            return f"Revenue typically converts to cash in ~{ar_days} days (expected)."

        ar = df_ar.copy()
        for c in ["issued_date", "payment_date", "amount", "status"]:
            if c in ar.columns:
                if c in ["issued_date", "payment_date"]:
                    ar[c] = pd.to_datetime(ar[c], errors="coerce")
                if c == "amount":
                    ar[c] = pd.to_numeric(ar[c], errors="coerce").fillna(0.0)

        if "payment_date" not in ar.columns or "amount" not in ar.columns:
            ar_days = int(settings.get("ar_default_days", 30) or 30)
            return f"Revenue typically converts to cash in ~{ar_days} days (expected)."

        paid = ar[(ar.get("status", "").astype(str).str.lower() == "paid")].copy()
        paid = paid[paid["payment_date"].notna()].copy()

        cash_this = 0.0
        cash_prev = 0.0
        if not paid.empty:
            paid["pay_month"] = paid["payment_date"].dt.to_period("M").dt.to_timestamp()
            cash_this = float(paid.loc[paid["pay_month"] == focus_start, "amount"].sum())
            cash_prev = float(paid.loc[paid["pay_month"] == prev_start, "amount"].sum())

        if "issued_date" in paid.columns and paid["issued_date"].notna().any():
            paid2 = paid[paid["issued_date"].notna()].copy()
            if not paid2.empty:
                lookback_start = focus_start - pd.DateOffset(months=6)
                paid2 = paid2[(paid2["payment_date"] >= lookback_start) & (paid2["payment_date"] < (focus_start + pd.DateOffset(months=1)))].copy()

                if not paid2.empty:
                    paid2["lag_days"] = (paid2["payment_date"].dt.normalize() - paid2["issued_date"].dt.normalize()).dt.days
                    paid2 = paid2[pd.to_numeric(paid2["lag_days"], errors="coerce").notna()].copy()

                    if not paid2.empty:
                        avg_lag = float(pd.to_numeric(paid2["lag_days"], errors="coerce").fillna(0.0).mean())

                        mid = focus_start - pd.DateOffset(months=3)
                        last3 = paid2[(paid2["payment_date"] >= mid)].copy()
                        prev3 = paid2[(paid2["payment_date"] < mid)].copy()

                        last3_lag = float(pd.to_numeric(last3["lag_days"], errors="coerce").fillna(0.0).mean()) if not last3.empty else avg_lag
                        prev3_lag = float(pd.to_numeric(prev3["lag_days"], errors="coerce").fillna(0.0).mean()) if not prev3.empty else last3_lag

                        if rev_this > rev_prev and cash_this < cash_prev:
                            return "🔴 Revenue up but cash down — collection risk rising"

                        if last3_lag > prev3_lag + 7:
                            return "⚠️ Revenue-to-cash lag increasing — expect delayed cash impact"

                        return f"ℹ️ Revenue typically converts to cash in ~{int(round(avg_lag))} days"

        ar_days = int(settings.get("ar_default_days", 30) or 30)
        return f"Revenue typically converts to cash in ~{ar_days} days (expected)."

    def _build_revenue_risk_section(
        df_pipe: pd.DataFrame,
        settings: dict,
        client_id: str,
        focus_month,
        secured_schedule_df: pd.DataFrame,
        currency_symbol: str,
    ) -> dict:
        risks = []
        data_issues = _data_issues(df_pipe)

        if _pipeline_ageing_risk_exists(df_pipe, settings):
            risks.append({"code": "A", "message": "Pipeline ageing detected — expected conversion may be delayed."})

        weighted_now = _compute_weighted_pipeline(df_pipe)
        weighted_prev = _weighted_pipeline_last_month(df_pipe, focus_month)

        rec_this = _get_actual_recognised_for_month_from_schedule_df(secured_schedule_df, focus_month)
        prev_month = (pd.to_datetime(focus_month).to_period("M").to_timestamp() - pd.DateOffset(months=1)).date()
        rec_prev = _get_actual_recognised_for_month_from_schedule_df(secured_schedule_df, prev_month)

        if weighted_now > weighted_prev and rec_this <= rec_prev:
            risks.append({"code": "B", "message": "Pipeline growth not converting into recognised revenue."})

        cash_this, cash_prev = _cash_received_this_prev_month(client_id, focus_month)
        rev_to_cash_line = _rev_to_cash_one_liner(client_id, focus_month, settings, secured_schedule_df)
        lag_deteriorating = ("lag increasing" in rev_to_cash_line.lower()) or ("collection risk" in rev_to_cash_line.lower())

        if (rec_this > rec_prev and cash_this < cash_prev) or lag_deteriorating:
            risks.append({"code": "C", "message": "Revenue recognised, but cash conversion is slower than expected."})

        top2_share, top2_names = _concentration_top2_weighted(df_pipe)
        conc_threshold = float(settings.get("pipeline_concentration_threshold", 0.50) or 0.50)
        if top2_share >= conc_threshold and top2_share > 0:
            risks.append({"code": "D", "message": f"Revenue outlook dependent on a small number of deals (Top 2 ~{top2_share*100:,.0f}%: {top2_names})."})

        bullets = []
        if rec_this >= rec_prev:
            bullets.append("Recognised revenue is stable vs last month, but forward growth depends on late-stage conversion.")
        else:
            bullets.append("Recognised revenue is down vs last month, increasing reliance on pipeline conversion to recover.")

        if any(r["code"] == "A" for r in risks):
            bullets.append("Aging pipeline suggests near-term revenue may slip if deals do not progress.")

        if any(r["code"] == "C" for r in risks):
            bullets.append("Cash impact from recognised revenue is deteriorating, implying higher working-capital pressure.")
        else:
            bullets.append(rev_to_cash_line.replace("ℹ️ ", "").replace("⚠️ ", "").replace("🔴 ", ""))

        if any(r["code"] == "D" for r in risks):
            bullets.append("Revenue outcomes are concentrated; slippage in one major deal will materially move results.")

        actions = []

        if any(r["code"] == "A" for r in risks):
            actions += [
                "Push stalled deals to decision or exit (book a buyer next-step within 72 hours).",
                "Revalidate deal status with the buyer and update stage_entry_date if it changed.",
            ]
        if any(r["code"] == "B" for r in risks):
            actions += [
                "Focus sales effort on late-stage (contract/demo) deals; deprioritise early-stage outreach for 7 days.",
                "Reassess conversion assumptions (win rate + cycle time) and adjust forecast expectations.",
            ]
        if any(r["code"] == "C" for r in risks):
            actions += [
                "Review payment terms on recent wins; tighten terms on any new contracts signed this week.",
                "Accelerate invoicing milestones and escalate collections on overdue invoices.",
            ]
        if any(r["code"] == "D" for r in risks):
            actions += ["Create a contingency plan: identify 3 next-best deals to advance if a top deal slips."]

        if data_issues:
            actions += [
                "Fix missing deal metadata (close month, method, value) and enforce pipeline hygiene immediately.",
                "Add a simple rule: no stage move without stage_entry_date update (and reason if moving backwards).",
            ]

        seen = set()
        deduped = []
        for a in actions:
            if a not in seen:
                deduped.append(a)
                seen.add(a)

        if len(deduped) < 3:
            deduped += [
                "Confirm next 7-day sales target is tied to contract-stage progression (not new leads).",
                "Remove dead deals from pipeline to avoid false confidence in weighted totals.",
            ]
            seen = set()
            tmp = []
            for a in deduped:
                if a not in seen:
                    tmp.append(a)
                    seen.add(a)
            deduped = tmp

        return {
            "risks": risks,
            "data_issues": data_issues,
            "what_this_means": bullets[:4],
            "actions": deduped[:5],
        }

    def _render_pipeline_waterfall(df_pipe: pd.DataFrame, currency_symbol: str):
        """
        Waterfall chart (WEIGHTED). ✅ TABLE REMOVED.
        Stages: Idea -> Demo -> Proposal -> Contract -> Closed Won -> Total Estimated
        """
        import matplotlib.pyplot as plt

        if df_pipe is None or df_pipe.empty:
            st.caption("No pipeline data to display.")
            return

        st.subheader("🧩 Pipeline Waterfall (Weighted)")

        dfw = df_pipe.copy()
        dfw["stage"] = dfw.get("stage", "").astype(str).str.lower().fillna("")
        dfw["value_total"] = pd.to_numeric(dfw.get("value_total"), errors="coerce").fillna(0.0)
        dfw["probability_pct"] = pd.to_numeric(dfw.get("probability_pct"), errors="coerce").fillna(0.0)

        stage_order = ["idea", "demo", "proposal", "contract", "closed_won"]
        dfw = dfw[dfw["stage"].isin(stage_order)].copy()
        if dfw.empty:
            st.caption("No staged pipeline rows to display.")
            return

        def _prob_factor(row) -> float:
            stg = str(row.get("stage", "")).lower()
            if stg == "closed_won":
                return 1.0
            p = float(row.get("probability_pct", 0.0) or 0.0)
            p = max(0.0, min(100.0, p))
            return p / 100.0

        dfw["weighted_value"] = dfw["value_total"] * dfw.apply(_prob_factor, axis=1)

        stage_sum = (
            dfw.groupby("stage", as_index=False)[["weighted_value"]]
            .sum()
            .set_index("stage")
            .reindex(stage_order)
            .fillna(0.0)
            .reset_index()
        )
        stage_sum["Stage"] = stage_sum["stage"].str.replace("_", " ").str.title()

        labels = stage_sum["Stage"].tolist()
        values = stage_sum["weighted_value"].astype(float).tolist()
        total_estimated = float(sum(values))

        labels_total = labels + ["Total Estimated"]
        step_heights = values + [total_estimated]

        bases = []
        cum = 0.0
        for v in values:
            bases.append(cum)
            cum += float(v)
        bases_total = bases + [0.0]

        fig, ax = plt.subplots(figsize=(9, 4))
        x = list(range(len(labels_total)))
        ax.bar(x, step_heights, bottom=bases_total)

        ax.set_xticks(x)
        ax.set_xticklabels(labels_total, rotation=0)
        ax.set_ylabel("Weighted pipeline value")
        ax.set_title("Pipeline Waterfall (Weighted)")

        for i, (b, h) in enumerate(zip(bases_total, step_heights)):
            ax.text(i, b + h, f"{currency_symbol}{h:,.0f}", ha="center", va="bottom", fontsize=9)

        ax.axhline(0, linewidth=1)
        ax.margins(x=0.02)
        st.pyplot(fig)

    # ============================================================
    # Pipeline overview (single-column now)
    # ============================================================
    st.subheader("📂 Pipeline overview")

    if df_pipeline.empty:
        st.info("No deals in the pipeline yet for this business.")
    else:
        k1, k2 = st.columns(2)
        with k1:
            st.metric("Total pipeline", f"{currency_symbol}{total_pipeline:,.0f}")
        with k2:
            st.metric("Open deals", open_deals_count)

        # Build secured schedule ONCE
        snap_df = build_revenue_schedule_for_client(
            selected_client_id,
            base_month=selected_month_start,
            n_months=12,
            mode="secured",
        )

        rec_this = _get_actual_recognised_for_month_from_schedule_df(snap_df, selected_month_start)
        rec_next3 = _get_next3_total_from_schedule_df(snap_df, selected_month_start)

        rk1, rk2 = st.columns(2)
        with rk1:
            st.metric("Recognised revenue (this month)", f"{currency_symbol}{rec_this:,.0f}")
        with rk2:
            st.metric("Recognised revenue (next 3 months)", f"{currency_symbol}{rec_next3:,.0f}")

        plan_rev = _get_plan_for_month(selected_client_id, selected_month_start)
        rev_vs_plan_pct = (rec_this - plan_rev) / plan_rev if plan_rev else None
        avg_deal = _avg_deal_size_active(df_pipeline)

        kk1, kk2 = st.columns(2)
        with kk1:
            if plan_rev <= 0:
                st.metric("Revenue vs Plan", "Plan not set")
            else:
                st.metric("Revenue vs Plan", f"{rev_vs_plan_pct*100:,.1f}%")
        with kk2:
            st.metric("Average deal size (Active only)", f"{currency_symbol}{avg_deal:,.0f}")

        ageing_line = _pipeline_ageing_one_liner(df_pipeline, settings, currency_symbol)
        if ageing_line:
            st.caption(ageing_line)

        cash_line = _rev_to_cash_one_liner(selected_client_id, selected_month_start, settings, snap_df)
        st.caption(cash_line)

        horizon_start = pd.to_datetime(selected_month_start).replace(day=1)
        horizon_end = horizon_start + pd.DateOffset(months=12)
        st.caption(
            f"Views below use deals closing between "
            f"{horizon_start.strftime('%b %Y')} and "
            f"{(horizon_end - pd.DateOffset(days=1)).strftime('%b %Y')} "
            "(next 12 months window)."
        )

        # Health panel (uses your updated build_pipeline_health_panel without the 90d KPI)
        sev, decisions, metrics = build_pipeline_health_panel(df_pipeline, selected_month_start, currency_symbol)

        if metrics:
            mcols = st.columns(len(metrics))
            for i, (label, val) in enumerate(metrics):
                mcols[i].metric(label, val)

        if sev == "error":
            st.error("\n".join([f"• {d}" for d in decisions]))
        elif sev == "warning":
            st.warning("\n".join([f"• {d}" for d in decisions]))
        else:
            st.success("\n".join([f"• {d}" for d in decisions]))

        # Data checks warning (external; unchanged)
        issues = build_pipeline_data_checks(df_pipeline)
        if issues:
            st.warning("**Data checks**\n" + "\n".join([f"• {i}" for i in issues]))

        # View full pipeline list expander (kept)
        display_cols = [
            c
            for c in [
                "deal_name",
                "customer_name",
                "stage",
                "value_total",
                "probability_pct",
                "start_month",
                "end_month",
                "method",
                "contract_months",
                "stage_entry_date",
            ]
            if c in df_pipeline.columns
        ]
        if display_cols:
            view = df_pipeline[display_cols].copy().rename(
                columns={
                    "deal_name": "Deal",
                    "customer_name": "Customer",
                    "stage": "Stage",
                    "value_total": "Value",
                    "probability_pct": "Win chance (%)",
                    "start_month": "Start month",
                    "end_month": "End month",
                    "method": "Revenue type",
                    "contract_months": "Contract months",
                    "stage_entry_date": "Stage entry date",
                }
            )
            with st.expander("View full pipeline list", expanded=False):
                st.dataframe(view, width="stretch")

    # ============================================================
    # Waterfall chart (TABLE REMOVED)
    # ============================================================
    _render_pipeline_waterfall(df_pipeline, currency_symbol)

    # ============================================================
    # Revenue / Risk / What this means / Action plan (NO tables)
    # ============================================================
    st.markdown("---")
    st.subheader("🚦 Revenue, Risk & data issues")

    try:
        snap_df
    except NameError:
        snap_df = build_revenue_schedule_for_client(
            selected_client_id,
            base_month=selected_month_start,
            n_months=12,
            mode="secured",
        )

    bundle = _build_revenue_risk_section(
        df_pipe=df_pipeline,
        settings=settings,
        client_id=selected_client_id,
        focus_month=selected_month_start,
        secured_schedule_df=snap_df,
        currency_symbol=currency_symbol,
    )

    if not bundle["risks"]:
        st.success("No revenue risk triggers fired for this month based on current data.")
    else:
        for r in bundle["risks"]:
            st.error(f"🔴 Risk {r['code']} — {r['message']}")

    if bundle["data_issues"]:
        st.warning("⚠️ Data issues detected (hygiene):")
        for di in bundle["data_issues"]:
            st.write(f"• {di}")
    else:
        st.caption("No data hygiene issues detected based on available fields.")

    st.markdown("#### What this means")
    for b in bundle["what_this_means"][:4]:
        st.write(f"• {b}")

    st.markdown("#### Action plan (next 7 days)")
    for i, a in enumerate(bundle["actions"][:5], start=1):
        st.write(f"{i}. {a}")



def create_pipeline_deal(
    client_id,
    deal_name,
    customer_name,
    value_total,
    probability_pct,
    method,
    stage,
    start_month,
    end_month,
    commentary,
    contract_months=None,
):
    """
    Create a new pipeline deal for this client.

    - start_month / end_month can be date, datetime or string
    - contract_months (int) optional (used by SaaS / straight-line / POC)
    - ✅ stage_entry_date is set on create (required for Pipeline Ageing Warning KPI)
    """

    def _to_month_iso(d):
        if d is None:
            return None
        dt = pd.to_datetime(d, errors="coerce")
        if pd.isna(dt):
            return None
        return dt.to_period("M").to_timestamp().date().isoformat()

    stage_norm = str(stage).lower().strip() if stage is not None else "proposal"

    payload = {
        "client_id": str(client_id),
        "deal_name": str(deal_name).strip(),
        "customer_name": str(customer_name).strip() if customer_name is not None else None,
        "value_total": float(value_total) if value_total is not None else 0.0,
        "probability_pct": float(probability_pct) if probability_pct is not None else 0.0,
        "method": str(method).strip() if method is not None else None,
        "stage": stage_norm,
        "start_month": _to_month_iso(start_month),
        "end_month": _to_month_iso(end_month),
        "commentary": str(commentary) if commentary is not None else None,
        # ✅ REQUIRED FOR AGEING KPI
        # Use "today" as stage entry date when creating.
        "stage_entry_date": datetime.now(timezone.utc).date().isoformat(),
        "updated_at": _utc_now_iso(),
    }

    # enforce stage rules
    if payload["stage"] == "closed_won":
        payload["probability_pct"] = 100.0
    elif payload["stage"] == "closed_lost":
        payload["probability_pct"] = 0.0

    if contract_months is not None:
        try:
            payload["contract_months"] = int(contract_months)
        except Exception:
            pass

    try:
        supabase.table("revenue_pipeline").insert(payload).execute()
        _clear_all_caches_safely()
        return True
    except Exception as e:
        print("Error creating deal:", e)
        return False


def create_investing_flow(
    client_id,
    month_date: date,
    amount: float,
    category: str | None = None,
    notes: str | None = None,
) -> bool:
    if client_id is None or month_date is None:
        return False

    try:
        payload = {
            "client_id": str(client_id),
            "month_date": month_date.isoformat(),
            "amount": float(amount),
            "category": category or None,
            "notes": notes or None,
        }
        resp = supabase.table("investing_flows").insert(payload).execute()

        error_obj = getattr(resp, "error", None)
        if error_obj:
            print("Error creating investing_flow:", error_obj)
            return False

        if isinstance(resp, dict) and resp.get("error"):
            print("Error creating investing_flow:", resp["error"])
            return False

        return True
    except Exception as e:
        print("Exception creating investing_flow:", e)
        return False


def create_financing_flow(
    client_id,
    month_date: date,
    amount: float,
    category: str | None = None,
    notes: str | None = None,
) -> bool:
    if client_id is None or month_date is None:
        return False

    try:
        payload = {
            "client_id": str(client_id),
            "month_date": month_date.isoformat(),
            "amount": float(amount),
            "category": category or None,
            "notes": notes or None,
        }
        resp = supabase.table("financing_flows").insert(payload).execute()

        error_obj = getattr(resp, "error", None)
        if error_obj:
            print("Error creating financing_flow:", error_obj)
            return False

        if isinstance(resp, dict) and resp.get("error"):
            print("Error creating financing_flow:", resp["error"])
            return False

        return True
    except Exception as e:
        print("Exception creating financing_flow:", e)
        return False

# ---------------------------------------------------------
# Investing & Financing flows – simple fetch helpers
# ---------------------------------------------------------

def fetch_investing_flows_for_client(client_id: str) -> pd.DataFrame | None:
    """
    Read investing_flows for this client.

    Returns:
        - None  -> Supabase connection problem
        - empty DataFrame -> call succeeded but no rows
        - non-empty DataFrame -> valid data
    """
    if not client_id:
        return pd.DataFrame()

    def _call():
        return (
            supabase
            .table("investing_flows")
            .select("*")
            .eq("client_id", str(client_id))
            .execute()
        )

    resp = sb_execute_with_retry(_call, label="fetch investing_flows")
    if resp is None:
        print("[ERROR] fetch_investing_flows_for_client -> connection problem (None)")
        return None

    data = getattr(resp, "data", None) if hasattr(resp, "data") else resp.get("data")
    df = pd.DataFrame(data or [])
    print(f"[DEBUG] fetch_investing_flows_for_client -> {len(df)} rows for {client_id}")
    return df


PAID_STATUSES = {"paid", "closed", "settled"}
NON_CASH_STATUSES = {"cancelled", "disputed"}


def normalise_ar_ap_for_cash(df: pd.DataFrame | None) -> pd.DataFrame:
    """
    Clean AR/AP rows for cash forecasting:

    - Drop fully paid / closed / settled / cancelled / disputed
    - Normalise `status` to lowercase
    - Ensure numeric 'amount' and 'amount_paid'
    - Compute 'effective_amount' = balance still due:
        * if partially_paid == True -> amount - amount_paid
        * otherwise -> amount
    - Add 'is_cash_relevant' for quick filtering (effective_amount > 0)
    """
    if df is None or df.empty:
        return pd.DataFrame()

    out = df.copy()

    # Normalise status and drop non-cash rows
    if "status" in out.columns:
        out["status"] = out["status"].astype(str).str.strip().str.lower()
        drop_statuses = PAID_STATUSES | NON_CASH_STATUSES
        out = out[~out["status"].isin(drop_statuses)]

    if out.empty:
        return out

    # Base amount
    out["amount"] = pd.to_numeric(out.get("amount"), errors="coerce").fillna(0.0)

    # Amount already paid
    if "amount_paid" in out.columns:
        out["amount_paid"] = pd.to_numeric(
            out.get("amount_paid"), errors="coerce"
        ).fillna(0.0)
    else:
        out["amount_paid"] = 0.0

    # Partially paid flag
    if "partially_paid" in out.columns:
        out["partially_paid"] = out["partially_paid"].fillna(False).astype(bool)
    else:
        out["partially_paid"] = False

    # Balance still due
    def _effective(row) -> float:
        amt = float(row.get("amount") or 0.0)
        amt_paid = float(row.get("amount_paid") or 0.0)
        if bool(row.get("partially_paid", False)):
            return max(0.0, amt - amt_paid)
        return amt

    out["effective_amount"] = out.apply(_effective, axis=1)

    # No negative balances
    out["effective_amount"] = out["effective_amount"].clip(lower=0.0)

    # Convenience flag
    out["is_cash_relevant"] = out["effective_amount"] > 0

    return out

def get_or_create_opening_cash(client_id: str, as_of_date: date) -> float:
    # 1) Try read from cash_opening_balance
    res = (
        supabase.table("cash_opening_balance")
        .select("opening_cash")
        .eq("client_id", str(client_id))
        .eq("as_of_month", as_of_date.isoformat())
        .limit(1)
        .execute()
    )
    if res.data:
        return float(res.data[0]["opening_cash"] or 0.0)

    # 2) If missing: seed from KPI (one time only)
    seed = 0.0
    kpi_row = fetch_kpis_for_client_month(client_id, as_of_date)
    if kpi_row and "cash_balance" in kpi_row:
        try:
            seed = float(kpi_row["cash_balance"] or 0.0)
        except Exception:
            seed = 0.0

    # 3) Insert seed row
    supabase.table("cash_opening_balance").upsert(
        {
            "client_id": str(client_id),
            "as_of_month": as_of_date.isoformat(),
            "opening_cash": round(seed, 2),
            "source": "kpi_seed",
        },
        on_conflict="client_id,as_of_month",
    ).execute()

    return seed


@st.cache_data(ttl=60)
def fetch_ar_ap_for_client(client_id):
    """
    Fetch AR and AP rows for the given client from ar_ap_tracker.
    Returns two DataFrames: df_ar, df_ap
    """
    if client_id is None:
        return pd.DataFrame(), pd.DataFrame()

    try:
        res = (
            supabase
            .table("ar_ap_tracker")
            .select("*")
            .eq("client_id", str(client_id))
            .execute()
        )
        rows = res.data or []
    except Exception:
        rows = []

    if not rows:
        return pd.DataFrame(), pd.DataFrame()

    df = pd.DataFrame(rows)

    # Format dates nicely if they exist
    for col in ["due_date", "expected_payment_date", "created_at", "updated_at", "issued_date"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce")

    # Split into AR and AP
    if "type" in df.columns:
        df_ar = df[df["type"] == "AR"].copy()
        df_ap = df[df["type"] == "AP"].copy()
    else:
        df_ar = pd.DataFrame()
        df_ap = pd.DataFrame()

    # Get client defaults
    settings = get_client_settings(client_id)
    client_ar_days = settings["ar_default_days"]

    # For AR: compute expected_date
    if not df_ar.empty:
        def _compute_expected(row):
            issued = row.get("issued_date")
            row_days = row.get("default_receipt_days")
            try:
                if issued is not None:
                    days = row_days if row_days is not None else client_ar_days
                    return issued + timedelta(days=int(days))
            except Exception:
                pass
            return row.get("due_date")

        df_ar["expected_date"] = df_ar.apply(_compute_expected, axis=1)

        # Keep everything as proper datetimes (not Python date)
        for col in ["issued_date", "due_date", "expected_date"]:
            if col in df_ar.columns:
                df_ar[col] = pd.to_datetime(df_ar[col], errors="coerce")


    return df_ar, df_ap


# ----------AR Aging ----------#

PAID_STATUSES = {"paid", "closed", "settled"}

def _to_date_safe(x):
    if x is None or (isinstance(x, float) and pd.isna(x)) or pd.isna(x):
        return None
    if isinstance(x, date) and not isinstance(x, datetime):
        return x
    try:
        return pd.to_datetime(x, errors="coerce").date()
    except Exception:
        return None

def _pick_first(row, candidates):
    for c in candidates:
        if c in row and row.get(c) not in (None, "", "nan"):
            return row.get(c)
    return None

def add_ar_aging(df_ar: pd.DataFrame, as_of: date, ar_default_days: int) -> pd.DataFrame:
    """
    Adds:
      - expected_date
      - days_past_expected
      - aging_bucket
      - is_over_default
    """
    if df_ar is None or df_ar.empty:
        return df_ar

    df = df_ar.copy()

    # Normalize common date columns if present
    for col in ["issue_date", "issued_date", "issue_Date", "due_date", "expected_date"]:
        if col in df.columns:
            df[col] = df[col].apply(_to_date_safe)

    # Normalize status
    if "status" in df.columns:
        df["status"] = df["status"].astype(str).str.strip().str.lower()
    else:
        df["status"] = ""

    # Compute expected_date if missing or mostly empty
    if "expected_date" not in df.columns:
        df["expected_date"] = None

    def _compute_expected(row):
        # If already has expected_date, keep it
        exp = row.get("expected_date")
        if exp:
            return exp

        # Prefer due_date if available (strongest explicit expectation)
        due = row.get("due_date")
        if due:
            return due

        issued = _pick_first(row, ["issue_date", "issued_date", "issue_Date"])
        if not issued:
            return None

        # Row-level override days (your schema)
        row_days = row.get("expected_payment_days", None)
        try:
            days = int(row_days) if row_days is not None and not pd.isna(row_days) else int(ar_default_days)
        except Exception:
            days = int(ar_default_days)

        return issued + timedelta(days=days)

    df["expected_date"] = df.apply(_compute_expected, axis=1)

    def _is_paid(row):
        return str(row.get("status", "")).strip().lower() in PAID_STATUSES

    def _days_past_expected(row):
        if _is_paid(row):
            return 0
        exp = row.get("expected_date")
        if not exp:
            return 0
        d = (as_of - exp).days
        return d

    df["days_past_expected"] = df.apply(_days_past_expected, axis=1).fillna(0).astype(int)

    def _bucket(days):
        if days <= 0:
            return "Not due / current"
        if days <= 30:
            return "0–30 days overdue"
        if days <= 60:
            return "30–60 days overdue"
        if days <= 90:
            return "60–90 days overdue"
        return "90+ days overdue"

    df["aging_bucket"] = df["days_past_expected"].apply(_bucket)
    df["is_over_default"] = df["days_past_expected"] > 0

    return df



def rebuild_cashflow_summary_for_client(client_id):
    """
    Core cashflow engine v1.

    Logic:
      - Read baseline_monthly for this client
      - For each month:
          * take closing_cash from baseline_monthly
          * compute net change vs previous month = operating_cf (for now)
          * investing_cf = 0, financing_cf = 0 (v1 placeholder)
          * free_cash_flow = operating_cf + investing_cf
          * cash_danger_flag = True if closing_cash <= 0
      - Upsert into cashflow_summary on (client_id, month_date)
    """
    if client_id is None:
        return False, "No client selected."

    df_base = fetch_baseline_monthly_for_client(client_id)
    if df_base.empty:
        return False, "No baseline_monthly data found for this client."

    # Make sure we have a month_date and closing_cash-like column
    if "month_date" not in df_base.columns:
        return False, "baseline_monthly has no month_date column after normalisation."

    # Try to find closing_cash / cash balance column
    cash_col = None
    for cand in ["closing_cash", "cash_balance", "closing_cash_balance"]:
        if cand in df_base.columns:
            cash_col = cand
            break

    if cash_col is None:
        return False, "baseline_monthly has no closing_cash / cash_balance column."

    # Sort by month
    df_base = df_base.sort_values("month_date").reset_index(drop=True)

    records = []
    prev_closing = None

    for _, row in df_base.iterrows():
        month_ts = row["month_date"]
        if pd.isna(month_ts):
            continue

        # Convert to python date
        if isinstance(month_ts, datetime):
            month_d = month_ts.date().replace(day=1)
        else:
            month_d = pd.to_datetime(month_ts).date().replace(day=1)

        raw_closing = row.get(cash_col)
        try:
            closing_cash = float(raw_closing) if raw_closing is not None else None
        except (TypeError, ValueError):
            closing_cash = None

        # Net change vs previous month (simple v1: treat as operating_CF)
        if prev_closing is None or closing_cash is None:
            operating_cf = None
        else:
            operating_cf = closing_cash - prev_closing

        prev_closing = closing_cash if closing_cash is not None else prev_closing

        investing_cf = 0.0  # placeholder – later you can feed real capex etc.
        financing_cf = 0.0  # placeholder – later you can feed raises, loans, repayments

        # Free cash flow (v1 = operating + investing, ignoring financing)
        free_cash_flow = None
        if operating_cf is not None:
            free_cash_flow = operating_cf + investing_cf

        cash_danger_flag = bool(
            (closing_cash is not None) and (closing_cash <= 0)
        )

        rec = {
            "client_id": str(client_id),
            "month_date": month_d.isoformat(),
            "operating_cf": operating_cf,
            "investing_cf": investing_cf,
            "financing_cf": financing_cf,
            "free_cash_flow": free_cash_flow,
            "closing_cash": closing_cash,
            "cash_danger_flag": cash_danger_flag,
        }
        records.append(rec)

    if not records:
        return False, "No valid months were found to write to cashflow_summary."

    try:
        # Upsert on (client_id, month_date) as per your unique constraint
        supabase.table("cashflow_summary").upsert(
            records,
            on_conflict="client_id,month_date",
        ).execute()
        return True, f"Rebuilt cashflow_summary for {len(records)} month(s)."
    except Exception as e:
        return False, f"Failed to upsert cashflow_summary: {e}"

def add_ap_aging(df_ap: pd.DataFrame, as_of: date, ap_default_days: int) -> pd.DataFrame:
    """
    Adds:
      - pay_expected_date
      - days_past_expected
      - aging_bucket
      - is_over_default
    """
    if df_ap is None or df_ap.empty:
        return df_ap

    df = df_ap.copy()

    # Normalize date columns
    for col in ["issue_date", "issued_date", "issue_Date", "due_date", "pay_expected_date", "expected_payment_date"]:
        if col in df.columns:
            df[col] = df[col].apply(_to_date_safe)

    # Normalize status
    if "status" in df.columns:
        df["status"] = df["status"].astype(str).str.strip().str.lower()
    else:
        df["status"] = ""

    # Ensure pay_expected_date exists
    if "pay_expected_date" not in df.columns:
        df["pay_expected_date"] = None

    def _compute_pay_expected(row):
        # If already has pay_expected_date, keep it
        pe = row.get("pay_expected_date")
        if pe:
            return pe

        # If you store an explicit expected_payment_date, use it
        epd = row.get("expected_payment_date")
        if epd:
            return epd

        # Then due_date
        due = row.get("due_date")
        if due:
            return due

        issued = _pick_first(row, ["issue_date", "issued_date", "issue_Date"])
        if not issued:
            return None

        # Row-level override days (same column name in your tracker)
        row_days = row.get("expected_payment_days", None)
        try:
            days = int(row_days) if row_days is not None and not pd.isna(row_days) else int(ap_default_days)
        except Exception:
            days = int(ap_default_days)

        return issued + timedelta(days=days)

    df["pay_expected_date"] = df.apply(_compute_pay_expected, axis=1)

    def _is_paid(row):
        return str(row.get("status", "")).strip().lower() in PAID_STATUSES

    def _days_past_expected(row):
        if _is_paid(row):
            return 0
        exp = row.get("pay_expected_date")
        if not exp:
            return 0
        return (as_of - exp).days

    df["days_past_expected"] = df.apply(_days_past_expected, axis=1).fillna(0).astype(int)

    df["aging_bucket"] = df["days_past_expected"].apply(lambda d: (
        "Not due / current" if d <= 0 else
        "0–30 days overdue" if d <= 30 else
        "30–60 days overdue" if d <= 60 else
        "60–90 days overdue" if d <= 90 else
        "90+ days overdue"
    ))

    df["is_over_default"] = df["days_past_expected"] > 0

    return df



def build_cash_commitments(
    df_ar: pd.DataFrame | None,
    df_ap: pd.DataFrame | None,
    as_of: date,
    horizon_days: int = 60,
    limit: int = 7,
) -> pd.DataFrame:
    """
    Build a short list of upcoming cash commitments around the Focus Month.

    Rules:
      - Only UNPAID AR/AP are included
        (paid / closed / settled / cancelled / disputed are dropped).
      - Timing is driven by DUE DATE (fallbacks only if due_date missing).
      - If an item is overdue (original_due < as_of) and still unpaid,
        we treat it as 'due now' → date = as_of.
      - Only commitments between as_of and as_of + horizon_days are shown.
      - Uses 'effective_amount' where available (balance due after partial payments).
      - Adds:
          * original_due (original invoice/bill due date)
          * overdue ("Yes"/"No") based on original_due < as_of
      - Sorted by date then amount (biggest first on same day).
    """
    frames: list[pd.DataFrame] = []
    horizon_end = as_of + timedelta(days=horizon_days)

    # ---------- AR: cash in ----------
    ar_norm = normalise_ar_ap_for_cash(df_ar)
    if ar_norm is not None and not ar_norm.empty:
        tmp = ar_norm.copy()
        tmp["direction"] = "Cash in"
        tmp["who"] = tmp.get("counterparty", "")

        # Use balance due if present
        amt_col = "effective_amount" if "effective_amount" in tmp.columns else "amount"
        tmp["amount"] = pd.to_numeric(tmp.get(amt_col, 0), errors="coerce").fillna(0.0)

        # Use due_date first, else expected_date
        if "due_date" in tmp.columns:
            base_dates = tmp["due_date"]
        elif "expected_date" in tmp.columns:
            base_dates = tmp["expected_date"]
        else:
            base_dates = pd.NaT

        # Original due = real invoice due
        tmp["original_due"] = pd.to_datetime(base_dates, errors="coerce").dt.date

        # Scheduling date used in this commitments grid
        tmp["date"] = tmp["original_due"]

        # Overdue but unpaid → treat as "due now" in this Focus Month
        overdue_mask = tmp["original_due"] < as_of
        tmp.loc[overdue_mask, "date"] = as_of

        # Keep commitments in [as_of, as_of + horizon_days]
        mask_window = (
            (tmp["date"] >= as_of) &
            (tmp["date"] <= horizon_end)
        )
        tmp = tmp[mask_window]

        if not tmp.empty:
            frames.append(
                tmp[["direction", "who", "amount", "date", "original_due"]]
            )

    # ---------- AP: cash out ----------
    ap_norm = normalise_ar_ap_for_cash(df_ap)
    if ap_norm is not None and not ap_norm.empty:
        tmp = ap_norm.copy()
        tmp["direction"] = "Cash out"
        tmp["who"] = tmp.get("counterparty", "")

        amt_col = "effective_amount" if "effective_amount" in tmp.columns else "amount"
        tmp["amount"] = pd.to_numeric(tmp.get(amt_col, 0), errors="coerce").fillna(0.0)

        # Use due_date first, then expected_payment_date, then pay_expected_date
        if "due_date" in tmp.columns:
            base_dates = tmp["due_date"]
        elif "expected_payment_date" in tmp.columns:
            base_dates = tmp["expected_payment_date"]
        elif "pay_expected_date" in tmp.columns:
            base_dates = tmp["pay_expected_date"]
        else:
            base_dates = pd.NaT

        tmp["original_due"] = pd.to_datetime(base_dates, errors="coerce").dt.date
        tmp["date"] = tmp["original_due"]

        overdue_mask = tmp["original_due"] < as_of
        tmp.loc[overdue_mask, "date"] = as_of

        mask_window = (
            (tmp["date"] >= as_of) &
            (tmp["date"] <= horizon_end)
        )
        tmp = tmp[mask_window]

        if not tmp.empty:
            frames.append(
                tmp[["direction", "who", "amount", "date", "original_due"]]
            )

    # ---------- Combine + Overdue flag ----------
    if not frames:
        return pd.DataFrame()

    combined = pd.concat(frames, ignore_index=True)

    # Overdue? badge (based on original_due vs as_of)
    combined["overdue"] = combined["original_due"].apply(
        lambda d: "Yes" if (pd.notna(d) and d < as_of) else "No"
    )

    # Normalise to python date for safety
    combined["date"] = pd.to_datetime(combined["date"], errors="coerce").dt.date

    # Sort by scheduling date, then amount (largest first on same day)
    combined = combined.sort_values(
        by=["date", "amount"],
        ascending=[True, False],
        ignore_index=True,
    )

    # Apply limit
    combined = combined.head(limit)

    return combined




def save_comment(client_id, page_name: str, body: str, month_start: date):
    """
    Insert a comment row into the comments table.
    """
    if not body or not body.strip():
        return False

    data = {
        "client_id": str(client_id) if client_id else None,
        "page_name": page_name,
        "context_type": "generic",
        "context_id": None,
        "month_date": month_start.isoformat() if month_start else None,
        "author_name": "Internal CFO",  # later replace with real user
        "author_role": "advisor",
        "body": body.strip(),
    }

    try:
        supabase.table("comments").insert(data).execute()
        return True
    except Exception:
        return False

@st.cache_data(ttl=60)
def fetch_recent_comments_for_client(client_id, limit: int = 50):
    """
    Fetch recent comments for this client across all pages/months.
    """
    if client_id is None:
        return []

    try:
        res = (
            supabase
            .table("comments")
            .select("*")
            .eq("client_id", str(client_id))
            .order("created_at", desc=True)
            .limit(limit)
            .execute()
        )
        return res.data or []
    except Exception:
        return []


def fetch_overdue_ar_for_client(client_id, as_of: date | None = None):
    """
    Fetch overdue AR rows for this client based on expected receipt date:
      expected_date = issued_date + default_receipt_days (row-level)
                      or issued_date + ar_default_days (client-level)
                      or fallback: due_date
    """
    if client_id is None:
        return pd.DataFrame()

    if as_of is None:
        as_of = date.today()

    # Get client-level defaults
    settings = get_client_settings(client_id)
    client_ar_days = settings["ar_default_days"]

    try:
        res = (
            supabase
            .table("ar_ap_tracker")
            .select("*")
            .eq("client_id", str(client_id))
            .eq("type", "AR")
            .execute()
        )
        rows = res.data or []
    except Exception:
        rows = []

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)

    # Parse dates
    for col in ["issued_date", "due_date"]:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce").dt.date

    # Compute expected_date
    def _compute_expected(row):
        issued = row.get("issued_date")
        row_days = row.get("default_receipt_days")
        try:
            if issued is not None:
                days = row_days if row_days is not None else client_ar_days
                return issued + timedelta(days=int(days))
        except Exception:
            pass
        # fallback
        return row.get("due_date")

    df["expected_date"] = df.apply(_compute_expected, axis=1)

    df = df[df["expected_date"].notna()]

   
    # Make sure expected_date is a date, not full datetime
    if "expected_date" in df.columns:
        df["expected_date"] = pd.to_datetime(df["expected_date"], errors="coerce").dt.date
    else:
        return pd.DataFrame()  # nothing to check

    # Make sure as_of is also a date object
    if isinstance(as_of, pd.Timestamp):
        as_of_date = as_of.date()
    else:
        as_of_date = as_of  # already a date

    # Keep only rows with a valid expected_date and overdue
    df = df[df["expected_date"].notna()]
    df = df[df["expected_date"] < as_of_date]


    # Exclude fully paid if we have a status column
    if "status" in df.columns:
        df = df[~df["status"].str.lower().isin(["paid", "closed", "settled"])]

    return df

@st.cache_data(ttl=60)
def build_12m_pnl_and_cash_view(
    client_id,
    base_month: date,
    n_months: int = 12,
) -> pd.DataFrame:
    """
    12-month view combining:
      - recognised revenue (from revenue engine)
      - payroll cost (from payroll_positions)
      - vendor cost (from AP)
      - cash in (AR receipts)
      - cash out (AP + payroll)
      - operating CF + closing cash

    Returns one row per month with:
      month_date, revenue_pnl, payroll_pnl, vendor_pnl,
      total_cost_pnl, net_burn_pnl,
      cash_in_ar, cash_out_ap, cash_out_payroll,
      operating_cf, closing_cash, cash_danger_flag
    """

    if client_id is None or base_month is None:
        return pd.DataFrame()

    # Month index (month starts)
    month_index = pd.date_range(
        start=pd.to_datetime(base_month),
        periods=n_months,
        freq="MS",
    )
    month_set = set(month_index)

    # ------------------------------
    # 1) Recognised revenue by month
    # ------------------------------
    rev_df = build_revenue_schedule_for_client(
        client_id,
        base_month=base_month,
        n_months=n_months,
    )

    if rev_df is None or rev_df.empty:
        rev_series = pd.Series(0.0, index=month_index)
    else:
        tmp = rev_df.copy()
        tmp["month_date"] = pd.to_datetime(tmp["month_date"], errors="coerce")
        tmp = tmp[tmp["month_date"].notna()]

        # pick revenue column
        value_col = None
        for cand in ["revenue_amount", "recognised_revenue", "amount"]:
            if cand in tmp.columns:
                value_col = cand
                break

        if value_col is None:
            rev_series = pd.Series(0.0, index=month_index)
        else:
            g = (
                tmp.groupby("month_date", as_index=True)[value_col]
                .sum()
            )
            # reindex so we have all months
            rev_series = g.reindex(month_index).fillna(0.0)

    # ------------------------------
    # 2) AR / AP by month (cash view)
    # ------------------------------
    df_ar, df_ap = fetch_ar_ap_for_client(client_id)

    # ---- AR ----
    if not df_ar.empty:
        df_ar = df_ar.copy()
        df_ar["amount"] = pd.to_numeric(df_ar.get("amount"), errors="coerce").fillna(0.0)

        if "expected_date" in df_ar.columns:
            df_ar["expected_date"] = pd.to_datetime(df_ar["expected_date"], errors="coerce")
        else:
            if "due_date" in df_ar.columns:
                df_ar["expected_date"] = pd.to_datetime(df_ar["due_date"], errors="coerce")
            else:
                df_ar["expected_date"] = pd.NaT

        df_ar["bucket_month"] = (
            df_ar["expected_date"]
            .dt.to_period("M")
            .dt.to_timestamp()
        )

        ar_monthly = (
            df_ar.groupby("bucket_month", as_index=True)["amount"]
            .sum()
        )
        ar_series = ar_monthly.reindex(month_index).fillna(0.0)
    else:
        ar_series = pd.Series(0.0, index=month_index)

    # ---- AP ----
    if not df_ap.empty:
        df_ap = df_ap.copy()
        df_ap["amount"] = pd.to_numeric(df_ap.get("amount"), errors="coerce").fillna(0.0)

        pay_col = "expected_payment_date"
        if pay_col not in df_ap.columns:
            pay_col = "due_date" if "due_date" in df_ap.columns else None

        if pay_col is not None:
            df_ap[pay_col] = pd.to_datetime(df_ap[pay_col], errors="coerce")
            df_ap["bucket_month"] = (
                df_ap[pay_col]
                .dt.to_period("M")
                .dt.to_timestamp()
            )
        else:
            df_ap["bucket_month"] = pd.NaT

        ap_monthly = (
            df_ap.groupby("bucket_month", as_index=True)["amount"]
            .sum()
        )
        ap_series = ap_monthly.reindex(month_index).fillna(0.0)
    else:
        ap_series = pd.Series(0.0, index=month_index)

    # For P&L view, treat vendor_pnl = AP amount in that month
    vendor_pnl_series = ap_series.copy()

    # ------------------------------
    # 3) Payroll by month
    # ------------------------------
    payroll_by_month = compute_payroll_by_month(client_id, month_index)
    payroll_series = pd.Series(
        [float(payroll_by_month.get(m, 0.0)) for m in month_index],
        index=month_index,
    )

    # ------------------------------
    # 4) Opening cash and monthly loop
    # ------------------------------
    # Opening cash from KPI for base month
    opening_cash = None
    kpi_row = fetch_kpis_for_client_month(client_id, base_month)
    if kpi_row and "cash_balance" in kpi_row:
        try:
            opening_cash = float(kpi_row["cash_balance"])
        except Exception:
            opening_cash = None

    if opening_cash is None:
        opening_cash = 120_000.0  # sensible fallback

    current_cash = opening_cash
    rows = []

    for m in month_index:
        revenue_pnl = float(rev_series.get(m, 0.0))
        payroll_pnl = float(payroll_series.get(m, 0.0))
        vendor_pnl = float(vendor_pnl_series.get(m, 0.0))
        total_cost_pnl = payroll_pnl + vendor_pnl
        net_burn_pnl = total_cost_pnl - revenue_pnl  # +ve = burning cash

        cash_in_ar = float(ar_series.get(m, 0.0))
        cash_out_ap = float(ap_series.get(m, 0.0))
        cash_out_payroll = payroll_pnl  # assume paid in same month for now

        operating_cf = cash_in_ar - cash_out_ap - cash_out_payroll
        investing_cf = 0.0
        financing_cf = 0.0

        closing_cash = current_cash + operating_cf + investing_cf + financing_cf
        danger = closing_cash <= 0

        rows.append(
            {
                "client_id": str(client_id),
                "month_date": m,
                "revenue_pnl": round(revenue_pnl, 2),
                "payroll_pnl": round(payroll_pnl, 2),
                "vendor_pnl": round(vendor_pnl, 2),
                "total_cost_pnl": round(total_cost_pnl, 2),
                "net_burn_pnl": round(net_burn_pnl, 2),
                "cash_in_ar": round(cash_in_ar, 2),
                "cash_out_ap": round(cash_out_ap, 2),
                "cash_out_payroll": round(cash_out_payroll, 2),
                "operating_cf": round(operating_cf, 2),
                "investing_cf": round(investing_cf, 2),
                "financing_cf": round(financing_cf, 2),
                "closing_cash": round(closing_cash, 2),
                "cash_danger_flag": danger,
            }
        )

        current_cash = closing_cash

    out_df = pd.DataFrame(rows)
    out_df["month_date"] = pd.to_datetime(out_df["month_date"], errors="coerce")
    out_df = out_df[out_df["month_date"].notna()]
    out_df = out_df.sort_values("month_date")
    out_df["month_label"] = out_df["month_date"].dt.strftime("%b %Y")

    return out_df


def fetch_opex_for_client(client_id: str) -> pd.DataFrame:
    """
    Load operating_opex rows for a client.
    One row = one month + category + amount (cash out).
    """
    if not client_id:
        return pd.DataFrame()

    try:
        resp = (
            supabase.table("operating_opex")
            .select("*")
            .eq("client_id", str(client_id))
            .execute()
        )
    except Exception as e:
        print("Error fetching operating_opex:", e)
        return pd.DataFrame()

    data = getattr(resp, "data", None) if hasattr(resp, "data") else resp.get("data")
    if not data:
        return pd.DataFrame()

    df = pd.DataFrame(data)
    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")

    if "amount" in df.columns:
        df["amount"] = pd.to_numeric(df["amount"], errors="coerce").fillna(0.0)

    return df



from datetime import date
import pandas as pd

def build_cashflow_recon_for_month(
    client_id,
    month_ref: date,
    opening_cash_hint: float | None = None,
) -> pd.DataFrame:
    if client_id is None or month_ref is None:
        return pd.DataFrame()

    # Normalise reference month to first-of-month Timestamp
    month_ts = pd.to_datetime(month_ref).to_period("M").to_timestamp()
    as_of_date = month_ts.date()

    # ------------------------
    # 1) Settings + AR/AP data
    # ------------------------
    settings = get_client_settings(client_id)
    ar_days = int(settings.get("ar_default_days", 30))
    ap_days = int(settings.get("ap_default_days", 30))

    df_ar, df_ap = fetch_ar_ap_for_client(client_id)

    # ----- AR -----
    if df_ar is not None and not df_ar.empty:
        df_ar = df_ar.copy()
        df_ar = add_ar_aging(df_ar, as_of=as_of_date, ar_default_days=ar_days)
        df_ar["expected_date"] = pd.to_datetime(df_ar.get("expected_date"), errors="coerce")
        df_ar["amount"] = pd.to_numeric(df_ar.get("amount"), errors="coerce").fillna(0.0)
    else:
        df_ar = pd.DataFrame(columns=["id", "expected_date", "amount"])

    # ----- AP -----
    if df_ap is not None and not df_ap.empty:
        df_ap = df_ap.copy()
        df_ap = add_ap_aging(df_ap, as_of=as_of_date, ap_default_days=ap_days)

        if "pay_expected_date" in df_ap.columns:
            pay_col = "pay_expected_date"
        elif "expected_payment_date" in df_ap.columns:
            pay_col = "expected_payment_date"
        else:
            pay_col = "due_date"

        df_ap[pay_col] = pd.to_datetime(df_ap.get(pay_col), errors="coerce")
        df_ap["amount"] = pd.to_numeric(df_ap.get("amount"), errors="coerce").fillna(0.0)
    else:
        df_ap = pd.DataFrame(columns=["id", "due_date", "amount"])
        pay_col = "due_date"

    # ---------- AR/AP payments map (SAME as engine) ----------
    month_index = pd.date_range(start=month_ts, periods=1, freq="MS")
    df_pay = fetch_ar_ap_payments_for_client(client_id)
    ar_cash_map, ap_cash_map = build_monthly_ar_ap_cash_maps(
        df_ar=df_ar,
        df_ap=df_ap,
        df_payments=df_pay,
        month_index=month_index,
    )

    cash_in_ar = float(ar_cash_map.get(month_ts, 0.0))
    cash_out_ap = float(ap_cash_map.get(month_ts, 0.0))

    # ------------------------
    # 2) Payroll (same engine logic)
    # ------------------------
    payroll_by_month = compute_payroll_by_month(client_id, month_index)
    payroll_cash = float(payroll_by_month.get(month_ts, 0.0))

    # ------------------------
    # 3) Opex (operating_opex)
    # ------------------------
    df_opex = fetch_opex_for_client(client_id)
    if df_opex is not None and not df_opex.empty:
        df_opex = df_opex.copy()
        df_opex["month_bucket"] = (
            pd.to_datetime(df_opex["month_date"], errors="coerce")
            .dt.to_period("M")
            .dt.to_timestamp()
        )
        mask_ox = df_opex["month_bucket"] == month_ts
        opex_cash = float(df_opex.loc[mask_ox, "amount"].sum())
    else:
        opex_cash = 0.0

    # ------------------------
    # 4) Other operating income (use same map as engine)
    # ------------------------
    try:
        df_oinc = fetch_operating_other_income_for_client(client_id)
    except Exception as e:
        print("Error fetching operating_other_income in recon:", e)
        df_oinc = pd.DataFrame()

    if df_oinc is not None and not df_oinc.empty:
        df_oinc = df_oinc.copy()

        if "month_bucket" not in df_oinc.columns:
            df_oinc["month_bucket"] = (
                pd.to_datetime(df_oinc["month_date"], errors="coerce")
                .dt.to_period("M")
                .dt.to_timestamp()
            )
        df_oinc["cash_in"] = pd.to_numeric(df_oinc.get("cash_in"), errors="coerce").fillna(0.0)

        mask_oi = df_oinc["month_bucket"] == month_ts
        other_income_cash = float(df_oinc.loc[mask_oi, "cash_in"].sum())
    else:
        other_income_cash = 0.0

    # ------------------------
    # 5) Investing & Financing (same grouping as engine)
    # ------------------------
    df_inv = fetch_investing_flows_for_client(client_id)
    if df_inv is not None and not df_inv.empty:
        df_inv = df_inv.copy()
        df_inv["month_bucket"] = (
            pd.to_datetime(df_inv["month_date"], errors="coerce")
            .dt.to_period("M")
            .dt.to_timestamp()
        )
        inv_agg = (
            df_inv.groupby("month_bucket", as_index=False)["amount"]
            .sum()
            .rename(columns={"month_bucket": "month_date"})
        )
        inv_by_month = {
            row["month_date"]: float(row["amount"] or 0.0)
            for _, row in inv_agg.iterrows()
        }
    else:
        inv_by_month = {}

    investing_cf_calc = float(inv_by_month.get(month_ts, 0.0))

    df_fin = fetch_financing_flows_for_client(client_id)
    if df_fin is not None and not df_fin.empty:
        df_fin = df_fin.copy()
        df_fin["month_bucket"] = (
            pd.to_datetime(df_fin["month_date"], errors="coerce")
            .dt.to_period("M")
            .dt.to_timestamp()
        )
        fin_agg = (
            df_fin.groupby("month_bucket", as_index=False)["amount"]
            .sum()
            .rename(columns={"month_bucket": "month_date"})
        )
        fin_by_month = {
            row["month_date"]: float(row["amount"] or 0.0)
            for _, row in fin_agg.iterrows()
        }
    else:
        fin_by_month = {}

    financing_cf_calc = float(fin_by_month.get(month_ts, 0.0))

    # ------------------------
    # 6) Engine row for this month
    # ------------------------
    engine_df = fetch_cashflow_summary_for_client(client_id)
    engine_df_row = None
    if engine_df is not None and not engine_df.empty:
        tmp = engine_df.copy()
        tmp["month_date"] = pd.to_datetime(tmp["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
        engine_df_row = tmp[tmp["month_date"] == month_ts]
        if engine_df_row.empty:
            engine_df_row = None

    opening_cash_engine = None
    operating_cf_engine = None
    investing_cf_engine = None
    financing_cf_engine = None
    free_cf_engine = None
    closing_cash_engine = None

    if engine_df_row is not None:
        row = engine_df_row.iloc(0)[0] if hasattr(engine_df_row, "iloc") else list(engine_df_row.to_dict("records"))[0]
        # safer, but you can stick with engine_df_row.iloc[0] in your code:
        row = engine_df_row.iloc[0]

        opening_cash_engine = float(row.get("opening_cash", 0.0) or 0.0)
        operating_cf_engine = float(row.get("operating_cf", 0.0) or 0.0)
        investing_cf_engine = float(row.get("investing_cf", 0.0) or 0.0)
        financing_cf_engine = float(row.get("financing_cf", 0.0) or 0.0)
        free_cf_engine = float(row.get("free_cash_flow", 0.0) or 0.0)
        closing_cash_engine = float(row.get("closing_cash", 0.0) or 0.0)

    # Opening cash for calc side
    if opening_cash_engine is not None:
        opening_cash_calc = opening_cash_engine
    else:
        if opening_cash_hint is not None:
            opening_cash_calc = float(opening_cash_hint)
        else:
            kpi_row = fetch_kpis_for_client_month(client_id, month_ts.date())
            opening_cash_calc = 0.0
            if kpi_row and "cash_balance" in kpi_row:
                try:
                    opening_cash_calc = float(kpi_row["cash_balance"])
                except Exception:
                    opening_cash_calc = 0.0

    # Our calculated CFs for this month (MUST match engine)
    operating_cf_calc = (
        cash_in_ar
        + other_income_cash
        - cash_out_ap
        - payroll_cash
        - opex_cash
    )
    free_cf_calc = operating_cf_calc + investing_cf_calc + financing_cf_calc
    closing_cash_calc = opening_cash_calc + free_cf_calc

    # ------------------------
    # 7) KPI cash balance
    # ------------------------
    kpi_row = fetch_kpis_for_client_month(client_id, month_ts.date())
    kpi_cash_balance = None
    if kpi_row and "cash_balance" in kpi_row:
        try:
            kpi_cash_balance = float(kpi_row["cash_balance"])
        except Exception:
            kpi_cash_balance = None

    # ------------------------
    # 8) Build reconciliation table
    # ------------------------
    rows = [
        ["AR cash-in", cash_in_ar, None, None, "Cash receipts from AR in this month (AR/AP payments map)"],
        ["AP cash-out", -cash_out_ap, None, None, "Cash payments to suppliers in this month (negative)"],
        ["Payroll cash-out", -payroll_cash, None, None, "Monthly payroll (negative)"],
        ["Opex cash-out", -opex_cash, None, None, "Monthly operating opex (negative)"],
        ["Other operating income", other_income_cash, None, None, "Other op. income cash in this month"],
        [
            "Operating CF (calc)",
            operating_cf_calc,
            operating_cf_engine,
            (operating_cf_engine - operating_cf_calc) if operating_cf_engine is not None else None,
            "Our calc vs engine.operating_cf",
        ],
        [
            "Investing CF (calc)",
            investing_cf_calc,
            investing_cf_engine,
            (investing_cf_engine - investing_cf_calc) if investing_cf_engine is not None else None,
            "Our calc vs engine.investing_cf",
        ],
        [
            "Financing CF (calc)",
            financing_cf_calc,
            financing_cf_engine,
            (financing_cf_engine - financing_cf_calc) if financing_cf_engine is not None else None,
            "Our calc vs engine.financing_cf",
        ],
        [
            "Free CF (calc)",
            free_cf_calc,
            free_cf_engine,
            (free_cf_engine - free_cf_calc) if free_cf_engine is not None else None,
            "Our calc vs engine.free_cash_flow",
        ],
        [
            "Opening cash (engine)",
            opening_cash_calc,
            opening_cash_engine,
            (opening_cash_engine - opening_cash_calc) if opening_cash_engine is not None else None,
            "Opening cash used for this month",
        ],
        [
            "Closing cash (calc)",
            closing_cash_calc,
            closing_cash_engine,
            (closing_cash_engine - closing_cash_calc) if closing_cash_engine is not None else None,
            "Our calc vs engine.closing_cash",
        ],
        [
            "Cash balance (KPI)",
            kpi_cash_balance,
            closing_cash_engine,
            (closing_cash_engine - kpi_cash_balance)
            if (kpi_cash_balance is not None and closing_cash_engine is not None)
            else None,
            "KPI cash_balance vs engine.closing_cash",
        ],
    ]

    recon_df = pd.DataFrame(
        rows,
        columns=["Metric", "Calculated", "Engine/DB", "Difference", "Notes"],
    )
    return recon_df


#.........Part payement to client...........................................

@st.cache_data(ttl=60)
def fetch_ar_ap_payments_for_client(client_id) -> pd.DataFrame:
    """
    Load all AR/AP payment events for this client.

    Expected table: public.ar_ap_payments
      - id           bigserial
      - client_id    uuid
      - ar_ap_id     <same type as ar_ap_tracker.id>
      - payment_date date
      - amount       numeric
      - status       text (optional, e.g. 'posted', 'reversed')

    Returns a DataFrame with parsed dates and numeric amounts.
    """
    if not client_id:
        return pd.DataFrame()

    try:
        resp = (
            supabase.table("ar_ap_payments")
            .select("*")
            .eq("client_id", str(client_id))
            .execute()
        )
    except Exception as e:
        print("Error fetching ar_ap_payments:", e)
        return pd.DataFrame()

    data = getattr(resp, "data", None) if hasattr(resp, "data") else resp.get("data", [])
    if not data:
        return pd.DataFrame()

    df = pd.DataFrame(data)

    if "payment_date" in df.columns:
        df["payment_date"] = pd.to_datetime(df["payment_date"], errors="coerce")

    if "amount" in df.columns:
        df["amount"] = pd.to_numeric(df["amount"], errors="coerce").fillna(0.0)

    # Optional: drop reversed / cancelled payments if you introduce such statuses later
    if "status" in df.columns:
        df["status"] = df["status"].astype(str).str.strip().str.lower()
        df = df[~df["status"].isin(["reversed", "cancelled", "canceled"])]

    return df



def build_monthly_ar_ap_cash_maps(
    df_ar: pd.DataFrame | None,
    df_ap: pd.DataFrame | None,
    df_payments: pd.DataFrame | None,
    month_index: pd.DatetimeIndex,
) -> tuple[dict[pd.Timestamp, float], dict[pd.Timestamp, float]]:
    """
    Build month -> cash maps for AR (cash-in) and AP (cash-out), using the
    *actual* amounts paid from the AR/AP tracker table.

    Rules (current schema):
      - Cash is taken from `amount_paid` (not `amount`).
      - Month bucket is based on the invoice `due_date` month
        (this matches your manual July'25 recon).
      - Invoices with status 'cancelled' / 'disputed' contribute 0
        *unless* there is some non-zero amount_paid on them.
      - Partial payments are naturally included because we sum `amount_paid`.
      - Any future extension to a separate payments table (with payment_date)
        can be wired in here later via df_payments; for now it is ignored.

    Returns:
      (ar_cash_map, ap_cash_map) where each is:
          dict[month_start_timestamp -> total_cash_for_that_month]
    """
    # Normalised month keys (first day of month)
    month_list = [
        pd.to_datetime(m).to_period("M").to_timestamp()
        for m in month_index
    ]
    ar_cash_map: dict[pd.Timestamp, float] = {m: 0.0 for m in month_list}
    ap_cash_map: dict[pd.Timestamp, float] = {m: 0.0 for m in month_list}

    def _from_invoices(df: pd.DataFrame | None) -> dict[pd.Timestamp, float]:
        if df is None or df.empty:
            return {}

        d = df.copy()

        # ---- Status handling: drop cancelled/disputed only if nothing paid ----
        if "status" in d.columns:
            d["status_norm"] = (
                d["status"]
                .astype(str)
                .str.strip()
                .str.lower()
            )
            amt_paid = pd.to_numeric(
                d.get("amount_paid", 0.0),
                errors="coerce"
            ).fillna(0.0)

            bad_statuses = {"cancelled", "canceled", "disputed"}
            mask_bad = d["status_norm"].isin(bad_statuses) & (amt_paid == 0.0)
            d = d[~mask_bad]

        # ---- Choose the date column used to bucket cash by month ----
        # We deliberately prefer `due_date`, because your July'25 recon
        # is based on "due date in July", not expected_payment_date.
        date_col = None
        for cand in [
            "due_date",
            "payment_date",          # future extension
            "expected_payment_date",
            "pay_expected_date",
            "expected_date",
            "issued_date",
        ]:
            if cand in d.columns:
                date_col = cand
                break

        if date_col is None:
            return {}

        # Your dates are in day-first format (e.g. 05-07-2025),
        # so we set dayfirst=True to avoid month/day swaps.
        d[date_col] = pd.to_datetime(
            d[date_col],
            errors="coerce",
            dayfirst=True,
        )
        d = d[d[date_col].notna()]
        if d.empty:
            return {}

        d["month_bucket"] = (
            d[date_col]
            .dt.to_period("M")
            .dt.to_timestamp()
        )

        # Use amount_paid if present; fall back to amount if not.
        amt_col = "amount_paid" if "amount_paid" in d.columns else "amount"
        d[amt_col] = pd.to_numeric(
            d.get(amt_col),
            errors="coerce",
        ).fillna(0.0)

        grp = (
            d.groupby("month_bucket", as_index=False)[amt_col]
            .sum()
        )

        out_map: dict[pd.Timestamp, float] = {m: 0.0 for m in month_list}
        for _, row in grp.iterrows():
            mb = row["month_bucket"]
            if mb in out_map:
                out_map[mb] += float(row[amt_col])

        return out_map

    # Build maps from the AR/AP invoice tables
    ar_map = _from_invoices(df_ar)
    ap_map = _from_invoices(df_ap)

    # Fill the final maps aligned to the month_index
    for m in month_list:
        if m in ar_map:
            ar_cash_map[m] = ar_map[m]
        if m in ap_map:
            ap_cash_map[m] = ap_map[m]

    # NOTE: df_payments is ignored for now. When you introduce a separate
    # ar_ap_payments table with (ar_ap_id, payment_date, amount, direction),
    # you can extend this function to OR in those entries as well.
    return ar_cash_map, ap_cash_map


def debug_cashflow_engine(
    client_id,
    base_month: date,
    n_months: int = 3,
    opening_cash_hint: float | None = None,
) -> pd.DataFrame:
    if client_id is None or base_month is None:
        return pd.DataFrame()

    base_month_ts = pd.to_datetime(base_month).replace(day=1)
    as_of_date = base_month_ts.date()
    month_index = pd.date_range(start=base_month_ts, periods=n_months, freq="MS")

    # ---------- Settings ----------
    settings = get_client_settings(client_id)
    ar_days = int(settings.get("ar_default_days", 30))
    ap_days = int(settings.get("ap_default_days", 30))

    # ---------- AR / AP ----------
    df_ar, df_ap = fetch_ar_ap_for_client(client_id)

    if df_ar is not None and not df_ar.empty:
        df_ar = df_ar.copy()
        df_ar = add_ar_aging(df_ar, as_of=as_of_date, ar_default_days=ar_days)
        df_ar["expected_date"] = pd.to_datetime(df_ar.get("expected_date"), errors="coerce")
        df_ar["amount"] = pd.to_numeric(df_ar.get("amount"), errors="coerce").fillna(0.0)
    else:
        df_ar = pd.DataFrame(columns=["id", "expected_date", "amount"])

    if df_ap is not None and not df_ap.empty:
        df_ap = df_ap.copy()
        df_ap = add_ap_aging(df_ap, as_of=as_of_date, ap_default_days=ap_days)
        if "pay_expected_date" in df_ap.columns:
            pay_col = "pay_expected_date"
        elif "expected_payment_date" in df_ap.columns:
            pay_col = "expected_payment_date"
        else:
            pay_col = "due_date"
        df_ap[pay_col] = pd.to_datetime(df_ap.get(pay_col), errors="coerce")
        df_ap["amount"] = pd.to_numeric(df_ap.get("amount"), errors="coerce").fillna(0.0)
    else:
        df_ap = pd.DataFrame(columns=["id", "due_date", "amount"])

    # ---------- Payments map ----------
    df_pay = fetch_ar_ap_payments_for_client(client_id)
    ar_cash_map, ap_cash_map = build_monthly_ar_ap_cash_maps(
        df_ar=df_ar,
        df_ap=df_ap,
        df_payments=df_pay,
        month_index=month_index,
    )

    # ---------- Opex ----------
    try:
        df_opex = fetch_opex_for_client(client_id)
        if df_opex is not None and not df_opex.empty:
            df_opex = df_opex.copy()
            df_opex["month_bucket"] = (
                pd.to_datetime(df_opex["month_date"], errors="coerce")
                .dt.to_period("M")
                .dt.to_timestamp()
            )
        else:
            df_opex = pd.DataFrame(columns=["month_bucket", "amount"])
    except Exception as e:
        print("Error fetching operating_opex in debug:", e)
        df_opex = pd.DataFrame(columns=["month_bucket", "amount"])

    # ---------- Other operating income ----------
    try:
        df_oinc = fetch_operating_other_income_for_client(client_id)
    except Exception as e:
        print("Error fetching operating_other_income in debug:", e)
        df_oinc = pd.DataFrame()

    if df_oinc is None or df_oinc.empty:
        df_oinc = pd.DataFrame(columns=["month_bucket", "cash_in"])
    else:
        df_oinc = df_oinc.copy()
        if "month_bucket" not in df_oinc.columns:
            df_oinc["month_bucket"] = (
                pd.to_datetime(df_oinc["month_date"], errors="coerce")
                .dt.to_period("M")
                .dt.to_timestamp()
            )
        df_oinc["cash_in"] = pd.to_numeric(df_oinc.get("cash_in"), errors="coerce").fillna(0.0)

    # ---------- Investing / Financing maps ----------
    df_inv = fetch_investing_flows_for_client(client_id)
    df_fin = fetch_financing_flows_for_client(client_id)

    inv_by_month: dict[pd.Timestamp, float] = {}
    if df_inv is not None and not df_inv.empty:
        # Keep a copy WITHOUT filtering for fallback
        orig_df_inv = df_inv.copy()
        df_inv = df_inv.copy()

        # Optional versioning by created_at – but only if it actually works
        if "created_at" in df_inv.columns:
            df_inv["created_at"] = pd.to_datetime(
                df_inv["created_at"], errors="coerce"
            )
            # Drop timezone if present
            try:
                df_inv["created_at"] = df_inv["created_at"].dt.tz_convert(None)
            except TypeError:
                pass  # already naive

            before_len = len(df_inv)
            df_inv = df_inv[df_inv["created_at"] <= as_of_cutoff]
            after_len = len(df_inv)
            print(
                f"[DEBUG] investing_flows created_at filter: "
                f"{before_len} -> {after_len} rows (as_of_cutoff={as_of_cutoff.date()})"
            )

            # 🔁 Fallback: if everything got filtered out, ignore created_at for now
            if after_len == 0:
                print(
                    "[WARN] investing_flows all filtered out by created_at; "
                    "falling back to using all investing rows without versioning."
                )
                df_inv = orig_df_inv

        # Now aggregate by month_date
        df_inv["month_bucket"] = (
            pd.to_datetime(df_inv["month_date"], errors="coerce")
            .dt.to_period("M")
            .dt.to_timestamp()
        )
        inv_agg = (
            df_inv.groupby("month_bucket", as_index=False)["amount"]
            .sum()
            .rename(columns={"month_bucket": "month_date"})
        )
        inv_by_month = {
            row["month_date"]: float(row["amount"] or 0.0)
            for _, row in inv_agg.iterrows()
        }


    fin_by_month: dict[pd.Timestamp, float] = {}
    if df_fin is not None and not df_fin.empty:
        df_fin = df_fin.copy()
        df_fin["month_bucket"] = (
            pd.to_datetime(df_fin["month_date"], errors="coerce")
            .dt.to_period("M")
            .dt.to_timestamp()
        )
        fin_agg = (
            df_fin.groupby("month_bucket", as_index=False)["amount"]
            .sum()
            .rename(columns={"month_bucket": "month_date"})
        )
        fin_by_month = {
            row["month_date"]: float(row["amount"] or 0.0)
            for _, row in fin_agg.iterrows()
        }

    # ---------- Payroll ----------
    payroll_by_month = compute_payroll_by_month(client_id, month_index)

    # ---------- Engine output ----------
    engine_df = fetch_cashflow_summary_for_client(client_id)
    if engine_df is None or engine_df.empty:
        st.warning("No rows in cashflow_summary yet – rebuild cashflow first.")
        return pd.DataFrame()

    engine_df = engine_df.copy()
    engine_df["month_date"] = pd.to_datetime(engine_df["month_date"], errors="coerce")
    engine_df["month_date"] = engine_df["month_date"].dt.to_period("M").dt.to_timestamp()

    # ---------- Opening cash ----------
    if opening_cash_hint is not None:
        opening_cash_0 = float(opening_cash_hint)
    else:
        kpi_row = fetch_kpis_for_client_month(client_id, as_of_date)
        opening_cash_0 = 0.0
        if kpi_row and "cash_balance" in kpi_row:
            try:
                opening_cash_0 = float(kpi_row["cash_balance"])
            except Exception:
                opening_cash_0 = 0.0

    rows = []
    current_closing = opening_cash_0

    for m in month_index:
        m_key = pd.to_datetime(m).to_period("M").to_timestamp()
        opening_cash = current_closing

        cash_in_ar = float(ar_cash_map.get(m_key, 0.0))
        cash_out_ap = float(ap_cash_map.get(m_key, 0.0))
        payroll_cash = float(payroll_by_month.get(m_key, 0.0))

        if df_opex is not None and not df_opex.empty:
            mask_ox = df_opex["month_bucket"] == m_key
            opex_cash = float(df_opex.loc[mask_ox, "amount"].sum())
        else:
            opex_cash = 0.0

        if df_oinc is not None and not df_oinc.empty:
            mask_oi = df_oinc["month_bucket"] == m_key
            other_income_cash = float(df_oinc.loc[mask_oi, "cash_in"].sum())
        else:
            other_income_cash = 0.0

        investing_cf_calc = float(inv_by_month.get(m_key, 0.0))
        financing_cf_calc = float(fin_by_month.get(m_key, 0.0))

        operating_cf_calc = (
            cash_in_ar
            + other_income_cash
            - cash_out_ap
            - payroll_cash
            - opex_cash
        )

        row_eng = engine_df[engine_df["month_date"] == m_key]
        if not row_eng.empty:
            operating_cf_engine = float(row_eng.get("operating_cf", 0.0).iloc[0] or 0.0)
            investing_cf_engine = float(row_eng.get("investing_cf", 0.0).iloc[0] or 0.0)
            financing_cf_engine = float(row_eng.get("financing_cf", 0.0).iloc[0] or 0.0)
            free_cf_engine = float(row_eng.get("free_cash_flow", 0.0).iloc[0] or 0.0)
            closing_cash_engine = float(row_eng.get("closing_cash", 0.0).iloc[0] or 0.0)
        else:
            operating_cf_engine = 0.0
            investing_cf_engine = 0.0
            financing_cf_engine = 0.0
            free_cf_engine = 0.0
            closing_cash_engine = float("nan")

        free_cf_calc = operating_cf_calc + investing_cf_calc + financing_cf_calc
        closing_cash_calc = opening_cash + free_cf_calc

        rows.append(
            {
                "month_date": m_key,
                "opening_cash_calc": round(opening_cash, 2),
                "cash_in_ar": round(cash_in_ar, 2),
                "cash_out_ap": round(cash_out_ap, 2),
                "payroll_cash": round(payroll_cash, 2),
                "opex_cash": round(opex_cash, 2),
                "other_income_cash": round(other_income_cash, 2),
                "operating_cf_calc": round(operating_cf_calc, 2),
                "operating_cf_engine": round(operating_cf_engine, 2),
                "investing_cf_calc": round(investing_cf_calc, 2),
                "investing_cf_engine": round(investing_cf_engine, 2),
                "financing_cf_calc": round(financing_cf_calc, 2),
                "financing_cf_engine": round(financing_cf_engine, 2),
                "free_cf_calc": round(free_cf_calc, 2),
                "free_cf_engine": round(free_cf_engine, 2),
                "closing_cash_calc": round(closing_cash_calc, 2),
                "closing_cash_engine": (
                    round(closing_cash_engine, 2)
                    if pd.notna(closing_cash_engine)
                    else None
                ),
            }
        )

        current_closing = closing_cash_calc

    debug_df = pd.DataFrame(rows)
    debug_df["diff_operating_cf"] = (
        debug_df["operating_cf_engine"] - debug_df["operating_cf_calc"]
    )
    debug_df["diff_investing_cf"] = (
        debug_df["investing_cf_engine"] - debug_df["investing_cf_calc"]
    )
    debug_df["diff_financing_cf"] = (
        debug_df["financing_cf_engine"] - debug_df["financing_cf_calc"]
    )
    debug_df["diff_free_cf"] = (
        debug_df["free_cf_engine"] - debug_df["free_cf_calc"]
    )
    debug_df["diff_closing"] = (
        debug_df["closing_cash_engine"] - debug_df["closing_cash_calc"]
    )

    return debug_df



def fetch_financing_flows_for_client(client_id) -> pd.DataFrame:
    if client_id is None:
        return pd.DataFrame()

    try:
        resp = (
            supabase.table("financing_flows")
            .select("*")
            .eq("client_id", str(client_id))
            .order("month_date", desc=False)
            .execute()
        )

        data = getattr(resp, "data", None)
        if data is None and isinstance(resp, dict):
            data = resp.get("data", [])

        if not data:
            return pd.DataFrame()

        df = pd.DataFrame(data)
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
        df["amount"] = pd.to_numeric(df["amount"], errors="coerce").fillna(0.0)
        return df

    except Exception as e:
        print("Error fetching financing_flows:", e)
        return pd.DataFrame()


def fetch_cash_curve_for_client(client_id, base_month: date | None, n_months: int = 12) -> pd.DataFrame:
    """
    Try to build a month-by-month cash curve for this client.

    Priority:
      1) cashflow_summary (if exists) with month_date + closing_cash
      2) baseline_monthly (if exists) with month_date + closing_cash
      3) fallback: simple synthetic curve for demo

    Returns DataFrame with columns:
      - month_date (datetime)
      - month_label (e.g. "Dec 2025")
      - closing_cash (float)
    """
    if client_id is None:
        return _build_synthetic_cash_curve(base_month, n_months)

    # Helper to filter & clean a table if present
    def _try_table(table_name: str, amount_col: str):
        try:
            res = (
                supabase
                .table(table_name)
                .select("*")
                .eq("client_id", str(client_id))
                .execute()
            )
            rows = res.data or []
        except Exception:
            rows = []

        if not rows:
            return pd.DataFrame()

        df = pd.DataFrame(rows)

        # Identify month column
        month_col = None
        for cand in ["month_date", "month", "period"]:
            if cand in df.columns:
                month_col = cand
                break
        if month_col is None or amount_col not in df.columns:
            return pd.DataFrame()

        df["month_date"] = pd.to_datetime(df[month_col], errors="coerce")
        df = df[df["month_date"].notna()]

        if base_month is not None:
            start = pd.to_datetime(base_month)
            end = start + pd.DateOffset(months=n_months)
            df = df[(df["month_date"] >= start) & (df["month_date"] < end)]

        if df.empty:
            return pd.DataFrame()

        out = df[["month_date", amount_col]].copy()
        out = out.rename(columns={amount_col: "closing_cash"})
        out = out.sort_values("month_date")
        out["month_label"] = out["month_date"].dt.strftime("%b %Y")
        return out

    # 1) Try cashflow_summary
    df = _try_table("cashflow_summary", "closing_cash")
    if not df.empty:
        return df

    # 2) Try baseline_monthly
    df = _try_table("baseline_monthly", "closing_cash")
    if not df.empty:
        return df

    # 3) Fallback synthetic curve
    return _build_synthetic_cash_curve(base_month, n_months)


def _build_synthetic_cash_curve(base_month: date | None, n_months: int = 12) -> pd.DataFrame:
    """
    Simple placeholder cash curve if DB has no data yet.
    Starts at 120k, burns 10k per month.
    """
    if base_month is None:
        base_month = date.today().replace(day=1)

    month_index = pd.date_range(
        start=pd.to_datetime(base_month),
        periods=n_months,
        freq="MS",
    )
    starting_cash = 120_000
    burn_per_month = 10_000

    cash_values = []
    for i in range(n_months):
        cash_values.append(starting_cash - burn_per_month * i)

    df = pd.DataFrame(
        {
            "month_date": month_index,
            "closing_cash": cash_values,
        }
    )
    df["month_label"] = df["month_date"].dt.strftime("%b %Y")
    return df



def compute_runway_and_effective_burn_from_df(
    engine_df: pd.DataFrame,
    month_ref: date | datetime,
) -> tuple[float | None, float | None]:
    """
    Forward-looking runway + effective burn from an existing cashflow_summary DataFrame.

    - Runway:
        first month where closing_cash <= 0, counted as:
        0 = this month, 1 = next month, etc.
        If no cash-out in the horizon -> runway_months stays None.

    - Effective burn:
        3-month average of negative cashflow (operating_cf/free_cash_flow/etc),
        returned as a positive number (e.g. 25k burn per month).
    """
    if engine_df is None or engine_df.empty:
        print("[DEBUG] compute_runway_and_effective_burn_from_df -> empty/no engine_df")
        return None, None

    df = engine_df.copy()
    df["month_date"] = pd.to_datetime(df.get("month_date"), errors="coerce")
    df = df[df["month_date"].notna()].sort_values("month_date").reset_index(drop=True)
    if df.empty:
        print("[DEBUG] compute_runway_and_effective_burn_from_df -> no valid month_date rows")
        return None, None

    # Normalise reference to first of month (Timestamp)
    start_month = pd.to_datetime(month_ref).to_period("M").to_timestamp()

    # --- Find base row (focus month) ---
    mask_exact = df["month_date"] == start_month
    if mask_exact.any():
        base_idx = int(mask_exact.idxmax())
    else:
        mask_future = df["month_date"] > start_month
        base_idx = int(mask_future.idxmax()) if mask_future.any() else len(df) - 1

    df_future = df.loc[base_idx:].reset_index(drop=True)

    # ======================
    # 1) Runway from closing_cash
    # ======================
    runway_months: float | None = None

    if "closing_cash" in df_future.columns:
        cash_series = pd.to_numeric(df_future["closing_cash"], errors="coerce")
        neg_mask = cash_series.notna() & (cash_series <= 0)

        if neg_mask.any():
            first_neg_idx = int(neg_mask.idxmax())
            runway_months = float(first_neg_idx)

    # ======================
    # 2) Effective burn (3-month forward avg)
    # ======================
    effective_burn: float | None = None

    burn_source_cols = [
        "operating_cf",      # from your sample engine logs
        "free_cash_flow",
        "net_cash_flow",
        "net_cashflow",
        "total_cash_flow",
    ]
    burn_col = next((c for c in burn_source_cols if c in df.columns), None)

    if burn_col is not None:
        window = df.loc[base_idx: base_idx + 2].copy()
        vals = pd.to_numeric(window[burn_col], errors="coerce")
        avg_cf = vals.mean()
        if pd.notna(avg_cf) and avg_cf < 0:
            effective_burn = float(-avg_cf)  # positive number

    # Fallback: derive burn from slope of closing_cash
    if effective_burn is None and "closing_cash" in df.columns:
        window = df.loc[base_idx: base_idx + 2].copy()
        cc = pd.to_numeric(window["closing_cash"], errors="coerce")
        delta = cc.diff().mean()   # negative = cash falling
        if pd.notna(delta) and delta < 0:
            effective_burn = float(-delta)

    # ======================
    # 3) Fallback runway = cash_now / burn
    # ======================
    if runway_months is None and effective_burn is not None and effective_burn > 0:
        try:
            cash_now = pd.to_numeric(df_future.loc[0, "closing_cash"], errors="coerce")
            cash_now = float(cash_now) if pd.notna(cash_now) else 0.0
            runway_months = (cash_now / effective_burn) if cash_now > 0 else 0.0
        except Exception as e:
            print("[WARN] runway fallback error:", e)

    print(
        f"[DEBUG] compute_runway_and_effective_burn_from_df -> "
        f"runway={runway_months}, eff_burn={effective_burn}"
    )
    return runway_months, effective_burn



def compute_danger_month_from_cash(df: pd.DataFrame, amount_col: str = "closing_cash"):
    """
    Given a DF with month_date, month_label, and closing_cash, find the first
    month where cash is <= 0. Returns (month_label, month_date) or (None, None).
    """
    if df is None or df.empty or amount_col not in df.columns:
        return None, None

    df_sorted = df.sort_values("month_date")
    for _, row in df_sorted.iterrows():
        value = row[amount_col]
        try:
            if float(value) <= 0:
                return row.get("month_label"), row.get("month_date")
        except Exception:
            continue

    return None, None




def fetch_comments_for_client_page(client_id: str, page_name: str):
    try:
        resp = (
            supabase.table("comments")
            .select("*")
            .eq("client_id", str(client_id))
            .eq("page_name", page_name)
            .order("created_at", desc=True)
            .execute()
        )
        return resp.data or []
    except Exception as e:
        print("Error fetching comments:", e)
        return []

def create_comment_for_client_page(
    client_id: str,
    page_name: str,
    comment_text: str,
    author_name: str = "Team member",
):
    try:
        data = {
            "client_id": str(client_id),
            "page_name": page_name,
            "comment_text": comment_text,
            "author_name": author_name,
            "created_at": datetime.now(timezone.utc).isoformat(),
        }
        supabase.table("comments").insert(data).execute()
        return True
    except Exception as e:
        print("Error creating comment:", e)
        return False



# ---------- Sidebar: global controls ----------

st.sidebar.title("Matfina")

# Client selector (multi-tenant ready)
clients_data = load_clients()
client_names = [c["name"] for c in clients_data]

selected_client_name = st.sidebar.selectbox("Select business", client_names)

# Find the matching client_id
selected_client_id = None
for c in clients_data:
    if c["name"] == selected_client_name:
        selected_client_id = c["id"]
        break

# Month selector (fully data-driven from Supabase, no hard upper limit)
month_options = get_client_focus_month_options(selected_client_id)

if not month_options:
    st.sidebar.warning(
        "No Focus Month available for this client (no dated data yet)."
    )
    st.stop()

selected_month_label = st.sidebar.selectbox(
    "Focus month",
    month_options,
    index=len(month_options) - 1,  # default to latest month with data
)

selected_month_start = parse_month_label_to_date(selected_month_label)




def save_task(client_id, page_name: str, title: str, month_start: date, kpi_key: str = "general"):
    if not title or not title.strip():
        return False

    month_bucket = _month_start(month_start)

    data = {
        "client_id": str(client_id),
        "page_name": page_name,
        "kpi_key": kpi_key,  # IMPORTANT (add this column in tasks table)
        "month_date": month_bucket_utc_iso(month_start),
        "title": title.strip(),
        "status": "open",
        "priority": "normal",
        "owner_name": "Internal CFO",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }

    try:
        supabase.table("tasks").insert(data).execute()
        return True
    except Exception as e:
        print("Error inserting task:", e)
        return False




def update_task_status(task_id, status: str = "done"):
    try:
        supabase.table("tasks").update({
            "status": status,
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }).eq("id", int(task_id)).execute()
        return True
    except Exception as e:
        print("Error updating task:", e)
        return False



# ---------------- Currency helpers ----------------

DEFAULT_CURRENCY_CODE = "AUD"

CURRENCY_SYMBOLS = {
    "AUD": "A$",
    "USD": "$",
    "EUR": "€",
    "GBP": "£",
    "NZD": "NZ$",
    "CAD": "C$",
}
def format_money(value, currency_symbol: str = "$", none_placeholder: str = "—") -> str:
    """
    Nicely format a numeric value as money with the right currency symbol.
    If value is None or a placeholder, returns an em dash.
    """
    if value is None or value == "—":
        return none_placeholder

    try:
        num = float(value)
    except (TypeError, ValueError):
        return none_placeholder

    return f"{currency_symbol}{num:,.0f}"

def get_client_currency(client_id: str) -> tuple[str, str]:
    """
    Return (currency_code, currency_symbol) for a client.

    Source of truth:
    - client_settings.currency_code

    Fallback = AUD
    """
    if not client_id:
        return "AUD", "A$"

    try:
        settings = get_client_settings(client_id) or {}
        code = settings.get("currency_code", "AUD")
        code = str(code).upper()
    except Exception as e:
        print("[WARN] Currency lookup failed, defaulting to AUD:", e)
        code = "AUD"

    symbol_map = {
        "AUD": "A$",
        "USD": "$",
        "EUR": "€",
        "GBP": "£",
        "NZD": "NZ$",
        "CAD": "C$",
        "INR": "₹",
    }

    symbol = symbol_map.get(code, f"{code} ")

    print(f"[DEBUG] get_client_currency -> code={code}, symbol={symbol}")
    return code, symbol


def format_money_for_client(client_id: str | None, value) -> str:
    """
    Format a numeric value into a currency string for this client.
    Handles None / dash cases gracefully.
    """
    if value is None or value == "—":
        return "—"

    try:
        num = float(value)
    except Exception:
        return str(value)

    _, symbol = get_client_currency(client_id)
    return f"{symbol}{num:,.0f}"

import streamlit as st

@st.cache_data(show_spinner=False, ttl=300)
def get_client_settings(client_id: str) -> dict:
    cid = str(client_id)

    # ---- fetch row (safe if missing) ----
    try:
        resp = (
            supabase.table("client_settings")
            .select("*")
            .eq("client_id", cid)
            .maybe_single()   # ✅ IMPORTANT: no crash if row missing
            .execute()
        )
        row = resp.data or {}
    except Exception:
        row = {}

    def _f(key: str, default: float = 0.0) -> float:
        try:
            v = row.get(key, default)
            return float(default if v is None else v)
        except Exception:
            return float(default)

    def _i(key: str, default: int = 0) -> int:
        try:
            v = row.get(key, default)
            return int(default if v is None else v)
        except Exception:
            return int(default)

    def _s(key: str, default: str = "") -> str:
        v = row.get(key, default)
        return default if v is None else str(v)

    # currency_code from DB (default AUD)
    currency_code = (_s("currency_code", "AUD") or "AUD").upper()

    return {
        # ---- Currency ----
        "currency_code": currency_code,

        # ---- Opening cash / buffer ----
        "opening_cash_start": _f("opening_cash_start", 0.0),
        "min_cash_buffer": _f("min_cash_buffer", 0.0),   # ✅ THIS WAS MISSING

        # ---- Payroll cash timing ----
        "super_lag_months": _i("super_lag_months", 1),
        "payroll_tax_lag_months": _i("payroll_tax_lag_months", 1),

        # ---- Employer on-costs ----
        "payroll_tax_threshold_annual": _f("payroll_tax_threshold_annual", 1_300_000.0),
        "payroll_tax_rate_pct": _f("payroll_tax_rate_pct", 0.0),
        "workers_comp_rate_pct": _f("workers_comp_rate_pct", 0.0),

        # ---- Next hire assumptions ----
        "next_hire_salary_annual": _f("next_hire_salary_annual", 0.0),
        "next_hire_fte": _f("next_hire_fte", 1.0),
        "next_hire_super_rate_pct": _f("next_hire_super_rate_pct", 11.0),
        "next_hire_benefits_rate_pct": _f("next_hire_benefits_rate_pct", 0.0),
        "next_hire_onboarding_monthly": _f("next_hire_onboarding_monthly", 0.0),

        # ---- AR/AP defaults ----
        "ar_default_days": _i("ar_default_days", 30),
        "ap_default_days": _i("ap_default_days", 30),

        # ---- Risk thresholds ----
        "runway_min_months": _f("runway_min_months", 4.0),
        "overspend_warn_pct": _f("overspend_warn_pct", 0.0),
        "overspend_high_pct": _f("overspend_high_pct", 0.0),

        # ---- Revenue recognition ----
        "revenue_recognition_method": _s("revenue_recognition_method", "saas") or "saas",
    }




def save_alert_thresholds(client_id, runway_min_months: float,
                          overspend_warn_pct: float, overspend_high_pct: float):
    """
    Upsert alert thresholds for this client into client_settings.
    (Keeps existing AR/AP defaults as they are.)
    """
    if client_id is None:
        return False

    try:
        current = get_client_settings(client_id)
        payload = {
            "client_id": str(client_id),
            # keep existing AR/AP defaults
            "ar_default_days": int(current["ar_default_days"]),
            "ap_default_days": int(current["ap_default_days"]),
            # update thresholds
            "runway_min_months": float(runway_min_months),
            "overspend_warn_pct": float(overspend_warn_pct),
            "overspend_high_pct": float(overspend_high_pct),
            "updated_at": datetime.utcnow().isoformat(),
        }
        supabase.table("client_settings").upsert(payload, on_conflict="client_id").execute()
         # Clear cached settings so new thresholds are picked up
        st.cache_data.clear()
        return True
    except Exception:
        return False


def bulk_update_deal_methods_for_client(client_id, method_display: str) -> bool:
    """
    Overwrite revenue_pipeline.method for ALL deals for this client.
    method_display is the human label stored in revenue_pipeline.method
    (e.g., 'SaaS subscription', 'Milestone project', etc.)
    """
    if client_id is None:
        return False

    try:
        payload = {
            "method": str(method_display),
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }

        # Update all rows matching client_id
        supabase.table("revenue_pipeline").update(payload).eq("client_id", str(client_id)).execute()

        # Clear cache so pipeline + schedule refresh
        _clear_all_caches_safely()
        return True

    except Exception as e:
        print("bulk_update_deal_methods_for_client error:", e)
        return False

# =========================
# SAVE BUSINESS DEFAULT METHOD (CLEAR CACHE + OPTIONAL APPLY-TO-ALL)
# =========================
from datetime import datetime, timezone

def _clear_all_caches_safely():
    """Clear Streamlit caches without crashing if unavailable."""
    try:
        st.cache_data.clear()
    except Exception:
        pass
    try:
        st.cache_resource.clear()
    except Exception:
        pass

def save_revenue_method(client_id, method: str) -> bool:
    """
    Save the default revenue recognition method for this client.
    Method is one of: 'saas', 'milestone', 'poc', 'straight_line', 'usage', 'point_in_time'.

    IMPORTANT:
    - Clears cache so build_revenue_schedule_for_client recalculates immediately.
    """
    if client_id is None:
        return False

    valid_methods = {"saas", "milestone", "poc", "straight_line", "usage", "point_in_time"}
    if method not in valid_methods:
        method = "saas"

    try:
        current = get_client_settings(client_id) or {}

        payload = {
            "client_id": str(client_id),
            "ar_default_days": int(current.get("ar_default_days", 30)),
            "ap_default_days": int(current.get("ap_default_days", 30)),
            "runway_min_months": float(current.get("runway_min_months", 3)),
            "overspend_warn_pct": float(current.get("overspend_warn_pct", 10)),
            "overspend_high_pct": float(current.get("overspend_high_pct", 20)),
            "revenue_recognition_method": method,
            "updated_at": datetime.now(timezone.utc).isoformat(),  # ✅ safe for older python
        }

        supabase.table("client_settings").upsert(payload, on_conflict="client_id").execute()

        # ✅ critical: clear cache so schedules update
        _clear_all_caches_safely()
        return True

    except Exception as e:
        print("save_revenue_method error:", e)
        return False




def save_client_settings(client_id, ar_days=None, ap_days=None, **extra_fields):
    if client_id is None:
        return False

    payload = {"client_id": str(client_id), "updated_at": datetime.utcnow().isoformat()}

    if ar_days is not None:
        payload["ar_default_days"] = int(ar_days)
    if ap_days is not None:
        payload["ap_default_days"] = int(ap_days)

    if extra_fields:
        payload.update(extra_fields)

    try:
        supabase.table("client_settings").upsert(payload, on_conflict="client_id").execute()

        try:
            get_client_settings.clear()
        except Exception:
            pass

        return True
    except Exception as e:
        print("Error in save_client_settings:", e)
        return False

def render_currency_settings_section(client_id: str) -> None:
    st.subheader("💱 Currency settings")

    if not client_id:
        st.info("Select a business first.")
        return

    settings = get_client_settings(client_id) or {}

    # --- Source of truth: currency_code ---
    current_code = settings.get("currency_code", "AUD")
    current_code = str(current_code).upper()

    currency_options = ["AUD", "USD", "EUR", "GBP", "NZD", "CAD", "INR"]

    try:
        idx = currency_options.index(current_code)
    except ValueError:
        idx = 0

    selected_code = st.selectbox(
        "Base currency for this business",
        options=currency_options,
        index=idx,
        help="All KPIs, dashboards, and reports will use this currency (no FX conversion yet).",
        key="client_currency_select",
    )

    if st.button("Save currency", key="save_client_currency_btn"):
        resp = supabase.table("client_settings").upsert(
            {
                "client_id": str(client_id),
                "currency_code": selected_code,
            }
        ).execute()

        if not resp or getattr(resp, "error", None):
            st.error("Could not update currency. Please try again.")
        else:
            st.success(f"Currency set to {selected_code}.")
            try:
                get_client_settings.clear()
            except Exception:
                pass
            st.rerun()

from datetime import datetime, timezone

def fetch_task_replies(task_id: int):
    try:
        resp = (
            supabase.table("task_replies")
            .select("*")
            .eq("task_id", int(task_id))
            .order("created_at", desc=False)
            .execute()
        )
        return resp.data or []
    except Exception as e:
        print("Error fetching task replies:", e)
        return []

def add_task_reply(client_id: str, task_id: int, author_name: str, message: str) -> bool:
    if not message or not message.strip():
        return False
    try:
        supabase.table("task_replies").insert({
            "client_id": str(client_id),
            "task_id": int(task_id),
            "author_name": (author_name or "").strip() or None,
            "message": message.strip(),
            "created_at": datetime.now(timezone.utc).isoformat(),
        }).execute()
        return True
    except Exception as e:
        print("Error adding task reply:", e)
        return False

def upsert_board_commentary(client_id: str, month_start: date, kpi_key: str, task_id: int, final_comment: str, author_name: str = None) -> bool:
    if not final_comment or not final_comment.strip():
        return False
    try:
        supabase.table("board_commentary").upsert({
            "client_id": str(client_id),
            "month_date": month_start.isoformat(),
            "kpi_key": str(kpi_key),
            "task_id": int(task_id) if task_id is not None else None,
            "final_comment": final_comment.strip(),
            "author_name": (author_name or "").strip() or None,
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }).execute()
        return True
    except Exception as e:
        print("Error saving board commentary:", e)
        return False



def fetch_tasks_for_page(client_id, page_name: str, month_start: date):
    """
    Fetch ALL tasks (any status) for this client + page + month.
    """
    if client_id is None or month_start is None:
        return []

    try:
        res = (
            supabase
            .table("tasks")
            .select("*")
            .eq("client_id", str(client_id))
            .eq("page_name", page_name)
            .eq("month_date", month_start.isoformat())
            .order("created_at", desc=True)
            .execute()
        )
        return res.data or []
    except Exception:
        return []

def sort_alerts_by_severity(alerts: list[dict]) -> list[dict]:
    """
    Sort alerts by severity (critical > high > medium > low),
    then by created_at (newest first).
    """
    if not alerts:
        return []

    severity_order = {
        "critical": 3,
        "high": 2,
        "medium": 1,
        "low": 0,
    }

    def alert_key(a):
        sev = str(a.get("severity", "medium")).lower()
        score = severity_order.get(sev, 1)
        created = a.get("created_at") or ""
        return (-score, created)  # higher severity first, then newest

    return sorted(alerts, key=alert_key)

def ensure_overdue_ar_alert(client_id, as_of: date | None = None, min_days_overdue: int = 0):
    """
    Create/resolve an 'ar_overdue' alert for the month of `as_of`.
    Overdue = unpaid invoices with days_past_expected > min_days_overdue
    where expected_date is derived from:
      expected_date (if present) else due_date else issue_date + expected_payment_days / client default
    """
    if client_id is None:
        return
    if as_of is None:
        as_of = date.today()

    # Pull settings so alert aligns to client rules
    settings = get_client_settings(client_id) or {}
    ar_default_days = int(settings.get("ar_default_days", 30))

    # IMPORTANT: use your existing fetch. If this already returns all AR, great.
    # If it returns only overdue, still fine; we compute expected anyway.
    df_ar = fetch_ar_ap_for_client(client_id) if "fetch_ar_for_client" in globals() else fetch_overdue_ar_for_client(client_id, as_of=as_of)

    if df_ar is None or df_ar.empty:
        _resolve_ar_alerts_for_month(client_id, as_of)
        return

    df_ar = add_ar_aging(df_ar, as_of=as_of, ar_default_days=ar_default_days)

    # unpaid + overdue threshold
    df_ar["status"] = df_ar.get("status", "").astype(str).str.strip().str.lower()
    unpaid = ~df_ar["status"].isin(list(PAID_STATUSES))
    overdue = df_ar["days_past_expected"] > int(min_days_overdue)

    df_overdue = df_ar[unpaid & overdue].copy()
    if df_overdue.empty:
        _resolve_ar_alerts_for_month(client_id, as_of)
        return

    count_invoices = int(len(df_overdue))
    total_amount = float(pd.to_numeric(df_overdue.get("amount", 0), errors="coerce").fillna(0).sum())
    max_days_overdue = int(df_overdue["days_past_expected"].max())

    # Severity
    if max_days_overdue > 60 or total_amount > 100_000:
        severity = "critical"
    elif max_days_overdue > 30 or total_amount > 50_000:
        severity = "high"
    else:
        severity = "medium"

    month_start = as_of.replace(day=1)

    # Prevent duplicate
    try:
        existing = (
            supabase.table("alerts")
            .select("id")
            .eq("client_id", str(client_id))
            .eq("alert_type", "ar_overdue")
            .eq("month_date", month_start.isoformat())
            .eq("is_active", True)
            .execute()
        )
        if existing.data:
            return
    except Exception:
        # if select fails, don't block insert attempt
        pass

    msg = (
        f"{count_invoices} customer invoice(s) are overdue more than {int(min_days_overdue)} days "
        f"past expected receipt date, totalling approx ${total_amount:,.0f}. "
        f"Oldest is about {max_days_overdue} days late."
    )

    data = {
        "client_id": str(client_id),
        "page_name": "cash_bills",
        "alert_type": "ar_overdue",
        "severity": severity,
        "message": msg,
        "month_date": month_start.isoformat(),
        "context_type": "ar_ap",
        "context_id": None,
        "is_active": True,
        "is_dismissed": False,
    }

    try:
        supabase.table("alerts").insert(data).execute()
    except Exception:
        pass


def _resolve_ar_alerts_for_month(client_id, as_of: date):
    """
    Internal helper: resolve existing ar_overdue alerts for the month of 'as_of'.
    """
    try:
        month_start = as_of.replace(day=1)
        supabase.table("alerts").update(
            {
                "is_active": False,
                "is_dismissed": True,
                "resolved_at": datetime.now(timezone.utc).isoformat(),
            }
        ).eq("client_id", str(client_id)) \
         .eq("alert_type", "ar_overdue") \
         .eq("month_date", month_start.isoformat()) \
         .eq("is_active", True) \
         .execute()
    except Exception:
        pass

def _resolve_cash_danger_alerts_for_month(client_id, month_start: date):
    """
    Internal helper: resolve existing cash_danger alerts for the given month.
    """
    if client_id is None or month_start is None:
        return

    try:
        supabase.table("alerts").update(
            {
                "is_active": False,
                "is_dismissed": True,
                "resolved_at": datetime.now(timezone.utc).isoformat(),
            }
        ).eq("client_id", str(client_id)) \
         .eq("alert_type", "cash_danger") \
         .eq("month_date", month_start.isoformat()) \
         .eq("is_active", True) \
         .execute()
    except Exception:
        pass

from datetime import datetime, timezone

def upsert_alert(
    client_id: str,
    alert_type: str,
    severity: str,
    title: str,
    message: str,
    month_date: date,
    page_name: str,
    context_type: str | None = None,
    context_id: str | None = None,
):
    """
    Insert or update an alert for a given client + type + month + page.

    If an active alert already exists for that combination, it is updated.
    Otherwise a new row is inserted.
    """

    if not client_id or month_date is None:
        return None

    try:
        # 1) Check if an active alert already exists for this "key"
        existing = (
            supabase
            .table("alerts")
            .select("id")
            .eq("client_id", str(client_id))
            .eq("alert_type", alert_type)
            .eq("month_date", month_date.isoformat())
            .eq("page_name", page_name)
            .eq("context_type", context_type)
            .eq("is_active", True)
            .limit(1)
            .execute()
        )

        base_payload = {
            "client_id": str(client_id),
            "alert_type": alert_type,
            "severity": severity,
            "title": title,
            "message": message,
            "month_date": month_date.isoformat(),
            "is_active": True,
            "page_name": page_name,
            "context_type": context_type,
            "context_id": context_id,
            "is_dismissed": False,
        }

        # 2) Update existing alert
        if existing.data:
            alert_id = existing.data[0]["id"]

            payload = {
                **base_payload,
                "resolved_at": None,            # keep it open
            }

            resp = (
                supabase
                .table("alerts")
                .update(payload)
                .eq("id", alert_id)
                .execute()
            )
            print(
                f"[ALERT] updated {alert_type} alert for "
                f"client={client_id}, month={month_date}"
            )
            return resp.data

        # 3) Insert new alert
        resp = (
            supabase
            .table("alerts")
            .insert(base_payload)
            .execute()
        )

        print(
            f"[ALERT] inserted {alert_type} alert for "
            f"client={client_id}, month={month_date}, severity={severity}"
        )
        return resp.data

    except Exception as e:
        print(
            f"[ALERT] upsert_alert FAILED for client={client_id}, "
            f"type={alert_type}, month={month_date}: {e}"
        )
        return None



def ensure_ar_ap_alerts_for_month(
    client_id: str,
    month_ref: date,
    ar_ageing_df: pd.DataFrame | None,
    ap_ageing_df: pd.DataFrame | None,
) -> None:
    """
    Create / update alerts for AR and AP ageing for a given client + month_ref.

    Expects:
      ar_ageing_df with columns: ["bucket", "total_amount"]
      ap_ageing_df with columns: ["bucket", "total_amount"]
      where 'bucket' includes values like:
        - "0–30 days overdue"
        - "30–60 days overdue"
        - "60–90 days overdue"
        - "90+ days overdue"
        - "Not due / current"
    """

    if not client_id or month_ref is None:
        return

    # ---------- Helper to compute overdue totals ----------
    def _compute_overdue(df: pd.DataFrame | None) -> float:
        if df is None or df.empty:
            return 0.0
        overdue = df[~df["bucket"].str.contains("Not due", case=False, na=False)]
        return float(overdue["total_amount"].fillna(0).sum())

    total_ar_overdue = _compute_overdue(ar_ageing_df)
    total_ap_overdue = _compute_overdue(ap_ageing_df)

    # ---------- AR overdue alert ----------
    if total_ar_overdue > 0:
        # choose thresholds to taste
        if total_ar_overdue >= 1_000_000:
            severity = "critical"
            title = "Large AR overdue"
            message = (
                f"Customer invoices overdue total about {total_ar_overdue:,.0f}. "
                "Prioritise collections on >60 day buckets."
            )
        elif total_ar_overdue >= 250_000:
            severity = "warning"
            title = "AR overdue building up"
            message = (
                f"Customer invoices overdue total about {total_ar_overdue:,.0f}. "
                "Review ageing and chase late payers."
            )
        else:
            severity = "info"
            title = "Some AR overdue"
            message = (
                f"Customer invoices overdue total about {total_ar_overdue:,.0f}. "
                "Keep an eye on collections."
            )

        upsert_alert(
            client_id=client_id,
            alert_type="ar_overdue",
            severity=severity,
            title=title,
            message=message,
            month_date=month_ref,
            page_name="ar_ap_ageing",   # or "cash_bills" if that’s your page
            context_type="kpi",
            context_id=None,
        )

    # ---------- AP overdue alert ----------
    if total_ap_overdue > 0:
        if total_ap_overdue >= 1_000_000:
            severity = "critical"
            title = "Large AP overdue"
            message = (
                f"Supplier bills overdue total about {total_ap_overdue:,.0f}. "
                "Risk of supplier pressure or credit holds."
            )
        elif total_ap_overdue >= 250_000:
            severity = "warning"
            title = "AP overdue building up"
            message = (
                f"Supplier bills overdue total about {total_ap_overdue:,.0f}. "
                "Plan payments and talk to key suppliers."
            )
        else:
            severity = "info"
            title = "Some AP overdue"
            message = (
                f"Supplier bills overdue total about {total_ap_overdue:,.0f}. "
                "Monitor cash and payment terms."
            )

        upsert_alert(
            client_id=client_id,
            alert_type="ap_overdue",
            severity=severity,
            title=title,
            message=message,
            month_date=month_ref,
            page_name="ar_ap_ageing",   # same page name as you’ll filter on
            context_type="kpi",
            context_id=None,
        )


def ensure_cash_danger_alert_for_month(client_id, month_start: date):
    """
    Use cashflow_summary for this client + month to manage a 'cash_danger' alert.

      - If cash_danger_flag is TRUE or closing_cash <= 0 -> ensure a critical alert.
      - If FALSE or no row                                -> resolve any existing alert.
    """
    if client_id is None or month_start is None:
        return

    df = fetch_cashflow_summary_for_client_as_of(client_id, month_start)
    if df is None or df.empty:
        _resolve_cash_danger_alerts_for_month(client_id, month_start)
        return

    # Make sure month_date is datetime-like for .dt
    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    else:
        _resolve_cash_danger_alerts_for_month(client_id, month_start)
        return

    if "cash_danger_flag" not in df.columns:
        _resolve_cash_danger_alerts_for_month(client_id, month_start)
        return

    # Compare by month (YYYY-MM)
    df["month_key"] = df["month_date"].dt.to_period("M")
    target_key = pd.to_datetime(month_start).to_period("M")

    row = df[df["month_key"] == target_key]
    if row.empty:
        _resolve_cash_danger_alerts_for_month(client_id, month_start)
        return

    row = row.iloc[0]
    flag = bool(row.get("cash_danger_flag", False))
    closing_cash = row.get("closing_cash", None)

    # Optional: treat <= 0 as danger even if flag is False
    try:
        closing_cash_val = float(closing_cash) if closing_cash is not None else None
    except Exception:
        closing_cash_val = None

    is_danger = flag or (closing_cash_val is not None and closing_cash_val <= 0)

    # If not a danger month -> resolve any existing alerts
    if not is_danger:
        _resolve_cash_danger_alerts_for_month(client_id, month_start)
        return

    # Avoid duplicate active alerts
    try:
        existing = (
            supabase
            .table("alerts")
            .select("id")
            .eq("client_id", str(client_id))
            .eq("alert_type", "cash_danger")
            .eq("month_date", month_start.isoformat())
            .eq("is_active", True)
            .execute()
        )
        if existing.data:
            return
    except Exception:
        return

    cash_str = ""
    try:
        if closing_cash_val is not None:
            cash_str = f" (projected cash ~${closing_cash_val:,.0f})"
    except Exception:
        pass

    msg = (
        f"Projected cash is at or below zero this month{cash_str}. "
        "This is your cliff month – plan funding or cuts before this point."
    )

    data = {
        "client_id": str(client_id),
        "page_name": "business_overview",
        "alert_type": "cash_danger",
        "severity": "critical",
        "message": msg,
        "month_date": month_start.isoformat(),
        "context_type": "cashflow_summary",
        "context_id": None,
        "is_active": True,
        "is_dismissed": False,
    }

    try:
        supabase.table("alerts").insert(data).execute()
    except Exception:
        pass


from datetime import datetime, date


def ensure_runway_alert_for_month(client_id, month_start: date):
    """
    Check runway for the given client + month and ensure we have the right runway alert:
      - If runway < client-specific runway_min_months -> create 'runway_low' alert.
      - If runway >= threshold -> resolve any existing runway_low alerts for that month.

    Uses cashflow_summary + compute_runway_and_effective_burn_from_df,
    so it does NOT rely on a separate KPIs table being populated.
    """
    if client_id is None or month_start is None:
        print("[ALERT] ensure_runway_alert_for_month: missing client_id or month_start")
        return

    # Per-client threshold
    settings = get_client_settings(client_id)
    min_months = float(settings["runway_min_months"])

    # --- derive runway from cashflow_summary ---
    df_cf = fetch_cashflow_summary_for_client(client_id)
    if df_cf is None or df_cf.empty:
        print(f"[ALERT] ensure_runway_alert_for_month: no cashflow_summary rows for client={client_id}")
        return

    try:
        # IMPORTANT: pass month_start as month_ref
        runway, eff_burn = compute_runway_and_effective_burn_from_df(df_cf, month_start)
    except Exception as e:
        print(f"[ALERT] ensure_runway_alert_for_month: error computing runway for client={client_id}: {e}")
        return

    if runway is None:
        print(f"[ALERT] ensure_runway_alert_for_month: runway is None for client={client_id}")
        return

    try:
        runway = float(runway)
    except (TypeError, ValueError):
        print(f"[ALERT] ensure_runway_alert_for_month: cannot cast runway='{runway}' to float for client={client_id}")
        return

    print(
        f"[ALERT] ensure_runway_alert_for_month: client={client_id}, "
        f"month={month_start}, runway={runway}, min_months={min_months}"
    )

    # Helper: resolve (deactivate) existing runway alerts if any
    def resolve_existing_runway_alerts():
        try:
            supabase.table("alerts").update(
                {
                    "is_active": False,
                    "is_dismissed": True,
                    "resolved_at": datetime.now(timezone.utc).isoformat(),
                }
            ).eq("client_id", str(client_id)) \
             .eq("alert_type", "runway_low") \
             .eq("month_date", month_start.isoformat()) \
             .eq("is_active", True) \
             .execute()
            print(f"[ALERT] resolved runway_low alerts for client={client_id}, month={month_start}")
        except Exception as e:
            print(
                f"[ALERT] failed to resolve runway_low alerts for "
                f"client={client_id}, month={month_start}: {e}"
            )

    # If runway is healthy, just resolve any old alerts and return
    if runway >= min_months:
        print(
            f"[ALERT] runway OK for client={client_id}, "
            f"month={month_start}, runway={runway} >= min={min_months}"
        )
        resolve_existing_runway_alerts()
        return

    # Decide severity based on how low it is
    if runway < 2:
        severity = "critical"
        title = "Runway critically low"
    elif runway < min_months:
        severity = "high"
        title = "Runway below target"
    else:
        severity = "medium"
        title = "Runway risk"

    # Check if an active alert for this month already exists
    try:
        existing = (
            supabase
            .table("alerts")
            .select("id")
            .eq("client_id", str(client_id))
            .eq("alert_type", "runway_low")
            .eq("month_date", month_start.isoformat())
            .eq("is_active", True)
            .execute()
        )
        if getattr(existing, "data", None):
            print(
                f"[ALERT] existing active runway_low alert found for "
                f"client={client_id}, month={month_start} – skipping insert"
            )
            return
    except Exception as e:
        print(
            f"[ALERT] failed to check existing runway_low alerts for "
            f"client={client_id}, month={month_start}: {e}"
        )
        return

    # Founder-friendly message
    msg = f"Runway is around {runway:.1f} months. Consider reducing spend or planning your next raise."

    data = {
        "client_id": str(client_id),
        "page_name": "business_overview",
        "alert_type": "runway_low",
        "severity": severity,
        "title": title,  # NOT NULL in Supabase
        "message": msg,
        "month_date": month_start.isoformat(),
        "context_type": "kpi",
        "context_id": None,
        "is_active": True,
        "is_dismissed": False,
    }

    try:
        resp = supabase.table("alerts").insert(data).execute()
        print(
            f"[ALERT] runway_low INSERTED for client={client_id}, "
            f"month={month_start}, severity={severity}, runway={runway}"
        )
        print(f"[ALERT] insert response: {getattr(resp, 'data', resp)}")
    except Exception as e:
        print(
            f"[ALERT] runway_low INSERT FAILED for client={client_id}, "
            f"month={month_start}: {e}"
        )
        pass



def fetch_alerts_for_client(
    client_id,
    only_active: bool = True,
    limit: int = 50,
    page_name: str | None = None,
    month_start: date | None = None,
    alert_types: list[str] | None = None,
):
    """
    Fetch alerts for this client.
    Options:
      - only_active      -> return only active alerts
      - page_name        -> filter alerts belonging to a specific page (eg: 'ar_ap_ageing')
      - month_start      -> filter alert month
      - alert_types      -> list of specific alert types eg ["ar_overdue","ap_overdue"]

    Sorted newest first.
    """

    if client_id is None:
        return []

    try:
        query = (
            supabase
            .table("alerts")
            .select("*")
            .eq("client_id", str(client_id))
        )

        if only_active:
            query = query.eq("is_active", True)

        if page_name:
            query = query.eq("page_name", page_name)

        if month_start:
            query = query.eq("month_date", month_start.isoformat())

        if alert_types:
            query = query.in_("alert_type", alert_types)

        query = query.order("created_at", desc=True).limit(limit)

        res = query.execute()
        return res.data or []

    except Exception:
        return []


def dismiss_alert(alert_id: int):
    """
    Mark an alert as dismissed (inactive) in the alerts table.
    """
    try:
        supabase.table("alerts").update(
            {
                "is_active": False,
                "is_dismissed": True,
                "resolved_at": datetime.now(timezone.utc).isoformat(),
            }
        ).eq("id", alert_id).execute()
        return True
    except Exception:
        return False


@st.cache_data(ttl=60)
def fetch_all_open_tasks(client_id):
    """
    Fetch all open tasks for this client (any page/month).
    """
    if client_id is None:
        return []

    try:
        res = (
            supabase
            .table("tasks")
            .select("*")
            .eq("client_id", str(client_id))
            .neq("status", "done")
            .order("created_at", desc=True)
            .limit(50)
            .execute()
        )
        return res.data or []
    except Exception:
        return []

# ---------- Reusable UI pieces ----------

def comments_block(page_name: str):
    """Simple shared comments UI, scoped by page_name + client."""
    st.markdown("---")
    st.subheader("💬 Team comments")

    if not selected_client_id:
        st.caption("Pick a business first to see comments.")
        return

    key_prefix = f"{selected_client_id}_{page_name}"

    # Load existing comments (page + client)
    comments = fetch_comments_for_client_page(selected_client_id, page_name)

    if not comments:
        st.caption("No comments yet. Start the discussion for this page.")
    else:
        for c in comments:
            author = c.get("author_name", "Someone on your team")
            created = c.get("created_at", "")
            text = c.get("comment_text", "")

            st.markdown(f"**{author}** · _{created}_")
            st.write(text)
            st.markdown("---")

    st.markdown("#### Add a new comment")

    new_comment = st.text_area(
        "What do you want others to know about this view?",
        key=f"comment_input_{key_prefix}",
        placeholder="E.g. 'December payroll is higher because of bonuses.'",
    )

    if st.button("Post comment", key=f"comment_btn_{key_prefix}"):
        if new_comment.strip():
            ok = create_comment_for_client_page(
                selected_client_id,
                page_name=page_name,
                comment_text=new_comment.strip(),
            )
            if ok:
                st.success("Comment posted.")
                st.rerun()
            else:
                st.error("Could not save comment. Please try again.")
        else:
            st.warning("Type something before posting.")


def tasks_block(page_name: str):
    """
    Lightweight task list per client + page.

    Expects Supabase table 'tasks' with columns:
      id, client_id, page_name, title, owner, status, priority, due_date, created_at
    """
    st.markdown("### ✅ Tasks for this page")

    if not selected_client_id:
        st.info("Select a business first to manage tasks.")
        return

    # ---- Load existing tasks ----
    try:
        resp = (
            supabase.table("tasks")
            .select("*")
            .eq("client_id", str(selected_client_id))
            .eq("page_name", page_name)
            .order("status", desc=False)
            .order("due_date", desc=False)
            .order("created_at", desc=False)
            .execute()
        )
        tasks = resp.data or []
    except Exception as e:
        st.error("Couldn't load tasks right now.")
        print("Error fetching tasks:", e)
        tasks = []

    # ---- Show tasks with quick status updates ----
    if not tasks:
        st.caption("No tasks yet for this page.")
    else:
        for t in tasks:
            tid = t.get("id")
            title = t.get("title") or "(no title)"
            owner = t.get("owner") or "Unassigned"
            status = t.get("status") or "Open"
            priority = t.get("priority") or "Medium"
            due = t.get("due_date") or ""

            row_cols = st.columns([4, 2, 2, 2, 1])
            with row_cols[0]:
                st.markdown(f"**{title}**")
                st.caption(f"👤 {owner}")
            with row_cols[1]:
                st.caption(f"Priority: **{priority}**")
            with row_cols[2]:
                st.caption(f"Due: {due or '—'}")
            with row_cols[3]:
                new_status = st.selectbox(
                    "Status",
                    options=["Open", "In progress", "Done"],
                    index=["Open", "In progress", "Done"].index(status)
                    if status in ["Open", "In progress", "Done"]
                    else 0,
                    key=f"task_status_{tid}",
                )
            with row_cols[4]:
                if st.button("💾", key=f"task_save_{tid}", help="Save status"):
                    try:
                        supabase.table("tasks").update(
                            {"status": new_status}
                        ).eq("id", tid).execute()
                        st.success("Updated.")
                        st.rerun()
                    except Exception as e:
                        st.error("Could not update task.")
                        print("Error updating task:", e)

    st.markdown("#### ➕ Add a new task")

    with st.form(key=f"task_form_{page_name}"):
        title = st.text_input("Task title")
        owner = st.text_input("Owner (name or role)", value="Founder")
        col1, col2 = st.columns(2)
        with col1:
            priority = st.selectbox(
                "Priority",
                options=["Low", "Medium", "High"],
                index=1,
            )
        with col2:
            due_date = st.date_input(
                "Due date (optional)",
                value=None,
                key=f"task_due_{page_name}",
            )

        submitted = st.form_submit_button("Create task")

        if submitted:
            if not title.strip():
                st.warning("Task title is required.")
            else:
                try:
                    # date_input may return None or a date; convert to iso if non-empty
                    due_iso = None
                    if due_date:
                        try:
                            due_iso = due_date.isoformat()
                        except Exception:
                            due_iso = None

                    data = {
                        "client_id": str(selected_client_id),
                        "page_name": page_name,
                        "title": title.strip(),
                        "owner": owner.strip() or "Founder",
                        "status": "Open",
                        "priority": priority,
                        "due_date": due_iso,
                        "created_at": datetime.utcnow().isoformat(),
                    }
                    supabase.table("tasks").insert(data).execute()
                    st.success("Task created.")
                    st.rerun()
                except Exception as e:
                    st.error("Could not create task.")
                    print("Error inserting task:", e)



def top_header(title: str):
    st.title(title)
    st.caption(f"{selected_client_name} • Focus month: {selected_month_label}")

def _safe_float(val, default: float = 0.0) -> float:
    try:
        if val is None or val == "—":
            return default
        return float(val)
    except Exception:
        return default


def build_kpi_change_explanations(client_id, month_date: date) -> dict:
    """
    Compare this month vs previous month for KPI values and
    return:
      {
        "deltas": {
            "money_in": float | None,
            "money_out": float | None,
            "cash": float | None,
            "runway": float | None,
        },
        "explanations": {
            "money_in": str,
            "money_out": str,
            "cash": str,
            "runway": str,
        },
    }
    If previous month is missing, deltas will be None and text will say
    it's the first tracked month.
    """
    result = {
        "deltas": {
            "money_in": None,
            "money_out": None,
            "cash": None,
            "runway": None,
        },
        "explanations": {
            "money_in": "First month tracked – no prior month to compare.",
            "money_out": "First month tracked – no prior month to compare.",
            "cash": "First month tracked – no prior month to compare.",
            "runway": "First month tracked – no prior month to compare.",
        },
    }

    if client_id is None or month_date is None:
        return result

    this_month_ts = pd.to_datetime(month_date).replace(day=1)
    prev_month_ts = (this_month_ts - pd.DateOffset(months=1)).date()

    curr = fetch_kpis_for_client_month(client_id, this_month_ts.date())
    prev = fetch_kpis_for_client_month(client_id, prev_month_ts)

    if not curr or not prev:
        # No previous row – keep defaults
        return result

    # Extract numeric values
    curr_in = _safe_float(curr.get("revenue", 0.0))
    prev_in = _safe_float(prev.get("revenue", 0.0))

    curr_out = _safe_float(curr.get("burn", 0.0))
    prev_out = _safe_float(prev.get("burn", 0.0))

    curr_cash = _safe_float(curr.get("cash_balance", 0.0))
    prev_cash = _safe_float(prev.get("cash_balance", 0.0))

    curr_runway = _safe_float(curr.get("runway_months", 0.0))
    prev_runway = _safe_float(prev.get("runway_months", 0.0))

    di = curr_in - prev_in
    do = curr_out - prev_out
    dc = curr_cash - prev_cash
    dr = curr_runway - prev_runway

    result["deltas"]["money_in"] = di
    result["deltas"]["money_out"] = do
    result["deltas"]["cash"] = dc
    result["deltas"]["runway"] = dr

    # Helper to build short change sentence
    def _change_sentence(delta: float, label: str) -> str:
        if abs(delta) < 1e-6:
            return f"{label} is roughly flat vs last month."
        direction = "up" if delta > 0 else "down"
        return f"{label} is {direction} by {abs(delta):,.0f} vs last month."

    result["explanations"]["money_in"] = _change_sentence(di, "Money in")
    result["explanations"]["money_out"] = _change_sentence(do, "Money out")

    if abs(dc) < 1e-6:
        result["explanations"]["cash"] = (
            "Cash in bank is roughly flat vs last month."
        )
    else:
        direction = "higher" if dc > 0 else "lower"
        result["explanations"]["cash"] = (
            f"Cash in bank is {direction} by {abs(dc):,.0f} vs last month."
        )

    if abs(dr) < 1e-3:
        result["explanations"]["runway"] = (
            "Runway is about the same as last month."
        )
    else:
        direction = "longer" if dr > 0 else "shorter"
        result["explanations"]["runway"] = (
            f"Runway is {direction} by {abs(dr):.1f} months vs last month."
        )

    return result


def build_month_summary(
    kpis: dict | None,
    deltas: dict | None,
    alerts_for_month: list[dict] | None = None,
) -> str:
    """
    Build a short, founder-friendly summary of the month using:
      - KPIs (money in, out, cash, runway)
      - Deltas vs last month
      - Key alerts (optional)
    """
    if not kpis:
        return (
            "We don't have enough data to summarise this month yet. "
            "Once AR/AP, payroll and cash are loaded, this will turn into a "
            "simple story about what happened to your money."
        )

    money_in = float(kpis.get("revenue") or 0.0)
    money_out = float(kpis.get("burn") or 0.0)
    cash = float(kpis.get("cash_balance") or 0.0)
    runway = kpis.get("runway_months")

    d_money_in = (deltas or {}).get("money_in")
    d_money_out = (deltas or {}).get("money_out")
    d_cash = (deltas or {}).get("cash")
    d_runway = (deltas or {}).get("runway")

    parts = []

    # 1) Headline on money in / out
    if money_in == 0 and money_out == 0:
        parts.append("No meaningful cash movement recorded this month yet.")
    else:
        if money_in >= money_out:
            parts.append(
                f"You brought in about ${money_in:,.0f} and spent about ${money_out:,.0f}, "
                "so cashflow from operations was **positive** this month."
            )
        else:
            parts.append(
                f"You brought in about ${money_in:,.0f} and spent about ${money_out:,.0f}, "
                "so cashflow from operations was **negative** this month."
            )

    # 2) Changes vs last month (if we have deltas)
    change_bits = []
    if d_money_in is not None and d_money_in != 0:
        direction = "up" if d_money_in > 0 else "down"
        change_bits.append(f"Money in is {direction} by ${abs(d_money_in):,.0f} vs last month")
    if d_money_out is not None and d_money_out != 0:
        direction = "up" if d_money_out > 0 else "down"
        change_bits.append(f"Money out is {direction} by ${abs(d_money_out):,.0f}")
    if d_cash is not None and d_cash != 0:
        direction = "up" if d_cash > 0 else "down"
        change_bits.append(f"Cash in bank is {direction} by ${abs(d_cash):,.0f}")

    if change_bits:
        parts.append(" ".join(change_bits) + ".")

    # 3) Cash + runway
    if cash > 0:
        if isinstance(runway, (int, float)) and runway is not None:
            parts.append(
                f"You finished the month with about ${cash:,.0f} in the bank, "
                f"which gives roughly **{runway:.1f} months of runway** at the recent burn rate."
            )
        else:
            parts.append(
                f"You finished the month with about ${cash:,.0f} in the bank. "
                "We don't have enough history yet to compute a stable runway."
            )
    else:
        parts.append(
            "Cash in bank is effectively at or below zero. "
            "You are in a critical cash position and should treat this as urgent."
        )

    # 4) Alerts (pull just 1–2 key ones)
    if alerts_for_month:
        # Sort so highest severity first
        sev_order = {"critical": 3, "high": 2, "medium": 1, "low": 0}
        alerts_sorted = sorted(
            alerts_for_month,
            key=lambda a: sev_order.get(str(a.get("severity", "medium")).lower(), 1),
            reverse=True,
        )

        top = alerts_sorted[:2]
        alert_msgs = []
        for a in top:
            msg = a.get("message", "")
            if msg:
                alert_msgs.append(msg)

        if alert_msgs:
            parts.append(
                "Top issues to pay attention to: " + " | ".join(alert_msgs)
            )

    return " ".join(parts)


def build_month_summary_insight(
    client_id: str,
    focus_month: date,
) -> str:
    """
    Generate a short, opinionated 'CFO-style' summary for the focus month.

    Uses cashflow_summary + computed runway/effective burn to answer:
      - What kind of month was this? (strong / controlled burn / tight / critical)
      - What mainly drove it? (operating vs investing vs financing)
      - What should the founder pay attention to next?
    """
    try:
        engine_df = fetch_cashflow_summary_for_client_as_of(client_id, focus_month)
        if engine_df is None or engine_df.empty:
            return (
                "We don't have any cashflow engine rows yet for this business. "
                "Once AR/AP, payroll and investing/financing flows are in and you've rebuilt "
                "the cashflow, this summary will explain what drove the month."
            )

        df = engine_df.copy()
        df["month_date"] = pd.to_datetime(
            df["month_date"], errors="coerce"
        ).dt.date
        df = df.sort_values("month_date")

        focus = pd.to_datetime(focus_month).date().replace(day=1)

        # ---- This month row ----
        this_row = df[df["month_date"] == focus]
        if this_row.empty:
            return (
                "We don't yet have a cashflow row for this month. "
                "Rebuild the cashflow engine, or pick a month that has data."
            )

        closing_cash = float(this_row.get("closing_cash", 0.0).iloc[0] or 0.0)
        free_cf_this = float(this_row.get("free_cash_flow", 0.0).iloc[0] or 0.0)
        investing_cf_this = float(
            this_row.get("investing_cf", 0.0).iloc[0] or 0.0
        )
        financing_cf_this = float(
            this_row.get("financing_cf", 0.0).iloc[0] or 0.0
        )
        operating_cf_this = float(
            this_row.get("operating_cf", 0.0).iloc[0] or 0.0
        )

        # ---- Previous month (for comparison) ----
        prev_row = df[df["month_date"] < focus].tail(1)
        closing_prev = None
        free_cf_prev = None
        if not prev_row.empty:
            closing_prev = float(
                prev_row.get("closing_cash", 0.0).iloc[0] or 0.0
            )
            free_cf_prev = float(
                prev_row.get("free_cash_flow", 0.0).iloc[0] or 0.0
            )

        # ---- Runway & effective burn ----
        runway_months, effective_burn = compute_runway_and_effective_burn_from_df(
        engine_df,
        focus,
    )

        # ------------------------------------------------------------------
        # 1) Classify the month
        # ------------------------------------------------------------------
        headline_parts: list[str] = []

        # Cash movement this month
        if free_cf_this > 0:
            headline_parts.append("Cash **increased** this month.")
        elif free_cf_this < 0:
            headline_parts.append("Cash **decreased** this month.")
        else:
            headline_parts.append("Cash was **roughly flat** this month.")

        # Runway view
        if runway_months is not None:
            rm = runway_months
            if rm < 3:
                headline_parts.append(
                    f"Runway is **critical** at about {rm:.1f} months."
                )
            elif rm < 6:
                headline_parts.append(
                    f"Runway is **tight**, around {rm:.1f} months."
                )
            elif rm < 12:
                headline_parts.append(
                    f"Runway is **comfortable** at about {rm:.1f} months, "
                    "with a controlled burn."
                )
            else:
                headline_parts.append(
                    f"Runway is **strong** at roughly {rm:.1f} months."
                )
        else:
            if effective_burn and effective_burn > 0:
                headline_parts.append(
                    "You're burning cash, but there isn't enough history yet "
                    "to estimate runway."
                )
            else:
                headline_parts.append(
                    "There isn't enough history yet to estimate runway."
                )

        headline = " ".join(headline_parts)

        # ------------------------------------------------------------------
        # 2) Identify the main driver of the month
        # ------------------------------------------------------------------
        driver_parts: list[str] = []

        # Magnitudes for attribution
        total_movement = abs(free_cf_this)
        inv_abs = abs(investing_cf_this)
        fin_abs = abs(financing_cf_this)
        op_abs = abs(operating_cf_this)

        # If there's very little movement, keep this simple
        if total_movement < 1e-6:
            driver_parts.append(
                "This month’s net cash movement was very small, so there "
                "were no major cash drivers."
            )
        else:
            # Decide which component is the biggest contributor
            # (operating vs investing vs financing)
            biggest = "operating"
            biggest_val = op_abs

            if inv_abs > biggest_val:
                biggest = "investing"
                biggest_val = inv_abs
            if fin_abs > biggest_val:
                biggest = "financing"
                biggest_val = fin_abs

            share = biggest_val / total_movement if total_movement > 0 else 0.0

            if biggest == "operating":
                if share > 0.6:
                    driver_parts.append(
                        "Most of the movement came from **operating cashflow** "
                        "(revenue vs operating spend)."
                    )
                else:
                    driver_parts.append(
                        "Operating cashflow mattered, but **investing/financing "
                        "moves also played a role**."
                    )
            elif biggest == "investing":
                driver_parts.append(
                    "The main driver was **investing activity** – things like "
                    "capex, asset purchases or disposals."
                )
            else:  # financing
                driver_parts.append(
                    "The main driver was **financing activity** – equity raises, "
                    "loan drawdowns/repayments or dividends."
                )

        # Also mention whether movement is a one-off vs recurring
        if abs(investing_cf_this) > 0 and abs(investing_cf_this) > 0.3 * total_movement:
            driver_parts.append(
                "Part of this looks like a **one-off investing move** "
                "(not a recurring monthly pattern)."
            )
        if abs(financing_cf_this) > 0 and abs(financing_cf_this) > 0.3 * total_movement:
            driver_parts.append(
                "There is also a **financing impact** this month "
                "(funding or repayments)."
            )

        driver_sentence = " ".join(driver_parts)

        # ------------------------------------------------------------------
        # 3) Suggest next actions / focus areas
        # ------------------------------------------------------------------
        action_parts: list[str] = []

        # If runway is short
        if runway_months is not None and runway_months < 6:
            action_parts.append(
                "Priority should be **extending runway** – "
                "review hiring pace, contractor spend and large supplier costs."
            )
        elif runway_months is not None and runway_months >= 12:
            action_parts.append(
                "Given the healthy runway, focus can stay on **growth and ROI** "
                "rather than immediate cost cutting."
            )

        # If burn is increasing vs last month
        if free_cf_prev is not None and free_cf_prev < 0 and free_cf_this < free_cf_prev:
            action_parts.append(
                "Burn is **higher than last month** – it’s worth checking "
                "for permanent step-ups in payroll or recurring vendor spend."
            )

        # If very little movement and strong cash
        if total_movement < 1000 and closing_cash > 0:
            action_parts.append(
                "Cash movement is small relative to balance – it may be time "
                "to decide whether to **deploy excess cash** or keep it as buffer."
            )

        if not action_parts:
            # Fallback generic guidance
            action_parts.append(
                "Use this view to confirm whether the month behaved as planned – "
                "if not, start with operating spend, then any large "
                "investing or financing moves."
            )

        actions_sentence = " ".join(action_parts)

        # ------------------------------------------------------------------
        # Final combined summary (single paragraph)
        # ------------------------------------------------------------------
        return f"{headline} {driver_sentence} {actions_sentence}"

    except Exception as e:
        print("Error building month summary insight:", e)
        return (
            "We hit an issue while building this month's summary. "
            "Check that cashflow_summary has data for this client and month."
        )


def render_cash_engine_debug_tile(client_id: str, base_month: date):
    """
    Debug tile: shows how opening cash, net free cashflow and closing cash
    relate over the 12-month window currently in view.

    Helps catch double-counting or weird jumps in the engine.
    """
    if not client_id or not base_month:
        st.caption("No business selected for cash debug.")
        return

    engine_df = fetch_cashflow_summary_for_client(client_id)
    if engine_df is None or engine_df.empty:
        st.caption("No cashflow engine data yet – rebuild cashflow first.")
        return

    df = engine_df.copy()
    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    df = df.dropna(subset=["month_date"])

    start = pd.to_datetime(base_month).replace(day=1)
    end = start + pd.DateOffset(months=12)

    window = df[(df["month_date"] >= start) & (df["month_date"] < end)].copy()
    window = window.sort_values("month_date")

    if window.empty:
        st.caption("No engine rows in this 12-month window yet.")
        return

    # Make sure columns exist
    for col in ["free_cash_flow", "closing_cash"]:
        if col not in window.columns:
            window[col] = 0.0

    first_row = window.iloc[0]
    last_row = window.iloc[-1]

    closing_0 = float(first_row.get("closing_cash", 0.0) or 0.0)
    free_cf_0 = float(first_row.get("free_cash_flow", 0.0) or 0.0)

    # Opening cash is implied by the first month: closing = opening + free_cf
    opening_0 = closing_0 - free_cf_0

    net_free_cf = float(window["free_cash_flow"].sum() or 0.0)
    expected_closing = opening_0 + net_free_cf

    closing_actual = float(last_row.get("closing_cash", 0.0) or 0.0)
    diff = closing_actual - expected_closing

    # Currency handling
    curr_code, curr_symbol = get_client_currency(client_id)

    try:
        fmt = format_money  # use your existing helper if present
    except NameError:
        def fmt(value, currency_symbol=curr_symbol or "$"):
            if value is None:
                return "—"
            try:
                return f"{currency_symbol}{float(value):,.0f}"
            except Exception:
                return str(value)

    st.markdown("#### 🧪 Cash engine sanity check")

    c1, c2, c3, c4 = st.columns(4)
    with c1:
        st.metric("Opening cash (start of window)", fmt(opening_0, curr_symbol))
    with c2:
        st.metric("Net free cash over 12 months", fmt(net_free_cf, curr_symbol))
    with c3:
        st.metric("Expected closing (open + net)", fmt(expected_closing, curr_symbol))
    with c4:
        st.metric(
            "Engine closing (last row)",
            fmt(closing_actual, curr_symbol),
            delta=fmt(diff, curr_symbol) if diff != 0 else None,
            help="Delta = engine closing minus recalculated closing.",
        )

    st.caption(
        "If the expected closing (open + net free CF) and engine closing differ a lot, "
        "it may indicate double counting or an off opening balance."
    )


def render_quick_scenario_block_for_overview():
    """
    Lightweight scenario playground for the Business Overview page.
    Uses the same engine as page_scenarios but without saved-scenario management.
    """

    if not selected_client_id or not selected_month_start:
        st.info("Pick a business and focus month at the top before running scenarios.")
        return

    currency_code, currency_symbol = get_client_currency(selected_client_id)

    base_df = get_engine_window_for_scenarios(
        selected_client_id,
        base_month=selected_month_start,
        n_months=12,
    )

    if base_df is None or base_df.empty:
        st.caption(
            "No cashflow engine data yet – go to **Business at a Glance → Rebuild cashflow** first."
        )
        return

    base_df = base_df.sort_values("month_date").reset_index(drop=True)

    # Shared month selector for scenario start
    month_opts = base_df["month_date"].dt.to_period("M").drop_duplicates()
    month_label_map = {p.strftime("%b %Y"): p.to_timestamp() for p in month_opts}
    month_labels = list(month_label_map.keys())

    col_output, col_controls = st.columns([2, 1])

    # ------------------------------------
    # RIGHT: presets + levers
    # ------------------------------------
    with col_controls:
        st.markdown("#### Scenario type")

        preset = st.selectbox(
            "Quick scenario preset",
            options=[
                "Revenue push (+20%)",
                "Hire 1 key role (+$15k/month)",
                "Cut supplier/operating spend by 20%",
                "Raise $250k now",
                "Raise $250k + monthly top-up",
                "Custom mix",
            ],
        )

        controls = {
            "rev_change_pct": 0.0,
            "collections_boost_pct": 0.0,
            "extra_payroll_per_month": 0.0,
            "spend_change_pct": 0.0,
            "extra_capex_one_off": 0.0,
            "extra_capex_recurring": 0.0,
            "equity_raise": 0.0,
            "recurring_funding": 0.0,
        }

        start_month_label = st.selectbox(
            "Scenario starts from month",
            options=month_labels,
        )
        start_month_ts = month_label_map[start_month_label]

        is_custom = preset == "Custom mix"
        disabled = not is_custom

        tab_rev, tab_hire, tab_spend, tab_capex, tab_fund = st.tabs(
            [
                "Revenue & customers",
                "Hiring & team",
                "Bills & spend",
                "Capex & assets",
                "Funding & runway",
            ]
        )

        # Revenue & collections
        with tab_rev:
            rev_change_pct = st.slider(
                "Change in cash from customers (%)",
                min_value=-50,
                max_value=100,
                value=20 if preset == "Revenue push (+20%)" else 0,
                step=5,
                disabled=disabled,
            )
            collections_boost_pct = st.slider(
                "Change from collecting invoices earlier / later (%)",
                min_value=-20,
                max_value=20,
                value=0,
                step=5,
                disabled=disabled,
            )
            if is_custom:
                controls["rev_change_pct"] = rev_change_pct
                controls["collections_boost_pct"] = collections_boost_pct

        # Hiring
        with tab_hire:
            extra_headcount_cost = st.number_input(
                f"Extra monthly payroll from that month ({currency_code})",
                value=15_000.0 if preset == "Hire 1 key role (+$15k/month)" else 0.0,
                step=1000.0,
                disabled=disabled,
            )
            if is_custom:
                controls["extra_payroll_per_month"] = extra_headcount_cost

        # Spend
        with tab_spend:
            spend_change_pct = st.slider(
                "Change in supplier / operating spend (%)",
                min_value=-50,
                max_value=50,
                value=-20 if preset == "Cut supplier/operating spend by 20%" else 0,
                step=5,
                disabled=disabled,
            )
            if is_custom:
                controls["spend_change_pct"] = spend_change_pct

        # Capex (kept simple)
        with tab_capex:
            extra_capex_one_off = st.number_input(
                f"One-off capex / asset sale in start month ({currency_code})",
                value=0.0,
                step=5000.0,
                disabled=disabled,
            )
            if is_custom:
                controls["extra_capex_one_off"] = extra_capex_one_off

        # Funding
        with tab_fund:
            equity_raise = st.number_input(
                f"One-off equity / loan cash-in in start month ({currency_code})",
                value=250_000.0 if preset in ["Raise $250k now", "Raise $250k + monthly top-up"] else 0.0,
                step=25_000.0,
                disabled=disabled,
            )
            recurring_funding = st.number_input(
                f"Recurring funding each month from start ({currency_code})",
                value=10_000.0 if preset == "Raise $250k + monthly top-up" else 0.0,
                step=5000.0,
                disabled=disabled,
            )
            if is_custom:
                controls["equity_raise"] = equity_raise
                controls["recurring_funding"] = recurring_funding

        # Apply non-custom presets
        if not is_custom:
            if preset == "Revenue push (+20%)":
                controls["rev_change_pct"] = 20.0
            elif preset == "Hire 1 key role (+$15k/month)":
                controls["extra_payroll_per_month"] = 15_000.0
            elif preset == "Cut supplier/operating spend by 20%":
                controls["spend_change_pct"] = -20.0
            elif preset == "Raise $250k now":
                controls["equity_raise"] = 250_000.0
            elif preset == "Raise $250k + monthly top-up":
                controls["equity_raise"] = 250_000.0
                controls["recurring_funding"] = 10_000.0

        run_clicked = st.button("Run this scenario", key="bo_run_scenario")

    # ------------------------------------
    # LEFT: simple outputs
    # ------------------------------------
    with col_output:
        if not run_clicked:
            st.info("Choose a preset on the right and click **Run this scenario**.")
            return

        if not any(abs(v) > 0 for v in controls.values()):
            st.info("Nothing changed yet – try a preset or adjust the sliders.")
            return

        scen_df, (base_label, base_date), (scen_label, scen_date) = run_founder_scenario(
            base_df=base_df,
            start_month=start_month_ts,
            controls=controls,
        )

        st.caption(f"Scenario run: **{preset} starting {start_month_label}**")

        st.markdown("#### 📌 Runway & danger month")
        col_m1, col_m2 = st.columns(2)
        with col_m1:
            st.metric("Baseline danger month", base_label or "None in 12 months")
        with col_m2:
            delta_txt = None
            if base_date is not None and scen_date is not None:
                month_diff = (scen_date.to_period("M") - base_date.to_period("M")).n
                if month_diff != 0:
                    sign_word = "later" if month_diff > 0 else "earlier"
                    delta_txt = f"{abs(month_diff)} month(s) {sign_word}"
            st.metric(
                "Scenario danger month",
                scen_label or "None in 12 months",
                delta=delta_txt,
            )

        # Plain-language summary
        st.markdown("#### 🧠 What this means in plain language")

        min_base = float(base_df["closing_cash"].min())
        min_scen = float(scen_df["closing_cash_scenario"].min())
        diff_min = min_scen - min_base

        lines = []
        if diff_min > 0:
            lines.append(
                f"- Your **lowest cash balance improves by {currency_symbol}{diff_min:,.0f}**."
            )
        elif diff_min < 0:
            lines.append(
                f"- Your **lowest cash balance drops by {currency_symbol}{abs(diff_min):,.0f}**."
            )

        if base_label and scen_label and base_label != scen_label:
            lines.append(
                f"- The month where cash gets tight moves from **{base_label} → {scen_label}**."
            )

        if controls.get("extra_payroll_per_month", 0.0):
            lines.append("- This includes extra monthly payroll for new hires.")

        if controls.get("equity_raise", 0.0) or controls.get("recurring_funding", 0.0):
            lines.append("- This includes extra cash in from funding (equity / loans / top-ups).")

        if controls.get("spend_change_pct", 0.0) < 0:
            lines.append("- You are **cutting operating spend**, which improves runway.")
        elif controls.get("spend_change_pct", 0.0) > 0:
            lines.append("- You are **increasing operating spend**, which reduces runway.")

        if not lines:
            lines.append("- This scenario only makes a small change in the next 12 months.")

        st.markdown("\n".join(lines))

        # Simple cash curve chart
        st.markdown("---")
        st.markdown(f"#### 📉 Cash curve – baseline vs scenario ({currency_code})")

        chart_df = pd.DataFrame(
            {
                "month_date": scen_df["month_date"],
                "Baseline closing cash": base_df["closing_cash"].values,
                "Scenario closing cash": scen_df["closing_cash_scenario"].values,
            }
        )
        chart_long = chart_df.melt(
            id_vars="month_date", var_name="Series", value_name="Cash"
        )

        cash_chart = (
            alt.Chart(chart_long)
            .mark_line()
            .encode(
                x=alt.X("month_date:T", title="Month", axis=alt.Axis(format="%b %Y")),
                y=alt.Y("Cash:Q", title=f"Closing cash ({currency_code})"),
                color=alt.Color("Series:N", title=""),
                tooltip=[
                    alt.Tooltip("month_date:T", title="Month", format="%b %Y"),
                    "Series:N",
                    alt.Tooltip("Cash:Q", title="Cash", format=",.0f"),
                ],
            )
        )

        st.altair_chart(cash_chart.interactive(), width="stretch")



def _safe_float(x, default=0.0):
    try:
        if x is None:
            return default
        return float(x)
    except Exception:
        return default


def _compute_health_verdict(
    runway_months: float | None,
    cash_balance: float | None,
    money_in: float | None,
    money_out: float | None,
    active_alerts: list | None,
):
    """
    Returns: (status, headline, message)
    status: "good" | "warn" | "risk"
    """
    runway = _safe_float(runway_months, default=None) if runway_months is not None else None
    cash = _safe_float(cash_balance, default=None) if cash_balance is not None else None
    rev = _safe_float(money_in, 0.0)
    burn = _safe_float(money_out, 0.0)

    # alert severity count
    critical = 0
    high = 0
    if active_alerts:
        for a in active_alerts:
            sev = str(a.get("severity", "")).lower()
            if sev == "critical":
                critical += 1
            elif sev == "high":
                high += 1

    # Basic health logic (simple + founder-friendly)
    if runway is not None and runway <= 2.0:
        status = "risk"
        headline = "🔴 Business health: At risk"
        msg = "Runway is very tight. You should act immediately on cash preservation and start funding prep now."
        return status, headline, msg

    if runway is not None and runway <= 4.0:
        status = "warn"
        headline = "🟡 Business health: Tight runway"
        msg = "Cash is stable short-term, but runway is tight. Tighten costs/collections and plan funding early."
        # If also many severe alerts, escalate tone
        if critical > 0 or high >= 2:
            msg = "Runway is tight and risk signals are elevated. Prioritise cash control this month and prepare funding."
        return status, headline, msg

    # If runway looks fine, but burn > revenue (or revenue is zero), still caution
    if burn > 0 and rev >= 0 and burn > rev and (runway is None or runway <= 8.0):
        status = "warn"
        headline = "🟡 Business health: Watching burn"
        msg = "Runway is okay for now, but spend is outpacing inflows. Keep burn controlled and improve collections."
        return status, headline, msg

    status = "good"
    headline = "🟢 Business health: Stable"
    msg = "Cash position looks stable in the near term. Focus on efficient growth and keep an eye on upcoming commitments."
    return status, headline, msg


def _compute_recovery_vs_funding_signal(
    runway_months: float | None,
    money_in: float | None,
    money_out: float | None,
):
    """
    Returns: (mode, message)
    mode: "recover" | "mixed" | "funding"
    """
    runway = _safe_float(runway_months, default=None) if runway_months is not None else None
    rev = _safe_float(money_in, 0.0)
    burn = _safe_float(money_out, 0.0)

    # Heuristics: founders need a clear fork
    if runway is not None and runway <= 3.0:
        return "funding", "Funding likely required unless you can cut burn fast. Start funding prep now (2–3 months lead time)."

    if runway is not None and runway <= 6.0 and burn > rev:
        return "mixed", "You may recover operationally (reduce burn / speed up collections), but also start light funding prep as a backup."

    return "recover", "You likely can recover through operating levers (collections, costs, pricing) without immediate funding—monitor runway monthly."


# ---------- Shared helpers: 6-month trend data for overview ----------

def _six_month_window(focus_month) -> pd.DatetimeIndex:
    """
    Return a 6-month monthly index ending at focus_month (inclusive),
    anchored to first of each month.
    """
    focus = pd.to_datetime(focus_month).replace(day=1)
    start = focus - pd.DateOffset(months=5)
    return pd.date_range(start=start, periods=6, freq="MS")


def _detect_revenue_value_column(df: pd.DataFrame) -> str | None:
    """
    Pick the best column to use for recognised revenue from a revenue schedule.
    """
    for cand in ["recognised_revenue", "revenue_amount", "amount", "Total revenue"]:
        if cand in df.columns:
            return cand
    return None




def _build_revenue_trend_df(client_id, focus_month) -> pd.DataFrame:
    """
    Build a 6-month recognised revenue trend ending at focus_month.

    Source: build_revenue_schedule_for_client (same as Sales & Deals page).
    Returns columns: ["month_date", "Month", "Revenue"].
    """
    if not client_id or focus_month is None:
        return pd.DataFrame()

    window = _six_month_window(focus_month)
    base_month = window[0].date()

    rev_df = build_revenue_schedule_for_client(
        client_id,
        base_month=base_month,
        n_months=len(window),
    )

    if rev_df is None or rev_df.empty:
        return pd.DataFrame()

    rev_df = rev_df.copy()
    rev_df["month_date"] = pd.to_datetime(rev_df["month_date"], errors="coerce")
    rev_df = rev_df[rev_df["month_date"].notna()]

    value_col = _detect_revenue_value_column(rev_df)
    if value_col is None:
        return pd.DataFrame()

    # Keep only our 6-month window
    rev_window = rev_df[
        rev_df["month_date"].isin(window)
    ][["month_date", value_col]].copy()

    # Aggregate by month and ensure all 6 months exist
    monthly = (
        rev_window.groupby("month_date", as_index=False)[value_col]
        .sum()
        .rename(columns={value_col: "Revenue"})
    )

    grid = pd.DataFrame({"month_date": window})
    monthly = grid.merge(monthly, on="month_date", how="left")
    monthly["Revenue"] = pd.to_numeric(
        monthly["Revenue"], errors="coerce"
    ).fillna(0.0)

    monthly["Month"] = monthly["month_date"].dt.strftime("%b %Y")
    return monthly[["month_date", "Month", "Revenue"]]

def _build_burn_and_cash_trend_df(
    engine_df: pd.DataFrame | None,
    focus_month,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """
    From cashflow_summary (engine_df), build 6-month burn + cash balance trends
    ending at focus_month.

    Burn logic:
      - Prefer an explicit 'burn' / 'net_burn' column if present.
      - Else, prefer 'operating_cf' (core operating burn, excluding financing).
      - Else, fall back to 'free_cash_flow'.

    Returns:
      burn_df:  ["month_date", "Month", "Burn"]
      cash_df:  ["month_date", "Month", "Closing cash"]
    """
    if engine_df is None or engine_df.empty or focus_month is None:
        return pd.DataFrame(), pd.DataFrame()

    # 6-month window ending at focus_month
    window = _six_month_window(focus_month)

    df = engine_df.copy()

    # Normalise month_date to month-start Timestamp
    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    df["month_date"] = df["month_date"].dt.to_period("M").dt.to_timestamp()
    df = df[df["month_date"].notna()]

    # NEW: if engine_df carries multiple plan versions per month_date,
    # keep only the latest as_of_month snapshot for each month.
    if "as_of_month" in df.columns:
        df["as_of_month"] = pd.to_datetime(df["as_of_month"], errors="coerce")
        df = (
            df.sort_values(["month_date", "as_of_month"])
              .groupby("month_date", as_index=False)
              .tail(1)
        )

    # Filter to the 6-month window
    df_window = df[df["month_date"].isin(window)].copy()

    # Ensure we have one row for every month in the window (even if no data)
    grid = pd.DataFrame({"month_date": window})
    df_window = grid.merge(df_window, on="month_date", how="left")

    # ---------- Burn series ----------
    # Prefer operating_cf over free_cash_flow so chart shows *operating burn*.
    burn_source_col = None
    for cand in ["burn", "net_burn", "operating_cf", "free_cash_flow"]:
        if cand in df_window.columns:
            burn_source_col = cand
            break

    burn_df = pd.DataFrame(columns=["month_date", "Month", "Burn"])
    if burn_source_col is not None:
        vals = pd.to_numeric(df_window[burn_source_col], errors="coerce").fillna(0.0)

        if burn_source_col in ["burn", "net_burn"]:
            # Already in "burn" convention (positive = burn)
            burn_vals = vals
        else:
            # Convert negative CF into positive burn magnitude; positive CF -> 0 burn
            burn_vals = vals.apply(lambda v: -v if v < 0 else 0.0)

        burn_df = pd.DataFrame(
            {
                "month_date": df_window["month_date"],
                "Month": df_window["month_date"].dt.strftime("%b %Y"),
                "Burn": burn_vals,
            }
        )

    # ---------- Cash balance series ----------
    cash_df = pd.DataFrame(columns=["month_date", "Month", "Closing cash"])
    if "closing_cash" in df_window.columns:
        cash_vals = pd.to_numeric(df_window["closing_cash"], errors="coerce").fillna(0.0)
        cash_df = pd.DataFrame(
            {
                "month_date": df_window["month_date"],
                "Month": df_window["month_date"].dt.strftime("%b %Y"),
                "Closing cash": cash_vals,
            }
        )

    return burn_df, cash_df



def fetch_cashflow_summary_for_client_as_of(
    client_id: str,
    as_of_month: date,
) -> pd.DataFrame | None:
    """
    Fetch cashflow_summary rows for a client that belong to the plan
    as-of a given Focus Month (as_of_month).
    """
    if client_id is None or as_of_month is None:
        return pd.DataFrame()

    def _call():
        return (
            supabase
            .table("cashflow_summary")
            .select("*")
            .eq("client_id", str(client_id))
            .eq("as_of_month", str(as_of_month))
            .order("month_date")
            .execute()
        )

    res = sb_execute_with_retry(_call, label="fetch cashflow_summary as_of")
    if res is None:
        print("[ERROR] fetch_cashflow_summary_for_client_as_of -> connection problem (None)")
        return None

    rows = res.data or []
    if not rows:
        print(
            "[DEBUG] fetch_cashflow_summary_for_client_as_of -> 0 rows for",
            client_id, "as_of", as_of_month,
        )
        return pd.DataFrame()

    df = pd.DataFrame(rows)

    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
        df = df[df["month_date"].notna()].sort_values("month_date")

    print(
        f"[DEBUG] fetch_cashflow_summary_for_client_as_of -> {len(df)} rows for",
        client_id, "as_of", as_of_month,
    )
    return df


HEALTH_LEVEL = {"STABLE": 0, "WATCH": 1, "AT RISK": 2}

def _clamp_one_level_change(prev_state: str | None, new_state: str) -> str:
    if not prev_state or prev_state not in HEALTH_LEVEL:
        return new_state
    p = HEALTH_LEVEL[prev_state]
    n = HEALTH_LEVEL[new_state]
    if n > p + 1:
        # only allow +1 step worsening
        inv = {v: k for k, v in HEALTH_LEVEL.items()}
        return inv[p + 1]
    if n < p - 1:
        # only allow -1 step improving
        inv = {v: k for k, v in HEALTH_LEVEL.items()}
        return inv[p - 1]
    return new_state


def _compute_cash_cliff_within_days(*, cliff_ts_exact: pd.Timestamp | None, focus_month: pd.Timestamp, days: int = 90) -> bool:
    """
    cliff_ts_exact can be a month timestamp or an exact date (weekly). We measure from focus_month start.
    """
    if cliff_ts_exact is None:
        return False
    try:
        focus = pd.to_datetime(focus_month).to_period("M").to_timestamp()
        cliff = pd.to_datetime(cliff_ts_exact)
        return (cliff - focus).days <= days
    except Exception:
        return False


def _compute_business_health_state(
    *,
    focus_month: pd.Timestamp,
    runway_months: float | None,
    cash_balance: float | None,
    revenue_this_month: float,
    revenue_prev_month: float | None,
    burn_this_month: float,
    burn_prev_month: float | None,
    cash_delta_mom: float | None,            # from your deltas['cash']
    burn_delta_mom: float | None,            # from your deltas['money_out']
    rev_delta_mom: float | None,             # from your deltas['money_in'] or rev delta if you prefer
    planned_burn: float | None = None,       # optional, use settings if you have it
    cash_cliff_within_90d: bool = False,
    one_off_outflow_flag: bool = False,
    forecast_runway_down_2m: bool = False,
) -> tuple[str, str]:
    """
    Returns (state, message) where state is ONLY one of:
      - "STABLE"
      - "WATCH"
      - "AT RISK"
    """

    # ---------------------------
    # 🔴 AT RISK (any one)
    # ---------------------------
    if runway_months is not None and runway_months < 6:
        return "AT RISK", "Business is at risk. Immediate corrective action required."
    if cash_cliff_within_90d:
        return "AT RISK", "Business is at risk. Immediate corrective action required."
    if (burn_delta_mom is not None and burn_delta_mom > 0) and (rev_delta_mom is not None and rev_delta_mom < 0):
        return "AT RISK", "Business is at risk. Immediate corrective action required."
    if cash_delta_mom is not None and cash_delta_mom < 0:
        # “cash ↓ materially with no clear operational driver” → we can only be mechanical.
        # Mark “material” conservatively: >10% of cash balance or >1 month of burn if available.
        try:
            material = False
            if cash_balance and cash_balance > 0:
                material = abs(cash_delta_mom) >= 0.10 * float(cash_balance)
            if not material and burn_this_month and burn_this_month > 0:
                material = abs(cash_delta_mom) >= 1.0 * float(burn_this_month)
            if material:
                return "AT RISK", "Business is at risk. Immediate corrective action required."
        except Exception:
            pass
    if forecast_runway_down_2m:
        return "AT RISK", "Business is at risk. Immediate corrective action required."

    # ---------------------------
    # 🟢 STABLE (all must be true)
    # ---------------------------
    stable_ok = True

    if runway_months is None or runway_months < 9:
        stable_ok = False
    if burn_delta_mom is not None and burn_delta_mom > 0:
        stable_ok = False
    if cash_cliff_within_90d:
        stable_ok = False

    # Cash MoM decline ≤ planned burn (or improving)
    if planned_burn is not None and cash_delta_mom is not None:
        # cash_delta_mom negative means cash fell; compare to planned burn magnitude
        try:
            if cash_delta_mom < 0 and abs(float(cash_delta_mom)) > abs(float(planned_burn)):
                stable_ok = False
        except Exception:
            stable_ok = False

    if stable_ok:
        return "STABLE", "Business is stable. Cash position and runway are within acceptable bounds."

    # ---------------------------
    # 🟠 WATCH (any one)
    # ---------------------------
    if runway_months is not None and 6 <= runway_months < 9:
        return "WATCH", "Business requires close attention. Trends are moving against plan."
    if burn_delta_mom is not None and burn_delta_mom > 0:
        return "WATCH", "Business requires close attention. Trends are moving against plan."
    if planned_burn is not None and cash_delta_mom is not None:
        try:
            if cash_delta_mom < 0 and abs(float(cash_delta_mom)) > abs(float(planned_burn)):
                return "WATCH", "Business requires close attention. Trends are moving against plan."
        except Exception:
            pass
    if (revenue_prev_month is not None) and (revenue_this_month < revenue_prev_month) and (burn_delta_mom is None or burn_delta_mom >= 0):
        return "WATCH", "Business requires close attention. Trends are moving against plan."
    if one_off_outflow_flag:
        return "WATCH", "Business requires close attention. Trends are moving against plan."

    # Default if not stable and not risk by triggers
    return "WATCH", "Business requires close attention. Trends are moving against plan."


def _compute_primary_focus(
    *,
    health_state: str,
    runway_months: float | None,
    cash_cliff_within_90d: bool,
    cash_delta_mom: float | None,
    burn_delta_mom: float | None,
    revenue_delta_mom: float | None,
    revenue_this_month: float,
    burn_this_month: float,
    payroll_share_up: bool = False,   # optional if you have it
    ar_collections_issue: bool = False,  # optional if you can detect
) -> str:
    # 🔐 Rule 1 — Health overrides everything
    if health_state == "AT RISK":
        return "Extend runway and protect cash"

    # 1️⃣ CASH SURVIVAL
    if (runway_months is not None and runway_months < 6) or cash_cliff_within_90d:
        return "Extend runway and protect cash"
    if cash_delta_mom is not None and burn_delta_mom is not None:
        # “Cash ↓ faster than burn ↓”
        if cash_delta_mom < 0 and burn_delta_mom < 0 and abs(cash_delta_mom) > abs(burn_delta_mom):
            return "Extend runway and protect cash"

    # 2️⃣ COST CONTROL
    if burn_delta_mom is not None and burn_delta_mom > 0:
        return "Control costs, especially payroll and discretionary spend"
    if payroll_share_up:
        return "Control costs, especially payroll and discretionary spend"
    if revenue_this_month <= 0 and burn_this_month > 0:
        return "Control costs, especially payroll and discretionary spend"
    if revenue_delta_mom is not None and revenue_delta_mom <= 0 and burn_delta_mom is not None and burn_delta_mom >= 0:
        return "Control costs, especially payroll and discretionary spend"

    # 3️⃣ REVENUE & COLLECTIONS
    if revenue_delta_mom is not None and revenue_delta_mom < 0:
        return "Stabilise revenue and accelerate cash collections"
    if ar_collections_issue:
        return "Stabilise revenue and accelerate cash collections"
    if (revenue_this_month > 0) and (cash_delta_mom is not None and cash_delta_mom < 0):
        return "Stabilise revenue and accelerate cash collections"

    # 4️⃣ OPTIMISATION / DISCIPLINE
    return "Maintain discipline and selectively optimise growth"

def _build_main_actions(*, primary_focus: str) -> list[dict]:
    """
    Returns actions with:
      - action: specific instruction
      - kpi_trace: list of KPI labels shown on this page
      - why: explicit justification tied to the KPI(s)
    HARD RULE:
      - If no KPI trace exists, the action must not exist.
    """
    lib = {
        "Extend runway and protect cash": [
            {
                "action": "Freeze all non-essential spend immediately",
                "kpi_trace": ["Runway (months)", "🏦 Cash in bank"],
                "why": "Runway is constrained relative to cash in bank; freezing spend protects runway by reducing cash-out pressure.",
            },
            {
                "action": "Delay or renegotiate upcoming cash outflows",
                "kpi_trace": ["💸 Expect Cash out", "🧯 Runway (months)"],
                "why": "Expected cash out is the controllable driver; reducing or shifting outflows extends runway.",
            },
            {
                "action": "Prioritise invoices & collections due in the next 30 days",
                "kpi_trace": ["💰 Expect Cash in", "🏦 Cash in bank"],
                "why": "Expected cash in is the fastest lever to improve cash in bank; pulling collections forward reduces cash stress risk.",
            },
            {
                "action": "Model funding scenarios (raise timing/amount) against the next 12 months",
                "kpi_trace": ["🔍 Cashflow engine – 12-month breakdown", "🧯 Runway (months)"],
                "why": "The 12-month engine is your forward runway driver; scenarios test whether funding timing/size prevents runway deterioration.",
            },
            {
                "action": "Prepare a downside cash plan for the next 90 days",
                "kpi_trace": ["🔎 Cashflow reconciliation – Focus Month", "🧯 Runway (months)"],
                "why": "Reconciliation reveals what is actually driving cash movement; a 90-day downside plan protects runway with concrete cash levers.",
            },
        ],

        "Control costs, especially payroll and discretionary spend": [
            {
                "action": "Pause hiring and review open roles against revenue impact",
                "kpi_trace": ["💸 Expect Cash out", "🧯 Runway (months)"],
                "why": "Hiring increases expected cash out; pausing protects runway while revenue signal is validated.",
            },
            {
                "action": "Cut or defer discretionary opex immediately",
                "kpi_trace": ["💸 Expect Cash out"],
                "why": "Expected cash out is elevated; discretionary opex is the fastest controllable component to reduce cash-out.",
            },
            {
                "action": "Reforecast burn for the next 3 months and lock a burn ceiling",
                "kpi_trace": ["🔍 Cashflow engine – 12-month breakdown", "🧯 Runway (months)"],
                "why": "Runway depends on the burn path; tightening the next 3 months improves the engine’s runway profile.",
            },
            {
                "action": "Set spend approval thresholds for all new commitments",
                "kpi_trace": ["💸 Expect Cash out"],
                "why": "Expected cash out is a governance-controlled KPI; approval thresholds prevent unplanned cost creep.",
            },
        ],

        "Stabilise revenue and accelerate cash collections": [
            {
                "action": "Focus on closing near-term revenue with clear cash dates",
                "kpi_trace": ["📊 This month (recognised revenue)", "💰 Expect Cash in"],
                "why": "Recognised revenue and expected cash in are diverging risk drivers; near-term closes with cash dates improves cash-in timing.",
            },
            {
                "action": "Chase overdue invoices aggressively (weekly cadence)",
                "kpi_trace": ["💰 Expect Cash in", "🏦 Cash in bank"],
                "why": "Expected cash in is the cash lever; stronger collections directly lifts cash in bank and reduces stress risk.",
            },
            {
                "action": "Tighten terms to pull forward collections on new invoices",
                "kpi_trace": ["💰 Expect Cash in"],
                "why": "Expected cash in reflects collection timing; tightening terms accelerates cash-in without relying on cost cuts.",
            },
        ],

        "Maintain discipline and selectively optimise growth": [
            {
                "action": "Validate forecast assumptions used in the cashflow engine",
                "kpi_trace": ["🔍 Cashflow engine – 12-month breakdown"],
                "why": "Your forward view is only as strong as the engine assumptions; validating them reduces false runway confidence.",
            },
            {
                "action": "Stress-test growth options with scenarios before committing spend",
                "kpi_trace": ["🧪 Scenarios & What-if", "🧯 Runway (months)"],
                "why": "Scenarios test whether growth spend preserves runway; runway is the binding constraint even when stable.",
            },
            {
                "action": "Preserve optionality: keep burn flat unless revenue confirms",
                "kpi_trace": ["📈 6-month trends – revenue, Operating Cash Flow, cash", "🧯 Runway (months)"],
                "why": "Trends show whether revenue is consistent; keeping burn flat protects runway until revenue trend confirms.",
            },
        ],
    }

    actions = lib.get(primary_focus, [])

    # HARD FILTER: action must have trace + why
    clean: list[dict] = []
    for a in actions:
        if not isinstance(a, dict):
            continue
        if not str(a.get("action") or "").strip():
            continue
        trace = a.get("kpi_trace")
        why = str(a.get("why") or "").strip()
        if not trace or not isinstance(trace, list) or len(trace) == 0:
            continue
        if not why:
            continue
        clean.append(a)

    return clean[:5]



def page_business_overview():
    top_header("Business at a Glance")

    if not selected_client_id:
        st.info("Select a client from the sidebar to view this page.")
        return

    settings = get_client_settings(selected_client_id) or {}
    currency_code, currency_symbol = get_client_currency(selected_client_id)

    # --- Recognised revenue (this month & last month) ---
    this_month_rev = get_recognised_revenue_for_month(selected_client_id, selected_month_start)

    prev_month_date = (
        pd.to_datetime(selected_month_start) - pd.DateOffset(months=1)
    ).to_period("M").to_timestamp().date()

    prev_month_rev = get_recognised_revenue_for_month(selected_client_id, prev_month_date)

    rev_delta_label = None
    if prev_month_rev is not None:
        diff = this_month_rev - prev_month_rev
        rev_delta_label = f"{diff:+,.0f}"

    # Pull KPIs from Supabase
    kpis = fetch_kpis_for_client_month(selected_client_id, selected_month_start) or {}

    # Alerts
    ensure_runway_alert_for_month(selected_client_id, selected_month_start)
    ensure_cash_danger_alert_for_month(selected_client_id, selected_month_start)

    # ------------------------------------------------------------------
    # Engine df + runway (FETCH ONCE)
    # ------------------------------------------------------------------
    engine_df = fetch_cashflow_summary_for_client_as_of(selected_client_id, selected_month_start)

    runway_from_engine = None
    effective_burn = None
    if engine_df is not None and not engine_df.empty:
        runway_from_engine, effective_burn = compute_runway_and_effective_burn_from_df(
            engine_df=engine_df,
            month_ref=selected_month_start,
        )

    # ---------- KPI deltas & explanations ----------
    kpi_change = build_kpi_change_explanations(selected_client_id, selected_month_start)
    deltas = kpi_change.get("deltas", {}) or {}
    expl = kpi_change.get("explanations", {}) or {}

    # Defaults for display
    money_in_val = kpis.get("revenue", "—")
    money_out_val = kpis.get("burn", "—")
    cash_in_bank_val = kpis.get("cash_balance", "—")
    runway_months_val = kpis.get("runway_months", "—")

    money_in_delta_label = f"{deltas['money_in']:+,.0f}" if deltas.get("money_in") is not None else None
    money_out_delta_label = f"{deltas['money_out']:+,.0f}" if deltas.get("money_out") is not None else None
    cash_delta_label = f"{deltas['cash']:+,.0f}" if deltas.get("cash") is not None else None
    runway_delta_label = f"{deltas['runway']:+.1f} mo" if deltas.get("runway") is not None else None

    # Prefer engine runway if available
    if runway_from_engine is not None:
        runway_months_val = f"{runway_from_engine:.1f}"

    # ------------------------------------------------------------------
    # ✅ Define runway_num ONCE (fixes your NameError + used by Health/Focus/Actions)
    # ------------------------------------------------------------------
    runway_num = None
    try:
        runway_num = float(runway_months_val) if runway_months_val not in (None, "—") else None
    except Exception:
        runway_num = None

    cash_num = None
    try:
        cash_num = _safe_float(cash_in_bank_val, default=None) if cash_in_bank_val != "—" else None
    except Exception:
        cash_num = None

    rev_num = _safe_float(money_in_val, 0.0) if money_in_val != "—" else 0.0
    burn_num = _safe_float(money_out_val, 0.0) if money_out_val != "—" else 0.0

    # ---------- Top KPI row ----------
    kpi_col1, kpi_col0, kpi_col2, kpi_col3, kpi_col4 = st.columns(5)

    with kpi_col1:
        st.metric(
            f"💰 Expect Cash in ({currency_code})",
            format_money(money_in_val, currency_symbol),
            delta=money_in_delta_label,
            help=(
                f"Expected cash **in** this month from customers "
                f"(unpaid AR with expected_date in this month, in {currency_code})."
            ),
        )
        st.caption(expl.get("money_in", ""))

    with kpi_col0:
        st.metric(
            f"📊 This month (recognised revenue, {currency_code})",
            f"{currency_symbol}{this_month_rev:,.0f}",
            delta=rev_delta_label,
            help=(
                "Recognised revenue for this month based on your Sales & Deals "
                "revenue recognition schedule (SaaS / milestone / POC / etc)."
            ),
        )

    with kpi_col2:
        st.metric(
            f"💸 Expect Cash out ({currency_code})",
            format_money(money_out_val, currency_symbol),
            delta=money_out_delta_label,
            help=(
                "Expected cash **out** this month: unpaid supplier bills due/expected this month "
                "+ payroll cash for this month + opex for this month "
                f"(in {currency_code})."
            ),
        )
        st.caption(expl.get("money_out", ""))

    with kpi_col3:
        st.metric(
            f"🏦 Cash in bank ({currency_code})",
            format_money(cash_in_bank_val, currency_symbol),
            delta=cash_delta_label,
        )
        st.caption(expl.get("cash", ""))

    with kpi_col4:
        st.metric(
            "🧯 Runway (months)",
            "—" if runway_months_val in (None, "—") else f"{float(runway_months_val):.1f}",
            delta=runway_delta_label,
        )
        st.caption(expl.get("runway", ""))

    # ---------- This month in 20 seconds ----------
    st.subheader("📌 This month in 20 seconds")
    summary_text = build_month_summary_insight(selected_client_id, selected_month_start)
    st.info(summary_text)

    # ---------- 6-month trends (revenue only) ----------
    st.markdown("---")
    st.subheader("📈 6-month trends – revenue")

    rev_trend_df = _build_revenue_trend_df(selected_client_id, selected_month_start)
    if rev_trend_df is None or rev_trend_df.empty:
        st.caption("No recognised revenue data available for this 6-month window.")
    else:
        rev_chart = (
            alt.Chart(rev_trend_df)
            .mark_line(point=True)
            .encode(
                x=alt.X("month_date:T", title="Month", axis=alt.Axis(format="%b %y")),
                y=alt.Y("Revenue:Q", title=f"Revenue ({currency_code})"),
                tooltip=[
                    alt.Tooltip("Month:N", title="Month"),
                    alt.Tooltip("Revenue:Q", title="Revenue", format=",.0f"),
                ],
            )
            .properties(height=220)
        )
        st.altair_chart(rev_chart.interactive(), width="stretch")

    # ---------- Cashflow engine controls ----------
    st.markdown("---")
    st.subheader("🔁 Cashflow engine controls")

    if st.button("Rebuild cashflow for next 12 months"):
        opening_hint = settings.get("opening_cash_start")
        ok = recompute_cashflow_from_ar_ap(
            selected_client_id,
            base_month=selected_month_start,
            n_months=12,
            opening_cash_hint=opening_hint,
        )
        if ok:
            st.success("Cashflow updated from AR/AP, payroll, opex and investing/financing moves.")
            st.rerun()
        else:
            st.error("Could not rebuild cashflow. Check AR/AP data and try again.")

    # ---------------------------------------------
    # ✅ Cashflow engine – 12-month breakdown (COMBO CHART)
    # ---------------------------------------------
    st.markdown("---")
    st.subheader("🔍 Cashflow engine – 12-month breakdown (visual)")

    if engine_df is None or engine_df.empty:
        st.caption(
            "No cashflow engine data yet for this Focus Month. "
            "Hit **'Rebuild cashflow for next 12 months'** above to generate it."
        )
    else:
        start = pd.to_datetime(selected_month_start).replace(day=1)
        month_index = pd.date_range(start=start, periods=12, freq="MS")
        grid = pd.DataFrame({"month_date": month_index})

        df = engine_df.copy()
        df["month_date"] = pd.to_datetime(df.get("month_date"), errors="coerce")
        df = df[df["month_date"].notna()].copy()
        df["month_date"] = df["month_date"].dt.to_period("M").dt.to_timestamp()

        engine_window = df[df["month_date"].isin(month_index)].copy()
        merged = grid.merge(engine_window, on="month_date", how="left")

        for c in ["operating_cf", "investing_cf", "financing_cf", "closing_cash"]:
            if c not in merged.columns:
                merged[c] = 0.0
            merged[c] = pd.to_numeric(merged[c], errors="coerce").fillna(0.0)

        bar_df = merged[["month_date", "operating_cf", "investing_cf", "financing_cf"]].melt(
            id_vars=["month_date"],
            value_vars=["operating_cf", "investing_cf", "financing_cf"],
            var_name="Component",
            value_name="Amount",
        )
        bar_df["Component"] = bar_df["Component"].map({
            "operating_cf": "Operating CF",
            "investing_cf": "Investing CF",
            "financing_cf": "Financing CF",
        })

        line_df = merged[["month_date", "closing_cash"]].copy()

        x_enc = alt.X("month_date:T", title="Month", axis=alt.Axis(format="%b %Y"))

        bars = alt.Chart(bar_df).mark_bar().encode(
            x=x_enc,
            y=alt.Y("Amount:Q", title=f"Cashflow movement ({currency_code})", axis=alt.Axis(format=",.0f")),
            color=alt.Color("Component:N", title=""),
            tooltip=[
                alt.Tooltip("month_date:T", title="Month", format="%b %Y"),
                "Component:N",
                alt.Tooltip("Amount:Q", title="Amount", format=",.0f"),
            ],
        )

        bar_labels = alt.Chart(bar_df).mark_text(dy=-6).encode(
            x=x_enc,
            y="Amount:Q",
            detail="Component:N",
            text=alt.Text("Amount:Q", format=",.0f"),
        )

        line = alt.Chart(line_df).mark_line(point=True).encode(
            x=x_enc,
            y=alt.Y("closing_cash:Q", title=None, axis=alt.Axis(format=",.0f")),
            tooltip=[
                alt.Tooltip("month_date:T", title="Month", format="%b %Y"),
                alt.Tooltip("closing_cash:Q", title="Closing cash", format=",.0f"),
            ],
        )

        line_labels = alt.Chart(line_df).mark_text(dy=-10).encode(
            x=x_enc,
            y="closing_cash:Q",
            text=alt.Text("closing_cash:Q", format=",.0f"),
        )

        combo = alt.layer(bars, bar_labels, line, line_labels).resolve_scale(y="independent").properties(height=320)
        st.altair_chart(combo.interactive(), width="stretch")
        st.caption("Bars = Operating/Investing/Financing CF. Line = Closing cash. Hover for exact values.")

    # ------------------------------------------------------------------
    # Payroll vs bills vs Net operating cash (6-month view) — keep as-is
    # ------------------------------------------------------------------
    st.markdown("---")
    st.subheader("💼 Payroll vs bills vs Net operating cash")

    window_months = pd.date_range(start=pd.to_datetime(selected_month_start), periods=6, freq="MS")
    engine_df_all = fetch_cashflow_summary_for_client(selected_client_id)
    df_ar, df_ap = fetch_ar_ap_for_client(selected_client_id)
    payroll_by_month = compute_payroll_by_month(selected_client_id, window_months)

    if engine_df_all is None or engine_df_all.empty:
        st.caption("No cashflow engine data yet – rebuild cashflow first.")
    else:
        engine_df_all = engine_df_all.copy()
        engine_df_all["month_date"] = pd.to_datetime(engine_df_all["month_date"], errors="coerce")

        engine_window = engine_df_all[engine_df_all["month_date"].isin(window_months)][["month_date", "operating_cf"]].copy()

        ap_monthly = pd.DataFrame({"month_date": window_months})
        if df_ap is not None and not df_ap.empty:
            df_ap = df_ap.copy()
            df_ap["amount"] = pd.to_numeric(df_ap.get("amount"), errors="coerce").fillna(0.0)

            pay_col = "expected_payment_date"
            if pay_col not in df_ap.columns:
                pay_col = "due_date" if "due_date" in df_ap.columns else None

            if pay_col is not None:
                df_ap[pay_col] = pd.to_datetime(df_ap[pay_col], errors="coerce")
                df_ap["bucket_month"] = df_ap[pay_col].dt.to_period("M").dt.to_timestamp()

                ap_agg = (
                    df_ap.groupby("bucket_month", as_index=False)["amount"]
                    .sum()
                    .rename(columns={"bucket_month": "month_date", "amount": "ap_cash_out"})
                )
                ap_monthly = ap_monthly.merge(ap_agg, on="month_date", how="left")

        ap_monthly["ap_cash_out"] = pd.to_numeric(ap_monthly.get("ap_cash_out"), errors="coerce").fillna(0.0)

        payroll_series = (
            pd.Series([float(payroll_by_month.get(m, 0.0)) for m in window_months], index=window_months, name="payroll_cash")
            .reset_index()
            .rename(columns={"index": "month_date"})
        )

        merged2 = pd.DataFrame({"month_date": window_months})
        merged2 = merged2.merge(engine_window, on="month_date", how="left")
        merged2 = merged2.merge(ap_monthly, on="month_date", how="left")
        merged2 = merged2.merge(payroll_series, on="month_date", how="left")

        merged2["operating_cf"] = pd.to_numeric(merged2.get("operating_cf"), errors="coerce").fillna(0.0)
        merged2["ap_cash_out"] = pd.to_numeric(merged2.get("ap_cash_out"), errors="coerce").fillna(0.0)
        merged2["payroll_cash"] = pd.to_numeric(merged2.get("payroll_cash"), errors="coerce").fillna(0.0)

        merged2["month_date"] = pd.to_datetime(merged2["month_date"], errors="coerce")
        merged2 = merged2.sort_values("month_date")

        plot_df = merged2.melt(
            id_vars="month_date",
            value_vars=["operating_cf", "payroll_cash", "ap_cash_out"],
            var_name="Series",
            value_name="Amount",
        )

        series_labels = {
            "operating_cf": "Net operating CF",
            "payroll_cash": "Payroll (cash out)",
            "ap_cash_out": "Bills / AP (cash out)",
        }
        plot_df["Series"] = plot_df["Series"].map(series_labels)

        chart = alt.Chart(plot_df).mark_line().encode(
            x=alt.X("month_date:T", title="Month", axis=alt.Axis(format="%b %Y")),
            y=alt.Y("Amount:Q", title=f"Amount (cash, {currency_code})"),
            color=alt.Color("Series:N", title=""),
            tooltip=[
                alt.Tooltip("month_date:T", title="Month", format="%b %Y"),
                "Series:N",
                alt.Tooltip("Amount:Q", title="Amount", format=",.0f"),
            ],
        )
        st.altair_chart(chart.interactive(), width="stretch")

    # ------------------------------------------------------------------
    # ✅ Business Health + Primary Focus + Main Actions (NO governance)
    # ------------------------------------------------------------------
    st.markdown("---")
    st.subheader("🧭 Business summary")

    focus_ts = pd.to_datetime(selected_month_start).to_period("M").to_timestamp()

    # Simple mechanical flags (kept conservative)
    cash_cliff_90d = bool(runway_num is not None and runway_num < 3.0)  # ~90 days proxy

    one_off_outflow_flag = False
    try:
        if deltas.get("cash") is not None and deltas.get("money_out") is not None:
            if float(deltas["cash"]) < 0 and float(deltas["money_out"]) <= 0:
                one_off_outflow_flag = True
    except Exception:
        one_off_outflow_flag = False

    planned_burn = None
    try:
        planned_burn = _safe_float(settings.get("planned_burn_monthly"), default=None)
    except Exception:
        planned_burn = None

    # 1) BUSINESS HEALTH (3 states only) — uses your existing function
    health_state, health_msg = _compute_business_health_state(
        focus_month=focus_ts,
        runway_months=runway_num,
        cash_balance=cash_num,
        revenue_this_month=float(this_month_rev or 0.0),
        revenue_prev_month=float(prev_month_rev) if prev_month_rev is not None else None,
        burn_this_month=float(burn_num or 0.0),
        burn_prev_month=None,  # optional: wire if you want
        cash_delta_mom=deltas.get("cash"),
        burn_delta_mom=deltas.get("money_out"),
        rev_delta_mom=deltas.get("money_in"),
        planned_burn=planned_burn,
        cash_cliff_within_90d=cash_cliff_90d,
        one_off_outflow_flag=one_off_outflow_flag,
        forecast_runway_down_2m=False,
    )

    st.markdown("### 1) Business health")
    if health_state == "AT RISK":
        st.error(f"🔴 **AT RISK** — {health_msg}")
    elif health_state == "WATCH":
        st.warning(f"🟠 **WATCH** — {health_msg}")
    else:
        st.success(f"🟢 **STABLE** — {health_msg}")

    # 2) PRIMARY FOCUS (single) — uses your existing function
    primary_focus = _compute_primary_focus(
        health_state=health_state,
        runway_months=runway_num,
        cash_cliff_within_90d=cash_cliff_90d,
        cash_delta_mom=deltas.get("cash"),
        burn_delta_mom=deltas.get("money_out"),
        revenue_delta_mom=deltas.get("money_in"),
        revenue_this_month=float(this_month_rev or 0.0),
        burn_this_month=float(burn_num or 0.0),
        payroll_share_up=False,
        ar_collections_issue=False,
    )

    st.markdown("### 2) Primary focus")
    st.info(f"**{primary_focus}**")

    # 3) MAIN ACTIONS — must be KPI traceable
    st.markdown("### 3) Main actions (this month)")

    # Your _build_main_actions(primary_focus=...) should return list of dicts:
    #   {"action": "...", "why": "...", "kpi_trace": ["Runway < 6m", "Burn ↑ MoM", ...]}
    actions_raw = _build_main_actions(primary_focus=primary_focus) or []

    # HARD FILTER: drop anything with no KPI trace
    actions = []
    for a in actions_raw:
        if not isinstance(a, dict):
            continue
        kpi_trace = a.get("kpi_trace") or []
        kpi_trace = [str(x).strip() for x in kpi_trace if str(x).strip()]
        if not kpi_trace:
            continue
        # also ensure why exists (since you requested)
        why = str(a.get("why") or "").strip()
        if not why:
            continue
        actions.append({**a, "kpi_trace": kpi_trace})

    actions = actions[:5]

    if not actions:
        st.caption("No actions shown — no CFO-grade action could be defensibly traced to a KPI on this page.")
    else:
        for i, a in enumerate(actions, start=1):
            st.markdown(f"**{i}. {a['action']}**")
            st.caption(f"Why: {a['why']}")
            st.caption("KPI trace: " + " • ".join(a["kpi_trace"]))


def load_engine_for_focus(client_id: str, focus_month: date) -> pd.DataFrame | None:
    """
    Canonical engine loader for ALL UI pages.
    Always returns the ONE plan version for the focus month.
    """
    if not client_id or not focus_month:
        return pd.DataFrame()

    focus = pd.to_datetime(focus_month).to_period("M").to_timestamp().date()
    return fetch_cashflow_summary_for_client_as_of(client_id, focus)

def apply_monthly_scenario_to_weekly(
    week_df: pd.DataFrame,
    monthly_df: pd.DataFrame,
    *,
    month_col: str = "month_date",
    monthly_delta_col: str = "delta_free_cf",
) -> pd.DataFrame:
    """
    Convert monthly scenario deltas (scenario - base) into a weekly scenario cash curve.

    Inputs:
      - week_df: output of build_14_week_cash_table (must have week_start, opening_cash, closing_cash, net_cash)
      - monthly_df: must have month_col (month start) and monthly_delta_col (delta to free cash flow)
    Output:
      - same rows as week_df with added *_scenario columns including:
          opening_cash_scenario, closing_cash_scenario, net_cash_scenario,
          plus scenario copies of the component lines where available.
    """
    if week_df is None or week_df.empty:
        return pd.DataFrame()

    w = week_df.copy()

    # ---- Ensure week_start is datetime ----
    w["week_start"] = pd.to_datetime(w.get("week_start"), errors="coerce")
    w = w.dropna(subset=["week_start"]).sort_values("week_start").reset_index(drop=True)

    # ---- Ensure required base numeric columns ----
    for c in ["opening_cash", "closing_cash", "net_cash"]:
        if c not in w.columns:
            w[c] = 0.0
        w[c] = pd.to_numeric(w[c], errors="coerce").fillna(0.0)

    # If no monthly deltas, scenario == base
    if monthly_df is None or monthly_df.empty or monthly_delta_col not in monthly_df.columns:
        for c in [
            "opening_cash", "cash_in_ar", "cash_in_other_income",
            "cash_out_ap", "cash_out_payroll", "cash_out_opex",
            "operating_cf", "investing_cf", "financing_cf", "net_cash", "closing_cash"
        ]:
            if c in w.columns:
                w[f"{c}_scenario"] = w[c]
        w["opening_cash_scenario"] = w["opening_cash"]
        w["closing_cash_scenario"] = w["closing_cash"]
        w["net_cash_scenario"] = w["net_cash"]
        return w

    # ---- Monthly deltas prep ----
    m = monthly_df.copy()
    m[month_col] = pd.to_datetime(m.get(month_col), errors="coerce")
    m = m[m[month_col].notna()].copy()
    m[month_col] = m[month_col].dt.to_period("M").dt.to_timestamp()
    m[monthly_delta_col] = pd.to_numeric(m.get(monthly_delta_col), errors="coerce").fillna(0.0)

    # map each week -> its month bucket
    w["month_ts"] = w["week_start"].dt.to_period("M").dt.to_timestamp()

    # number of weeks per month in the horizon
    weeks_per_month = w.groupby("month_ts").size().to_dict()

    # delta per week for each month
    delta_per_week_map = {}
    for _, r in m.iterrows():
        mo = r[month_col]
        if mo in weeks_per_month and weeks_per_month[mo] > 0:
            delta_per_week_map[mo] = float(r[monthly_delta_col]) / float(weeks_per_month[mo])

    w["delta_week"] = w["month_ts"].map(delta_per_week_map).fillna(0.0)

    # ---- Build scenario net cash + running scenario opening/closing ----
    # We apply delta to net cash each week, then recompute running balance.
    w["net_cash_scenario"] = w["net_cash"] + w["delta_week"]

    opening_vals = []
    closing_vals = []

    running_open = float(w.loc[0, "opening_cash"])
    for i in range(len(w)):
        opening_vals.append(running_open)
        close = running_open + float(w.loc[i, "net_cash_scenario"])
        closing_vals.append(close)
        running_open = close

    w["opening_cash_scenario"] = opening_vals
    w["closing_cash_scenario"] = closing_vals

    # ---- Scenario copies of components (optional; keeps table display consistent) ----
    # NOTE: We only truly changed net cash, not line items, because scenario deltas are monthly-level.
    # If you later want to allocate deltas into specific lines, do it in run_founder_scenario and pass weekly deltas per line.
    component_cols = [
        "cash_in_ar", "cash_in_other_income",
        "cash_out_ap", "cash_out_payroll", "cash_out_opex",
        "operating_cf", "investing_cf", "financing_cf",
    ]
    for c in component_cols:
        if c in w.columns:
            w[f"{c}_scenario"] = w[c]

    # But update the total cashflow lines to match scenario cash
    w["closing_cash_scenario"] = pd.to_numeric(w["closing_cash_scenario"], errors="coerce").fillna(0.0)

    return w

def run_multi_lever_scenario(
    base_df: pd.DataFrame,
    start_month_ts: pd.Timestamp,
    op_delta_pct: float = 0.0,
    extra_payroll_per_month: float = 0.0,
    dept_burn_delta_pct: float = 0.0,
    extra_capex_per_month: float = 0.0,
    one_off_capex: float = 0.0,
    one_off_financing: float = 0.0,
    recurring_financing: float = 0.0,
) -> pd.DataFrame:
    """
    Generic scenario engine that takes the 12-month baseline cashflow engine
    (base_df) and applies different levers to produce a scenario cash curve.

    All inputs are incremental changes on top of the baseline:
      - op_delta_pct: % change to operating_cf (net operating cash) from start month
      - extra_payroll_per_month: additional (negative) payroll cost per month
      - dept_burn_delta_pct: extra % change applied only to negative operating_cf
      - extra_capex_per_month: extra investing CF per month (− = more capex, + = asset sales)
      - one_off_capex: one-off investing CF in the start month
      - one_off_financing: one-off financing CF in the start month (equity / loan inflow)
      - recurring_financing: recurring financing CF each month from start month
    """
    df = base_df.copy()
    df = df.sort_values("month_date").reset_index(drop=True)

    # Safety: make sure required columns exist & numeric
    for col in ["operating_cf", "investing_cf", "financing_cf", "free_cash_flow", "closing_cash"]:
        if col not in df.columns:
            df[col] = 0.0
        df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0.0)

    mask = df["month_date"] >= start_month_ts

    # ---------- Operating CF scenario ----------
    df["operating_cf_scenario"] = df["operating_cf"]

    # 1) Global % change on net operating cash from start month
    if op_delta_pct != 0.0:
        factor = 1.0 + op_delta_pct / 100.0
        df.loc[mask, "operating_cf_scenario"] = df.loc[mask, "operating_cf_scenario"] * factor

    # 2) Extra payroll (always cash out => negative impact)
    if extra_payroll_per_month != 0.0:
        df.loc[mask, "operating_cf_scenario"] = (
            df.loc[mask, "operating_cf_scenario"] - float(extra_payroll_per_month)
        )

    # 3) Extra burn % on the negative portion only (e.g. department-level increases)
    if dept_burn_delta_pct != 0.0:
        burn_factor = 1.0 + dept_burn_delta_pct / 100.0
        neg_mask = (df["operating_cf_scenario"] < 0.0) & mask
        df.loc[neg_mask, "operating_cf_scenario"] = (
            df.loc[neg_mask, "operating_cf_scenario"] * burn_factor
        )

    # ---------- Investing CF scenario ----------
    df["investing_cf_scenario"] = df["investing_cf"]

    # Extra monthly capex or asset sale from start month
    if extra_capex_per_month != 0.0:
        df.loc[mask, "investing_cf_scenario"] = (
            df.loc[mask, "investing_cf_scenario"] + float(extra_capex_per_month)
        )

    # One-off capex / asset sale in the start month
    if one_off_capex != 0.0:
        df.loc[df["month_date"] == start_month_ts, "investing_cf_scenario"] = (
            df.loc[df["month_date"] == start_month_ts, "investing_cf_scenario"]
            + float(one_off_capex)
        )

    # ---------- Financing CF scenario ----------
    df["financing_cf_scenario"] = df["financing_cf"]

    # One-off financing in start month (equity / new loan)
    if one_off_financing != 0.0:
        df.loc[df["month_date"] == start_month_ts, "financing_cf_scenario"] = (
            df.loc[df["month_date"] == start_month_ts, "financing_cf_scenario"]
            + float(one_off_financing)
        )

    # Recurring financing each month from start month (e.g. monthly loan drawdown)
    if recurring_financing != 0.0:
        df.loc[mask, "financing_cf_scenario"] = (
            df.loc[mask, "financing_cf_scenario"] + float(recurring_financing)
        )

    # ---------- Rebuild free CF & scenario closing cash ----------
    df["free_cash_flow_scenario"] = (
        df["operating_cf_scenario"]
        + df["investing_cf_scenario"]
        + df["financing_cf_scenario"]
    )

    # Derive the implied baseline opening cash for month 0
    first_row = df.iloc[0]
    base_opening_0 = float(first_row["closing_cash"]) - float(first_row["free_cash_flow"])

    closing_vals = []
    opening = base_opening_0
    for _, r in df.iterrows():
        free_cf_s = float(r["free_cash_flow_scenario"])
        closing = opening + free_cf_s
        closing_vals.append(closing)
        opening = closing

    df["closing_cash_scenario"] = closing_vals

    return df

def run_founder_scenario(
    base_df: pd.DataFrame,
    start_month: pd.Timestamp,
    controls: dict,
    *,
    week_df_base: pd.DataFrame | None,
    cash_buffer: float,
    compute_cash_stress_kpis_fn,
    apply_monthly_scenario_to_weekly_fn,
    today: pd.Timestamp | None = None,
    client_id: str | None = None,
    as_of_month: pd.Timestamp | None = None,
    fetch_cashflow_summary_for_client_as_of=None,   # ✅ allow passing same as Page_cash_bills
):
    """
    Scenario engine aligned to Cash & Bills:
      - monthly scenario simulated on base_df
      - weekly scenario produced from the SAME 14-week base table
      - runway/cliff computed using scenario weekly closing cash (and fallback to monthly scenario if needed)

    Returns:
      scen_week_df, base_stress_kpis, scen_stress_kpis
    """

    # -------------------------
    # Helpers
    # -------------------------
    def _safe_month_start(x) -> pd.Timestamp:
        ts = pd.to_datetime(x, errors="coerce")
        if ts is pd.NaT:
            return pd.Timestamp.utcnow().tz_localize(None).to_period("M").to_timestamp()
        return ts.to_period("M").to_timestamp()

    def _fallback_stress_from_monthly(
        monthly_df: pd.DataFrame,
        cash_buffer: float,
        today: pd.Timestamp,
    ) -> dict:
        """
        If scenario does not breach buffer inside 14 weeks, approximate from monthly scenario.
        Uses first month where closing_cash_scenario <= buffer.
        """
        out = {}
        if monthly_df is None or monthly_df.empty:
            return out
        if cash_buffer is None or float(cash_buffer) <= 0:
            return out

        m = monthly_df.copy()
        m["month_date"] = pd.to_datetime(m.get("month_date"), errors="coerce")
        m = m[m["month_date"].notna()].sort_values("month_date").reset_index(drop=True)
        if m.empty or "closing_cash_scenario" not in m.columns:
            return out

        cc = pd.to_numeric(m["closing_cash_scenario"], errors="coerce")
        breach = (cc.notna()) & (cc <= float(cash_buffer))

        if not breach.any():
            # no breach in the monthly horizon either
            out["stress_date_exact"] = None
            out["stress_week_label"] = "No cliff"
            out["runway_weeks"] = None
            out["runway_months"] = None
            out["stress_source"] = "monthly_scenario"
            return out

        first_idx = int(breach.idxmax())
        stress_month = pd.to_datetime(m.loc[first_idx, "month_date"]).to_period("M").to_timestamp()
        stress_date_exact = stress_month  # month-start approximation (conservative)

        runway_weeks = int(np.ceil((stress_date_exact - today).days / 7.0))
        runway_weeks = max(runway_weeks, 0)
        runway_months = float(runway_weeks) / 4.345  # approx

        out["stress_date_exact"] = stress_date_exact
        out["stress_week_label"] = f"Week of {stress_date_exact.strftime('%d %b %Y')} (est.)"
        out["runway_weeks"] = runway_weeks
        out["runway_months"] = runway_months
        out["stress_source"] = "monthly_scenario_estimate"
        return out

    # -------------------------
    # Monthly scenario (same as your existing logic)
    # -------------------------
    if base_df is None or base_df.empty:
        return pd.DataFrame(), {}, {}

    df = base_df.copy().sort_values("month_date").reset_index(drop=True)

    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce")
    df = df[df["month_date"].notna()].copy()
    df["month_date"] = df["month_date"].dt.to_period("M").dt.to_timestamp()

    for c in ["operating_cf", "investing_cf", "financing_cf", "free_cash_flow", "closing_cash"]:
        if c not in df.columns:
            df[c] = 0.0
        df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0.0)

    start_month_ts = _safe_month_start(start_month)
    mask = df["month_date"] >= start_month_ts

    rev_change_pct = float(controls.get("rev_change_pct", 0.0) or 0.0)
    collections_boost_pct = float(controls.get("collections_boost_pct", 0.0) or 0.0)
    extra_payroll_per_month = float(controls.get("extra_payroll_per_month", 0.0) or 0.0)
    spend_change_pct = float(controls.get("spend_change_pct", 0.0) or 0.0)
    extra_capex_one_off = float(controls.get("extra_capex_one_off", 0.0) or 0.0)
    extra_capex_recurring = float(controls.get("extra_capex_recurring", 0.0) or 0.0)
    equity_raise = float(controls.get("equity_raise", 0.0) or 0.0)
    recurring_funding = float(controls.get("recurring_funding", 0.0) or 0.0)

    # Operating
    op_base = df["operating_cf"].astype(float)
    pos = op_base.clip(lower=0.0)
    neg = op_base.where(op_base < 0.0, 0.0)

    pos_s = pos.copy()
    if rev_change_pct != 0.0:
        pos_s[mask] = pos_s[mask] * (1.0 + rev_change_pct / 100.0)

    neg_s = neg.copy()
    if spend_change_pct != 0.0:
        neg_s[mask] = neg_s[mask] * (1.0 + spend_change_pct / 100.0)

    op_s = pos_s + neg_s
    if collections_boost_pct != 0.0:
        op_s = op_s * (1.0 + collections_boost_pct / 100.0)

    if extra_payroll_per_month != 0.0:
        op_s[mask] = op_s[mask] - abs(extra_payroll_per_month)

    df["operating_cf_scenario"] = pd.to_numeric(op_s, errors="coerce").fillna(0.0)

    # Investing
    inv_s = df["investing_cf"].astype(float).copy()
    if extra_capex_recurring != 0.0:
        inv_s[mask] = inv_s[mask] + extra_capex_recurring
    if extra_capex_one_off != 0.0:
        inv_s[df["month_date"] == start_month_ts] = inv_s[df["month_date"] == start_month_ts] + extra_capex_one_off
    df["investing_cf_scenario"] = pd.to_numeric(inv_s, errors="coerce").fillna(0.0)

    # Financing
    fin_s = df["financing_cf"].astype(float).copy()
    if recurring_funding != 0.0:
        fin_s[mask] = fin_s[mask] + recurring_funding
    if equity_raise != 0.0:
        fin_s[df["month_date"] == start_month_ts] = fin_s[df["month_date"] == start_month_ts] + equity_raise
    df["financing_cf_scenario"] = pd.to_numeric(fin_s, errors="coerce").fillna(0.0)

    df["free_cash_flow_scenario"] = df["operating_cf_scenario"] + df["investing_cf_scenario"] + df["financing_cf_scenario"]

    opening_0 = float(df.iloc[0]["closing_cash"]) - float(df.iloc[0]["free_cash_flow"])
    closing_vals = []
    running = opening_0
    for _, r in df.iterrows():
        running = running + float(r["free_cash_flow_scenario"])
        closing_vals.append(running)
    df["closing_cash_scenario"] = pd.to_numeric(closing_vals, errors="coerce")

    df["delta_free_cf"] = df["free_cash_flow_scenario"] - df["free_cash_flow"]

    # -------------------------
    # Weekly scenario
    # -------------------------
    today = today or pd.Timestamp.utcnow().tz_localize(None)

    week_df_base = week_df_base if week_df_base is not None else pd.DataFrame()
    if week_df_base.empty:
        # still return something sensible
        return pd.DataFrame(), {}, {}

    base_week = week_df_base.copy()
    # IMPORTANT: build_14_week_cash_table returns week_start as ISO string
    base_week["week_start"] = pd.to_datetime(base_week.get("week_start"), errors="coerce")
    base_week = base_week.dropna(subset=["week_start"]).sort_values("week_start").reset_index(drop=True)
    base_week["closing_cash"] = pd.to_numeric(base_week.get("closing_cash"), errors="coerce").fillna(0.0)

    # Apply deltas to weekly closing cash
    week_scen = apply_monthly_scenario_to_weekly_fn(
        week_df=base_week,
        monthly_df=df[["month_date", "delta_free_cf"]],
        month_col="month_date",
        monthly_delta_col="delta_free_cf",
    )

    # ✅ Ensure scenario KPI input uses scenario cash, not base cash
    scen_week_for_kpis = week_scen.copy()
    scen_week_for_kpis["week_start"] = pd.to_datetime(scen_week_for_kpis.get("week_start"), errors="coerce")
    scen_week_for_kpis = scen_week_for_kpis.dropna(subset=["week_start"]).sort_values("week_start").reset_index(drop=True)

    if "closing_cash_scenario" in scen_week_for_kpis.columns:
        scen_week_for_kpis["closing_cash"] = pd.to_numeric(
            scen_week_for_kpis["closing_cash_scenario"], errors="coerce"
        ).fillna(0.0)
    else:
        # fallback (shouldn't happen)
        scen_week_for_kpis["closing_cash"] = pd.to_numeric(
            scen_week_for_kpis.get("closing_cash"), errors="coerce"
        ).fillna(0.0)

    # -------------------------
    # Stress KPIs (Base + Scenario)
    # -------------------------
    base_stress = compute_cash_stress_kpis_fn(
        week_df=base_week,
        cash_buffer=float(cash_buffer or 0.0),
        today=today,
        client_id=client_id,
        as_of_month=as_of_month,
        fetch_cashflow_summary_for_client_as_of=fetch_cashflow_summary_for_client_as_of,
    )

    scen_stress = compute_cash_stress_kpis_fn(
        week_df=scen_week_for_kpis,
        cash_buffer=float(cash_buffer or 0.0),
        today=today,
        client_id=client_id,
        as_of_month=as_of_month,
        fetch_cashflow_summary_for_client_as_of=fetch_cashflow_summary_for_client_as_of,
    )

    # ✅ If no breach inside 14w, but monthly scenario shows breach later, update runway/cliff
    # (This is the missing piece that makes runway move when improvements push breach beyond 14 weeks)
    if float(cash_buffer or 0.0) > 0:
        scen_has_weekly_breach = False
        try:
            scen_has_weekly_breach = (scen_week_for_kpis["closing_cash"] <= float(cash_buffer)).any()
        except Exception:
            scen_has_weekly_breach = False

        if not scen_has_weekly_breach:
            fallback = _fallback_stress_from_monthly(df, float(cash_buffer), today)
            # only override if fallback gives a better answer
            if fallback.get("stress_week_label") is not None:
                for k, v in fallback.items():
                    scen_stress[k] = v

    return week_scen, base_stress, scen_stress


def render_scenarios_what_if(*, selected_client_id, selected_month_start):
    st.markdown("---")
    st.subheader("🧪 Scenarios & What-if")

    st.markdown(
        "Answer simple questions like **“Can I hire now?”**, "
        "**“What if revenue grows 20%?”**, or **“When do I need to raise?”** "
        "without touching your actual data."
    )

    # ✅ default return value (so caller always gets a list)
    scenario_summaries: list[dict] = []

    if not selected_client_id or not selected_month_start:
        st.info("Pick a business and focus month at the top before running scenarios.")
        st.session_state["scenario_summaries"] = scenario_summaries
        return scenario_summaries

    currency_code, currency_symbol = get_client_currency(selected_client_id)

    base_df = get_engine_window_for_scenarios(
        selected_client_id,
        base_month=selected_month_start,
        n_months=12,
    )

    if base_df is None or base_df.empty:
        st.warning(
            "No cashflow engine data yet for this business. "
            "Go to **Business at a Glance → Rebuild cashflow for next 12 months** first."
        )
        st.session_state["scenario_summaries"] = scenario_summaries
        return scenario_summaries

    base_df = base_df.sort_values("month_date").reset_index(drop=True)
    base_df["month_date"] = pd.to_datetime(base_df["month_date"], errors="coerce")
    base_df = base_df[base_df["month_date"].notna()].copy()
    base_df["month_date"] = base_df["month_date"].dt.to_period("M").dt.to_timestamp()

    month_opts = base_df["month_date"].dt.to_period("M").drop_duplicates()
    month_label_map = {p.strftime("%b %Y"): p.to_timestamp() for p in month_opts}
    month_labels = list(month_label_map.keys())

    st.markdown("---")
    st.subheader("🎛 Choose your scenario")

    col_output, col_controls = st.columns([2, 1])

    # ----------------------------
    # RIGHT: CONTROLS + SAVE
    # ----------------------------
    with col_controls:
        st.markdown("#### Scenario type")

        preset = st.selectbox(
            "Quick scenario preset",
            options=[
                "Custom",
                "Revenue push (+20%)",
                "Hire 1 key role (+$15k/month)",
                "Cut supplier/operating spend by 20%",
                "Raise $250k now",
                "Raise $250k + monthly top-up",
            ],
            help="Presets override the sliders below when you run the scenario.",
            key="scen_preset",
        )

        controls = {
            "rev_change_pct": 0.0,
            "collections_boost_pct": 0.0,
            "extra_payroll_per_month": 0.0,
            "spend_change_pct": 0.0,
            "extra_capex_one_off": 0.0,
            "extra_capex_recurring": 0.0,
            "equity_raise": 0.0,
            "recurring_funding": 0.0,
        }

        start_month_label = st.selectbox(
            "Scenario starts from month",
            options=month_labels,
            key="scen_start_month",
        )
        start_month_ts = month_label_map[start_month_label]

        widgets_disabled = preset != "Custom"

        tab_rev, tab_hire, tab_spend, tab_capex, tab_fund = st.tabs(
            ["Revenue & customers", "Hiring & team", "Bills & spend", "Capex & assets", "Funding & runway"]
        )

        with tab_rev:
            st.caption("Change customer cash-in (simple multiplier).")
            rev_change_pct = st.slider(
                "Change in cash from customers (%)",
                min_value=-50, max_value=100, value=0, step=5,
                disabled=widgets_disabled,
                key="scen_rev_change",
            )
            collections_boost_pct = st.slider(
                "Collections speed effect (%)",
                min_value=-20, max_value=20, value=0, step=5,
                disabled=widgets_disabled,
                key="scen_collections_boost",
            )
            if preset == "Custom":
                controls["rev_change_pct"] = float(rev_change_pct)
                controls["collections_boost_pct"] = float(collections_boost_pct)

        with tab_hire:
            st.caption("Extra payroll (spread across weeks in each month).")
            extra_headcount_cost = st.number_input(
                f"Extra monthly payroll from that month ({currency_code})",
                value=0.0, step=1000.0,
                disabled=widgets_disabled,
                key="scen_extra_payroll",
            )
            if preset == "Custom":
                controls["extra_payroll_per_month"] = float(extra_headcount_cost)

        with tab_spend:
            st.caption("Adjust operating spend (opex).")
            spend_change_pct = st.slider(
                "Change in operating spend (opex) (%)",
                min_value=-50, max_value=50, value=0, step=5,
                disabled=widgets_disabled,
                key="scen_spend_change",
            )
            if preset == "Custom":
                controls["spend_change_pct"] = float(spend_change_pct)

        with tab_capex:
            st.caption("Adjust investing cashflow (usually negative for purchases).")
            extra_capex_one_off = st.number_input(
                f"One-off capex cash impact in start month ({currency_code})",
                value=0.0, step=5000.0,
                disabled=widgets_disabled,
                key="scen_capex_oneoff",
            )
            extra_capex_recurring = st.number_input(
                f"Recurring monthly capex cash impact from start ({currency_code})",
                value=0.0, step=1000.0,
                disabled=widgets_disabled,
                key="scen_capex_recurring",
            )
            if preset == "Custom":
                controls["extra_capex_one_off"] = float(extra_capex_one_off)
                controls["extra_capex_recurring"] = float(extra_capex_recurring)

        with tab_fund:
            st.caption("Adjust financing cashflow (positive = cash-in).")
            equity_raise = st.number_input(
                f"One-off equity/loan cash-in in start month ({currency_code})",
                value=0.0, step=25000.0,
                disabled=widgets_disabled,
                key="scen_equity_raise",
            )
            recurring_funding = st.number_input(
                f"Recurring monthly funding cash-in from start ({currency_code})",
                value=0.0, step=5000.0,
                disabled=widgets_disabled,
                key="scen_recurring_funding",
            )
            if preset == "Custom":
                controls["equity_raise"] = float(equity_raise)
                controls["recurring_funding"] = float(recurring_funding)

        # Apply preset overrides
        final_controls = controls.copy()
        if preset == "Revenue push (+20%)":
            final_controls["rev_change_pct"] = 20.0
        elif preset == "Hire 1 key role (+$15k/month)":
            final_controls["extra_payroll_per_month"] = 15_000.0
        elif preset == "Cut supplier/operating spend by 20%":
            final_controls["spend_change_pct"] = -20.0
        elif preset == "Raise $250k now":
            final_controls["equity_raise"] = 250_000.0
        elif preset == "Raise $250k + monthly top-up":
            final_controls["equity_raise"] = 250_000.0
            final_controls["recurring_funding"] = 10_000.0

        st.markdown("—")
        st.markdown("#### 💾 Save this scenario")

        scen_name = st.text_input(
            "Scenario name",
            value=f"{preset} @ {start_month_label}",
            key="scen_name",
        )
        scen_desc = st.text_area(
            "Short note (optional)",
            value="",
            key="scen_desc",
        )

        save_clicked = st.button("Save scenario", key="save_scenario_btn")
        run_clicked = st.button("Run this scenario", key="run_scenario_btn")

        if save_clicked:
            if not scen_name.strip():
                st.warning("Give this scenario a name before saving.")
            else:
                ok = save_scenario_for_client(
                    client_id=selected_client_id,
                    name=scen_name,
                    controls=final_controls,
                    base_month=selected_month_start,
                    description=scen_desc,
                )
                if ok:
                    st.success("Scenario saved.")
                    st.rerun()
                else:
                    st.error("Could not save scenario. Check logs.")

    # ----------------------------
    # LEFT: SAVED SCENARIOS + OUTPUT
    # ----------------------------
    with col_output:
        saved_scenarios = fetch_scenarios_for_client(selected_client_id) or []

        run_saved = False
        scenario_label = f"{preset} @ {start_month_label}"
        chosen_controls_from_saved = None
        chosen_base_month = None
        chosen_id = None

        if saved_scenarios:
            st.markdown("#### 📂 Run a saved scenario")

            labels = []
            id_map = {}
            for s in saved_scenarios:
                nm = s.get("name") or "Scenario"
                created = (s.get("created_at") or "")[:10]
                lab = f"{nm} (saved {created})"
                labels.append(lab)
                id_map[lab] = s.get("id")

            sel_label = st.selectbox(
                "Saved scenarios",
                options=["(None)"] + labels,
                key="saved_scenario_select",
            )

            colA, colB = st.columns([1, 1])
            with colA:
                run_saved = st.button("Run saved scenario", key="run_saved_btn")
            with colB:
                delete_saved = st.button("Delete selected", key="delete_saved_btn")

            if sel_label != "(None)":
                chosen_id = id_map.get(sel_label)
                scen_row = fetch_scenario_by_id(chosen_id) if chosen_id is not None else None
                if scen_row:
                    chosen_controls_from_saved = scen_row.get("controls") or {}
                    chosen_base_month = scen_row.get("base_month")
                    scenario_label = f"Saved: {scen_row.get('name') or 'Scenario'}"

            if sel_label != "(None)" and delete_saved and chosen_id is not None:
                if delete_scenario_by_id(chosen_id):
                    st.success("Scenario deleted.")
                    st.rerun()
                else:
                    st.error("Could not delete scenario. Check logs.")

        # Decide which scenario to run
        use_controls = final_controls
        use_start_month_ts = start_month_ts

        if run_saved and chosen_controls_from_saved:
            use_controls = chosen_controls_from_saved
            if chosen_base_month:
                try:
                    use_start_month_ts = pd.to_datetime(chosen_base_month).to_period("M").to_timestamp()
                except Exception:
                    use_start_month_ts = start_month_ts

        # Gate: must have at least one change
        try:
            has_any_change = any(abs(float(v)) > 0 for v in (use_controls or {}).values())
        except Exception:
            has_any_change = False

        # ✅ Build weekly base once (also used to compute saved scenario summaries)
        week_df_base = build_14_week_cash_table(selected_client_id, selected_month_start)
        if week_df_base is None or week_df_base.empty:
            st.warning("14-week cash table not available yet — cannot run runway/cliff scenario.")
            st.session_state["scenario_summaries"] = scenario_summaries
            return scenario_summaries

        settings = get_client_settings(selected_client_id) or {}
        cash_buffer = float(settings.get("min_cash_buffer") or 0.0)

        # ✅ NEW: compute summaries for up to 2 saved scenarios (for Board narrative)
        def _compute_saved_summaries(max_n: int = 2) -> list[dict]:
            if not saved_scenarios:
                return []

            # Prefer most recent by created_at if available
            def _sort_key(row):
                return str(row.get("created_at") or "")
            sorted_rows = sorted(saved_scenarios, key=_sort_key, reverse=True)

            out: list[dict] = []
            for row in sorted_rows:
                if len(out) >= max_n:
                    break

                sid = row.get("id")
                nm = (row.get("name") or "").strip() or "Scenario"
                scen_row = fetch_scenario_by_id(sid) if sid is not None else None
                if not scen_row:
                    continue

                ctrls = scen_row.get("controls") or {}
                base_m = scen_row.get("base_month")

                # Determine scenario start month
                use_ts = start_month_ts
                if base_m:
                    try:
                        use_ts = pd.to_datetime(base_m).to_period("M").to_timestamp()
                    except Exception:
                        use_ts = start_month_ts

                try:
                    _, base_stress, scen_stress = run_founder_scenario(
                        base_df=base_df,
                        start_month=use_ts,
                        controls=ctrls,
                        week_df_base=week_df_base,
                        cash_buffer=cash_buffer,
                        compute_cash_stress_kpis_fn=compute_cash_stress_kpis,
                        apply_monthly_scenario_to_weekly_fn=apply_monthly_scenario_to_weekly,
                        today=pd.to_datetime(selected_month_start).to_period("M").to_timestamp(),
                        client_id=selected_client_id,
                        as_of_month=selected_month_start,
                        fetch_cashflow_summary_for_client_as_of=fetch_cashflow_summary_for_client_as_of,
                    )
                except Exception:
                    continue

                out.append({
                    "name": nm,
                    "runway_months": (scen_stress or {}).get("runway_months"),
                    "stress_date_exact": (scen_stress or {}).get("stress_date_exact"),
                })

            return out

        # Cache in session_state to avoid re-running if nothing changed
        cache_key = (
            str(selected_client_id),
            str(pd.to_datetime(selected_month_start).to_period("M")),
            float(cash_buffer),
            tuple((s.get("id"), str(s.get("created_at") or "")) for s in (saved_scenarios or [])),
        )
        prev_key = st.session_state.get("_scenario_summaries_cache_key")
        if prev_key != cache_key:
            scenario_summaries = _compute_saved_summaries(max_n=2)
            st.session_state["_scenario_summaries_cache_key"] = cache_key
            st.session_state["scenario_summaries"] = scenario_summaries
        else:
            scenario_summaries = st.session_state.get("scenario_summaries", []) or []

        # If user hasn't run anything, we still return summaries for board narrative
        if not has_any_change and not run_saved:
            st.info("Pick a preset or change at least one control, then click **Run this scenario**.")
            return scenario_summaries

        if not (run_clicked or run_saved):
            st.info("Adjust controls or pick a saved scenario, then click **Run**.")
            return scenario_summaries

        # Run scenario (existing behavior)
        scen_week_df, base_stress, scen_stress = run_founder_scenario(
            base_df=base_df,
            start_month=use_start_month_ts,
            controls=use_controls,
            week_df_base=week_df_base,
            cash_buffer=cash_buffer,
            compute_cash_stress_kpis_fn=compute_cash_stress_kpis,
            apply_monthly_scenario_to_weekly_fn=apply_monthly_scenario_to_weekly,
            today=pd.to_datetime(selected_month_start).to_period("M").to_timestamp(),
            client_id=selected_client_id,
            as_of_month=selected_month_start,
            fetch_cashflow_summary_for_client_as_of=fetch_cashflow_summary_for_client_as_of,
        )

        st.caption(f"Scenario run: **{scenario_label}**")

        st.subheader("📌 Base case vs Scenario (aligned to Cash & Bills)")

        summary_tbl = build_base_vs_scenario_summary_table(
            currency_symbol=currency_symbol,
            base_week_df=week_df_base,
            scen_week_df=scen_week_df,
            base_stress_kpis=base_stress,
            scen_stress_kpis=scen_stress,
        )
        st.dataframe(summary_tbl, width="stretch")

        with st.expander("View scenario 14-week cash table", expanded=False):
            if scen_week_df is None or scen_week_df.empty:
                st.caption("Scenario weekly table not available.")
            else:
                view = scen_week_df.copy()
                view["Week"] = pd.to_datetime(view["week_start"], errors="coerce").dt.strftime("%d %b %Y")

                cols = [
                    "Week",
                    "cash_in_ar_scenario",
                    "cash_in_other_income_scenario",
                    "cash_out_ap_scenario",
                    "cash_out_payroll_scenario",
                    "cash_out_opex_scenario",
                    "operating_cf_scenario",
                    "investing_cf_scenario",
                    "financing_cf_scenario",
                    "net_cash_scenario",
                    "closing_cash_scenario",
                ]
                cols = [c for c in cols if c in view.columns]
                if not cols:
                    cols = ["Week", "closing_cash", "closing_cash_scenario"]
                    cols = [c for c in cols if c in view.columns]

                st.dataframe(view[cols], width="stretch")

    # ✅ MUST return summaries so Funding page can use them
    return scenario_summaries


def _to_month_start(x) -> pd.Timestamp:
    ts = pd.to_datetime(x, errors="coerce")
    if ts is pd.NaT:
        return pd.NaT
    return ts.to_period("M").to_timestamp()

def _build_weeks_per_month_map(week_df: pd.DataFrame) -> dict[pd.Timestamp, int]:
    w = week_df.copy()
    w["week_start"] = pd.to_datetime(w["week_start"], errors="coerce")
    w = w.dropna(subset=["week_start"]).sort_values("week_start")
    w["month_ts"] = w["week_start"].dt.to_period("M").dt.to_timestamp()
    return w.groupby("month_ts").size().to_dict()

def apply_scenario_controls_to_weekly_cash_table(
    week_df_base: pd.DataFrame,
    *,
    start_month: pd.Timestamp,
    controls: dict,
) -> pd.DataFrame:
    """
    Takes the *base* 14-week cash table (from build_14_week_cash_table)
    and produces a scenario version by adjusting weekly components, then
    recomputing running opening/closing cash.

    Returns week_df_scenario with:
      - all base columns
      - *_scenario component columns
      - operating_cf_scenario, investing_cf_scenario, financing_cf_scenario
      - net_cash_scenario, opening_cash_scenario, closing_cash_scenario
    """
    if week_df_base is None or week_df_base.empty:
        return pd.DataFrame()

    w = week_df_base.copy()
    w["week_start"] = pd.to_datetime(w["week_start"], errors="coerce")
    w = w.dropna(subset=["week_start"]).sort_values("week_start").reset_index(drop=True)

    # Ensure numeric for required columns
    numeric_cols = [
        "opening_cash", "closing_cash",
        "cash_in_ar", "cash_in_other_income",
        "cash_out_ap", "cash_out_payroll", "cash_out_opex",
        "operating_cf", "investing_cf", "financing_cf", "net_cash",
    ]
    for c in numeric_cols:
        if c not in w.columns:
            w[c] = 0.0
        w[c] = pd.to_numeric(w[c], errors="coerce").fillna(0.0)

    start_month_ts = _to_month_start(start_month)
    if start_month_ts is pd.NaT:
        start_month_ts = w["week_start"].min().to_period("M").to_timestamp()

    # Week->month mapping
    w["month_ts"] = w["week_start"].dt.to_period("M").dt.to_timestamp()
    mask = w["month_ts"] >= start_month_ts

    # ---- controls ----
    rev_change_pct = float(controls.get("rev_change_pct", 0.0) or 0.0)
    collections_boost_pct = float(controls.get("collections_boost_pct", 0.0) or 0.0)
    extra_payroll_per_month = float(controls.get("extra_payroll_per_month", 0.0) or 0.0)
    spend_change_pct = float(controls.get("spend_change_pct", 0.0) or 0.0)
    extra_capex_one_off = float(controls.get("extra_capex_one_off", 0.0) or 0.0)
    extra_capex_recurring = float(controls.get("extra_capex_recurring", 0.0) or 0.0)
    equity_raise = float(controls.get("equity_raise", 0.0) or 0.0)
    recurring_funding = float(controls.get("recurring_funding", 0.0) or 0.0)

    # Allocate monthly “one-off / recurring” across weeks of each month in horizon
    weeks_per_month = _build_weeks_per_month_map(w)

    def _per_week(month_ts: pd.Timestamp, monthly_amount: float) -> float:
        n = max(int(weeks_per_month.get(month_ts, 0) or 0), 1)
        return float(monthly_amount) / float(n)

    # ----------------------------
    # Scenario component columns
    # ----------------------------
    # 1) AR cash-in adjustment: revenue + collections (simple multiplier)
    #    (If you later want “collections timing shift”, you’ll need invoice-level rescheduling.)
    ar_factor = (1.0 + rev_change_pct / 100.0) * (1.0 + collections_boost_pct / 100.0)
    w["cash_in_ar_scenario"] = w["cash_in_ar"]
    if ar_factor != 1.0:
        w.loc[mask, "cash_in_ar_scenario"] = w.loc[mask, "cash_in_ar"] * ar_factor

    # 2) Opex adjustment
    spend_factor = (1.0 + spend_change_pct / 100.0)
    w["cash_out_opex_scenario"] = w["cash_out_opex"]
    if spend_factor != 1.0:
        w.loc[mask, "cash_out_opex_scenario"] = w.loc[mask, "cash_out_opex"] * spend_factor

    # 3) Payroll adjustment (monthly amount spread across weeks)
    w["cash_out_payroll_scenario"] = w["cash_out_payroll"]
    if extra_payroll_per_month != 0.0:
        for i in range(len(w)):
            if w.loc[i, "month_ts"] >= start_month_ts:
                w.loc[i, "cash_out_payroll_scenario"] += _per_week(w.loc[i, "month_ts"], abs(extra_payroll_per_month))

    # 4) Investing adjustment (monthly recurring + one-off in start month)
    w["investing_cf_scenario"] = w["investing_cf"]
    if extra_capex_recurring != 0.0:
        for i in range(len(w)):
            if w.loc[i, "month_ts"] >= start_month_ts:
                w.loc[i, "investing_cf_scenario"] += _per_week(w.loc[i, "month_ts"], extra_capex_recurring)

    if extra_capex_one_off != 0.0:
        # allocate one-off only within start_month weeks
        for i in range(len(w)):
            if w.loc[i, "month_ts"] == start_month_ts:
                w.loc[i, "investing_cf_scenario"] += _per_week(start_month_ts, extra_capex_one_off)

    # 5) Financing adjustment (monthly recurring + one-off in start month)
    w["financing_cf_scenario"] = w["financing_cf"]
    if recurring_funding != 0.0:
        for i in range(len(w)):
            if w.loc[i, "month_ts"] >= start_month_ts:
                w.loc[i, "financing_cf_scenario"] += _per_week(w.loc[i, "month_ts"], recurring_funding)

    if equity_raise != 0.0:
        for i in range(len(w)):
            if w.loc[i, "month_ts"] == start_month_ts:
                w.loc[i, "financing_cf_scenario"] += _per_week(start_month_ts, equity_raise)

    # Keep unchanged components
    w["cash_in_other_income_scenario"] = w["cash_in_other_income"]
    w["cash_out_ap_scenario"] = w["cash_out_ap"]

    # Recompute operating CF scenario from components (same as build_14_week_cash_table logic)
    w["operating_cf_scenario"] = (
        w["cash_in_ar_scenario"]
        + w["cash_in_other_income_scenario"]
        - w["cash_out_ap_scenario"]
        - w["cash_out_payroll_scenario"]
        - w["cash_out_opex_scenario"]
    )

    # net cash scenario and running balance
    w["net_cash_scenario"] = (
        w["operating_cf_scenario"] + w["investing_cf_scenario"] + w["financing_cf_scenario"]
    )

    opening0 = float(w.loc[0, "opening_cash"])
    open_list, close_list = [], []
    running_open = opening0
    for i in range(len(w)):
        open_list.append(float(running_open))
        close = float(running_open) + float(w.loc[i, "net_cash_scenario"])
        close_list.append(close)
        running_open = close

    w["opening_cash_scenario"] = open_list
    w["closing_cash_scenario"] = close_list
    w["cash_danger_flag_scenario"] = (pd.to_numeric(w["closing_cash_scenario"], errors="coerce").fillna(0.0) <= 0.0)

    return w

def build_base_vs_scenario_summary_table(
    *,
    currency_symbol: str,
    base_week_df: pd.DataFrame,
    scen_week_df: pd.DataFrame,
    base_stress_kpis: dict,
    scen_stress_kpis: dict,
) -> pd.DataFrame:
    import numpy as np
    import pandas as pd

    def _money(x):
        if x is None or (isinstance(x, float) and np.isnan(x)):
            return "—"
        return f"{currency_symbol}{float(x):,.0f}"

    def _months(x):
        return "—" if x is None else f"{float(x):.1f} months"

    def _weeks(x):
        return "—" if x is None else f"{int(x)} weeks"

    base_runway_m = base_stress_kpis.get("runway_months")
    base_runway_w = base_stress_kpis.get("runway_weeks")
    scen_runway_m = scen_stress_kpis.get("runway_months")
    scen_runway_w = scen_stress_kpis.get("runway_weeks")

    base_cliff = base_stress_kpis.get("stress_week_label") or "No cliff"
    scen_cliff = scen_stress_kpis.get("stress_week_label") or "No cliff"

    base_lowest = (
        float(pd.to_numeric(base_week_df["closing_cash"], errors="coerce").min())
        if base_week_df is not None and not base_week_df.empty and "closing_cash" in base_week_df.columns
        else np.nan
    )

    scen_lowest = (
        float(pd.to_numeric(scen_week_df["closing_cash_scenario"], errors="coerce").min())
        if scen_week_df is not None and not scen_week_df.empty and "closing_cash_scenario" in scen_week_df.columns
        else np.nan
    )

    base_burn = compute_effective_burn_from_weekly(
        base_week_df,
        net_source_cols=["net_cash"],
        window_weeks=13
    ) if base_week_df is not None and not base_week_df.empty else None

    scen_burn = compute_effective_burn_from_weekly(
        scen_week_df,
        net_source_cols=["net_cash_scenario", "net_cash"],
        window_weeks=13
    ) if scen_week_df is not None and not scen_week_df.empty else None

    rows = [
        ("Runway (to buffer breach)", f"{_months(base_runway_m)} ({_weeks(base_runway_w)})", f"{_months(scen_runway_m)} ({_weeks(scen_runway_w)})"),
        ("Cash cliff week (buffer breach)", base_cliff, scen_cliff),
        ("Lowest cash (next 14 weeks)", _money(base_lowest), _money(scen_lowest)),
        ("Net burn / month (3-mo avg from weekly)", _money(base_burn), _money(scen_burn)),
    ]
    return pd.DataFrame(rows, columns=["Metric", "Base case", "Scenario"])


def compute_effective_burn_from_weekly(
    week_df: pd.DataFrame,
    *,
    anchor_week: pd.Timestamp | None = None,
    net_source_cols: list[str] | None = None,
    window_weeks: int = 13,              # ~3 months
    weeks_per_month: float = 52.0 / 12.0 # 4.333...
) -> float | None:
    """
    Weekly version of compute_runway_and_effective_burn_from_df().

    Returns:
      effective_burn_monthly (positive number) or None

    Burn logic:
      1) Take ~3 months forward avg of a cashflow column (weekly net cash)
      2) If avg is negative => burn = -avg * weeks_per_month
      3) Fallback: slope of closing_cash over the same window
    """
    if week_df is None or week_df.empty:
        return None

    w = week_df.copy()

    # --- normalize dates ---
    if "week_start" in w.columns:
        w["week_start"] = pd.to_datetime(w["week_start"], errors="coerce")
        w = w[w["week_start"].notna()].sort_values("week_start").reset_index(drop=True)

    if w.empty:
        return None

    # --- choose anchor: default first row (matches your Cash & Bills anchor logic) ---
    if anchor_week is not None and "week_start" in w.columns:
        anchor_week = pd.to_datetime(anchor_week, errors="coerce")
        if anchor_week is not pd.NaT:
            w = w[w["week_start"] >= anchor_week].reset_index(drop=True)

    if w.empty:
        return None

    # --- pick burn source column ---
    if net_source_cols is None:
        net_source_cols = [
            "net_cash_scenario",
            "net_cash",
            "operating_cf_scenario",
            "operating_cf",
            "free_cash_flow_scenario",
            "free_cash_flow",
        ]
    burn_col = next((c for c in net_source_cols if c in w.columns), None)

    # --- 1) 3-month forward avg on chosen column ---
    effective_burn_monthly = None

    window = w.head(window_weeks).copy()

    if burn_col is not None:
        vals = pd.to_numeric(window[burn_col], errors="coerce")
        avg_cf = vals.mean()
        if pd.notna(avg_cf) and avg_cf < 0:
            effective_burn_monthly = float(-avg_cf) * float(weeks_per_month)

    # --- 2) fallback: slope of closing_cash ---
    if effective_burn_monthly is None and "closing_cash" in window.columns:
        cc = pd.to_numeric(window["closing_cash"], errors="coerce")
        delta = cc.diff().mean()  # negative => falling cash
        if pd.notna(delta) and delta < 0:
            effective_burn_monthly = float(-delta) * float(weeks_per_month)

    return effective_burn_monthly


def recompute_monthly_cash_path_from_actual(base_df: pd.DataFrame, opening_cash_actual: float) -> pd.DataFrame:
    df = base_df.copy().sort_values("month_date").reset_index(drop=True)
    df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()

    # ensure net_cash exists
    if "net_cash" not in df.columns or df["net_cash"].isna().all():
        df["net_cash"] = (
            pd.to_numeric(df.get("operating_cf", 0), errors="coerce").fillna(0.0)
            + pd.to_numeric(df.get("investing_cf", 0), errors="coerce").fillna(0.0)
            + pd.to_numeric(df.get("financing_cf", 0), errors="coerce").fillna(0.0)
        )

    opens, closes = [], []
    o = float(opening_cash_actual)
    for _, r in df.iterrows():
        opens.append(o)
        c = o + float(r["net_cash"] or 0.0)
        closes.append(c)
        o = c

    df["opening_cash_sim"] = opens
    df["closing_cash_sim"] = closes
    return df


def get_monthly_targets_for_month(client_id, month_ts: pd.Timestamp) -> dict:
    """
    Fetch one row from client_monthly_targets for the given month (month start).
    Returns {} if not found.
    """
    if not client_id or month_ts is None:
        return {}

    m = pd.to_datetime(month_ts).to_period("M").to_timestamp().date().isoformat()

    try:
        res = (
            supabase
            .table("client_monthly_targets")
            .select("*")
            .eq("client_id", str(client_id))
            .eq("month_date", m)
            .limit(1)
            .execute()
        )
        rows = res.data or []
        return rows[0] if rows else {}
    except Exception:
        return {}



def compute_headcount_month_split(df_positions: pd.DataFrame, month_ts: pd.Timestamp) -> dict:
    """
    Headcount split for the focus month.
    Definitions:
      - headcount_fte_eom: FTE active as at END of month (exits before month end removed)
      - hires_fte_in_month: start_date in [month_start, month_end]
      - exits_fte_in_month: end_date in [month_start, month_end)  (strictly before month_end)
    """
    out = {
        "headcount_fte_eom": 0.0,
        "hires_fte_in_month": 0.0,
        "exits_fte_in_month": 0.0,
        "net_change_fte_in_month": 0.0,
    }

    if df_positions is None or df_positions.empty or month_ts is None:
        return out

    pos = df_positions.copy()
    pos["start_date"] = pd.to_datetime(pos.get("start_date"), errors="coerce", dayfirst=True)
    pos["end_date"] = pd.to_datetime(pos.get("end_date"), errors="coerce", dayfirst=True)

    if "fte" not in pos.columns:
        pos["fte"] = 0.0
    pos["fte"] = pd.to_numeric(pos["fte"], errors="coerce").fillna(0.0)

    month_start = pd.to_datetime(month_ts).to_period("M").to_timestamp()
    month_end = month_start + pd.offsets.MonthEnd(0)  # last day of month (Timestamp)

    # ✅ EOM headcount: must still be employed at month_end
    eom_mask = (
        (pos["start_date"].notna()) &
        (pos["start_date"] <= month_end) &
        (pos["end_date"].isna() | (pos["end_date"] >= month_end))
    )
    out["headcount_fte_eom"] = float(pos.loc[eom_mask, "fte"].sum()) if eom_mask.any() else 0.0

    # ✅ Hires within focus month (inclusive)
    hires_mask = pos["start_date"].notna() & (pos["start_date"] >= month_start) & (pos["start_date"] <= month_end)
    out["hires_fte_in_month"] = float(pos.loc[hires_mask, "fte"].sum()) if hires_mask.any() else 0.0

    # ✅ Exits within focus month, but BEFORE month_end (strict)
    exits_mask = pos["end_date"].notna() & (pos["end_date"] >= month_start) & (pos["end_date"] < month_end)
    out["exits_fte_in_month"] = float(pos.loc[exits_mask, "fte"].sum()) if exits_mask.any() else 0.0

    out["net_change_fte_in_month"] = out["hires_fte_in_month"] - out["exits_fte_in_month"]
    return out





def compute_headcount_movement_last_30_days_basic(df_positions: pd.DataFrame, month_ts: pd.Timestamp) -> dict:
    """
    Movement in the 30 days leading up to month_ts (inclusive).
    - hires_fte_30d: sum FTE where start_date is within last 30 days
    - exits_fte_30d: sum FTE where end_date is within last 30 days
    - headcount_change_fte_30d = hires - exits
    """
    out = {"hires_fte_30d": 0.0, "exits_fte_30d": 0.0, "headcount_change_fte_30d": 0.0}

    if df_positions is None or df_positions.empty:
        return out

    pos = df_positions.copy()
    pos["start_date"] = pd.to_datetime(pos.get("start_date"), errors="coerce", dayfirst=True)
    pos["end_date"] = pd.to_datetime(pos.get("end_date"), errors="coerce", dayfirst=True)

    if "fte" not in pos.columns:
        pos["fte"] = 0.0
    pos["fte"] = pd.to_numeric(pos["fte"], errors="coerce").fillna(0.0)

    end_window = pd.to_datetime(month_ts)
    start_window = end_window - pd.Timedelta(days=30)

    hires_mask = pos["start_date"].notna() & (pos["start_date"] >= start_window) & (pos["start_date"] <= end_window)
    exits_mask = pos["end_date"].notna() & (pos["end_date"] >= start_window) & (pos["end_date"] <= end_window)

    hires = float(pos.loc[hires_mask, "fte"].sum()) if hires_mask.any() else 0.0
    exits = float(pos.loc[exits_mask, "fte"].sum()) if exits_mask.any() else 0.0

    out["hires_fte_30d"] = hires
    out["exits_fte_30d"] = exits
    out["headcount_change_fte_30d"] = hires - exits
    return out

import pandas as pd
import altair as alt
import streamlit as st


def compute_payroll_cash_components_by_month(client_id, month_index) -> pd.DataFrame:
    """
    Per-month payroll CASH components aligned to your existing payroll engine logic.

    CASH (paid month):
      - headcount_fte_eom: headcount as at END of month (exits before month end removed)
      - net_salary_cash: cash to employees this month (gross_work - payg_generated)
      - payg_cash: PAYG cash remitted this month (lagged by payroll_tax_lag_months)
      - super_cash: super cash paid this month (lagged by super_lag_months)
      - total_payroll_cash: net_salary_cash + payg_cash + super_cash (must match compute_payroll_by_month)

    GENERATED (work month):
      - gross_work
      - payg_generated
      - super_generated
      - workers_comp_generated (employer on-cost % of wages, always)
      - payroll_tax_generated (state payroll tax % of wages, ONLY if annualised wages (EOM) > threshold)
      - monthly_payroll_cost = gross_work + super_generated + workers_comp_generated + payroll_tax_generated
    """

    month_list = sorted({pd.to_datetime(m).to_period("M").to_timestamp() for m in month_index})
    base_cols = [
        "month_date",
        "headcount_fte_eom",
        "annualised_wages_eom",
        "payroll_tax_applicable",
        "gross_work", "payg_generated", "super_generated",
        "workers_comp_generated", "payroll_tax_generated", "monthly_payroll_cost",
        "net_salary_cash", "payg_cash", "super_cash", "total_payroll_cash",
        "engine_total", "diff_vs_engine",
    ]
    if not month_list:
        return pd.DataFrame(columns=base_cols)

    settings = get_client_settings(client_id) or {}

    # cash timing (existing)
    super_lag = int(settings.get("super_lag_months", 1) or 1)
    payroll_tax_lag = int(settings.get("payroll_tax_lag_months", 1) or 1)  # used as PAYG lag here
    lookback = max(super_lag, payroll_tax_lag, 0)

    # employer on-cost settings (NEW)
    payroll_tax_threshold_annual = float(settings.get("payroll_tax_threshold_annual", 1_300_000) or 1_300_000)
    payroll_tax_rate_pct = float(settings.get("payroll_tax_rate_pct", 0.0) or 0.0) / 100.0
    workers_comp_rate_pct = float(settings.get("workers_comp_rate_pct", 0.0) or 0.0) / 100.0

    df_positions = fetch_payroll_positions_for_client(client_id)
    if df_positions is None or df_positions.empty:
        df0 = pd.DataFrame({"month_date": month_list})
        for c in base_cols:
            if c not in df0.columns and c != "month_date":
                # bool column default to False, everything else 0.0
                df0[c] = False if c == "payroll_tax_applicable" else 0.0
        engine_totals = compute_payroll_by_month(client_id, month_list) or {}
        df0["engine_total"] = df0["month_date"].apply(lambda m: float(engine_totals.get(m, 0.0)))
        df0["diff_vs_engine"] = (df0["total_payroll_cash"] - df0["engine_total"]).abs()
        return df0[base_cols]

    pos = df_positions.copy()
    pos["start_date"] = pd.to_datetime(pos.get("start_date"), errors="coerce", dayfirst=True)
    pos["end_date"] = pd.to_datetime(pos.get("end_date"), errors="coerce", dayfirst=True)
    pos = pos[pos["start_date"].notna()]
    if pos.empty:
        df0 = pd.DataFrame({"month_date": month_list})
        for c in base_cols:
            if c not in df0.columns and c != "month_date":
                df0[c] = False if c == "payroll_tax_applicable" else 0.0
        engine_totals = compute_payroll_by_month(client_id, month_list) or {}
        df0["engine_total"] = df0["month_date"].apply(lambda m: float(engine_totals.get(m, 0.0)))
        df0["diff_vs_engine"] = (df0["total_payroll_cash"] - df0["engine_total"]).abs()
        return df0[base_cols]

    # IMPORTANT: remove payroll_tax_pct from roles. Employer rates come from client_settings now.
    for col in ["base_salary_annual", "fte", "super_rate_pct"]:
        if col not in pos.columns:
            pos[col] = 0.0
        pos[col] = pd.to_numeric(pos[col], errors="coerce").fillna(0.0)

    # Annual base salary (fte)
    pos["salary_annual_fte"] = pos["base_salary_annual"] * pos["fte"]

    # Annual PAYG (slab function)
    pos["payg_annual"] = pos["salary_annual_fte"].apply(compute_annual_income_tax_au)

    # Annual super
    pos["super_annual"] = pos["salary_annual_fte"] * (pos["super_rate_pct"] / 100.0)

    # Full-month generated values (work month)
    pos["gross_month_full"] = pos["salary_annual_fte"] / 12.0
    pos["payg_month_full"] = pos["payg_annual"] / 12.0
    pos["super_month_full"] = pos["super_annual"] / 12.0

    def _month_fraction(row, month_start: pd.Timestamp, month_end: pd.Timestamp) -> float:
        role_start = row["start_date"]
        role_end = row["end_date"]
        if pd.isna(role_end):
            role_end = month_end

        start_eff = max(role_start, month_start)
        end_eff = min(role_end, month_end)

        if pd.isna(start_eff) or pd.isna(end_eff):
            return 0.0
        if end_eff < month_start or start_eff > month_end:
            return 0.0

        days_in_month = (month_end - month_start).days + 1
        days_active = (end_eff - start_eff).days + 1
        if days_in_month <= 0 or days_active <= 0:
            return 0.0
        return float(days_active / days_in_month)

    # Extended months for lagged cash
    if lookback > 0:
        first_month = min(month_list)
        extended_months = [(first_month - pd.DateOffset(months=k)) for k in range(lookback, 0, -1)] + month_list
    else:
        extended_months = month_list

    gross_work = {m: 0.0 for m in extended_months}
    payg_work = {m: 0.0 for m in extended_months}
    super_work = {m: 0.0 for m in extended_months}

    for m in extended_months:
        month_start = m
        month_end = m + pd.offsets.MonthEnd(0)
        fractions = pos.apply(lambda row: _month_fraction(row, month_start, month_end), axis=1)

        gross_work[m] = float((pos["gross_month_full"] * fractions).sum())
        payg_work[m] = float((pos["payg_month_full"] * fractions).sum())
        super_work[m] = float((pos["super_month_full"] * fractions).sum())

    rows = []
    for m in month_list:
        month_end = m + pd.offsets.MonthEnd(0)

        # Headcount EOM
        eom_mask = (
            (pos["start_date"] <= month_end) &
            (pos["end_date"].isna() | (pos["end_date"] >= month_end))
        )
        headcount_fte_eom = float(pos.loc[eom_mask, "fte"].sum()) if eom_mask.any() else 0.0

        # Annualised wages at EOM (your threshold test)
        annualised_wages_eom = float(pos.loc[eom_mask, "salary_annual_fte"].sum()) if eom_mask.any() else 0.0
        payroll_tax_applicable = bool(annualised_wages_eom > payroll_tax_threshold_annual)

        gross_this = float(gross_work.get(m, 0.0))
        payg_generated = float(payg_work.get(m, 0.0))
        super_generated = float(super_work.get(m, 0.0))

        # Employer on-costs (generated, work month)
        workers_comp_generated = float(gross_this * workers_comp_rate_pct)
        payroll_tax_generated = float(gross_this * payroll_tax_rate_pct) if payroll_tax_applicable else 0.0

        # Cash to employees this month
        net_salary_cash = gross_this - payg_generated

        # PAYG cash paid this month (lagged)
        if payroll_tax_lag > 0:
            prior_payg_month = m - pd.DateOffset(months=payroll_tax_lag)
            payg_cash = float(payg_work.get(prior_payg_month, 0.0))
        else:
            payg_cash = float(payg_generated)

        # Super cash paid this month (lagged)
        if super_lag > 0:
            prior_super_month = m - pd.DateOffset(months=super_lag)
            super_cash = float(super_work.get(prior_super_month, 0.0))
        else:
            super_cash = float(super_generated)

        total_payroll_cash = float(net_salary_cash + payg_cash + super_cash)

        # Economic monthly payroll cost (work month)
        monthly_payroll_cost = float(
            gross_this
            + super_generated
            + workers_comp_generated
            + payroll_tax_generated
        )

        rows.append({
            "month_date": m,
            "headcount_fte_eom": headcount_fte_eom,
            "annualised_wages_eom": annualised_wages_eom,
            "payroll_tax_applicable": payroll_tax_applicable,
            "gross_work": gross_this,
            "payg_generated": payg_generated,
            "super_generated": super_generated,
            "workers_comp_generated": workers_comp_generated,
            "payroll_tax_generated": payroll_tax_generated,
            "monthly_payroll_cost": monthly_payroll_cost,
            "net_salary_cash": float(net_salary_cash),
            "payg_cash": float(payg_cash),
            "super_cash": float(super_cash),
            "total_payroll_cash": float(total_payroll_cash),
        })

    df = pd.DataFrame(rows)

    # Reconcile vs engine (cash engine remains unchanged)
    engine_totals = compute_payroll_by_month(client_id, month_list) or {}
    df["engine_total"] = df["month_date"].apply(lambda mm: float(engine_totals.get(mm, 0.0)))
    df["diff_vs_engine"] = (df["total_payroll_cash"] - df["engine_total"]).abs()

    # Ensure all base_cols exist (safety)
    for c in base_cols:
        if c not in df.columns:
            df[c] = False if c == "payroll_tax_applicable" else 0.0

    return df[base_cols]



# ============================================================
# NEW HELPER (add this once in a shared utils/module)
# ============================================================

def compute_next_hire_and_burn_kpis(
    client_id,
    month_ts: pd.Timestamp,
    monthly_payroll_cost: float,
) -> dict:
    """
    3 KPIs:
      1) Runway sensitivity to 1 hire (months)
      2) Fully loaded cost of next hire (monthly + annualised, incremental)
      3) Payroll % of burn (monthly payroll cost / effective monthly burn) + band
    """
    settings = get_client_settings(client_id) or {}

    # -----------------------
    # Effective burn (monthly)
    # -----------------------
    effective_burn_monthly = 0.0
    try:
        df_cf = fetch_cashflow_summary_for_client(client_id)
        _, eff_burn = compute_runway_and_effective_burn_from_df(df_cf, month_ts)
        effective_burn_monthly = float(eff_burn or 0.0)
    except Exception:
        effective_burn_monthly = 0.0

    payroll_pct_of_burn = None
    payroll_burn_band = None
    if effective_burn_monthly > 0:
        payroll_pct_of_burn = float(monthly_payroll_cost / effective_burn_monthly * 100.0)
        if payroll_pct_of_burn < 50:
            payroll_burn_band = "Healthy"
        elif payroll_pct_of_burn <= 65:
            payroll_burn_band = "Watch"
        else:
            payroll_burn_band = "High risk"

    # -----------------------
    # Next hire assumptions
    # -----------------------
    next_hire_salary_annual = float(settings.get("next_hire_salary_annual", 0.0) or 0.0)
    next_hire_fte = float(settings.get("next_hire_fte", 1.0) or 1.0)
    next_hire_super_rate = float(settings.get("next_hire_super_rate_pct", 11.0) or 11.0) / 100.0
    next_hire_benefits_rate = float(settings.get("next_hire_benefits_rate_pct", 0.0) or 0.0) / 100.0
    next_hire_onboarding_monthly = float(settings.get("next_hire_onboarding_monthly", 0.0) or 0.0)

    # Employer on-cost settings
    payroll_tax_threshold_annual = float(settings.get("payroll_tax_threshold_annual", 1_300_000) or 1_300_000)
    payroll_tax_rate = float(settings.get("payroll_tax_rate_pct", 0.0) or 0.0) / 100.0
    workers_comp_rate = float(settings.get("workers_comp_rate_pct", 0.0) or 0.0) / 100.0

    # Current annualised wages EOM to test payroll tax after hire
    annualised_wages_eom = 0.0
    try:
        comp_df = compute_payroll_cash_components_by_month(client_id, [month_ts])
        comp = comp_df.iloc[0].to_dict() if comp_df is not None and not comp_df.empty else {}
        annualised_wages_eom = float(comp.get("annualised_wages_eom", 0.0) or 0.0)
    except Exception:
        annualised_wages_eom = 0.0

    # Fully loaded annual cost of next hire (incremental)
    salary_annual_fte = next_hire_salary_annual * next_hire_fte
    super_annual = salary_annual_fte * next_hire_super_rate
    workers_comp_annual = salary_annual_fte * workers_comp_rate
    benefits_annual = salary_annual_fte * next_hire_benefits_rate

    payroll_tax_applicable_after_hire = bool((annualised_wages_eom + salary_annual_fte) > payroll_tax_threshold_annual)
    payroll_tax_annual = (salary_annual_fte * payroll_tax_rate) if payroll_tax_applicable_after_hire else 0.0

    next_hire_annual_cost = float(
        salary_annual_fte
        + super_annual
        + workers_comp_annual
        + payroll_tax_annual
        + benefits_annual
        + (next_hire_onboarding_monthly * 12.0)
    )
    next_hire_monthly_cost = float(next_hire_annual_cost / 12.0)

    # Runway sensitivity (months)
    runway_impact_months = None
    if effective_burn_monthly > 0:
        annual_burn = effective_burn_monthly * 12.0
        runway_impact_months = float(next_hire_annual_cost / annual_burn * 12.0)

    return {
        "effective_burn_monthly": effective_burn_monthly,
        "payroll_pct_of_burn": payroll_pct_of_burn,
        "payroll_burn_band": payroll_burn_band,
        "next_hire_monthly_cost": next_hire_monthly_cost,
        "next_hire_annual_cost": next_hire_annual_cost,
        "runway_impact_months": runway_impact_months,
        "payroll_tax_applicable_after_hire": payroll_tax_applicable_after_hire,
    }


# ============================================================
# UPDATED: page_team_spending (ONLY changes are the new KPIs)
# ============================================================

# ============================================================
# NEW HELPERS (add once in a shared utils/module)
# ============================================================

def _pct_change(curr: float, prev: float) -> float | None:
    try:
        curr = float(curr)
        prev = float(prev)
    except Exception:
        return None
    if prev == 0:
        return None
    return (curr - prev) / prev * 100.0

def compute_team_risk_panels(
    client_id,
    month_ts: pd.Timestamp,
    df_positions: pd.DataFrame,
    monthly_payroll_cost: float,
    payroll_pct_of_burn: float | None,
    runway_impact_months: float | None,
    next_hire_monthly_cost: float,
    month_revenue: float,
) -> dict:
    """
    Produces the 3-column panel content:
      - risks: list[(title, msg)]
      - obligations: list[str]
      - what_this_means: list[str]  (max 5)
      - actions_7d: list[str]       (max 2)
    """

    # ---------- Helpers ----------
    def pct_change(curr: float, prev: float) -> float | None:
        if prev is None or prev == 0:
            return None
        return (curr - prev) / prev * 100.0

    # ---------- Prior month metrics ----------
    prev_month = (month_ts - pd.DateOffset(months=1)).to_period("M").to_timestamp()
    prev2_month = (month_ts - pd.DateOffset(months=2)).to_period("M").to_timestamp()

    # Payroll cost trend (economic)
    prev_payroll_cost = None
    prev2_payroll_cost = None
    try:
        df_prev = compute_payroll_cash_components_by_month(client_id, [prev_month])
        if df_prev is not None and not df_prev.empty:
            prev_payroll_cost = float(df_prev.iloc[0].get("monthly_payroll_cost", 0.0) or 0.0)

        df_prev2 = compute_payroll_cash_components_by_month(client_id, [prev2_month])
        if df_prev2 is not None and not df_prev2.empty:
            prev2_payroll_cost = float(df_prev2.iloc[0].get("monthly_payroll_cost", 0.0) or 0.0)
    except Exception:
        pass

    payroll_mom = pct_change(monthly_payroll_cost, prev_payroll_cost) if prev_payroll_cost is not None else None
    payroll_2m = pct_change(monthly_payroll_cost, prev2_payroll_cost) if prev2_payroll_cost is not None else None

    # Revenue trend
    prev_revenue = None
    prev2_revenue = None
    try:
        k_prev = fetch_kpis_for_client_month(client_id, prev_month.date())
        k_prev2 = fetch_kpis_for_client_month(client_id, prev2_month.date())
        prev_revenue = float(k_prev.get("revenue") or 0.0) if k_prev else 0.0
        prev2_revenue = float(k_prev2.get("revenue") or 0.0) if k_prev2 else 0.0
    except Exception:
        prev_revenue = None
        prev2_revenue = None

    revenue_mom = pct_change(month_revenue, prev_revenue) if prev_revenue is not None else None
    revenue_2m = pct_change(month_revenue, prev2_revenue) if prev2_revenue is not None else None

    # Headcount trend
    headcount_now = 0.0
    headcount_prev = None
    try:
        df_now = compute_payroll_cash_components_by_month(client_id, [month_ts])
        headcount_now = float(df_now.iloc[0].get("headcount_fte_eom", 0.0) or 0.0)

        df_prev_hc = compute_payroll_cash_components_by_month(client_id, [prev_month])
        if df_prev_hc is not None and not df_prev_hc.empty:
            headcount_prev = float(df_prev_hc.iloc[0].get("headcount_fte_eom", 0.0) or 0.0)
    except Exception:
        pass

    hc_delta = (headcount_now - headcount_prev) if headcount_prev is not None else None

    # ---------- Upcoming payroll obligations (NEXT MONTH cash amounts) ----------
    next_month = (month_ts + pd.DateOffset(months=1)).to_period("M").to_timestamp()

    next_net_pay = None
    next_super_payable = None
    next_payg_payable = None

    try:
        df_next = compute_payroll_cash_components_by_month(client_id, [next_month])
        if df_next is not None and not df_next.empty:
            row = df_next.iloc[0].to_dict()
            # cash in next month
            next_net_pay = float(row.get("net_salary_cash", 0.0) or 0.0)
            next_super_payable = float(row.get("super_cash", 0.0) or 0.0)
            next_payg_payable = float(row.get("payg_cash", 0.0) or 0.0)
    except Exception:
        pass

    # ---------- Risk triggers ----------
    risks = []

    # Risk 1: Hiring Without Revenue Cover
    risk1 = False
    if hc_delta is not None and hc_delta > 0:
        # revenue flat or down: <= +1% MoM treated as flat
        if (revenue_mom is None) or (revenue_mom <= 1.0):
            risk1 = True
    if risk1:
        risks.append((
            "🔴 Risk 1: Hiring Without Revenue Cover",
            f"Headcount increased (ΔFTE={hc_delta:+.1f}) while revenue is flat/down (MoM={revenue_mom if revenue_mom is not None else '—'}%)."
        ))
    else:
        risks.append(("🟢 Risk 1: Hiring Without Revenue Cover", "No signal triggered for this month."))

    # Risk 2: Payroll Growth > Revenue Growth (over 2 months)
    risk2 = False
    if payroll_2m is not None and revenue_2m is not None:
        if payroll_2m > revenue_2m:
            risk2 = True
    if risk2:
        risks.append((
            "🔴 Risk 2: Payroll Growth > Revenue Growth",
            f"Over ~2 months, payroll growth ({payroll_2m:+.1f}%) is exceeding revenue growth ({revenue_2m:+.1f}%)."
        ))
    else:
        risks.append(("🟢 Risk 2: Payroll Growth > Revenue Growth", "No signal triggered for this month."))

    # Risk 3: Payroll Dominating Burn
    risk3 = bool(payroll_pct_of_burn is not None and payroll_pct_of_burn > 65.0)
    if risk3:
        risks.append((
            "🔴 Risk 3: Payroll Dominating Burn",
            f"Payroll is {payroll_pct_of_burn:.0f}% of burn (>65%). This increases failure risk if revenue slips."
        ))
    else:
        risks.append(("🟢 Risk 3: Payroll Dominating Burn", "No signal triggered for this month."))

    # Risk 4: Silent Hiring (contractors added but not counted as FTEs)
    risk4 = False
    contractor_count = 0
    contractor_fte = 0.0
    if df_positions is not None and not df_positions.empty:
        wt_col = None
        if "worker_type" in df_positions.columns:
            wt_col = "worker_type"
        elif "employment_type" in df_positions.columns:
            wt_col = "employment_type"

        if wt_col:
            tmp = df_positions.copy()
            tmp[wt_col] = tmp[wt_col].astype(str).str.lower()
            tmp["fte"] = pd.to_numeric(tmp.get("fte"), errors="coerce").fillna(0.0)
            contractors = tmp[tmp[wt_col].str.contains("contract", na=False)]
            contractor_count = int(len(contractors))
            contractor_fte = float(contractors["fte"].sum()) if not contractors.empty else 0.0

            # "silent": contractors exist but are effectively excluded from FTE totals
            if contractor_count > 0 and contractor_fte <= 0.01:
                risk4 = True

    if risk4:
        risks.append((
            "🔴 Risk 4: Silent Hiring",
            f"{contractor_count} contractor role(s) detected but contractor FTE≈{contractor_fte:.2f}. Contractor cost may be hiding outside headcount."
        ))
    else:
        risks.append(("🟢 Risk 4: Silent Hiring", "No signal triggered (or contractor classification not available)."))

    # ---------- Obligations text ----------
    obligations = []
    if next_net_pay is not None:
        obligations.append(f"Next month net pay (existing employees): {next_net_pay:,.0f}")
    else:
        obligations.append("Next month net pay (existing employees): —")

    if next_super_payable is not None:
        obligations.append(f"Super payable next month (cash): {next_super_payable:,.0f}")
    else:
        obligations.append("Super payable next month (cash): —")

    if next_payg_payable is not None:
        obligations.append(f"PAYG payable next month (cash): {next_payg_payable:,.0f}")
    else:
        obligations.append("PAYG payable next month (cash): —")

    # ---------- What this means (max 5 bullets) ----------
    what = []

    # a) Current payroll sustainability
    if payroll_pct_of_burn is not None:
        what.append(f"Payroll sustainability: **Payroll % of burn = {payroll_pct_of_burn:.0f}%** → tighter margin for error if revenue dips.")
    else:
        what.append("Payroll sustainability: Payroll % of burn is **not available** → burn source missing, risk may be understated.")

    # b) Direction of payroll vs revenue
    if payroll_mom is not None and revenue_mom is not None:
        what.append(f"Direction: **Payroll MoM {payroll_mom:+.1f}% vs Revenue MoM {revenue_mom:+.1f}%** → divergence increases cash pressure.")
    else:
        what.append("Direction: Payroll vs revenue trend is **partially missing** → ensure revenue KPI exists for both months.")

    # c) Runway sensitivity
    if runway_impact_months is not None:
        what.append(f"Runway sensitivity: **1 hire reduces runway by ~{runway_impact_months:.2f} months** → hiring mistakes compress reaction time.")
    else:
        what.append("Runway sensitivity: **not available** (needs valid effective burn + next-hire assumptions).")

    # d) Hiring risk
    if risk1:
        what.append(f"Hiring risk: **Headcount up (ΔFTE {hc_delta:+.1f})** while **revenue not growing** → risk of hiring into a flat top-line.")
    else:
        # keep concise
        what.append("Hiring risk: No headcount-vs-revenue mismatch triggered this month.")

    # e) Immediate implication (using obligations)
    if next_net_pay is not None:
        what.append(f"Immediate implication: next month cash obligations include **net pay {next_net_pay:,.0f}** + statutory remittances → protect liquidity window now.")
    else:
        what.append("Immediate implication: next month payroll obligations unavailable → check payroll engine inputs and role dates.")

    what_this_means = what[:5]

    # ---------- Action plan (max 2 bullets; must map to risks) ----------
    actions = []

    # Priority 1: Payroll dominating burn
    if risk3:
        actions.append("Risk 3 → **Freeze new hires** and **review contractor spend** to pull Payroll % of Burn below 65%.")

    # Priority 2: Hiring without revenue cover OR payroll>revenue growth
    if risk1 or risk2:
        actions.append("Risk 1/2 → **Pause hiring pipeline** and **revalidate hiring assumptions vs revenue reality** (update the next-hire assumptions + hiring plan).")

    # Priority 3: Runway sensitivity high (if still space)
    if len(actions) < 2 and runway_impact_months is not None and runway_impact_months >= 0.50:
        actions.append("Runway sensitivity → **Model a downside hiring scenario** and **extend reaction window via cost control** (stop discretionary roles).")

    # If nothing triggered, still give 1 decisive action
    if not actions:
        actions.append("No critical payroll risks triggered → **confirm next-month statutory payments** (Super/PAYG) are funded and scheduled.")

    actions_7d = actions[:2]

    return {
        "risks": risks,
        "obligations": obligations,
        "what_this_means": what_this_means,
        "actions_7d": actions_7d,
    }



# ============================================================
# UPDATED: page_team_spending (removes charts + decision panel,
# adds the 3-column risk/meaning/action layout)
#
import math
import pandas as pd
from pandas.tseries.offsets import MonthEnd

def page_team_spending():
    top_header("Team & Spending")

    if not selected_client_id:
        st.info("Select a business from the sidebar to view this page.")
        return

    currency_code, currency_symbol = get_client_currency(selected_client_id)

    # ------------------------------------------------------------------
    # Load data
    # ------------------------------------------------------------------
    with st.spinner("Loading team and spending data..."):
        df_positions = fetch_payroll_positions_for_client(selected_client_id)
        df_ar_ap = fetch_ar_ap_for_client(selected_client_id)
        kpis = fetch_kpis_for_client_month(selected_client_id, selected_month_start)

    if isinstance(df_ar_ap, tuple):
        df_ar, df_ap = df_ar_ap
    else:
        df_ar, df_ap = None, None

    # ------------------------------------------------------------------
    # Month setup
    # ------------------------------------------------------------------
    month_ts = pd.to_datetime(selected_month_start).to_period("M").to_timestamp()

    # ------------------------------------------------------------------
    # Payroll cash engine (source of truth)
    # ------------------------------------------------------------------
    month_index_1 = pd.date_range(start=month_ts, periods=1, freq="MS")
    payroll_by_month = compute_payroll_by_month(selected_client_id, month_index_1)
    month_payroll_cash_engine = float(payroll_by_month.get(month_ts, 0.0))

    # ------------------------------------------------------------------
    # Payroll components + economic payroll cost
    # ------------------------------------------------------------------
    comp_df = compute_payroll_cash_components_by_month(selected_client_id, month_index_1)
    comp = comp_df.iloc[0].to_dict() if comp_df is not None and not comp_df.empty else {}

    headcount_fte_eom = float(comp.get("headcount_fte_eom", 0.0) or 0.0)

    # Cash KPIs
    salaries_paid_cash = float(comp.get("net_salary_cash", 0.0) or 0.0)
    super_paid_cash = float(comp.get("super_cash", 0.0) or 0.0)
    tax_paid_cash = float(comp.get("payg_cash", 0.0) or 0.0)
    total_payroll_cash_components = float(comp.get("total_payroll_cash", 0.0) or 0.0)

    # Economic payroll cost (work month)
    monthly_payroll_cost = float(comp.get("monthly_payroll_cost", 0.0) or 0.0)
    avg_cost_per_fte = (monthly_payroll_cost / headcount_fte_eom) if headcount_fte_eom else 0.0

    # Employer on-costs
    workers_comp_generated = float(comp.get("workers_comp_generated", 0.0) or 0.0)
    payroll_tax_generated = float(comp.get("payroll_tax_generated", 0.0) or 0.0)
    annualised_wages_eom = float(comp.get("annualised_wages_eom", 0.0) or 0.0)
    payroll_tax_applicable = bool(comp.get("payroll_tax_applicable", False))

    # Tie-out
    tie_diff = float(total_payroll_cash_components - month_payroll_cash_engine)

    # ------------------------------------------------------------------
    # Headcount split for the month (focus month)
    # ------------------------------------------------------------------
    hc = compute_headcount_month_split(df_positions, month_ts)
    hires_fte_in_month = float(hc.get("hires_fte_in_month", 0.0))
    exits_fte_in_month = float(hc.get("exits_fte_in_month", 0.0))
    net_change_in_month = float(hc.get("net_change_fte_in_month", 0.0))

    # ------------------------------------------------------------------
    # Headcount trend (past 6 months)  ✅ NEW CHART
    # ------------------------------------------------------------------
    hc_trend_df = None
    try:
        # Past 6 months incl current focus month
        month_index_6_past = pd.date_range(end=month_ts, periods=6, freq="MS")
        comp_6_df = compute_payroll_cash_components_by_month(selected_client_id, month_index_6_past)

        if comp_6_df is not None and not comp_6_df.empty:
            hc_trend_df = comp_6_df[["month_date", "headcount_fte_eom"]].copy()
            hc_trend_df["month_date"] = pd.to_datetime(hc_trend_df["month_date"], errors="coerce")
            hc_trend_df["headcount_fte_eom"] = pd.to_numeric(
                hc_trend_df["headcount_fte_eom"], errors="coerce"
            ).fillna(0.0)
            hc_trend_df = hc_trend_df.sort_values("month_date")
    except Exception:
        hc_trend_df = None

    # ------------------------------------------------------------------
    # Payroll target (client_monthly_targets)
    # ------------------------------------------------------------------
    targets = get_monthly_targets_for_month(selected_client_id, month_ts)
    payroll_target = float(targets.get("payroll_target", 0.0) or 0.0)
    payroll_var = monthly_payroll_cost - payroll_target
    payroll_var_pct = (payroll_var / payroll_target * 100.0) if payroll_target else None

    # ------------------------------------------------------------------
    # Revenue (from KPIs)
    # ------------------------------------------------------------------
    month_revenue = 0.0
    if kpis:
        try:
            month_revenue = float(kpis.get("revenue") or 0.0)
        except Exception:
            month_revenue = 0.0

    # ------------------------------------------------------------------
    # 3 new KPIs: burn + next hire + sensitivity (source of truth)
    # ------------------------------------------------------------------
    hire_kpis = compute_next_hire_and_burn_kpis(
        selected_client_id,
        month_ts=month_ts,
        monthly_payroll_cost=monthly_payroll_cost,
    )
    effective_burn_monthly = float(hire_kpis.get("effective_burn_monthly", 0.0) or 0.0)
    payroll_pct_of_burn = hire_kpis.get("payroll_pct_of_burn", None)
    payroll_burn_band = hire_kpis.get("payroll_burn_band", None)
    next_hire_monthly_cost = float(hire_kpis.get("next_hire_monthly_cost", 0.0) or 0.0)
    next_hire_annual_cost = float(hire_kpis.get("next_hire_annual_cost", 0.0) or 0.0)
    runway_impact_months = hire_kpis.get("runway_impact_months", None)

    # ------------------------------------------------------------------
    # KPI HEADER (row-style sections)
    # ------------------------------------------------------------------
    st.subheader("👥 Team & Payroll KPIs (focus month)")

    # -------------------- Headcount (ROW) --------------------
    st.markdown("### 👥 Headcount")
    c1, c2, c3, c4 = st.columns([1.2, 1, 1, 1.2])
    with c1:
        st.metric("Total HC (FTE)", f"{headcount_fte_eom:,.1f}")
    with c2:
        st.metric("New hires", f"+{hires_fte_in_month:,.1f}")
    with c3:
        st.metric("Exits", f"-{exits_fte_in_month:,.1f}")
    with c4:
        st.metric("Net change", f"{net_change_in_month:+.1f} FTE")

    # ✅ NEW: headcount trend line (past 6 months) under headcount KPIs
    if hc_trend_df is not None and not hc_trend_df.empty:
        st.markdown("#### Headcount trend (past 6 months)")
        try:
            import altair as alt

            hc_chart = (
                alt.Chart(hc_trend_df)
                .mark_line()
                .encode(
                    x=alt.X("month_date:T", title="Month", axis=alt.Axis(format="%b %Y")),
                    y=alt.Y("headcount_fte_eom:Q", title="FTE (EOM)"),
                    tooltip=[
                        alt.Tooltip("month_date:T", title="Month", format="%b %Y"),
                        alt.Tooltip("headcount_fte_eom:Q", title="Headcount (FTE)", format=",.1f"),
                    ],
                )
            )
            st.altair_chart(hc_chart, width="stretch")
        except Exception:
            # fallback to Streamlit line_chart if Altair fails
            tmp = hc_trend_df.set_index("month_date")[["headcount_fte_eom"]]
            st.line_chart(tmp)
    else:
        st.caption("Headcount trend unavailable (no payroll roles / insufficient history).")

    st.markdown("---")

    # -------------------- Payroll (work month) (ROW) --------------------
    st.markdown("### 💰 Payroll (work month)")
    p1, p2, p3, p4 = st.columns([1.3, 1.0, 1.1, 1.4])
    with p1:
        st.metric("Monthly payroll cost", f"{currency_symbol}{monthly_payroll_cost:,.0f}")
    with p2:
        st.metric("Avg cost per FTE", f"{currency_symbol}{avg_cost_per_fte:,.0f}")
    with p3:
        st.metric("Workers comp (gen.)", f"{currency_symbol}{workers_comp_generated:,.0f}")
    with p4:
        st.metric("Payroll tax (gen.)", f"{currency_symbol}{payroll_tax_generated:,.0f}")

    st.caption(
        f"Payroll tax applicable: **{'Yes' if payroll_tax_applicable else 'No'}** · "
        f"Annualised wages (EOM): **{currency_symbol}{annualised_wages_eom:,.0f}**"
    )

    st.markdown("#### 🎯 Payroll vs target")

    t1, t2, t3 = st.columns([1.2, 1.2, 1.6])

    with t1:
        # ✅ Payroll target KPI removed intentionally
        st.caption("Payroll target is used only to calculate variance (not displayed as a KPI).")

    with t2:
        if payroll_target:
            delta_txt = f"{currency_symbol}{payroll_var:,.0f}"
            if payroll_var_pct is not None:
                delta_txt += f" ({payroll_var_pct:+.0f}%)"
            st.metric("Variance vs target", delta_txt)
        else:
            st.metric("Variance vs target", "—")

    with t3:
        if payroll_target:
            st.caption("Variance is compared to **work-month payroll cost** (economic), not cash timing.")
        else:
            st.caption("No payroll target set for this month in client_monthly_targets — variance cannot be computed.")


    st.markdown("#### 🧪 Hiring sensitivity & burn mix")
    h2, h3, h4 = st.columns([ 1.6, 1.2, 1.2])
    with h2:
        st.metric(
            "Fully loaded cost of next hire",
            f"{currency_symbol}{next_hire_monthly_cost:,.0f} / month",
            help="Incremental cost: salary + super + workers comp + payroll tax (if applicable) + benefits + onboarding (if recurring).",
        )
        st.caption(f"Annualised: **{currency_symbol}{next_hire_annual_cost:,.0f} / year**")
    with h3:
        st.metric(
            "Runway sensitivity (1 hire)",
            f"-{float(runway_impact_months):.2f} months" if runway_impact_months is not None else "—",
        )
    with h4:
        st.metric(
            "Payroll % of burn",
            f"{float(payroll_pct_of_burn):.0f}%" if payroll_pct_of_burn is not None else "—",
        )

    if payroll_pct_of_burn is not None:
        if payroll_burn_band == "High risk":
            st.error("Payroll % of burn: **High risk** (>65%)")
        elif payroll_burn_band == "Watch":
            st.warning("Payroll % of burn: **Watch** (50–65%)")
        else:
            st.success("Payroll % of burn: **Healthy** (<50%)")

    st.markdown("---")

    # -------------------- Payroll (cash paid) (ROW) --------------------
    st.markdown("### 💳 Payroll (cash paid)")
    cc1, cc2, cc3, cc4 = st.columns([1.2, 1.2, 1.2, 1.4])
    with cc1:
        st.metric("Base salaries (net)", f"{currency_symbol}{salaries_paid_cash:,.0f}")
    with cc2:
        st.metric("Super paid (lagged)", f"{currency_symbol}{super_paid_cash:,.0f}")
    with cc3:
        st.metric("PAYG paid (lagged)", f"{currency_symbol}{tax_paid_cash:,.0f}")
    with cc4:
        st.metric(
            "Total payroll cash",
            f"{currency_symbol}{month_payroll_cash_engine:,.0f}",
            delta=f"{currency_symbol}{tie_diff:,.0f}",
        )

    if abs(tie_diff) > 0.01:
        st.warning(
            f"Payroll engine tie-out mismatch: components={currency_symbol}{total_payroll_cash_components:,.2f} "
            f"vs engine={currency_symbol}{month_payroll_cash_engine:,.2f} "
            f"(diff={currency_symbol}{tie_diff:,.2f})."
        )

    st.markdown("---")

    # ------------------------------------------------------------------
    # ROW-STYLE PANELS (no 3 columns)
    # ------------------------------------------------------------------
    panel = compute_team_risk_panels(
        selected_client_id,
        month_ts,
        df_positions=df_positions,
        monthly_payroll_cost=monthly_payroll_cost,
        payroll_pct_of_burn=payroll_pct_of_burn,
        runway_impact_months=runway_impact_months,
        next_hire_monthly_cost=next_hire_monthly_cost,
        month_revenue=month_revenue,
    )

    st.markdown("## 🧾 Payroll risk & compliance")
    for title, msg in panel.get("risks", []):
        if title.startswith("🔴"):
            st.error(f"**{title}**\n\n{msg}")
        elif title.startswith("🟢"):
            st.success(f"**{title}**\n\n{msg}")
        else:
            st.info(f"**{title}**\n\n{msg}")

    st.markdown("## 📅 Upcoming payroll obligations (next month)")
    for ob in panel.get("obligations", [])[:3]:
        st.markdown(f"- {ob}")

    if df_positions is not None and not df_positions.empty:
        if ("worker_type" not in df_positions.columns) and ("employment_type" not in df_positions.columns):
            st.caption("Silent hiring detection note: add `worker_type` or `employment_type` to classify contractors.")

    st.markdown("---")

    st.markdown("## 🧠 What this means")
    for b in panel.get("what_this_means", [])[:5]:
        st.markdown(f"- {b}")

    st.markdown("---")

    st.markdown("## ✅ Next 7 days action plan")
    for a in panel.get("actions_7d", [])[:2]:
        st.markdown(f"- {a}")

    # Manage roles section intentionally removed.


# ============================================================
# UPDATED: page_client_settings (adds Next Hire assumptions UI)
# ============================================================


def page_client_settings():
    st.title("⚙️ Client Settings")

    if not selected_client_id:
        st.info("Select a client from the sidebar to configure settings.")
        return

    client_id_str = str(selected_client_id)

    # Always load settings once at top
    settings = get_client_settings(selected_client_id) or {}

    # ----------------- Currency -----------------
    st.subheader("💱 Currency Settings")
    render_currency_settings_section(selected_client_id)

    st.markdown("---")

# ----------------- NEW: Pipeline ageing thresholds -----------------
    st.subheader("⏳ Pipeline ageing thresholds (locked defaults)")

    st.caption(
        "Used for the **Pipeline ageing warning** on Sales & Deals.\n\n"
        "Hard defaults: Proposal 30d, Demo 45d, Contract 60d. "
        "Idea is excluded. Closed_Won/Lost excluded."
    )

    cA, cB, cC = st.columns(3)
    with cA:
        stage_max_days_proposal = st.number_input(
            "Proposal max healthy days",
            min_value=1,
            max_value=365,
            value=int(settings.get("stage_max_days_proposal", 30) or 30),
        )
    with cB:
        stage_max_days_demo = st.number_input(
            "Demo max healthy days",
            min_value=1,
            max_value=365,
            value=int(settings.get("stage_max_days_demo", 45) or 45),
        )
    with cC:
        stage_max_days_contract = st.number_input(
            "Contract max healthy days",
            min_value=1,
            max_value=365,
            value=int(settings.get("stage_max_days_contract", 60) or 60),
        )

    

    if st.button("Save pipeline ageing thresholds", key="save_stage_ageing"):
        payload = {
            "client_id": client_id_str,
            "stage_max_days_proposal": int(stage_max_days_proposal),
            "stage_max_days_demo": int(stage_max_days_demo),
            "stage_max_days_contract": int(stage_max_days_contract),
        }

        try:
            supabase.table("client_settings").upsert(payload, on_conflict="client_id").execute()

            try:
                st.cache_data.clear()
            except Exception:
                pass

            st.success("Pipeline ageing thresholds saved.")
            st.rerun()

        except APIError as e:
            msg = str(e)
            if "PGRST204" in msg and "stage_max_days_contract" in msg:
                st.error(
                    "Your Supabase schema is missing the pipeline ageing columns on `client_settings` "
                    "(or PostgREST schema cache hasn’t reloaded)."
                )
                st.markdown("### Run this SQL in Supabase, then reload schema:")
                st.code(
                    """alter table public.client_settings
    add column if not exists stage_max_days_proposal integer,
    add column if not exists stage_max_days_demo integer,
    add column if not exists stage_max_days_contract integer;""",
                    language="sql",
                )
            else:
                st.error("Could not save pipeline ageing thresholds.")
                st.code(msg, language="text")

    st.markdown("---")


    # ----------------- Opening cash -----------------
    st.subheader("💰 Opening cash")

    opening_cash = st.number_input(
        "Opening cash balance at start of focus month",
        min_value=0.0,
        value=float(settings.get("opening_cash_start", 0.0) or 0.0),
        step=1000.0,
        help="Used as the starting cash when you rebuild the 12-month cashflow engine.",
    )

    if st.button("Save opening cash", key="save_opening_cash"):
        supabase.table("client_settings").upsert(
            {"client_id": client_id_str, "opening_cash_start": float(opening_cash)},
            on_conflict="client_id",
        ).execute()
        st.success("Opening cash saved.")
        st.rerun()

    st.markdown("---")

    # ----------------- NEW: Cash safety buffer -----------------
    st.subheader("🛟 Cash safety buffer")

    min_cash_buffer = st.number_input(
        "Minimum cash buffer (safety cushion)",
        min_value=0.0,
        value=float(settings.get("min_cash_buffer", 0.0) or 0.0),
        step=5000.0,
        help=(
            "Used in Cash & Bills → 'Cash danger summary'.\n\n"
            "If your lowest projected cash drops below this number, the dashboard warns you "
            "even if cash never goes negative."
        ),
    )

    if st.button("Save cash safety buffer", key="save_min_cash_buffer"):
        try:
            supabase.table("client_settings").upsert(
                {"client_id": client_id_str, "min_cash_buffer": float(min_cash_buffer)},
                on_conflict="client_id",
            ).execute()
            st.success("Cash safety buffer saved.")
            st.rerun()
        except APIError as e:
            st.error(
                "Supabase schema is missing `min_cash_buffer` on `client_settings`.\n\n"
                "Run the SQL migration below to add the column, then refresh PostgREST schema."
            )
            st.code(str(e), language="text")

    st.markdown("---")

    # ----------------- Payroll cash timing (super + PAYG cash lag) -----------------
    st.subheader("🕒 Payroll cash timing (Super & PAYG)")

    col_s1, col_s2 = st.columns(2)

    with col_s1:
        super_lag = st.number_input(
            "Superannuation payment lag (months)",
            min_value=0,
            max_value=12,
            value=int(settings.get("super_lag_months", 1) or 1),
            help="How many months after salary month you actually pay super.",
        )

    with col_s2:
        payroll_tax_lag = st.number_input(
            "PAYG payment lag (months)",
            min_value=0,
            max_value=12,
            value=int(settings.get("payroll_tax_lag_months", 1) or 1),
            help="How many months after salary month you remit PAYG withholding.",
        )

    if st.button("Save payroll cash timing", key="save_payroll_cash_timing"):
        supabase.table("client_settings").upsert(
            {
                "client_id": client_id_str,
                "super_lag_months": int(super_lag),
                "payroll_tax_lag_months": int(payroll_tax_lag),
            },
            on_conflict="client_id",
        ).execute()
        st.success("Payroll cash timing saved.")
        st.rerun()

    st.markdown("---")

    # ----------------- Employer on-costs -----------------
    st.subheader("🏛️ Employer on-costs (Payroll tax + Workers comp)")

    col1, col2, col3 = st.columns(3)
    with col1:
        payroll_tax_threshold_annual = st.number_input(
            "Payroll tax threshold (annual wages)",
            min_value=0.0,
            step=50000.0,
            value=float(settings.get("payroll_tax_threshold_annual", 1_300_000) or 1_300_000),
        )
    with col2:
        payroll_tax_rate_pct = st.number_input(
            "Payroll tax rate (%)",
            min_value=0.0,
            max_value=20.0,
            step=0.1,
            value=float(settings.get("payroll_tax_rate_pct", 0.0) or 0.0),
        )
    with col3:
        workers_comp_rate_pct = st.number_input(
            "Workers comp rate (%)",
            min_value=0.0,
            max_value=10.0,
            step=0.1,
            value=float(settings.get("workers_comp_rate_pct", 0.0) or 0.0),
        )

    if st.button("Save employer on-costs", key="save_employer_on_costs"):
        supabase.table("client_settings").upsert(
            {
                "client_id": client_id_str,
                "payroll_tax_threshold_annual": float(payroll_tax_threshold_annual),
                "payroll_tax_rate_pct": float(payroll_tax_rate_pct),
                "workers_comp_rate_pct": float(workers_comp_rate_pct),
            },
            on_conflict="client_id",
        ).execute()
        st.success("Employer on-costs saved.")
        st.rerun()

    st.markdown("---")

    # ----------------- Next hire assumptions -----------------
    st.subheader("🧑‍💼 Next hire assumptions (used for hiring sensitivity KPIs)")

    nh1, nh2, nh3 = st.columns(3)

    with nh1:
        next_hire_salary_annual = st.number_input(
            "Next hire base salary (annual)",
            min_value=0.0,
            step=5000.0,
            value=float(settings.get("next_hire_salary_annual", 0.0) or 0.0),
        )
        next_hire_fte = st.number_input(
            "Next hire FTE",
            min_value=0.1,
            max_value=2.0,
            step=0.1,
            value=float(settings.get("next_hire_fte", 1.0) or 1.0),
        )

    with nh2:
        next_hire_super_rate_pct = st.number_input(
            "Next hire super rate (%)",
            min_value=0.0,
            max_value=30.0,
            step=0.5,
            value=float(settings.get("next_hire_super_rate_pct", 11.0) or 11.0),
        )
        next_hire_benefits_rate_pct = st.number_input(
            "Next hire benefits rate (%)",
            min_value=0.0,
            max_value=30.0,
            step=0.5,
            value=float(settings.get("next_hire_benefits_rate_pct", 0.0) or 0.0),
            help="Incremental benefits as % of base salary (allowances, tools, health, etc.).",
        )

    with nh3:
        next_hire_onboarding_monthly = st.number_input(
            "Onboarding cost (monthly, if recurring)",
            min_value=0.0,
            step=250.0,
            value=float(settings.get("next_hire_onboarding_monthly", 0.0) or 0.0),
            help="If it’s a one-off cost, model it elsewhere; this field is for recurring onboarding/enablement cost.",
        )

    st.caption(
        "If you see an error like **PGRST204 could not find column in schema cache**, "
        "it means the column does not exist yet OR PostgREST needs a schema reload."
    )

    if st.button("Save next hire assumptions", key="save_next_hire"):
        payload = {
            "client_id": client_id_str,
            "next_hire_salary_annual": float(next_hire_salary_annual),
            "next_hire_fte": float(next_hire_fte),
            "next_hire_super_rate_pct": float(next_hire_super_rate_pct),
            "next_hire_benefits_rate_pct": float(next_hire_benefits_rate_pct),
            "next_hire_onboarding_monthly": float(next_hire_onboarding_monthly),
        }

        try:
            supabase.table("client_settings").upsert(payload, on_conflict="client_id").execute()
            st.success("Next hire assumptions saved.")
            st.rerun()

        except APIError as e:
            # Your exact failure was: PGRST204 missing column in schema cache
            msg = str(e)
            if "PGRST204" in msg:
                st.error(
                    "Supabase says one or more next-hire columns are missing from `client_settings` "
                    "(or PostgREST schema cache hasn’t reloaded).\n\n"
                    "Run the SQL migration below, then reload PostgREST schema."
                )
            else:
                st.error("Failed to save next hire assumptions.")
            st.code(msg, language="text")

    st.markdown("---")

    # ----------------- AR / AP defaults -----------------
    st.subheader("📌 AR / AP Default Rules")
    col1, col2 = st.columns(2)

    with col1:
        ar_days = st.number_input(
            "Default AR collection days",
            min_value=0,
            max_value=180,
            value=int(settings.get("ar_default_days", 0) or 0),
        )
    with col2:
        ap_days = st.number_input(
            "Default AP payment days",
            min_value=0,
            max_value=180,
            value=int(settings.get("ap_default_days", 0) or 0),
        )

    if st.button("Save AR/AP defaults", key="save_ar_ap_defaults"):
        supabase.table("client_settings").upsert(
            {"client_id": client_id_str, "ar_default_days": int(ar_days), "ap_default_days": int(ap_days)},
            on_conflict="client_id",
        ).execute()
        st.success("AR/AP defaults saved.")
        st.rerun()

    st.markdown("---")

    # ----------------- Runway / overspend thresholds -----------------
    st.subheader("🚨 Runway & Overspend Thresholds")
    col1, col2, col3 = st.columns(3)

    with col1:
        runway_min = st.number_input(
            "Minimum runway months",
            min_value=0.0,
            max_value=24.0,
            value=float(settings.get("runway_min_months", 4.0) or 4.0),
            step=0.5,
        )
    with col2:
        warn_pct = st.number_input(
            "Overspend warning %",
            min_value=0.0,
            max_value=100.0,
            value=float(settings.get("overspend_warn_pct", 0.0) or 0.0),
        )
    with col3:
        danger_pct = st.number_input(
            "Overspend danger %",
            min_value=0.0,
            max_value=100.0,
            value=float(settings.get("overspend_high_pct", 0.0) or 0.0),
        )

    if st.button("Save runway & overspend thresholds", key="save_risk_thresholds"):
        supabase.table("client_settings").upsert(
            {
                "client_id": client_id_str,
                "runway_min_months": float(runway_min),
                "overspend_warn_pct": float(warn_pct),
                "overspend_high_pct": float(danger_pct),
            },
            on_conflict="client_id",
        ).execute()
        st.success("Risk thresholds saved.")
        st.rerun()

    st.markdown("---")

    # ----------------- Revenue recognition -----------------
    st.subheader("📊 Revenue Recognition Method")

    method = st.selectbox(
        "Select default revenue recognition method",
        ["saas", "milestone", "straight_line", "instant", "deferred"],
        index=["saas", "milestone", "straight_line", "instant", "deferred"]
        .index(settings.get("revenue_recognition_method", "saas")),
    )

    if st.button("Save revenue recognition method", key="save_revrec"):
        supabase.table("client_settings").upsert(
            {"client_id": client_id_str, "revenue_recognition_method": method},
            on_conflict="client_id",
        ).execute()
        st.success("Revenue recognition method updated.")
        st.rerun()
# ----------------- Monthly targets / budgets -----------------
    st.markdown("---")
    st.subheader("🎯 Monthly targets & budgets")
    st.caption(
        "Set internal targets for each month. These are used to compare **actual vs plan** "
        "across revenue, payroll, and headcount. (Always editable — no locking.)"
    )

    targets_df = fetch_client_monthly_targets(selected_client_id)

    start = pd.to_datetime(selected_month_start).replace(day=1)
    months = pd.date_range(start=start, periods=12, freq="MS")
    base = pd.DataFrame({"month_date": months})

    if targets_df is not None and not targets_df.empty:
        # ensure month_date is comparable (month bucket)
        targets_df = targets_df.copy()
        targets_df["month_date"] = pd.to_datetime(targets_df["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()

        targets_df = targets_df[[
            "month_date",
            "revenue_target",
            "payroll_target",
            "opex_target",
            "headcount_fte_target"
        ]]
        base = base.merge(targets_df, on="month_date", how="left")

    for c in ["revenue_target", "payroll_target", "opex_target", "headcount_fte_target"]:
        if c not in base.columns:
            base[c] = 0.0
        base[c] = pd.to_numeric(base[c], errors="coerce").fillna(0.0)

    base["Month"] = base["month_date"].dt.strftime("%b %Y")

    editable = base[
        ["Month", "month_date", "revenue_target", "payroll_target", "opex_target", "headcount_fte_target"]
    ].rename(
        columns={
            "revenue_target": "Revenue target",
            "payroll_target": "Payroll target",
            "opex_target": "Opex target",
            "headcount_fte_target": "Headcount FTE target",
        }
    )

    edited = st.data_editor(
        editable,
        num_rows="fixed",
        width="stretch",
        key="client_monthly_targets_editor",
    )

    def _r2(x):
        """Stable compare: round to cents (or 0.01 FTE for headcount)."""
        try:
            return float(x)
        except Exception:
            return 0.0

    if st.button("💾 Save monthly targets", key="save_monthly_targets"):
        # Build existing map for change detection
        existing = {}
        if targets_df is not None and not targets_df.empty:
            for _, rr in targets_df.iterrows():
                key = pd.to_datetime(rr["month_date"]).date().isoformat()
                existing[key] = {
                    "revenue_target": float(rr.get("revenue_target") or 0.0),
                    "payroll_target": float(rr.get("payroll_target") or 0.0),
                    "opex_target": float(rr.get("opex_target") or 0.0),
                    "headcount_fte_target": float(rr.get("headcount_fte_target") or 0.0),
                }

        rows_to_upsert = []

        for _, r in edited.iterrows():
            m = pd.to_datetime(r["month_date"], errors="coerce").to_period("M").to_timestamp()
            if pd.isna(m):
                continue

            month_key = m.date().isoformat()

            new_row = {
                "client_id": client_id_str,
                "month_date": month_key,
                "revenue_target": float(_r2(r.get("Revenue target"))),
                "payroll_target": float(_r2(r.get("Payroll target"))),
                "opex_target": float(_r2(r.get("Opex target"))),
                "headcount_fte_target": float(_r2(r.get("Headcount FTE target"))),
            }

            # Round for change detection (prevents “changed” due to float noise)
            new_cmp = {
                "revenue_target": round(new_row["revenue_target"], 2),
                "payroll_target": round(new_row["payroll_target"], 2),
                "opex_target": round(new_row["opex_target"], 2),
                "headcount_fte_target": round(new_row["headcount_fte_target"], 2),
            }

            old = existing.get(month_key)
            if old is None:
                rows_to_upsert.append(new_row)
            else:
                old_cmp = {
                    "revenue_target": round(float(old["revenue_target"]), 2),
                    "payroll_target": round(float(old["payroll_target"]), 2),
                    "opex_target": round(float(old["opex_target"]), 2),
                    "headcount_fte_target": round(float(old["headcount_fte_target"]), 2),
                }

                changed = (
                    new_cmp["revenue_target"] != old_cmp["revenue_target"]
                    or new_cmp["payroll_target"] != old_cmp["payroll_target"]
                    or new_cmp["opex_target"] != old_cmp["opex_target"]
                    or new_cmp["headcount_fte_target"] != old_cmp["headcount_fte_target"]
                )
                if changed:
                    rows_to_upsert.append(new_row)

        if not rows_to_upsert:
            st.info("No changes detected — nothing to save.")
        else:
            try:
                supabase.table("client_monthly_targets").upsert(
                    rows_to_upsert,
                    on_conflict="client_id,month_date",
                ).execute()

                try:
                    st.cache_data.clear()
                except Exception:
                    pass

                st.success(f"Monthly targets saved ({len(rows_to_upsert)} updated).")
                st.rerun()

            except Exception as e:
                st.error(f"Could not save monthly targets: {e}")





def build_team_decision_panel(
    runway_months: float | None,
    runway_min: float | None,
    payroll_vs_rev_pct: float | None,
    delta_payroll_6m: float,
):
    decisions = []
    severity = "info"  # info / warning / error

    # Runway rule
    if runway_months is not None and runway_min is not None:
        if runway_months < runway_min:
            decisions.append(
                f"Runway is **{runway_months:.1f} months**, below your safety threshold "
                f"of **{runway_min:.1f} months**."
            )
            decisions.append("Hiring freeze recommended until runway improves.")
            severity = "error"

    # Payroll vs revenue rule
    if payroll_vs_rev_pct is not None:
        if payroll_vs_rev_pct > 100:
            decisions.append(
                f"Payroll is **{payroll_vs_rev_pct:.0f}% of revenue**, which is not sustainable "
                f"without near-term revenue growth."
            )
            severity = max(severity, "warning")

    # Forward cost rule
    if delta_payroll_6m > 0:
        decisions.append(
            f"Payroll is expected to increase by **${delta_payroll_6m:,.0f}** over the next 6 months."
        )

    if not decisions:
        decisions.append(
            "Team cost is currently aligned with runway and revenue assumptions."
        )

    return severity, decisions

def build_team_spending_story(currency_symbol, month_label, month_payroll, month_ap, payroll_share, payroll_vs_rev_pct, delta_payroll_6m):
    lines = []
    lines.append(f"Payroll is {currency_symbol}{month_payroll:,.0f} in {month_label}.")
    if month_ap > 0:
        lines.append(f"Supplier bills expected this month are {currency_symbol}{month_ap:,.0f}.")
    lines.append(f"Payroll is {payroll_share:.0f}% of tracked cash-out on this page.")

    if payroll_vs_rev_pct is None:
        lines.append("Revenue for this month is not available yet, so payroll-to-revenue can’t be assessed.")
    else:
        if payroll_vs_rev_pct >= 150:
            lines.append(f"Payroll is {payroll_vs_rev_pct:.0f}% of revenue → **high burn risk** unless revenue ramps fast.")
        elif payroll_vs_rev_pct >= 80:
            lines.append(f"Payroll is {payroll_vs_rev_pct:.0f}% of revenue → monitor hiring pace and collections closely.")
        else:
            lines.append(f"Payroll is {payroll_vs_rev_pct:.0f}% of revenue → broadly healthy for an early-stage team (context matters).")

    # Forward trend
    if delta_payroll_6m > 0:
        lines.append(f"Forward view: payroll rises by {currency_symbol}{delta_payroll_6m:,.0f} over the next 6 months (new roles / raises / FTE changes).")
    elif delta_payroll_6m < 0:
        lines.append(f"Forward view: payroll falls by {currency_symbol}{abs(delta_payroll_6m):,.0f} over the next 6 months (roles ending / cost reduction).")
    else:
        lines.append("Forward view: payroll is broadly flat over the next 6 months.")

    # Actions
    actions = [
        "Confirm every planned role maps to a revenue milestone or delivery deadline.",
        "If runway is tight: pause non-critical hiring and review contractors.",
        "Check role end-dates and FTE assumptions are accurate (these drive the cash engine).",
    ]
    return lines, actions





def build_14_week_cash_table(
    client_id: str,
    focus_month: date,
    opening_cash_hint: float | None = None,
) -> pd.DataFrame:
    """
    14-week forward cash view with a TRUE running balance:
      closing_cash[w] = opening_cash[w] + net_cash[w]

    FIXES:
      - Anchors week 0 opening cash to the SAME plan anchor (bank truth / override / seed)
      - Includes Investing + Financing in weekly net cash (distributed from monthly totals)
      - Returns both opening_cash and closing_cash (so UI isn't "just a calculated net" number)

    Recommended UI:
      - Show "Opening cash (bank/anchor)" and "Closing cash (engine)" clearly.
    """
    if not client_id or focus_month is None:
        return pd.DataFrame()

    focus_start = focus_month.replace(day=1)
    focus_start_ts = pd.to_datetime(focus_start)
    as_of_date = focus_start_ts.date()
    as_of_cutoff = focus_start_ts + MonthEnd(0)

    horizon_weeks = 14
    week_starts = [focus_start_ts + timedelta(days=7 * i) for i in range(horizon_weeks)]
    horizon_end_ts = week_starts[-1] + timedelta(days=7)

    # -------------------------------
    # Opening cash anchor (same as monthly engine)
    # -------------------------------
    opening_cash_0 = resolve_opening_cash(
        client_id=str(client_id),
        base_month_ts=focus_start_ts,
        as_of_date=as_of_date,
        opening_cash_hint=opening_cash_hint,
    )

    # -------------------------------
    # Settings + AR/AP data
    # -------------------------------
    settings = get_client_settings(client_id) or {}
    ar_days = int(settings.get("ar_default_days", 30) or 30)
    ap_days = int(settings.get("ap_default_days", 30) or 30)

    df_ar_raw, df_ap_raw = fetch_ar_ap_for_client(client_id)
    df_ar = normalise_ar_ap_for_cash(df_ar_raw) if df_ar_raw is not None else pd.DataFrame()
    df_ap = normalise_ar_ap_for_cash(df_ap_raw) if df_ap_raw is not None else pd.DataFrame()

    # Payments table
    try:
        pay_df = fetch_ar_ap_payments_for_client(client_id)
        if pay_df is None:
            pay_df = pd.DataFrame()
    except Exception:
        pay_df = pd.DataFrame()

    def _week_index(dts: pd.Series) -> pd.Series:
        return ((dts - focus_start_ts).dt.days // 7).astype("int64")

    def _clip_week_index(ix: pd.Series) -> pd.Series:
        return ix.where((ix >= 0) & (ix < horizon_weeks), other=pd.NA)

    def _safe_to_dt(df: pd.DataFrame, col: str) -> None:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce")

    def _safe_to_num(df: pd.DataFrame, col: str) -> None:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0.0)

    # -------------------------------
    # 1) AR cash-in by week
    # -------------------------------
    cash_in_ar = {i: 0.0 for i in range(horizon_weeks)}

    # A) payments table AR
    if not pay_df.empty and {"ar_ap_id", "payment_date"}.issubset(pay_df.columns) and not df_ar.empty and "id" in df_ar.columns:
        ar_pay = pay_df.copy()
        _safe_to_dt(ar_pay, "payment_date")
        _safe_to_num(ar_pay, "amount")

        ar_ids = set(df_ar["id"].astype(str).tolist())
        ar_pay["ar_ap_id_str"] = ar_pay["ar_ap_id"].astype(str)
        ar_pay = ar_pay[ar_pay["ar_ap_id_str"].isin(ar_ids)]
        ar_pay = ar_pay[ar_pay["payment_date"].notna()].copy()

        ar_pay["effective_date"] = ar_pay["payment_date"].where(ar_pay["payment_date"] >= focus_start_ts, focus_start_ts)
        ar_pay["week_index"] = _clip_week_index(_week_index(ar_pay["effective_date"]))
        ar_pay = ar_pay[ar_pay["week_index"].notna()].copy()

        if not ar_pay.empty:
            g = ar_pay.groupby("week_index", as_index=False)["amount"].sum()
            for _, r in g.iterrows():
                cash_in_ar[int(r["week_index"])] += float(r["amount"])

    # B) tracker-based AR (if your tracker stores amount_paid/payment_date on the row)
    if not df_ar.empty and {"payment_date", "amount_paid"}.issubset(df_ar.columns):
        ar_tr = df_ar.copy()
        _safe_to_dt(ar_tr, "payment_date")
        _safe_to_num(ar_tr, "amount_paid")

        paid = ar_tr[(ar_tr["amount_paid"] > 0) & ar_tr["payment_date"].notna()].copy()
        paid = paid[(paid["payment_date"] >= focus_start_ts) & (paid["payment_date"] < horizon_end_ts)].copy()

        if not paid.empty:
            paid["week_index"] = _clip_week_index(_week_index(paid["payment_date"]))
            paid = paid[paid["week_index"].notna()].copy()
            g = paid.groupby("week_index", as_index=False)["amount_paid"].sum()
            for _, r in g.iterrows():
                cash_in_ar[int(r["week_index"])] += float(r["amount_paid"])

    # C) unpaid scheduled AR
    if not df_ar.empty:
        ar = df_ar.copy()
        try:
            ar = add_ar_aging(ar, as_of=(focus_start_ts + MonthEnd(0)).date(), ar_default_days=ar_days)
        except Exception:
            pass

        date_col = "due_date" if "due_date" in ar.columns else ("expected_date" if "expected_date" in ar.columns else None)
        if date_col:
            _safe_to_dt(ar, date_col)
            value_col = "effective_amount" if "effective_amount" in ar.columns else "amount"
            _safe_to_num(ar, value_col)

            ar = ar[ar[date_col].notna()].copy()
            ar = ar[ar[value_col] > 0].copy()

            # reduce double-count (payments table)
            if not pay_df.empty and "ar_ap_id" in pay_df.columns and "id" in ar.columns:
                p = pay_df.copy()
                _safe_to_dt(p, "payment_date")
                _safe_to_num(p, "amount")
                p["ar_ap_id_str"] = p["ar_ap_id"].astype(str)

                p = p[
                    (p["payment_date"].notna())
                    & (p["payment_date"] >= focus_start_ts)
                    & (p["payment_date"] < horizon_end_ts)
                ].copy()

                if not p.empty:
                    paid_map = p.groupby("ar_ap_id_str", as_index=False)["amount"].sum().rename(columns={"amount": "__paid__"})
                    ar["id_str"] = ar["id"].astype(str)
                    ar = ar.merge(paid_map, left_on="id_str", right_on="ar_ap_id_str", how="left")
                    ar["__paid__"] = ar["__paid__"].fillna(0.0)
                    ar[value_col] = (ar[value_col] - ar["__paid__"]).clip(lower=0.0)

            # overdue => week 0
            ar["effective_date"] = ar[date_col].where(ar[date_col] >= focus_start_ts, focus_start_ts)
            ar["week_index"] = _clip_week_index(_week_index(ar["effective_date"]))
            ar = ar[ar["week_index"].notna()].copy()
            ar = ar[ar[value_col] > 0].copy()

            if not ar.empty:
                g = ar.groupby("week_index", as_index=False)[value_col].sum()
                for _, r in g.iterrows():
                    cash_in_ar[int(r["week_index"])] += float(r[value_col])

    # -------------------------------
    # 2) AP cash-out by week
    # -------------------------------
    cash_out_ap = {i: 0.0 for i in range(horizon_weeks)}

    # A) payments table AP
    if not pay_df.empty and {"ar_ap_id", "payment_date"}.issubset(pay_df.columns) and not df_ap.empty and "id" in df_ap.columns:
        ap_pay = pay_df.copy()
        _safe_to_dt(ap_pay, "payment_date")
        _safe_to_num(ap_pay, "amount")

        ap_ids = set(df_ap["id"].astype(str).tolist())
        ap_pay["ar_ap_id_str"] = ap_pay["ar_ap_id"].astype(str)
        ap_pay = ap_pay[ap_pay["ar_ap_id_str"].isin(ap_ids)]
        ap_pay = ap_pay[ap_pay["payment_date"].notna()].copy()

        ap_pay["effective_date"] = ap_pay["payment_date"].where(ap_pay["payment_date"] >= focus_start_ts, focus_start_ts)
        ap_pay["week_index"] = _clip_week_index(_week_index(ap_pay["effective_date"]))
        ap_pay = ap_pay[ap_pay["week_index"].notna()].copy()

        if not ap_pay.empty:
            g = ap_pay.groupby("week_index", as_index=False)["amount"].sum()
            for _, r in g.iterrows():
                cash_out_ap[int(r["week_index"])] += float(r["amount"])

    # B) tracker-based AP
    if not df_ap.empty and {"payment_date", "amount_paid"}.issubset(df_ap.columns):
        ap_tr = df_ap.copy()
        _safe_to_dt(ap_tr, "payment_date")
        _safe_to_num(ap_tr, "amount_paid")

        paid = ap_tr[(ap_tr["amount_paid"] > 0) & ap_tr["payment_date"].notna()].copy()
        paid = paid[(paid["payment_date"] >= focus_start_ts) & (paid["payment_date"] < horizon_end_ts)].copy()

        if not paid.empty:
            paid["week_index"] = _clip_week_index(_week_index(paid["payment_date"]))
            paid = paid[paid["week_index"].notna()].copy()
            g = paid.groupby("week_index", as_index=False)["amount_paid"].sum()
            for _, r in g.iterrows():
                cash_out_ap[int(r["week_index"])] += float(r["amount_paid"])

    # C) unpaid scheduled AP
    if not df_ap.empty:
        ap = df_ap.copy()
        try:
            ap = add_ap_aging(ap, as_of=(focus_start_ts + MonthEnd(0)).date(), ap_default_days=ap_days)
        except Exception:
            pass

        date_col_ap = None
        for c in ["expected_payment_date", "pay_expected_date", "due_date"]:
            if c in ap.columns:
                date_col_ap = c
                break

        if date_col_ap:
            _safe_to_dt(ap, date_col_ap)
            value_col_ap = "effective_amount" if "effective_amount" in ap.columns else "amount"
            _safe_to_num(ap, value_col_ap)

            ap = ap[ap[date_col_ap].notna()].copy()
            ap = ap[ap[value_col_ap] > 0].copy()

            # overdue => week 0
            ap["effective_date"] = ap[date_col_ap].where(ap[date_col_ap] >= focus_start_ts, focus_start_ts)
            ap["week_index"] = _clip_week_index(_week_index(ap["effective_date"]))
            ap = ap[ap["week_index"].notna()].copy()
            ap = ap[ap[value_col_ap] > 0].copy()

            if not ap.empty:
                g = ap.groupby("week_index", as_index=False)[value_col_ap].sum()
                for _, r in g.iterrows():
                    cash_out_ap[int(r["week_index"])] += float(r[value_col_ap])

    # -------------------------------
    # 3) Payroll + Opex + Other income weekly (spread evenly)
    # -------------------------------
    # We'll compute monthly totals for months intersecting horizon, then distribute across weeks that fall in that month.
    months_in_horizon = pd.date_range(start=focus_start_ts, end=horizon_end_ts, freq="MS")

    # Payroll monthly totals -> weekly
    payroll_by_month = compute_payroll_by_month(client_id, months_in_horizon)
    cash_out_payroll = {i: 0.0 for i in range(horizon_weeks)}

    # Opex monthly totals -> weekly
    try:
        df_opex = fetch_opex_for_client(client_id)
        if df_opex is None or df_opex.empty:
            opex_by_month = {}
        else:
            df_opex = df_opex.copy()
            df_opex["month_bucket"] = (
                pd.to_datetime(df_opex["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
            )
            df_opex["amount"] = pd.to_numeric(df_opex.get("amount"), errors="coerce").fillna(0.0)
            opex_agg = df_opex.groupby("month_bucket", as_index=False)["amount"].sum()
            opex_by_month = {r["month_bucket"]: float(r["amount"] or 0.0) for _, r in opex_agg.iterrows()}
    except Exception as e:
        print("[WARN] weekly opex fetch failed:", e)
        opex_by_month = {}

    cash_out_opex = {i: 0.0 for i in range(horizon_weeks)}

    # Other income monthly totals -> weekly
    try:
        df_oinc = fetch_operating_other_income_for_client(client_id)
        if df_oinc is None or df_oinc.empty:
            oinc_by_month = {}
        else:
            df_oinc = df_oinc.copy()
            df_oinc["month_bucket"] = (
                pd.to_datetime(df_oinc["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
            )
            df_oinc["cash_in"] = pd.to_numeric(df_oinc.get("cash_in"), errors="coerce").fillna(0.0)
            oinc_agg = df_oinc.groupby("month_bucket", as_index=False)["cash_in"].sum()
            oinc_by_month = {r["month_bucket"]: float(r["cash_in"] or 0.0) for _, r in oinc_agg.iterrows()}
    except Exception as e:
        print("[WARN] weekly other income fetch failed:", e)
        oinc_by_month = {}

    cash_in_other = {i: 0.0 for i in range(horizon_weeks)}

    # Helper: weeks per month inside horizon
    def _month_bucket(ts: pd.Timestamp) -> pd.Timestamp:
        return ts.to_period("M").to_timestamp()

    weeks_in_month = {}
    for i, ws in enumerate(week_starts):
        mb = _month_bucket(ws)
        weeks_in_month.setdefault(mb, []).append(i)

    for mb, week_list in weeks_in_month.items():
        denom = max(len(week_list), 1)

        # payroll
        p_month = float(payroll_by_month.get(mb, 0.0))
        for wi in week_list:
            cash_out_payroll[wi] += p_month / denom

        # opex
        ox_month = float(opex_by_month.get(mb, 0.0))
        for wi in week_list:
            cash_out_opex[wi] += ox_month / denom

        # other income
        oi_month = float(oinc_by_month.get(mb, 0.0))
        for wi in week_list:
            cash_in_other[wi] += oi_month / denom

    # -------------------------------
    # 4) Investing + Financing weekly (distributed from monthly totals)
    # -------------------------------
    inv_week = {i: 0.0 for i in range(horizon_weeks)}
    fin_week = {i: 0.0 for i in range(horizon_weeks)}

    # Load investing + financing (versioned same way as monthly)
    def _monthly_flow_map(df: pd.DataFrame, label: str) -> dict[pd.Timestamp, float]:
        if df is None or df.empty:
            return {}

        df0 = df.copy()
        orig = df0.copy()

        if "created_at" in df0.columns:
            df0["created_at"] = pd.to_datetime(df0["created_at"], errors="coerce")
            try:
                df0["created_at"] = df0["created_at"].dt.tz_convert(None)
            except TypeError:
                pass

            before = len(df0)
            df0 = df0[df0["created_at"] <= as_of_cutoff]
            after = len(df0)
            print(f"[DEBUG] {label} created_at filter: {before} -> {after} rows (as_of_cutoff={as_of_cutoff.date()})")
            if after == 0:
                print(f"[WARN] {label} all filtered out by created_at; falling back to all rows (no versioning).")
                df0 = orig

        df0["month_bucket"] = (
            pd.to_datetime(df0["month_date"], errors="coerce").dt.to_period("M").dt.to_timestamp()
        )
        df0["amount"] = pd.to_numeric(df0.get("amount"), errors="coerce").fillna(0.0)

        agg = df0.groupby("month_bucket", as_index=False)["amount"].sum()
        return {r["month_bucket"]: float(r["amount"] or 0.0) for _, r in agg.iterrows()}

    try:
        df_inv = fetch_investing_flows_for_client(client_id)
    except Exception:
        df_inv = pd.DataFrame()

    try:
        df_fin = fetch_financing_flows_for_client(client_id)
    except Exception:
        df_fin = pd.DataFrame()

    inv_by_month = _monthly_flow_map(df_inv, "investing_flows")
    fin_by_month = _monthly_flow_map(df_fin, "financing_flows")

    for mb, week_list in weeks_in_month.items():
        denom = max(len(week_list), 1)
        inv_m = float(inv_by_month.get(mb, 0.0))
        fin_m = float(fin_by_month.get(mb, 0.0))
        for wi in week_list:
            inv_week[wi] += inv_m / denom
            fin_week[wi] += fin_m / denom

    # -------------------------------
    # 5) Build the table with OPENING + CLOSING (running)
    # -------------------------------
    rows = []
    running_close = float(opening_cash_0)

    for i, ws in enumerate(week_starts):
        open_cash = float(running_close)

        op_cf = float(cash_in_ar[i] + cash_in_other[i] - cash_out_ap[i] - cash_out_payroll[i] - cash_out_opex[i])
        inv_cf = float(inv_week[i])
        fin_cf = float(fin_week[i])
        net_cash = op_cf + inv_cf + fin_cf
        close_cash = open_cash + net_cash

        rows.append(
            {
                "week_start": ws.date().isoformat(),
                "opening_cash": round(open_cash, 2),
                "cash_in_ar": round(float(cash_in_ar[i]), 2),
                "cash_in_other_income": round(float(cash_in_other[i]), 2),
                "cash_out_ap": round(float(cash_out_ap[i]), 2),
                "cash_out_payroll": round(float(cash_out_payroll[i]), 2),
                "cash_out_opex": round(float(cash_out_opex[i]), 2),
                "operating_cf": round(op_cf, 2),
                "investing_cf": round(inv_cf, 2),
                "financing_cf": round(fin_cf, 2),
                "net_cash": round(net_cash, 2),
                "closing_cash": round(close_cash, 2),
                "cash_danger_flag": (close_cash <= 0.0),
            }
        )

        running_close = float(close_cash)

    return pd.DataFrame(rows)



def _to_naive(ts) -> pd.Timestamp:
    """Convert any datetime-like to tz-naive pandas Timestamp (UTC kept as clock time)."""
    t = pd.to_datetime(ts, errors="coerce")
    if t is pd.NaT:
        return pd.NaT
    # If tz-aware, drop tz info
    try:
        if getattr(t, "tzinfo", None) is not None:
            t = t.tz_localize(None)
    except Exception:
        pass
    return t


def _days_in_month(month_start: pd.Timestamp) -> int:
    month_start = pd.to_datetime(month_start).to_period("M").to_timestamp()
    month_end = month_start + MonthEnd(0)
    return int((month_end - month_start).days) + 1

def _week_floor(dt) -> pd.Timestamp:
    """Return Monday of the week for a date/timestamp (week starts Monday)."""
    ts = _to_naive(dt)
    if ts is pd.NaT:
        return pd.NaT
    return (ts - pd.to_timedelta(ts.weekday(), unit="D")).normalize()

def fetch_opening_cash_override(client_id: str, month_date: date) -> float | None:
    """
    Returns a manual opening cash override for a given client and month (first-of-month).
    Use this when you want 'bank truth' to win over computed / seeded values.
    """
    if not client_id or not month_date:
        return None

    m = pd.to_datetime(month_date).replace(day=1).date()

    try:
        res = (
            supabase.table("opening_cash_override")
            .select("opening_cash")
            .eq("client_id", str(client_id))
            .eq("month_date", m.isoformat())
            .order("created_at", desc=True)
            .limit(1)
            .execute()
        )
        data = getattr(res, "data", None)
        if data:
            v = data[0].get("opening_cash")
            if v is None:
                return None
            return float(v)
    except Exception as e:
        print("[WARN] fetch_opening_cash_override failed:", e)

    return None



# ---------- Page: Cash & Bills ----------

def _estimate_breach_date_in_month(
    month_start: pd.Timestamp,
    opening_cash: float,
    closing_cash: float,
    buffer: float,
) -> pd.Timestamp | None:
    """
    Estimate the date within the month where cash crosses below buffer.
    Assumes linear change from opening_cash -> closing_cash across the month.
    """
    month_start = pd.to_datetime(month_start).to_period("M").to_timestamp()
    if opening_cash < buffer:
        return month_start.normalize()

    # if it never crosses this month, return None
    if closing_cash >= buffer:
        return None

    delta = float(closing_cash) - float(opening_cash)  # likely negative
    if abs(delta) < 1e-9:
        # flat but somehow below at close; treat as month end
        return (month_start + MonthEnd(0)).normalize()

    # fraction of month when it hits buffer
    frac = (float(buffer) - float(opening_cash)) / delta  # delta negative -> frac between 0..1
    frac = max(0.0, min(1.0, frac))

    dim = _days_in_month(month_start)
    # Use day offset inside month (0..dim-1)
    day_offset = int(round(frac * (dim - 1)))
    breach_date = (month_start + pd.Timedelta(days=day_offset)).normalize()
    return breach_date

def _estimate_breach_date_in_week(
    week_start: pd.Timestamp,
    opening_cash: float,
    closing_cash: float,
    threshold: float,
) -> pd.Timestamp | None:
    """
    Estimate the first calendar date (within that week) when cash falls below threshold.
    Uses linear interpolation between opening_cash and closing_cash across 7 days.
    Returns a tz-naive pd.Timestamp (date at midnight), or None.
    """
    if week_start is None or pd.isna(week_start):
        return None

    ws = _to_naive(pd.to_datetime(week_start, errors="coerce"))
    if ws is pd.NaT:
        return None

    try:
        o = float(opening_cash or 0.0)
        c = float(closing_cash or 0.0)
        th = float(threshold or 0.0)
    except Exception:
        return None

    # Already below at start of week
    if o < th:
        return ws.normalize()

    # Never breaches within this week
    if c >= th:
        return None

    denom = (c - o)
    if denom == 0:
        # Flat line but c < th implies o < th would have been true; safe fallback
        return ws.normalize()

    # Solve o + f*(c-o) = th  => f = (th - o) / (c - o)
    f = (th - o) / denom
    f = max(0.0, min(1.0, float(f)))

    # Convert fraction of week to day index (0..6)
    day_index = int(math.floor(f * 7.0))
    day_index = max(0, min(6, day_index))

    return (ws + pd.Timedelta(days=day_index)).normalize()




def compute_cash_stress_kpis(
    week_df: pd.DataFrame,
    cash_buffer: float,
    today=None,
    *,
    client_id: str | None = None,
    as_of_month=None,
    fetch_cashflow_summary_for_client_as_of=None,
    horizon_months: int = 12,
) -> dict:
    """
    Base-case only.

    Primary (weekly):
      - Stress week comes from first week_start where closing_cash < buffer.

    Fallback (monthly plan engine):
      - If NO buffer breach in the next 14 weeks, estimate exact stress date from
        cashflow_summary (filtered by as_of_month plan version).

    KEY FIX:
      - Runway / reaction window must be measured relative to the *projection anchor*
        (focus month / first week in week_df), NOT relative to "today".
    """
    out = {
        "cash_buffer": float(cash_buffer or 0.0),
        "stress_week": None,               # backward compatible (week-start)
        "runway_weeks": None,
        "runway_months": None,
        "reaction_window_weeks": None,
        "reaction_band": None,
        "immediate_attention": False,

        # new
        "stress_date_exact": None,
        "stress_week_start": None,
        "stress_week_label": None,
        "stress_source": None,

        # extra debug-friendly anchors
        "anchor_week": None,               # week we measure runway from
        "anchor_source": None,             # "weekly_projection" / "focus_month" / "today"
    }

    buffer = out["cash_buffer"]
    if buffer <= 0:
        return out

    # ---------- helpers ----------
    def _naive_ts(x):
        x = pd.to_datetime(x, errors="coerce")
        if x is pd.NaT:
            return pd.NaT
        try:
            return x.tz_localize(None)
        except Exception:
            return x

    # Anchor week preference:
    # 1) week_df first week (projection start)  ✅ best
    # 2) focus month start week                ✅ good
    # 3) today week                            last resort
    anchor_week = pd.NaT
    anchor_source = None

    # Try to anchor from week_df
    if week_df is not None and not week_df.empty and "week_start" in week_df.columns:
        w0 = week_df.copy()
        w0["week_start"] = w0["week_start"].apply(_naive_ts)
        w0 = w0.dropna(subset=["week_start"]).sort_values("week_start")
        if not w0.empty:
            anchor_week = _week_floor(w0.iloc[0]["week_start"]).normalize()
            anchor_source = "weekly_projection"

    # If not available, anchor from focus month
    if anchor_week is pd.NaT and as_of_month is not None:
        m = pd.to_datetime(as_of_month, errors="coerce")
        if m is not pd.NaT:
            anchor_week = _week_floor(m.to_period("M").to_timestamp()).normalize()
            anchor_source = "focus_month"

    # Fallback to today
    if anchor_week is pd.NaT:
        if today is None:
            today = pd.Timestamp.today()
        anchor_week = _week_floor(_naive_ts(today)).normalize()
        anchor_source = "today"

    out["anchor_week"] = anchor_week
    out["anchor_source"] = anchor_source

    # -------------------------
    # 1) Weekly (14-week) logic
    # -------------------------
    if week_df is not None and not week_df.empty and {"week_start", "closing_cash"}.issubset(set(week_df.columns)):
        w = week_df.copy()
        w["week_start"] = w["week_start"].apply(_naive_ts)
        w["closing_cash"] = pd.to_numeric(w["closing_cash"], errors="coerce").fillna(0.0)
        w = w.dropna(subset=["week_start"]).sort_values("week_start")

        if not w.empty:
            stress_rows = w[w["closing_cash"] < buffer]
            if not stress_rows.empty:
                stress_week_start = _week_floor(stress_rows.iloc[0]["week_start"]).normalize()

                out["stress_week"] = stress_week_start
                out["stress_date_exact"] = stress_week_start  # weekly precision
                out["stress_week_start"] = stress_week_start
                out["stress_week_label"] = f"Wk of {stress_week_start.strftime('%d %b %Y')}"
                out["stress_source"] = "weekly"

                # ✅ FIX: runway is from anchor_week (projection start), not today
                delta_days = (stress_week_start - anchor_week).days
                runway_weeks = int(max(0, math.ceil(delta_days / 7)))

                out["runway_weeks"] = runway_weeks
                out["runway_months"] = round(runway_weeks / 4.345, 1)

                # reaction window = same as runway for base-case
                out["reaction_window_weeks"] = runway_weeks

                rw = runway_weeks
                if rw >= 12:
                    band = "Comfortable"
                elif 8 <= rw < 12:
                    band = "Watch"
                elif 4 <= rw < 8:
                    band = "Tight"
                else:
                    band = "Critical"
                out["reaction_band"] = band

                # ✅ "Immediate" should be relative to runway from anchor
                out["immediate_attention"] = (rw <= 8)

                return out

    # ---------------------------------------
    # 2) Fallback: monthly plan engine estimate
    # ---------------------------------------
    if client_id and as_of_month is not None and callable(fetch_cashflow_summary_for_client_as_of):
        try:
            as_of_month_ts = pd.to_datetime(as_of_month).to_period("M").to_timestamp()
            dfm = fetch_cashflow_summary_for_client_as_of(str(client_id), as_of_month_ts.date())
        except Exception:
            dfm = None

        if dfm is not None and not dfm.empty:
            dfm = dfm.copy()

            if "month_date" in dfm.columns:
                dfm["month_date"] = (
                    pd.to_datetime(dfm["month_date"], errors="coerce")
                    .dt.to_period("M")
                    .dt.to_timestamp()
                )
            if "opening_cash" in dfm.columns:
                dfm["opening_cash"] = pd.to_numeric(dfm["opening_cash"], errors="coerce").fillna(0.0)
            if "closing_cash" in dfm.columns:
                dfm["closing_cash"] = pd.to_numeric(dfm["closing_cash"], errors="coerce").fillna(0.0)

            dfm = dfm.dropna(subset=["month_date"]).sort_values("month_date")

            breach = dfm[dfm["closing_cash"] < buffer]
            if not breach.empty and {"month_date", "opening_cash", "closing_cash"}.issubset(breach.columns):
                r0 = breach.iloc[0]
                m0 = pd.to_datetime(r0["month_date"]).to_period("M").to_timestamp()
                open0 = float(r0.get("opening_cash", 0.0))
                close0 = float(r0.get("closing_cash", 0.0))

                stress_date_exact = _estimate_breach_date_in_month(
                    month_start=m0,
                    opening_cash=open0,
                    closing_cash=close0,
                    buffer=buffer,
                )

                if stress_date_exact is not None:
                    stress_date_exact = _naive_ts(stress_date_exact).normalize()
                    stress_week_start = _week_floor(stress_date_exact).normalize()

                    out["stress_date_exact"] = stress_date_exact
                    out["stress_week_start"] = stress_week_start
                    out["stress_week"] = stress_week_start
                    out["stress_week_label"] = f"Wk of {stress_week_start.strftime('%d %b %Y')}"
                    out["stress_source"] = "monthly_estimate"

                    # ✅ FIX: runway from anchor_week (projection/focus month), not today
                    delta_days = (stress_week_start - anchor_week).days
                    runway_weeks = int(max(0, math.ceil(delta_days / 7)))

                    out["runway_weeks"] = runway_weeks
                    out["runway_months"] = round(runway_weeks / 4.345, 1)
                    out["reaction_window_weeks"] = runway_weeks

                    rw = runway_weeks
                    if rw >= 12:
                        band = "Comfortable"
                    elif 8 <= rw < 12:
                        band = "Watch"
                    elif 4 <= rw < 8:
                        band = "Tight"
                    else:
                        band = "Critical"
                    out["reaction_band"] = band
                    out["immediate_attention"] = (rw <= 8)

                    return out

    return out



def page_cash_bills():
    top_header("Cash & Bills")

    if not selected_client_id:
        st.info("Select a business from the sidebar to view this page.")
        return

    from datetime import timedelta
    import math

    # ✅ Ensure these ALWAYS exist (prevents UnboundLocalError forever)
    week_df = pd.DataFrame()
    df_ar, df_ap = pd.DataFrame(), pd.DataFrame()
    runway_weeks = None
    runway_months = None
    reaction_weeks = None
    band = None
    stress_week = None
    immediate = False

    # NEW (safe defaults)
    stress_date_exact = None
    stress_source = None
    stress_week_label = None

    # Currency for formatting
    currency_code, currency_symbol = get_client_currency(selected_client_id)

    # ==========================
    # Anchor everything to Focus Month
    # ==========================
    focus_start = selected_month_start
    focus_end_ts = pd.to_datetime(focus_start) + MonthEnd(0)
    as_of = focus_end_ts.date()

    horizon_4w_start = focus_start
    horizon_4w_end = horizon_4w_start + timedelta(days=28)

    # --- Load settings once (buffer + defaults used across the page) ---
    settings = get_client_settings(selected_client_id) or {}
    cash_buffer = float(settings.get("min_cash_buffer") or 0.0)
    ar_days = int(settings.get("ar_default_days", 30) or 30)
    ap_days = int(settings.get("ap_default_days", 30) or 30)

    # Keep AR overdue alerts up to date
    ensure_overdue_ar_alert(
        selected_client_id,
        as_of=as_of,
        min_days_overdue=0,
    )

    # ==========================
    # Load AR/AP + 14-week cash table FIRST
    # ==========================
    with st.spinner("Loading cash, invoices and commitments..."):
        try:
            df_ar, df_ap = fetch_ar_ap_for_client(selected_client_id)
        except Exception as e:
            st.error(f"Failed to load AR/AP: {e}")
            df_ar, df_ap = pd.DataFrame(), pd.DataFrame()

        try:
            week_df = build_14_week_cash_table(selected_client_id, selected_month_start)
            if week_df is None:
                week_df = pd.DataFrame()
        except Exception as e:
            week_df = pd.DataFrame()
            st.error(f"Failed to build 14-week cash table: {e}")

    # Normalise
    try:
        df_ar = normalise_ar_ap_for_cash(df_ar)
    except Exception:
        pass
    try:
        df_ap = normalise_ar_ap_for_cash(df_ap)
    except Exception:
        pass

    # Add ageing info
    if df_ar is not None and not df_ar.empty:
        try:
            df_ar = add_ar_aging(df_ar, as_of=as_of, ar_default_days=ar_days)
        except Exception:
            pass

    if df_ap is not None and not df_ap.empty:
        try:
            df_ap = add_ap_aging(df_ap, as_of=as_of, ap_default_days=ap_days)
        except Exception:
            pass

    # ==========================
    # KPI STRIP: Runway / Stress Date / Reaction Window (Base-case)
    # ==========================
    st.markdown("### 🧭 Cash stress & runway (base-case)")

    if cash_buffer <= 0:
        st.warning("Cash Stress Buffer is not set (or is 0). Set it in Client Settings → Cash safety buffer.")
    elif week_df is None or week_df.empty:
        st.info("14-week cashflow not available yet — cannot calculate runway/stress date.")
    else:
        # ✅ UPDATED: pass focus-month plan version + fetcher so we can compute exact stress date > 14 weeks
        stress_kpis = compute_cash_stress_kpis(
            week_df=week_df,
            cash_buffer=cash_buffer,
            today=pd.Timestamp.utcnow().tz_localize(None),
            client_id=selected_client_id,
            as_of_month=focus_start,  # Focus Month start
            fetch_cashflow_summary_for_client_as_of=fetch_cashflow_summary_for_client_as_of,
        )

        # Backward-compatible fields
        stress_week = stress_kpis.get("stress_week")  # week-start (Monday) even if estimated from monthly
        runway_weeks = stress_kpis.get("runway_weeks")
        runway_months = stress_kpis.get("runway_months")
        reaction_weeks = stress_kpis.get("reaction_window_weeks")
        band = stress_kpis.get("reaction_band")
        immediate = bool(stress_kpis.get("immediate_attention", False))

        # NEW fields
        stress_date_exact = stress_kpis.get("stress_date_exact")
        stress_source = stress_kpis.get("stress_source")
        stress_week_label = stress_kpis.get("stress_week_label")

        c2, c3 = st.columns([1.3, 1.2])

        with c2:
            # ✅ UPDATED: show exact date if available (even >14 weeks)
            if stress_date_exact is not None:
                st.metric("Cash stress date", stress_date_exact.strftime("%d %b %Y"))
                # show week label + source
                extra = []
                if stress_week_label:
                    extra.append(stress_week_label)
                if stress_source:
                    extra.append(f"source: {stress_source}")
                if extra:
                    st.caption(" • ".join(extra))
            else:
                # fallback as before
                st.metric("Cash stress date", "Not in next 14 weeks")
                st.caption("No buffer breach found in weekly view or plan data.")

        with c3:
            st.metric("Reaction window (weeks)", f"{reaction_weeks} w" if reaction_weeks is not None else "—")

        # Status messages (unchanged logic, but now stress_week might be derived from monthly estimate)
        if stress_week is not None and immediate:
            st.error("Immediate Attention: cash stress is ≤ 8 weeks away under base-case.")
        elif stress_week is not None and band == "Tight":
            st.warning("Tight: reaction window 4–8 weeks (base-case).")
        elif stress_week is not None and band == "Watch":
            st.warning("Watch: reaction window 8–12 weeks (base-case).")
        elif stress_week is not None:
            st.info("Base-case shows cash remains above the buffer for now, but the runway metrics should be used for decisions.")
        else:
            st.info("Base-case shows no buffer breach in the next 14 weeks. Maintain discipline on collections and commitments.")

    # ------------------------------------------------------------------
    # 1) Top KPIs – Cash + AR/AP next 4 weeks
    # ------------------------------------------------------------------
    st.subheader("📌 Cash & working capital snapshot")

    # Helper to keep old overdue logic (do NOT remove)
    ar_overdue_amt = 0.0
    ar_overdue_count = 0
    ap_overdue_amt = 0.0
    ap_overdue_count = 0

    def _not_paid_status(s):
        return str(s or "").lower() not in ["paid", "closed", "settled"]

    # --- Keep overdue counts exactly as before (uses df_ar/df_ap) ---
    if df_ar is not None and not df_ar.empty:
        ar_tmp = df_ar.copy()
        if "days_past_expected" in ar_tmp.columns:
            mask_overdue = (ar_tmp["days_past_expected"] > 0) & ar_tmp["status"].apply(_not_paid_status)
            overdue_rows = ar_tmp[mask_overdue]
            amt_col_overdue_ar = "effective_amount" if "effective_amount" in overdue_rows.columns else "amount"
            ar_overdue_amt = float(overdue_rows[amt_col_overdue_ar].fillna(0.0).sum())
            ar_overdue_count = int(len(overdue_rows))

    if df_ap is not None and not df_ap.empty:
        ap_tmp = df_ap.copy()
        if "days_past_expected" in ap_tmp.columns:
            mask_overdue_ap = (ap_tmp["days_past_expected"] > 0) & ap_tmp["status"].apply(_not_paid_status)
            overdue_rows_ap = ap_tmp[mask_overdue_ap]
            amt_col_overdue_ap = "effective_amount" if "effective_amount" in overdue_rows_ap.columns else "amount"
            ap_overdue_amt = float(overdue_rows_ap[amt_col_overdue_ap].fillna(0.0).sum())
            ap_overdue_count = int(len(overdue_rows_ap))

    # ------------------------------------------------------------------
    # ✅ KPI totals must match the 14-week cashflow table
    # We define "next 4 weeks" as week_index 0..3 (first 4 rows in week_df)
    # ------------------------------------------------------------------
    ar_next_4w = 0.0
    ap_next_4w = 0.0

    if week_df is not None and not week_df.empty:
        w4 = week_df.head(4).copy()

        # Ensure numeric
        for c in ["cash_in_ar", "cash_out_ap"]:
            if c in w4.columns:
                w4[c] = pd.to_numeric(w4[c], errors="coerce").fillna(0.0)

        ar_next_4w = float(w4["cash_in_ar"].sum()) if "cash_in_ar" in w4.columns else 0.0
        ap_next_4w = float(w4["cash_out_ap"].sum()) if "cash_out_ap" in w4.columns else 0.0

    # Approximate "cash in bank now" from the 14-week table if available (keep as-is)
    cash_now_label = "—"
    opening_first = None
    if week_df is not None and not week_df.empty:
        first_row = week_df.iloc[0]
        closing_first = float(first_row.get("closing_cash") or 0.0)
        net_first = float(first_row.get("net_cash") or 0.0)
        opening_first = closing_first - net_first
        cash_now_label = f"{currency_symbol}{opening_first:,.0f}"

    net_4w = ar_next_4w - ap_next_4w

    k1, k2, k3, k4 = st.columns(4)
    with k1:
        st.metric("Projected opening cash (week 1)", cash_now_label)
    with k2:
        st.metric("Invoices to collect (next 4 weeks)", f"{currency_symbol}{ar_next_4w:,.0f}")
    with k3:
        st.metric("Bills to pay (next 4 weeks)", f"{currency_symbol}{ap_next_4w:,.0f}")
    with k4:
        st.metric("Net position (next 4 weeks)", f"{currency_symbol}{net_4w:,.0f}")

    # ==========================
    # Governance: Runway change > ±10% MoM
    # ==========================
    try:
        prev_month_start = (pd.to_datetime(selected_month_start).to_period("M") - 1).to_timestamp()
        prev_week_df = build_14_week_cash_table(selected_client_id, prev_month_start)

        if cash_buffer > 0 and prev_week_df is not None and not prev_week_df.empty and runway_weeks is not None:
            prev = compute_cash_stress_kpis(prev_week_df, cash_buffer=cash_buffer, today=pd.Timestamp.utcnow())
            prev_runway = prev.get("runway_weeks")

            if prev_runway is not None and prev_runway > 0:
                pct_change = (runway_weeks - prev_runway) / prev_runway * 100.0
                if abs(pct_change) > 10:
                    st.warning(f"Governance flag: runway changed {pct_change:+.0f}% vs last month (>{10}% threshold).")
    except Exception:
        pass

    # ------------------------------------------------------------------
    # Cash danger summary strip (from 14-week view)
    # ------------------------------------------------------------------
    st.markdown("### 🚨 Cash danger summary (next 14 weeks)")

    min_cash = None
    worst_week = None

    if week_df is None or week_df.empty or "closing_cash" not in week_df.columns:
        st.info("14-week cashflow not available yet to assess cash risk.")
    else:
        w = week_df.copy()
        if "week_start" in w.columns:
            w["week_start"] = pd.to_datetime(w["week_start"], errors="coerce")

        min_cash = float(w["closing_cash"].min())
        idx_min = int(w["closing_cash"].idxmin())
        worst_week = w.loc[idx_min, "week_start"] if "week_start" in w.columns else None

        buffer_shortfall = max(0.0, cash_buffer - min_cash)
        cash_deficit = max(0.0, -min_cash)

        if min_cash <= 0:
            headline = "Cash goes negative in the next 14 weeks."
            gap_label = "Cash deficit (to reach $0)"
            gap_value = cash_deficit
            st.error(f"• {headline}")
        elif min_cash < cash_buffer:
            headline = "Cash stays positive, but falls below your safety buffer."
            gap_label = "Gap to safety buffer"
            gap_value = buffer_shortfall
            st.warning(f"• {headline}")
        else:
            headline = "Cash stays above your safety buffer."
            gap_label = "Gap to safety buffer"
            gap_value = 0.0
            st.info(f"• {headline}")

        cA, cB, cC, cD = st.columns(4)
        cA.metric("Lowest projected cash", f"{currency_symbol}{min_cash:,.0f}")
        cB.metric("Week of lowest cash", worst_week.strftime("%d %b %Y") if worst_week is not None else "—")
        cC.metric("Safety buffer", f"{currency_symbol}{cash_buffer:,.0f}")
        cD.metric(gap_label, f"{currency_symbol}{gap_value:,.0f}")

    # Short caption for overdue situation
    overdue_bits = []
    if ar_overdue_amt > 0:
        overdue_bits.append(
            f"{ar_overdue_count} invoice(s) overdue from customers "
            f"≈ {currency_symbol}{ar_overdue_amt:,.0f}"
        )
    if ap_overdue_amt > 0:
        overdue_bits.append(
            f"{ap_overdue_count} bill(s) overdue to suppliers "
            f"≈ {currency_symbol}{ap_overdue_amt:,.0f}"
        )
    if overdue_bits:
        st.caption("Overdue snapshot: " + " · ".join(overdue_bits))

    st.markdown("---")

    # ==========================
    # WHAT THIS MEANS (CFO interpretation — HARD RULES)
    # Max 4 bullets. Mandatory coverage: (1) current cash, (2) runway direction,
    # (3) timing of risk, (4) implication for decisions. No optimism language.
    # ==========================
    st.subheader("What this means")

    what_means = []

    # 1) Current cash position
    if opening_first is None:
        what_means.append("• Current cash position cannot be inferred because the 14-week projection is unavailable.")
    else:
        if cash_buffer > 0 and min_cash is not None and min_cash < cash_buffer:
            what_means.append(f"• Current cash position is {currency_symbol}{opening_first:,.0f}, with projected cash falling below the safety buffer of {currency_symbol}{cash_buffer:,.0f}.")
        else:
            what_means.append(f"• Current cash position is {currency_symbol}{opening_first:,.0f}, based on week 1 of the projection (not a live bank feed).")

    # 2) Direction of runway
    if runway_weeks is None:
        what_means.append("• Runway is not quantified because the stress buffer or 14-week projection is missing.")
    else:
        # directional reference without optimism
        if runway_weeks <= 8:
            what_means.append(f"• Base-case runway is {runway_weeks} weeks, indicating limited time to protect the buffer under current trajectory.")
        else:
            what_means.append(f"• Base-case runway is {runway_weeks} weeks under current assumptions and will shorten if operating outflows exceed customer collections.")

    # 3) Timing of risk
    if stress_week is not None:
        what_means.append(f"• Cash stress is expected around {stress_week.strftime('%d %b %Y')} if current timing of receipts and payments holds.")
    else:
        # fallback: use worst week if available
        if worst_week is not None and min_cash is not None and cash_buffer > 0 and min_cash < cash_buffer:
            what_means.append(f"• Buffer breach risk concentrates around {worst_week.strftime('%d %b %Y')} based on the lowest projected cash point.")
        else:
            what_means.append("• No stress date is identified within the 14-week window using the current buffer and data available.")

    # 4) Implication for decisions
    # Tie to AR / AP / buffer logic explicitly
    if cash_buffer > 0 and runway_weeks is not None and runway_weeks < 8:
        what_means.append("• Near-term decisions must prioritise cash preservation: defer discretionary outflows and accelerate collections to extend runway.")
    else:
        # still decision-implication even if not critical
        if ar_next_4w > 0:
            what_means.append("• Near-term decisions should be constrained by collections timing: AR due in the next 4 weeks must land to avoid tightening the buffer.")
        else:
            what_means.append("• Near-term decisions must be gated by committed cash outflows because forward receipts are not visible in the next 4 weeks.")

    # Enforce max 4 bullets
    what_means = what_means[:4]
    st.markdown("\n".join(what_means))

    st.markdown("---")

    # ==========================
    # ACTION PLAN — NEXT 7–14 DAYS (HARD RULES)
    # 3–5 actions only. Executable in 7–14 days. Directly linked to KPI/risk.
    # No "Monitor"/"Review".
    # ==========================
    st.subheader("Action Plan – Next 7–14 days")

    actions = []

    reaction_is_critical = (reaction_weeks is not None and reaction_weeks < 8) or (immediate is True)

    # Build commitments signal (doesn't use a "review" action)
    commitments_drive = False
    try:
        commitments_df = build_cash_commitments(df_ar, df_ap, as_of=as_of, horizon_days=60, limit=20)
        commitments_drive = (commitments_df is not None and not commitments_df.empty)
    except Exception:
        commitments_drive = False

    # --- If reaction window < 8 weeks (mandatory-type actions) ---
    if reaction_is_critical:
        actions.append("• Freeze discretionary spend approvals for the next 14 days and require founder sign-off for any non-essential cash outflow.")
        actions.append("• Accelerate invoicing and collections: issue all outstanding invoices within 48 hours and confirm payment dates with top customers within 7 days.")
        actions.append("• Delay non-critical cash outflows: renegotiate supplier payment dates to push at least 30–45% of near-term payables beyond the next 14 days.")
        actions.append("• Build a downside cash plan within 7 days: identify immediate cost cuts and payment deferrals that extend runway by at least 4 weeks.")

    # --- If AR is a critical driver ---
    # Treat as critical if overdue AR exists OR AR due next 4 weeks is meaningful vs buffer
    ar_is_critical_driver = (ar_overdue_amt > 0) or (cash_buffer > 0 and ar_next_4w >= 0.5 * cash_buffer)
    if ar_is_critical_driver:
        actions.append(f"• Chase top overdue invoices: escalate the 3 largest overdue balances totaling up to {currency_symbol}{ar_overdue_amt:,.0f} and secure confirmed settlement dates within 7 days.")
        actions.append(f"• Lock collections timing: contact customers representing {currency_symbol}{ar_next_4w:,.0f} due in the next 4 weeks and obtain written payment commitments within 7 days.")

    # --- If commitments drive stress ---
    if commitments_drive:
        actions.append("• Renegotiate near-term commitments: secure revised dates for any large outflows in the next 30 days and defer discretionary commitments by 14–30 days.")

    # De-duplicate while preserving order
    seen = set()
    actions_dedup = []
    for a in actions:
        if a not in seen:
            actions_dedup.append(a)
            seen.add(a)

    # Enforce 3–5 actions:
    if len(actions_dedup) < 3:
        # add hard, executable defaults linked to cash/buffer even when data is thin
        actions_dedup.append("• Cut cash burn in 14 days: stop or pause at least 2 discretionary expense categories and implement immediately.")
        actions_dedup.append("• Pull forward cash-in: offer a small early-payment incentive (or remove holds) for priority customers to bring receipts into the next 14 days.")
        actions_dedup.append("• Push out cash-out: request extended terms from top suppliers and confirm updated due dates in writing within 7 days.")

        # de-dup again
        seen = set()
        tmp = []
        for a in actions_dedup:
            if a not in seen:
                tmp.append(a)
                seen.add(a)
        actions_dedup = tmp

    actions_dedup = actions_dedup[:5]
    if len(actions_dedup) < 3:
        actions_dedup = actions_dedup + ["• Freeze discretionary spend approvals for the next 14 days."] * (3 - len(actions_dedup))
        actions_dedup = actions_dedup[:3]

    st.markdown("\n".join(actions_dedup))

    st.markdown("---")

    # ------------------------------------------------------------------
    # 2) AR & AP ageing – side by side + drilldown
    # ------------------------------------------------------------------
    st.subheader("📊 AR & AP ageing – who pays late, who you pay")

    col_ar, col_ap = st.columns(2)

    # ---------- AR ageing summary ----------
    with col_ar:
        st.markdown("#### 📥 Invoices coming in (AR)")

        if df_ar is None or df_ar.empty:
            st.info("No AR records yet for this client.")
        else:
            if "aging_bucket" in df_ar.columns:

                def _is_cash_row(row):
                    if not bool(row.get("is_cash_relevant", True)):
                        return False
                    return float(row.get("effective_amount", 0.0)) > 0

                df_age = df_ar[df_ar.apply(_is_cash_row, axis=1)].copy()

                if "days_past_expected" in df_age.columns:
                    df_age = df_age[df_age["days_past_expected"] >= 0]

                if df_age.empty:
                    st.caption("No open customer invoices due as of this month-end.")
                else:
                    amt_col = "effective_amount" if "effective_amount" in df_age.columns else "amount"

                    ageing_summary = (
                        df_age.groupby("aging_bucket", as_index=False)[amt_col]
                        .sum()
                        .rename(columns={"aging_bucket": "Bucket", amt_col: "Total amount"})
                    )

                    st.dataframe(ageing_summary, width="stretch")
                    st.caption(
                        "Open customer invoices due on or before this month-end, "
                        "grouped by ageing bucket using **balance due** (after partial payments)."
                    )

            with st.expander("View AR invoice details (drilldown)", expanded=False):
                ar_view = df_ar.copy()
                cols = []
                if "counterparty" in ar_view.columns:
                    cols.append("counterparty")
                if "issued_date" in ar_view.columns:
                    cols.append("issued_date")
                if "expected_date" in ar_view.columns:
                    cols.append("expected_date")
                if "amount" in ar_view.columns:
                    cols.append("amount")
                if "amount_paid" in ar_view.columns:
                    cols.append("amount_paid")
                if "effective_amount" in ar_view.columns:
                    cols.append("effective_amount")
                if "due_date" in ar_view.columns:
                    cols.append("due_date")
                if "status" in ar_view.columns:
                    cols.append("status")
                if "days_past_expected" in ar_view.columns:
                    cols.append("days_past_expected")
                if "aging_bucket" in ar_view.columns:
                    cols.append("aging_bucket")
                if "is_over_default" in ar_view.columns:
                    cols.append("is_over_default")

                if cols:
                    ar_view = ar_view[cols].rename(
                        columns={
                            "counterparty": "Customer",
                            "issued_date": "Issued",
                            "expected_date": "Expected",
                            "amount": "Original amount",
                            "amount_paid": "Amount paid",
                            "effective_amount": "Balance due",
                            "due_date": "Original due",
                            "status": "Status",
                            "days_past_expected": "Days overdue",
                            "aging_bucket": "Aging bucket",
                            "is_over_default": "Over default terms?",
                        }
                    )

                    filter_text = st.text_input("Filter AR by customer or status", value="", key="ar_filter")
                    if filter_text:
                        ft = filter_text.lower()
                        mask = (
                            ar_view["Customer"].astype(str).str.lower().str.contains(ft)
                            | ar_view["Status"].astype(str).str.lower().str.contains(ft)
                        )
                        ar_view = ar_view[mask]

                    st.dataframe(ar_view, width="stretch")
                else:
                    st.caption("AR table missing expected columns for detailed view.")

    # 💳 Record AR payment
    with st.expander("💳 Record a payment on an AR invoice", expanded=False):
        ar_open = df_ar[df_ar["is_cash_relevant"] & (df_ar["effective_amount"] > 0)].copy() if df_ar is not None else pd.DataFrame()

        if ar_open.empty:
            st.caption("No open customer invoices to record payments against.")
        else:
            AR_TABLE_NAME = "ar_ap_tracker"  # 🔁 change to your real table name

            label_map = {}
            for _, r in ar_open.iterrows():
                inv_id = r.get("id")
                inv_no = r.get("invoice_number", inv_id)
                cust = r.get("counterparty", "")
                due = r.get("due_date", "")
                bal = float(r.get("effective_amount") or 0.0)

                try:
                    due_txt = pd.to_datetime(due).strftime("%d %b %Y")
                except Exception:
                    due_txt = str(due)

                label_map[inv_id] = f"{inv_no} – {cust} – due {due_txt} – bal {currency_symbol}{bal:,.0f}"

            selected_id = st.selectbox(
                "Select invoice to update",
                options=list(label_map.keys()),
                format_func=lambda x: label_map.get(x, str(x)),
                key="ar_payment_invoice",
            )

            payment_type = st.radio(
                "Payment type",
                ["Full payment", "Partial payment"],
                horizontal=True,
                key="ar_payment_type",
            )

            current_row = ar_open[ar_open["id"] == selected_id].iloc[0]
            current_balance = float(current_row.get("effective_amount") or 0.0)

            default_amount = current_balance if payment_type == "Full payment" else 0.0

            payment_amount = st.number_input(
                "Amount paid",
                min_value=0.0,
                value=default_amount,
                step=100.0,
                key="ar_payment_amount",
            )

            payment_date = st.date_input("Payment date", value=as_of, key="ar_payment_date")

            if st.button("Record AR payment", key="ar_payment_button"):
                try:
                    sb = get_supabase_client() or get_supabase_client()

                    prev_paid = float(current_row.get("amount_paid") or 0.0)
                    new_paid_total = prev_paid + float(payment_amount)

                    if new_paid_total >= float(current_row.get("amount") or 0.0) - 0.01:
                        new_status = "paid"
                        partially_paid_flag = False
                    else:
                        new_status = "open"
                        partially_paid_flag = True

                    update_data = {
                        "amount_paid": new_paid_total,
                        "partially_paid": partially_paid_flag,
                        "status": new_status,
                        "payment_date": payment_date.isoformat(),
                        "updated_at": pd.Timestamp.utcnow().isoformat(),
                    }

                    sb.table(AR_TABLE_NAME).update(update_data).eq("id", selected_id).execute()

                    st.success("Payment recorded. Rerun the page to refresh ageing & cash.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Failed to record payment: {e}")

    # ---------- AP ageing summary ----------
    with col_ap:
        st.markdown("#### 📤 Bills to pay (AP)")

        if df_ap is None or df_ap.empty:
            st.info("No AP records yet for this client.")
        else:
            if "aging_bucket" in df_ap.columns:

                def _is_cash_row_ap(row):
                    if not bool(row.get("is_cash_relevant", True)):
                        return False
                    return float(row.get("effective_amount", 0.0)) > 0

                df_age_ap = df_ap[df_ap.apply(_is_cash_row_ap, axis=1)].copy()

                if "days_past_expected" in df_age_ap.columns:
                    df_age_ap = df_age_ap[df_age_ap["days_past_expected"] >= 0]

                if df_age_ap.empty:
                    st.caption("No open supplier bills due as of this month-end.")
                else:
                    amt_col_ap = "effective_amount" if "effective_amount" in df_age_ap.columns else "amount"

                    ap_ageing_summary = (
                        df_age_ap.groupby("aging_bucket", as_index=False)[amt_col_ap]
                        .sum()
                        .rename(columns={"aging_bucket": "Bucket", amt_col_ap: "Total amount"})
                    )

                    st.dataframe(ap_ageing_summary, width="stretch")
                    st.caption(
                        "Open supplier bills due on or before this month-end, "
                        "grouped by ageing bucket using **balance due**."
                    )

            with st.expander("View AP bill details (drilldown)", expanded=False):
                ap_view = df_ap.copy()
                cols = []
                if "counterparty" in ap_view.columns:
                    cols.append("counterparty")
                if "amount" in ap_view.columns:
                    cols.append("amount")
                if "amount_paid" in ap_view.columns:
                    cols.append("amount_paid")
                if "effective_amount" in ap_view.columns:
                    cols.append("effective_amount")
                if "due_date" in ap_view.columns:
                    cols.append("due_date")
                if "pay_expected_date" in ap_view.columns:
                    cols.append("pay_expected_date")
                elif "expected_payment_date" in ap_view.columns:
                    cols.append("expected_payment_date")
                if "status" in ap_view.columns:
                    cols.append("status")
                if "days_past_expected" in ap_view.columns:
                    cols.append("days_past_expected")
                if "aging_bucket" in ap_view.columns:
                    cols.append("aging_bucket")
                if "is_over_default" in ap_view.columns:
                    cols.append("is_over_default")

                if cols:
                    ap_view = ap_view[cols].rename(
                        columns={
                            "counterparty": "Supplier",
                            "amount": "Original amount",
                            "amount_paid": "Amount paid",
                            "effective_amount": "Balance due",
                            "due_date": "Original due",
                            "pay_expected_date": "Expected to pay",
                            "expected_payment_date": "Expected to pay",
                            "status": "Status",
                            "days_past_expected": "Days overdue",
                            "aging_bucket": "Aging bucket",
                            "is_over_default": "Over default terms?",
                        }
                    )

                    filter_text_ap = st.text_input("Filter AP by supplier or status", value="", key="ap_filter")
                    if filter_text_ap:
                        ft = filter_text_ap.lower()
                        mask_ap = (
                            ap_view["Supplier"].astype(str).str.lower().str.contains(ft)
                            | ap_view["Status"].astype(str).str.lower().str.contains(ft)
                        )
                        ap_view = ap_view[mask_ap]

                    st.dataframe(ap_view, width="stretch")
                else:
                    st.caption("AP table missing expected columns for detailed view.")

    # 💳 Record AP payment
    with st.expander("💳 Record a payment on an AP bill", expanded=False):
        ap_open = df_ap[df_ap["is_cash_relevant"] & (df_ap["effective_amount"] > 0)].copy() if df_ap is not None else pd.DataFrame()

        if ap_open.empty:
            st.caption("No open supplier bills to record payments against.")
        else:
            AP_TABLE_NAME = "ar_ap_tracker"  # 🔁 change if AP lives in a different table

            label_map_ap = {}
            for _, r in ap_open.iterrows():
                bill_id = r.get("id")
                bill_no = r.get("invoice_number", bill_id)
                supp = r.get("counterparty", "")
                due = r.get("due_date", "")
                bal = float(r.get("effective_amount") or 0.0)

                try:
                    due_txt = pd.to_datetime(due).strftime("%d %b %Y")
                except Exception:
                    due_txt = str(due)

                label_map_ap[bill_id] = f"{bill_no} – {supp} – due {due_txt} – bal {currency_symbol}{bal:,.0f}"

            selected_bill_id = st.selectbox(
                "Select bill to update",
                options=list(label_map_ap.keys()),
                format_func=lambda x: label_map_ap.get(x, str(x)),
                key="ap_payment_bill",
            )

            payment_type_ap = st.radio(
                "Payment type",
                ["Full payment", "Partial payment"],
                horizontal=True,
                key="ap_payment_type",
            )

            current_ap_row = ap_open[ap_open["id"] == selected_bill_id].iloc[0]
            current_balance_ap = float(current_ap_row.get("effective_amount") or 0.0)

            default_amount_ap = current_balance_ap if payment_type_ap == "Full payment" else 0.0

            payment_amount_ap = st.number_input(
                "Amount paid",
                min_value=0.0,
                value=default_amount_ap,
                step=100.0,
                key="ap_payment_amount",
            )

            payment_date_ap = st.date_input("Payment date", value=as_of, key="ap_payment_date")

            if st.button("Record AP payment", key="ap_payment_button"):
                try:
                    sb = get_supabase_client() or get_supabase_client()

                    prev_paid_ap = float(current_ap_row.get("amount_paid") or 0.0)
                    new_paid_total_ap = prev_paid_ap + float(payment_amount_ap)

                    if new_paid_total_ap >= float(current_ap_row.get("amount") or 0.0) - 0.01:
                        new_status_ap = "paid"
                        partially_paid_flag_ap = False
                    else:
                        new_status_ap = "open"
                        partially_paid_flag_ap = True

                    update_data_ap = {
                        "amount_paid": new_paid_total_ap,
                        "partially_paid": partially_paid_flag_ap,
                        "status": new_status_ap,
                        "payment_date": payment_date_ap.isoformat(),
                        "updated_at": pd.Timestamp.utcnow().isoformat(),
                    }

                    sb.table(AP_TABLE_NAME).update(update_data_ap).eq("id", selected_bill_id).execute()

                    st.success("Payment recorded. Rerun the page to refresh ageing & cash.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Failed to record payment: {e}")

    # ------------------------------------------------------------------
    # 3) Cash commitments (near-term list)
    # ------------------------------------------------------------------
    st.markdown("---")
    st.subheader("🔍 Cash commitments coming up")

    commitments = build_cash_commitments(
        df_ar,
        df_ap,
        as_of=as_of,
        horizon_days=60,
        limit=7,
    )

    if commitments.empty:
        st.info("No upcoming cash movements found yet.")
    else:
        display_commitments = commitments.rename(
            columns={
                "direction": "Direction",
                "who": "Customer / Supplier",
                "amount": "Amount",
                "date": "Planned date",
                "original_due": "Original due",
                "overdue": "Overdue?",
            }
        )

        st.dataframe(display_commitments, width="stretch")
        st.caption(
            "Top upcoming customer receipts and supplier payments. "
            "`Overdue? = Yes` means the original due date was before this Focus Month and is still unpaid."
        )

    # ------------------------------------------------------------------
    # 4) 14-week cashflow – chart + table (NO 'What this 14-week view is telling you')
    # ------------------------------------------------------------------
    st.markdown("---")
    st.subheader("📆 14-week cashflow – can we sleep at night?")

    if week_df is None or week_df.empty:
        st.caption(
            "Not enough AR/AP + payroll data to build a 14-week view yet. "
            "Add invoices and team data to see near-term cash pressure."
        )
    else:
        view = week_df.copy()

        if "week_start" in view.columns:
            view["week_start"] = pd.to_datetime(view["week_start"], errors="coerce")

        view["Week"] = view["week_start"].dt.strftime("%d %b %Y")

        cols = [
            "Week",
            "cash_in_ar",
            "cash_in_other_income",     # ✅ add
            "cash_out_ap",
            "cash_out_payroll",
            "cash_out_opex",            # ✅ add
            "operating_cf",
            "investing_cf",
            "financing_cf",
            "net_cash",
            "closing_cash",
        ]
        cols = [c for c in cols if c in view.columns]

        view = view[cols].rename(
            columns={
                "cash_in_ar": "Cash in (customers)",
                "cash_in_other_income": "Cash in (other income)",   # ✅ add
                "cash_out_ap": "Cash out (bills/AP)",
                "cash_out_payroll": "Cash out (payroll)",
                "cash_out_opex": "Cash out (opex)",                 # ✅ add
                "operating_cf": "Operating CF",
                "investing_cf": "Investing CF",
                "financing_cf": "Financing CF",
                "net_cash": "Net cash",
                "closing_cash": "Closing cash",
            }
        )


        st.dataframe(view, width="stretch")
        st.caption(
            "Expected cash in/out by week for the next ~3 months, "
            "split by operating, investing and financing activity."
        )

        # Cash cliff highlight (first week where closing_cash <= 0)
        if "closing_cash" in week_df.columns:
            danger_weeks = week_df[week_df["closing_cash"] <= 0]
            if not danger_weeks.empty:
                first_danger = danger_weeks.iloc[0]
                st.error(
                    f"⚠️ Cash cliff around **week of "
                    f"{first_danger['week_start'].strftime('%d %b %Y')}** – "
                    f"closing cash projected at "
                    f"{currency_symbol}{first_danger['closing_cash']:,.0f}."
                )

        base_line = (
            alt.Chart(week_df)
            .mark_line()
            .encode(
                x=alt.X("week_start:T", title="Week starting", axis=alt.Axis(format="%d %b")),
                y=alt.Y("closing_cash:Q", title=f"Projected closing cash ({currency_code})"),
                tooltip=[
                    alt.Tooltip("week_start:T", title="Week of", format="%d %b %Y"),
                    alt.Tooltip("closing_cash:Q", title="Closing cash", format=",.0f"),
                    alt.Tooltip("operating_cf:Q", title="Operating CF", format=",.0f"),
                    alt.Tooltip("investing_cf:Q", title="Investing CF", format=",.0f"),
                    alt.Tooltip("financing_cf:Q", title="Financing CF", format=",.0f"),
                ],
            )
        )

        if cash_buffer and cash_buffer > 0:
            buffer_rule = (
                alt.Chart(pd.DataFrame({"y": [cash_buffer]}))
                .mark_rule(strokeDash=[6, 6])
                .encode(y="y:Q")
            )
            chart = base_line + buffer_rule
            st.caption(f"Dashed line = safety buffer ({currency_symbol}{cash_buffer:,.0f}).")
        else:
            chart = base_line

        st.altair_chart(chart.interactive(), width="stretch")


# -----------------------------------------------------------------------------


def resolve_opening_cash(
    client_id: str,
    base_month_ts: pd.Timestamp,
    as_of_date: date,
    opening_cash_hint: float | None,
) -> float:
    """
    Single source of truth for plan anchor opening cash.
    Priority:
      0) opening_cash_hint
      1) cash override for focus month (bank truth)
      2) cash_opening_balance table for this as_of_month
      3) seed from previous plan's focus-month closing cash
      4) fallback seed from KPI cash_balance
    """
    if not client_id or base_month_ts is None or as_of_date is None:
        return 0.0

    if opening_cash_hint is not None:
        return float(opening_cash_hint)

    base_month = base_month_ts.replace(day=1).date()

    # 1) manual override for the focus month (bank truth)
    ov = fetch_opening_cash_override(client_id, base_month)  # float|None
    if ov is not None:
        return float(ov)

    # 2) plan anchor table
    try:
        res = (
            supabase.table("cash_opening_balance")
            .select("opening_cash")
            .eq("client_id", str(client_id))
            .eq("as_of_month", as_of_date.isoformat())
            .limit(1)
            .execute()
        )
        data = getattr(res, "data", None)
        if data:
            return float(data[0].get("opening_cash") or 0.0)
    except Exception as e:
        print("[WARN] cash_opening_balance read failed:", e)

    # 3) seed from previous plan closing cash (roll plan-to-plan)
    try:
        prev_as_of = (base_month_ts.to_period("M") - 1).to_timestamp().date()
        df_prev_plan = fetch_cashflow_summary_for_client_as_of(client_id, prev_as_of)

        if df_prev_plan is not None and not df_prev_plan.empty:
            df_prev_plan = df_prev_plan.copy()
            df_prev_plan["month_date"] = pd.to_datetime(
                df_prev_plan["month_date"], errors="coerce"
            ).dt.date

            prev_row = df_prev_plan[df_prev_plan["month_date"] == prev_as_of]
            if not prev_row.empty and pd.notna(prev_row.iloc[0].get("closing_cash")):
                seeded = float(prev_row.iloc[0]["closing_cash"])

                try:
                    supabase.table("cash_opening_balance").upsert(
                        {
                            "client_id": str(client_id),
                            "as_of_month": as_of_date.isoformat(),
                            "opening_cash": round(seeded, 2),
                            "source": "prev_plan_closing_seed",
                            "notes": None,
                        },
                        on_conflict="client_id,as_of_month",
                    ).execute()
                except Exception as e:
                    print("[WARN] cash_opening_balance upsert (prev seed) failed:", e)

                return seeded
    except Exception as e:
        print("[WARN] prev plan seed failed:", e)

    # 4) fallback KPI seed
    seed = 0.0
    try:
        kpi_row = fetch_kpis_for_client_month(client_id, as_of_date)
        if kpi_row and kpi_row.get("cash_balance") is not None:
            seed = float(kpi_row["cash_balance"] or 0.0)
    except Exception as e:
        print("[WARN] KPI seed failed, defaulting to 0:", e)
        seed = 0.0

    try:
        supabase.table("cash_opening_balance").upsert(
            {
                "client_id": str(client_id),
                "as_of_month": as_of_date.isoformat(),
                "opening_cash": round(seed, 2),
                "source": "kpi_seed",
                "notes": None,
            },
            on_conflict="client_id,as_of_month",
        ).execute()
    except Exception as e:
        print("[WARN] cash_opening_balance seed upsert failed:", e)

    return float(seed)




def recompute_cashflow_from_ar_ap(
    client_id: str,
    base_month: date,
    n_months: int = 12,
    opening_cash_hint: float | None = None,
) -> bool:
    """
    Monthly plan engine (12-month by default).
    FIXES:
      - Financing CF is always included in free_cash_flow and closing_cash
      - Persists ONE plan version keyed by as_of_month (= focus month start date)
      - Uses created_at <= as_of_cutoff versioning for investing/financing, with safe fallback
      - closing_cash is deterministic: close = open + (op + inv + fin)

    IMPORTANT:
      Anything that needs "cash balance for focus month" should read:
        cashflow_summary WHERE as_of_month = focus_month AND month_date = focus_month (closing_cash)
      NOT from an unfiltered cashflow_summary fetch that mixes multiple versions.
    """
    if not client_id or base_month is None:
        return False

    try:
        base_month_ts = pd.to_datetime(base_month).replace(day=1)
        as_of_date = base_month_ts.date()            # plan version key
        as_of_cutoff = base_month_ts + MonthEnd(0)   # version cutoff

        month_index = pd.date_range(start=base_month_ts, periods=n_months, freq="MS")

        # ---------- Opening cash anchor ----------
        use_hint = (opening_cash_hint is not None and float(opening_cash_hint) != 0.0)

        if use_hint:
            opening_cash_0 = float(opening_cash_hint)
            try:
                supabase.table("cash_opening_balance").upsert(
                    {
                        "client_id": str(client_id),
                        "as_of_month": as_of_date.isoformat(),
                        "opening_cash": round(opening_cash_0, 2),
                        "source": "manual_hint",
                        "notes": None,
                    },
                    on_conflict="client_id,as_of_month",
                ).execute()
            except Exception as e:
                print("[WARN] Failed to persist opening_cash_hint:", e)
        else:
            opening_cash_0 = resolve_opening_cash(
                client_id=str(client_id),
                base_month_ts=base_month_ts,
                as_of_date=as_of_date,
                opening_cash_hint=None,
            )

        # ---------- Settings ----------
        settings = get_client_settings(client_id) or {}
        ar_days = int(settings.get("ar_default_days", 30) or 30)
        ap_days = int(settings.get("ap_default_days", 30) or 30)

        # ---------- AR/AP (timing + partials) ----------
        df_ar, df_ap = fetch_ar_ap_for_client(client_id)

        if df_ar is not None and not df_ar.empty:
            df_ar = df_ar.copy()
            df_ar = add_ar_aging(df_ar, as_of=as_of_date, ar_default_days=ar_days)
            df_ar["expected_date"] = pd.to_datetime(df_ar.get("expected_date"), errors="coerce")
            df_ar["amount"] = pd.to_numeric(df_ar.get("amount"), errors="coerce").fillna(0.0)
        else:
            df_ar = pd.DataFrame(columns=["id", "expected_date", "amount"])

        if df_ap is not None and not df_ap.empty:
            df_ap = df_ap.copy()
            df_ap = add_ap_aging(df_ap, as_of=as_of_date, ap_default_days=ap_days)

            if "pay_expected_date" in df_ap.columns:
                pay_col = "pay_expected_date"
            elif "expected_payment_date" in df_ap.columns:
                pay_col = "expected_payment_date"
            else:
                pay_col = "due_date"

            df_ap[pay_col] = pd.to_datetime(df_ap.get(pay_col), errors="coerce")
            df_ap["amount"] = pd.to_numeric(df_ap.get("amount"), errors="coerce").fillna(0.0)
        else:
            df_ap = pd.DataFrame(columns=["id", "due_date", "amount"])
            pay_col = "due_date"

        df_pay = fetch_ar_ap_payments_for_client(client_id)
        ar_cash_map, ap_cash_map = build_monthly_ar_ap_cash_maps(
            df_ar=df_ar, df_ap=df_ap, df_payments=df_pay, month_index=month_index
        )

        # ---------- Opex ----------
        try:
            df_opex = fetch_opex_for_client(client_id)
            if df_opex is not None and not df_opex.empty:
                df_opex = df_opex.copy()
                df_opex["month_bucket"] = (
                    pd.to_datetime(df_opex["month_date"], errors="coerce")
                    .dt.to_period("M")
                    .dt.to_timestamp()
                )
                df_opex["amount"] = pd.to_numeric(df_opex.get("amount"), errors="coerce").fillna(0.0)
            else:
                df_opex = pd.DataFrame(columns=["month_bucket", "amount"])
        except Exception as e:
            print("Error fetching operating_opex:", e)
            df_opex = pd.DataFrame(columns=["month_bucket", "amount"])

        # ---------- Other operating income ----------
        try:
            df_oinc = fetch_operating_other_income_for_client(client_id)
        except Exception as e:
            print("Error fetching operating_other_income:", e)
            df_oinc = pd.DataFrame()

        if df_oinc is None or df_oinc.empty:
            other_income_by_month = {}
        else:
            df_oinc = df_oinc.copy()
            if "month_bucket" not in df_oinc.columns:
                df_oinc["month_bucket"] = (
                    pd.to_datetime(df_oinc["month_date"], errors="coerce")
                    .dt.to_period("M")
                    .dt.to_timestamp()
                )
            df_oinc["cash_in"] = pd.to_numeric(df_oinc.get("cash_in"), errors="coerce").fillna(0.0)
            oinc_agg = (
                df_oinc.groupby("month_bucket", as_index=False)["cash_in"]
                .sum()
                .rename(columns={"month_bucket": "month_date"})
            )
            other_income_by_month = {
                row["month_date"]: float(row["cash_in"] or 0.0)
                for _, row in oinc_agg.iterrows()
            }

        # ---------- Investing & Financing (versioned by created_at <= as_of_cutoff) ----------
        inv_by_month: dict[pd.Timestamp, float] = {}
        fin_by_month: dict[pd.Timestamp, float] = {}

        # Investing
        try:
            df_inv = fetch_investing_flows_for_client(client_id)
        except Exception as e:
            print("Error fetching investing_flows:", e)
            df_inv = pd.DataFrame()

        if df_inv is not None and not df_inv.empty:
            orig = df_inv.copy()
            df_inv = df_inv.copy()

            if "created_at" in df_inv.columns:
                df_inv["created_at"] = pd.to_datetime(df_inv["created_at"], errors="coerce")
                try:
                    df_inv["created_at"] = df_inv["created_at"].dt.tz_convert(None)
                except TypeError:
                    pass

                before = len(df_inv)
                df_inv = df_inv[df_inv["created_at"] <= as_of_cutoff]
                after = len(df_inv)
                print(f"[DEBUG] investing_flows created_at filter: {before} -> {after} rows (as_of_cutoff={as_of_cutoff.date()})")
                if after == 0:
                    print("[WARN] investing_flows all filtered out by created_at; falling back to all investing rows (no versioning).")
                    df_inv = orig

            df_inv["month_bucket"] = (
                pd.to_datetime(df_inv["month_date"], errors="coerce")
                .dt.to_period("M")
                .dt.to_timestamp()
            )
            df_inv["amount"] = pd.to_numeric(df_inv.get("amount"), errors="coerce").fillna(0.0)

            inv_agg = df_inv.groupby("month_bucket", as_index=False)["amount"].sum().rename(columns={"month_bucket": "month_date"})
            inv_by_month = {row["month_date"]: float(row["amount"] or 0.0) for _, row in inv_agg.iterrows()}

        # Financing
        try:
            df_fin = fetch_financing_flows_for_client(client_id)
        except Exception as e:
            print("Error fetching financing_flows:", e)
            df_fin = pd.DataFrame()

        if df_fin is not None and not df_fin.empty:
            orig = df_fin.copy()
            df_fin = df_fin.copy()

            if "created_at" in df_fin.columns:
                df_fin["created_at"] = pd.to_datetime(df_fin["created_at"], errors="coerce")
                try:
                    df_fin["created_at"] = df_fin["created_at"].dt.tz_convert(None)
                except TypeError:
                    pass

                before = len(df_fin)
                df_fin = df_fin[df_fin["created_at"] <= as_of_cutoff]
                after = len(df_fin)
                print(f"[DEBUG] financing_flows created_at filter: {before} -> {after} rows (as_of_cutoff={as_of_cutoff.date()})")
                if after == 0:
                    print("[WARN] financing_flows all filtered out by created_at; falling back to all financing rows (no versioning).")
                    df_fin = orig

            df_fin["month_bucket"] = (
                pd.to_datetime(df_fin["month_date"], errors="coerce")
                .dt.to_period("M")
                .dt.to_timestamp()
            )
            df_fin["amount"] = pd.to_numeric(df_fin.get("amount"), errors="coerce").fillna(0.0)

            fin_agg = df_fin.groupby("month_bucket", as_index=False)["amount"].sum().rename(columns={"month_bucket": "month_date"})
            fin_by_month = {row["month_date"]: float(row["amount"] or 0.0) for _, row in fin_agg.iterrows()}

        # ---------- Payroll ----------
        payroll_by_month = compute_payroll_by_month(client_id, month_index)

        # ---------- Build rows (running cash) ----------
        rows = []
        current_closing = float(opening_cash_0)

        for m in month_index:
            m_key = pd.to_datetime(m).to_period("M").to_timestamp()
            opening_cash = float(current_closing)

            cash_in_ar = float(ar_cash_map.get(m_key, 0.0))
            cash_out_ap = float(ap_cash_map.get(m_key, 0.0))
            payroll_cash = float(payroll_by_month.get(m_key, 0.0))

            if df_opex is not None and not df_opex.empty:
                opex_cash = float(df_opex.loc[df_opex["month_bucket"] == m_key, "amount"].sum())
            else:
                opex_cash = 0.0

            other_income_cash = float(other_income_by_month.get(m_key, 0.0))
            investing_cf = float(inv_by_month.get(m_key, 0.0))
            financing_cf = float(fin_by_month.get(m_key, 0.0))   # ✅ ALWAYS included

            operating_cf = cash_in_ar + other_income_cash - cash_out_ap - payroll_cash - opex_cash
            free_cf = operating_cf + investing_cf + financing_cf
            closing_cash = opening_cash + free_cf

            print(
                f"[ENGINE] {m_key.strftime('%Y-%m')} | "
                f"open={opening_cash:.2f}, ar_in={cash_in_ar:.2f}, other_in={other_income_cash:.2f}, "
                f"ap_out={cash_out_ap:.2f}, payroll={payroll_cash:.2f}, opex={opex_cash:.2f}, "
                f"inv={investing_cf:.2f}, fin={financing_cf:.2f}, free_cf={free_cf:.2f}, close={closing_cash:.2f}"
            )

            rows.append(
                {
                    "client_id": str(client_id),
                    "month_date": m_key.date().isoformat(),
                    "as_of_month": as_of_date.isoformat(),
                    "opening_cash": round(opening_cash, 2),
                    "operating_cf": round(operating_cf, 2),
                    "investing_cf": round(investing_cf, 2),
                    "financing_cf": round(financing_cf, 2),
                    "free_cash_flow": round(free_cf, 2),
                    "closing_cash": round(closing_cash, 2),
                    "cash_danger_flag": (closing_cash <= 0.0),
                    "notes": None,
                }
            )

            current_closing = float(closing_cash)

        if not rows:
            print("No rows computed for cashflow_summary recompute.")
            return False

        # ---------- Persist only THIS plan version ----------
        try:
            (
                supabase.table("cashflow_summary")
                .delete()
                .eq("client_id", str(client_id))
                .eq("as_of_month", as_of_date.isoformat())
                .execute()
            )
        except Exception as e:
            print("Warning deleting old cashflow_summary rows (versioned):", e)

        try:
            resp = supabase.table("cashflow_summary").insert(rows).execute()
        except Exception as e:
            print("Error inserting cashflow_summary:", e)
            return False

        if isinstance(resp, dict) and resp.get("error"):
            print("Error inserting cashflow_summary:", resp["error"])
            return False

        print(f"Recomputed {len(rows)} rows into cashflow_summary for as_of_month={as_of_date}")

        # Clear caches (if you use st.cache_data)
        try:
            fetch_cashflow_summary_for_client.clear()
        except Exception:
            pass
        try:
            fetch_alerts_for_client.clear()
        except Exception:
            pass

        # Run alerts
        try:
            run_all_alerts_for_client_months(client_id, month_index)
        except Exception as e:
            print(f"Warning running alerts after cashflow recompute: {e}")

        return True

    except Exception as e:
        print("Error recomputing cashflow_summary:", e)
        return False



def cash_flows_editor_block():
    st.markdown("---")
    st.subheader("💹 Investing & Financing cash moves")

    col_info, col_btn = st.columns([3, 1])
    with col_info:
        st.caption(
            "Use this to record *big one-off moves* like buying equipment, "
            "raising equity, taking or repaying loans, or paying dividends. "
            "Positive = cash **in**, negative = cash **out**."
        )

    with col_btn:
        if st.button("Rebuild cashflow after changes", key="rebuild_after_inv_fin"):
            ok = recompute_cashflow_from_ar_ap(
                selected_client_id,
                base_month=selected_month_start,
                n_months=12,
                opening_cash_hint=None,
            )
            if ok:
                st.success("Cashflow rebuilt with latest investing/financing moves.")
            else:
                st.error("Could not rebuild cashflow. Check data and try again.")

    # Load existing flows
    df_inv = fetch_investing_flows_for_client(selected_client_id)
    df_fin = fetch_financing_flows_for_client(selected_client_id)

    tab_inv, tab_fin = st.tabs(["🏗 Investing flows", "🏦 Financing flows"])

    # ----------------- INVESTING TAB -----------------
    with tab_inv:
        st.markdown("#### Investing flows (capex, asset sales, etc.)")

        if df_inv is None or df_inv.empty:
            st.caption("No investing flows yet for this business.")
        else:
            df_show = df_inv.copy()
            df_show["month_date"] = pd.to_datetime(
                df_show["month_date"], errors="coerce"
            )
            df_show = df_show.sort_values("month_date")
            df_show["Month"] = df_show["month_date"].dt.strftime("%b %Y")
            st.dataframe(
                df_show[["id", "Month", "amount", "category", "notes"]],
                width="stretch",
            )

        st.markdown("##### ➕ Add a new investing flow")
        with st.form(key="new_investing_flow_form"):
            col1, col2 = st.columns(2)
            with col1:
                inv_date = st.date_input(
                    "Month (use any date in that month)",
                    value=selected_month_start,
                    key="inv_date",
                )
                inv_amount = st.number_input(
                    "Amount (positive = cash in, negative = cash out)",
                    value=0.0,
                    step=1000.0,
                    key="inv_amount",
                )
            with col2:
                inv_category = st.text_input(
                    "Category",
                    value="Capex",
                    key="inv_category",
                )
                inv_notes = st.text_area(
                    "Notes (optional)",
                    value="",
                    key="inv_notes",
                )

            submitted_inv = st.form_submit_button("Save investing flow")
            if submitted_inv:
                ok = create_investing_flow(
                    selected_client_id,
                    month_date=inv_date,
                    amount=inv_amount,
                    category=inv_category,
                    notes=inv_notes,
                )
                if ok:
                    st.success("Investing flow saved.")
                    st.rerun()
                else:
                    st.error("Could not save investing flow. Check inputs and try again.")

    # ----------------- FINANCING TAB -----------------
    with tab_fin:
        st.markdown("#### Financing flows (equity, loans, repayments, dividends)")

        if df_fin is None or df_fin.empty:
            st.caption("No financing flows yet for this business.")
        else:
            df_show = df_fin.copy()
            df_show["month_date"] = pd.to_datetime(
                df_show["month_date"], errors="coerce"
            )
            df_show = df_show.sort_values("month_date")
            df_show["Month"] = df_show["month_date"].dt.strftime("%b %Y")
            st.dataframe(
                df_show[["id", "Month", "amount", "category", "notes"]],
                width="stretch",
            )

        st.markdown("##### ➕ Add a new financing flow")
        with st.form(key="new_financing_flow_form"):
            col1, col2 = st.columns(2)
            with col1:
                fin_date = st.date_input(
                    "Month (use any date in that month)",
                    value=selected_month_start,
                    key="fin_date",
                )
                fin_amount = st.number_input(
                    "Amount (positive = cash in, negative = cash out)",
                    value=0.0,
                    step=1000.0,
                    key="fin_amount",
                )
            with col2:
                fin_category = st.text_input(
                    "Category",
                    value="Equity raise",
                    key="fin_category",
                )
                fin_notes = st.text_area(
                    "Notes (optional)",
                    value="",
                    key="fin_notes",
                )

            submitted_fin = st.form_submit_button("Save financing flow")
            if submitted_fin:
                ok = create_financing_flow(
                    selected_client_id,
                    month_date=fin_date,
                    amount=fin_amount,
                    category=fin_category,
                    notes=fin_notes,
                )
                if ok:
                    st.success("Financing flow saved.")
                    st.rerun()
                else:
                    st.error("Could not save financing flow. Check inputs and try again.")



def _resolve_dept_overspend_alerts_for_month(client_id, month_start: date):
    """
    Internal helper: resolve existing dept_overspend alerts for the given month.
    """
    try:
        supabase.table("alerts").update(
            {
                "is_active": False,
                "is_dismissed": True,
                "resolved_at": datetime.now(timezone.utc).isoformat(),
            }
        ).eq("client_id", str(client_id)) \
         .eq("alert_type", "dept_overspend") \
         .eq("month_date", month_start.isoformat()) \
         .eq("is_active", True) \
         .execute()
    except Exception:
        pass


def _add_title_only_slide(prs, title_text: str):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    return slide


def _add_text_box(slide, left_in, top_in, width_in, height_in, text, font_size=14, bold=False):
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return tx_box


def _add_chart_image(slide, fig, left_in, top_in, width_in):
    """Save a matplotlib fig to PNG in-memory and insert as picture."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(left_in), Inches(top_in), width=Inches(width_in))

def build_funding_strategy_snapshot(client_id: str, focus_month: datetime.date) -> dict:
    settings = get_client_settings(client_id) or {}
    currency_code, currency_symbol = get_client_currency(client_id)

    engine_df = fetch_cashflow_summary_for_client(client_id)

    runway, eff_burn = compute_runway_and_effective_burn_from_df(
        engine_df=engine_df,
        month_ref=focus_month,
    )

    # Closing cash at focus month (engine)
    closing_cash = 0.0
    if engine_df is not None and not engine_df.empty and "closing_cash" in engine_df.columns:
        tmp = engine_df.copy()
        tmp["month_date"] = pd.to_datetime(tmp["month_date"], errors="coerce").dt.date
        fm = pd.to_datetime(focus_month).to_period("M").to_timestamp().date()
        row = tmp[tmp["month_date"] == fm]
        if not row.empty:
            closing_cash = float(row["closing_cash"].iloc[0] or 0.0)

    # Strategy knobs (saved in client_settings or default)
    target_runway = float(settings.get("target_runway_months", 12) or 12)
    burn_change_pct = float(settings.get("funding_burn_change_pct", 0) or 0)

    adjusted_burn = (float(eff_burn) if eff_burn is not None else 0.0) * (1 + burn_change_pct / 100.0)

    required_cash = adjusted_burn * target_runway if adjusted_burn > 0 else 0.0
    funding_gap = max(0.0, required_cash - closing_cash)

    pre_money = float(settings.get("rough_pre_money", 0.0) or 0.0)
    dilution = None
    if pre_money > 0 and funding_gap > 0:
        dilution = funding_gap / (pre_money + funding_gap) * 100.0

    return {
        "currency_code": currency_code,
        "currency_symbol": currency_symbol,
        "runway_months": runway,
        "effective_burn": eff_burn,
        "target_runway": target_runway,
        "burn_change_pct": burn_change_pct,
        "adjusted_burn": adjusted_burn,
        "closing_cash": closing_cash,
        "funding_gap": funding_gap,
        "pre_money": pre_money,
        "dilution_pct": dilution,
        "engine_df": engine_df,  # keep for chart
    }


def build_funding_strategy_commentary(snapshot: dict) -> list[str]:
    cs = snapshot["currency_symbol"]
    runway = snapshot.get("runway_months")
    burn = snapshot.get("effective_burn")
    adj_burn = snapshot.get("adjusted_burn")
    burn_pct = snapshot.get("burn_change_pct")
    target = snapshot.get("target_runway")
    gap = snapshot.get("funding_gap")
    cash = snapshot.get("closing_cash")
    dilution = snapshot.get("dilution_pct")

    lines = []

    # What happened / where you stand
    if runway is None:
        lines.append("Runway does not hit zero in the forecast horizon (cash stays positive).")
    else:
        lines.append(f"Runway is ~{runway:.1f} months based on the forward cash forecast.")

    if burn is not None:
        lines.append(f"Effective burn is ~{cs}{float(burn):,.0f}/month (next 3-month average).")

    # What changes with the plan
    if burn_pct and burn is not None:
        lines.append(f"With planned burn change ({burn_pct:+.0f}%), adjusted burn is {cs}{float(adj_burn):,.0f}/month.")

    # What decision is required
    if gap and gap > 0:
        lines.append(f"To reach {target:.0f} months runway, estimated funding required is {cs}{float(gap):,.0f}.")
        if dilution is not None:
            lines.append(f"At current pre-money assumption, implied dilution is ~{float(dilution):.1f}% (equity-only).")
        lines.append("Recommendation: align spend to runway target and begin raise prep now (metrics pack + pipeline + milestones).")
    else:
        lines.append(f"No immediate funding gap to reach {target:.0f} months runway at current assumptions.")
        lines.append("Recommendation: raise opportunistically (better terms) and keep collections + hiring discipline tight.")

    return lines[:6]


def add_funding_strategy_slide(
    prs: Presentation,
    client_id: str,
    focus_month: datetime.date,
) -> None:
    snap = build_funding_strategy_snapshot(client_id, focus_month)
    currency_code = snap["currency_code"]
    cs = snap["currency_symbol"]
    engine_df = snap["engine_df"]

    slide = _add_title_only_slide(prs, "5. Funding Strategy – Investing & Financing")

    # ---------- KPI block ----------
    runway = snap.get("runway_months")
    eff_burn = snap.get("effective_burn")
    target = snap.get("target_runway")
    gap = snap.get("funding_gap")
    cash = snap.get("closing_cash")
    dilution = snap.get("dilution_pct")

    kpi_lines = [
        f"Current closing cash: {cs}{float(cash):,.0f}",
        f"Runway (forward): {('—' if runway is None else f'{float(runway):.1f} months')}",
        f"Effective burn: {('—' if eff_burn is None else f'{cs}{float(eff_burn):,.0f} / month')}",
        f"Target runway: {float(target):.0f} months",
        f"Funding required (gap): {cs}{float(gap):,.0f}",
    ]
    if dilution is not None:
        kpi_lines.append(f"Implied dilution (equity-only): ~{float(dilution):.1f}%")

    _add_text_box(
        slide,
        left_in=0.5,
        top_in=1.3,
        width_in=4.8,
        height_in=2.2,
        text="Funding snapshot\n" + "\n".join(f"• {x}" for x in kpi_lines),
        font_size=13,
    )

    # ---------- Chart: investing_cf vs financing_cf next 12 months ----------
    if engine_df is not None and not engine_df.empty and ("investing_cf" in engine_df.columns or "financing_cf" in engine_df.columns):
        df = engine_df.copy()
        df["month_date"] = pd.to_datetime(df.get("month_date"), errors="coerce")
        df = df[df["month_date"].notna()].sort_values("month_date")

        start = pd.to_datetime(focus_month).replace(day=1)
        end = start + pd.DateOffset(months=12)
        df = df[(df["month_date"] >= start) & (df["month_date"] < end)]

        if not df.empty:
            months = df["month_date"].dt.strftime("%b %y").tolist()
            inv = pd.to_numeric(df.get("investing_cf", 0), errors="coerce").fillna(0.0).tolist()
            fin = pd.to_numeric(df.get("financing_cf", 0), errors="coerce").fillna(0.0).tolist()

            x = np.arange(len(months))
            fig, ax = plt.subplots(figsize=(7.8, 3.0))
            ax.bar(x - 0.2, inv, width=0.4, label="Investing CF")
            ax.bar(x + 0.2, fin, width=0.4, label="Financing CF")
            ax.axhline(0, linewidth=1)
            ax.set_title("Investing vs Financing cashflows (next 12 months)")
            ax.set_ylabel(f"Cashflow ({currency_code})")
            ax.set_xticks(x)
            ax.set_xticklabels(months, rotation=45, ha="right")
            ax.legend()
            fig.tight_layout()

            _add_chart_image(slide, fig, left_in=5.4, top_in=1.3, width_in=4.7)
        else:
            _add_text_box(slide, 5.4, 1.3, 4.7, 2.0, "No engine rows in the next 12 months window.", font_size=13)
    else:
        _add_text_box(slide, 5.4, 1.3, 4.7, 2.0, "No investing/financing cashflow columns in cashflow_summary yet.", font_size=13)

    # ---------- Commentary ----------
    lines = build_funding_strategy_commentary(snap)
    _add_text_box(
        slide,
        left_in=0.5,
        top_in=3.8,
        width_in=9.6,
        height_in=2.4,
        text="Board commentary\n" + "\n".join(f"• {l}" for l in lines),
        font_size=13,
    )


def generate_board_pack_pptx(client_id: str, focus_month: datetime.date) -> str:
    settings = get_client_settings(client_id) or {}
    currency_code, currency_symbol = get_client_currency(client_id)

    kpis = fetch_kpis_for_client_month(client_id, focus_month) or {}
    kpi_change = build_kpi_change_explanations(client_id, focus_month)
    deltas = kpi_change.get("deltas", {})
    expl = kpi_change.get("explanations", {})

    month_label = pd.to_datetime(focus_month).strftime("%b %Y")

    rev_sched = build_revenue_schedule_for_client(client_id, base_month=focus_month, n_months=12)

    cash_df = fetch_cashflow_summary_for_client(client_id)
    df_pipeline = fetch_pipeline_for_client(client_id)
    df_positions = fetch_payroll_positions_for_client(client_id)
    df_ar, df_ap = fetch_ar_ap_for_client(client_id)

    # Runway + burn from engine (use function params!)
    engine_df = cash_df
    runway_from_engine, effective_burn = compute_runway_and_effective_burn_from_df(
        engine_df=engine_df,
        month_ref=focus_month,
    )

    month_summary = build_month_summary_insight(client_id, focus_month)

    prs = Presentation()

    # Cover
    cover_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(cover_layout)
    cover.shapes.title.text = "Board Pack"
    subtitle = cover.placeholders[1].text_frame
    subtitle.text = f"Client: {client_id}\nPeriod: {month_label}"

    # 1) Exec Summary
    slide = _add_title_only_slide(prs, "1. Executive Summary")

    revenue_this = float(kpis.get("revenue") or 0.0)
    burn_this = float(kpis.get("burn") or 0.0)
    cash_balance = float(kpis.get("cash_balance") or 0.0)
    runway_months = runway_from_engine if runway_from_engine is not None else kpis.get("runway_months")

    kpi_text = [
        f"Revenue ({month_label}): {currency_symbol}{revenue_this:,.0f}",
        f"Burn ({month_label}): {currency_symbol}{burn_this:,.0f}",
        f"Cash in bank: {currency_symbol}{cash_balance:,.0f}",
    ]
    if runway_months is not None:
        kpi_text.append(f"Runway (forward-looking): {float(runway_months):.1f} months")
    if effective_burn is not None:
        kpi_text.append(f"Effective burn (avg): {currency_symbol}{float(effective_burn):,.0f} / month")

    _add_text_box(slide, 0.5, 1.3, 4.5, 2.0, "Key KPIs\n" + "\n".join(f"• {line}" for line in kpi_text), font_size=14)

    alerts = fetch_alerts_for_client(client_id, only_active=True, limit=50) or []
    key_risks, key_opps = [], []
    for a in alerts:
        sev = str(a.get("severity", "medium")).lower()
        msg = a.get("message", "")
        if sev in ["high", "critical"]:
            key_risks.append(msg)
        else:
            key_opps.append(msg)
    if not key_risks:
        key_risks.append("No high / critical risks currently flagged.")
    if not key_opps:
        key_opps.append("Focus on executing pipeline and improving collections.")

    _add_text_box(slide, 5.2, 1.3, 4.5, 1.5, "Key risks\n" + "\n".join(f"• {r}" for r in key_risks[:5]), font_size=13)
    _add_text_box(slide, 5.2, 3.0, 4.5, 1.5, "Key opportunities\n" + "\n".join(f"• {o}" for o in key_opps[:5]), font_size=13)

    exec_lines = []
    if expl.get("money_in"):
        exec_lines.append(f"Revenue: {expl['money_in']}")
    if expl.get("money_out"):
        exec_lines.append(f"Burn: {expl['money_out']}")
    if expl.get("cash"):
        exec_lines.append(f"Cash: {expl['cash']}")
    if expl.get("runway"):
        exec_lines.append(f"Runway: {expl['runway']}")
    exec_lines.append(month_summary)

    _add_text_box(
        slide,
        0.5,
        3.7,
        9.2,
        2.5,
        "Executive commentary\n" + "\n".join(f"• {line}" for line in exec_lines[:6]),
        font_size=13,
    )
# ------------------------------------------------------------------
    # 3) Revenue Summary slide – with chart
    # ------------------------------------------------------------------
    slide = _add_title_only_slide(prs, "2. Revenue Summary – Sales & Deals")

    if rev_sched is not None and not rev_sched.empty:
        rev_df = rev_sched.copy()
        rev_df["month_date"] = pd.to_datetime(rev_df["month_date"], errors="coerce")
        rev_df = rev_df[rev_df["month_date"].notna()]

        rev_col = None
        for cand in ["recognised_revenue", "revenue_amount", "amount"]:
            if cand in rev_df.columns:
                rev_col = cand
                break

        if rev_col:
            monthly = (
                rev_df.groupby("month_date", as_index=False)[rev_col]
                .sum()
                .sort_values("month_date")
            )

            # Split last 6 / next 6 around focus_month
            focus_start = pd.to_datetime(focus_month).replace(day=1)
            past = monthly[monthly["month_date"] <= focus_start].tail(6)
            future = monthly[monthly["month_date"] > focus_start].head(6)

            # Build line chart (12-month view)
            fig, ax = plt.subplots(figsize=(7.5, 3))
            ax.plot(monthly["month_date"], monthly[rev_col], marker="o")
            ax.axvline(focus_start, linestyle="--", linewidth=1)
            ax.set_title("Recognised revenue – last 6 & next 6 months")
            ax.set_ylabel(f"Revenue ({currency_code})")
            ax.tick_params(axis="x", rotation=45)
            fig.tight_layout()

            _add_chart_image(slide, fig, left_in=0.4, top_in=1.3, width_in=6.5)

            # Text summary on right
            rev_text_lines = []
            if not past.empty:
                last_val = float(past[rev_col].iloc[-1])
                rev_text_lines.append(
                    f"Latest month recognised revenue: {currency_symbol}{last_val:,.0f}"
                )
            if not future.empty:
                future_total = float(future[rev_col].sum())
                rev_text_lines.append(
                    f"Next 6 months forecast total: {currency_symbol}{future_total:,.0f}"
                )

            # Pipeline summary
            total_pipeline = 0.0
            weighted_pipeline = 0.0
            if df_pipeline is not None and not df_pipeline.empty:
                tmp = df_pipeline.copy()
                tmp["value_total"] = pd.to_numeric(
                    tmp.get("value_total"), errors="coerce"
                ).fillna(0.0)
                tmp["probability_pct"] = pd.to_numeric(
                    tmp.get("probability_pct"), errors="coerce"
                ).fillna(0.0)
                total_pipeline = float(tmp["value_total"].sum())
                weighted_pipeline = float(
                    (tmp["value_total"] * tmp["probability_pct"] / 100.0).sum()
                )
                rev_text_lines.append(
                    f"Total pipeline: {currency_symbol}{total_pipeline:,.0f}"
                )
                rev_text_lines.append(
                    f"Weighted pipeline: {currency_symbol}{weighted_pipeline:,.0f}"
                )

            if total_pipeline > 0:
                qual = (weighted_pipeline / total_pipeline) * 100.0
                rev_text_lines.append(
                    f"Pipeline quality (weighted/gross): {qual:.1f}%"
                )

            rev_commentary = [
                "Revenue commentary",
                "• Recent revenue plus weighted pipeline gives visibility into the next 3–6 months.",
            ]
            if future is not None and not future.empty:
                rev_commentary.append(
                    "• Check if the forecast is driven by a few large deals – monitor conversion risk."
                )
            if total_pipeline > 0:
                rev_commentary.append(
                    "• Consider whether pipeline coverage is at least 3× your target revenue."
                )

            right_block = "\n".join(
                ["Key numbers"] + [f"• {l}" for l in rev_text_lines] + [""] + rev_commentary
            )

            _add_text_box(slide, 7.2, 1.3, 3.0, 3.5, right_block, font_size=12)

        else:
            _add_text_box(
                slide,
                0.5,
                1.5,
                9.0,
                3.0,
                "No recognised revenue schedule yet – add deals and revenue methods to see a full revenue summary.",
                font_size=14,
            )

    # ------------------------------------------------------------------
    # 4) Team & Spending slide – payroll chart + commentary
    # ------------------------------------------------------------------
    slide = _add_title_only_slide(prs, "3. Team & Spending")

    # Payroll trend for next 6 months (same logic as dashboard)
    months_6 = pd.date_range(start=pd.to_datetime(focus_month), periods=6, freq="MS")
    payroll_by_month = compute_payroll_by_month(client_id, months_6)

    if payroll_by_month:
        pay_vals = [float(payroll_by_month.get(m, 0.0)) for m in months_6]

        fig, ax = plt.subplots(figsize=(7.5, 3))
        ax.plot(months_6, pay_vals, marker="o")
        ax.set_title("Payroll cash-out – next 6 months")
        ax.set_ylabel(f"Payroll ({currency_code})")
        ax.tick_params(axis="x", rotation=45)
        fig.tight_layout()
        _add_chart_image(slide, fig, left_in=0.4, top_in=1.3, width_in=6.5)

        avg_payroll = sum(pay_vals) / len(pay_vals) if pay_vals else 0.0

        # Relative to current revenue
        revenue_this = revenue_this or 0.0
        payroll_vs_rev = (avg_payroll / revenue_this * 100.0) if revenue_this else None

        lines = [
            f"Average monthly payroll (next 6 months): {currency_symbol}{avg_payroll:,.0f}",
        ]
        if payroll_vs_rev is not None:
            lines.append(f"Payroll as % of current revenue: {payroll_vs_rev:.0f}%")

        # Headcount from positions
        headcount_fte = 0.0
        if df_positions is not None and not df_positions.empty:
            pos = df_positions.copy()
            pos["start_date"] = pd.to_datetime(pos.get("start_date"), errors="coerce")
            pos["end_date"] = pd.to_datetime(pos.get("end_date"), errors="coerce")
            month_start = pd.to_datetime(focus_month)

            active_mask = (pos["start_date"] <= month_start) & (
                pos["end_date"].isna() | (pos["end_date"] >= month_start)
            )
            pos_active = pos[active_mask]
            if not pos_active.empty:
                headcount_fte = float(pos_active.get("fte", 0.0).fillna(0.0).sum())

        if headcount_fte > 0:
            lines.append(f"Active headcount this month: {headcount_fte:.1f} FTE")

        commentary = [
            "Team & opex commentary",
            "• Check whether payroll growth is in line with revenue growth.",
            "• Identify any high-cost roles that are not clearly linked to revenue or product milestones.",
            "• Review departments that are consistently over budget.",
        ]

        text_block = "\n".join(
            ["Key numbers"] + [f"• {l}" for l in lines] + [""] + commentary
        )

        _add_text_box(slide, 7.2, 1.3, 3.0, 3.5, text_block, font_size=12)
    else:
        _add_text_box(
            slide,
            0.5,
            1.5,
            9.0,
            3.0,
            "No payroll data yet – add roles in the Team & Spending page to see headcount and payroll trends.",
            font_size=14,
        )

    # ------------------------------------------------------------------
    # 5) Cash Position & Forecast – cash engine chart + AR/AP snapshot
    # ------------------------------------------------------------------
    slide = _add_title_only_slide(prs, "4. Cash Position & Cash Forecast")

    if cash_df is not None and not cash_df.empty:
        df_cash = cash_df.copy()
        df_cash["month_date"] = pd.to_datetime(df_cash["month_date"], errors="coerce")
        df_cash = df_cash[df_cash["month_date"].notna()].sort_values("month_date")

        fig, ax = plt.subplots(figsize=(7.5, 3))
        if "closing_cash" in df_cash.columns:
            ax.plot(df_cash["month_date"], df_cash["closing_cash"], marker="o")
            ax.set_ylabel(f"Closing cash ({currency_code})")
        if "free_cash_flow" in df_cash.columns:
            ax.bar(
                df_cash["month_date"],
                df_cash["free_cash_flow"],
                alpha=0.3,
            )
        ax.set_title("Monthly cash forecast (engine)")
        ax.tick_params(axis="x", rotation=45)
        fig.tight_layout()
        _add_chart_image(slide, fig, left_in=0.4, top_in=1.3, width_in=6.5)

        cash_today = cash_balance
        cliff = None
        if "closing_cash" in df_cash.columns:
            below_zero = df_cash[df_cash["closing_cash"] <= 0]
            if not below_zero.empty:
                cliff = below_zero["month_date"].iloc[0].date()

        wc_lines = [f"Cash today: {currency_symbol}{cash_today:,.0f}"]
        if cliff:
            wc_lines.append(f"Cash cliff (engine): {cliff.strftime('%d %b %Y')}")

        # Working capital from AR / AP
        ar_out = 0.0
        ap_out = 0.0
        if df_ar is not None and not df_ar.empty:
            df_ar2 = df_ar.copy()
            df_ar2["amount"] = pd.to_numeric(df_ar2.get("amount"), errors="coerce").fillna(0.0)
            ar_out = float(df_ar2["amount"].sum())
        if df_ap is not None and not df_ap.empty:
            df_ap2 = df_ap.copy()
            df_ap2["amount"] = pd.to_numeric(df_ap2.get("amount"), errors="coerce").fillna(0.0)
            ap_out = float(df_ap2["amount"].sum())

        wc = ar_out - ap_out
        wc_lines.append(f"AR outstanding: {currency_symbol}{ar_out:,.0f}")
        wc_lines.append(f"AP outstanding: {currency_symbol}{ap_out:,.0f}")
        wc_lines.append(f"Net working capital (AR – AP): {currency_symbol}{wc:,.0f}")

        cash_comment = [
            "Cash commentary",
            "• If cash cliff is within 6–9 months, begin funding or cost-reduction planning now.",
            "• Use AR and AP levers (collections, payment terms) before cutting critical growth spend.",
        ]
        if cliff:
            cash_comment.insert(
                1,
                f"• Cash cliff expected around {cliff.strftime('%b %Y')} if no changes are made.",
            )

        text_block = "\n".join(
            ["Key numbers"] + [f"• {l}" for l in wc_lines] + [""] + cash_comment
        )
        _add_text_box(slide, 7.2, 1.3, 3.0, 3.5, text_block, font_size=12)
    else:
        _add_text_box(
            slide,
            0.5,
            1.5,
            9.0,
            3.0,
            "No cashflow engine data yet – rebuild cashflow from the Business at a Glance page.",
            font_size=14,
        )



    # ✅ NEW 5) Funding Strategy – replaces Scenario + old I&F slide
    add_funding_strategy_slide(prs, client_id=client_id, focus_month=focus_month)

    # 6) Risks, Issues & Notes (renumbered)
    slide = _add_title_only_slide(prs, "6. Risks, Issues & Notes")

    risk_lines = []
    for a in alerts:
        sev = str(a.get("severity", "medium")).title()
        msg = a.get("message", "")
        owner = a.get("owner") or a.get("owner_name") or "Unassigned"
        risk_lines.append(f"• [{sev}] {msg}  (Owner: {owner})")

    if not risk_lines:
        risk_lines = ["• No major risks recorded – confirm this with the founder."]

    _add_text_box(slide, 0.5, 1.3, 9.0, 4.5, "Top risks & issues\n" + "\n".join(risk_lines[:8]), font_size=13)

    safe_month = pd.to_datetime(focus_month).strftime("%Y%m")
    filename = f"BoardPack_{client_id}_{safe_month}.pptx"
    file_path = os.path.join("outputs", filename)
    os.makedirs("outputs", exist_ok=True)
    prs.save(file_path)
    return file_path


def run_all_alerts_for_client_months(client_id, month_index):
    """
    For each month in month_index, run all alert checks:
      - runway_low      (KPI/cashflow based)
      - cash_danger     (cashflow_summary-based)
      - ar_overdue      (AR aging-based)
    """
    if client_id is None or month_index is None:
        print("[ALERT] run_all_alerts_for_client_months called with missing client_id or month_index")
        return

    # Normalise to unique first-of-month dates
    month_list = sorted(
        {pd.to_datetime(m).to_period("M").to_timestamp().date() for m in month_index}
    )

    print(f"[ALERT] running alerts for client={client_id}, months={month_list}")

    for m in month_list:
        # 1) runway alert
        try:
            print(f"[ALERT] ensure_runway_alert_for_month -> {m}")
            ensure_runway_alert_for_month(client_id, m)
        except Exception as e:
            print(f"[ALERT] runway_low failed for {m}: {e}")

        # 2) cash danger alert
        try:
            print(f"[ALERT] ensure_cash_danger_alert_for_month -> {m}")
            ensure_cash_danger_alert_for_month(client_id, m)
        except Exception as e:
            print(f"[ALERT] cash_danger failed for {m}: {e}")

        # 3) AR overdue alert
        try:
            print(f"[ALERT] ensure_overdue_ar_alert -> {m}")
            ensure_overdue_ar_alert(client_id, as_of=m)
        except Exception as e:
            print(f"[ALERT] ar_overdue failed for {m}: {e}")


def fetch_board_commentary(client_id: str, month_start: date):
    try:
        resp = (
            supabase.table("board_commentary")
            .select("*")
            .eq("client_id", str(client_id))
            .eq("month_date", month_start.isoformat())
            .execute()
        )
        return resp.data or []
    except Exception as e:
        print("Error fetching board commentary:", e)
        return []

def upsert_board_commentary(
    client_id: str,
    month_start: date,
    page_name: str,
    kpi_key: str,
    final_commentary: str,
    updated_by: str = "Team member",
) -> bool:
    try:
        payload = {
            "client_id": str(client_id),
            "month_date": month_start.isoformat(),
            "page_name": page_name,
            "kpi_key": kpi_key,
            "final_commentary": final_commentary.strip(),
            "updated_by": updated_by,
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }
        supabase.table("board_commentary").upsert(payload).execute()
        return True
    except Exception as e:
        print("Error saving board commentary:", e)
        return False

from datetime import datetime, timezone

def render_final_commentary_box(client_id: str, month_start: date, df_tasks: pd.DataFrame):
    st.markdown("---")
    st.subheader("🧾 Final commentary for board report")
    st.caption(
        "Finalise commentary in KPI → Task order. This is what goes into the board pack."
    )

    # Load saved final commentary (existing)
    saved = fetch_board_commentary(client_id, month_start)
    saved_map = {}
    for r in saved:
        saved_map[(r.get("page_name"), r.get("kpi_key"))] = r.get("final_commentary", "")

    if df_tasks is None or df_tasks.empty:
        st.info("No tasks yet. Create tasks above to generate board commentary.")
        return

    # Only include relevant tasks for this month
    df = df_tasks.copy()
    if "month_date" in df.columns:
        df["month_date"] = pd.to_datetime(df["month_date"], errors="coerce").dt.date
        df = df[df["month_date"] == month_start]

    # Require KPI tagging
    if "kpi_key" not in df.columns:
        st.warning("Tasks are missing `kpi_key`. Add it to your tasks table + create-task form.")
        return

    # Group by page -> kpi
    df["page_name"] = df["page_name"].fillna("general")
    df["kpi_key"] = df["kpi_key"].fillna("general")

    grouped = df.groupby(["page_name", "kpi_key"], as_index=False).size()

    # Render each KPI group
    for _, g in grouped.iterrows():
        page_name = g["page_name"]
        kpi_key = g["kpi_key"]

        page_label = PAGE_KPI_MAP.get(page_name, {}).get("label", page_name)
        kpi_label = PAGE_KPI_MAP.get(page_name, {}).get("kpis", {}).get(kpi_key, kpi_key)

        st.markdown(f"### {page_label} → **{kpi_label}**")

        # Show tasks under this KPI
        subset = df[(df["page_name"] == page_name) & (df["kpi_key"] == kpi_key)].copy()
        sort_cols = [c for c in ["status", "created_at"] if c in subset.columns]
        if sort_cols:
            subset = subset.sort_values(by=sort_cols, ascending=[True] * len(sort_cols))


        for _, t in subset.iterrows():
            st.markdown(f"- **{t.get('title','Untitled')}**  _(status: {t.get('status','open')})_")

        # Final commentary editor
        default_text = saved_map.get((page_name, kpi_key), "")
        key_txt = f"final_commentary_{page_name}_{kpi_key}_{month_start.isoformat()}"

        final_text = st.text_area(
            "Final board commentary (what will appear in the board pack)",
            value=default_text,
            height=120,
            key=key_txt,
        )

        c1, c2 = st.columns([1, 3])
        with c1:
            if st.button("Save", key=f"save_final_{page_name}_{kpi_key}_{month_start.isoformat()}"):
                if not final_text.strip():
                    st.warning("Commentary is empty.")
                else:
                    ok = upsert_board_commentary(
                        client_id=client_id,
                        month_start=month_start,
                        page_name=page_name,
                        kpi_key=kpi_key,
                        final_commentary=final_text,
                        updated_by="Team member",
                    )
                    if ok:
                        st.success("Saved.")
                    else:
                        st.error("Could not save. Check Supabase table + permissions.")

        st.markdown("---")




# ---------- Page: Alerts & To-Dos ----------
def page_alerts_todos():
    """
    Collaboration Hub for the founder & team:
    - Snapshot of key alerts (runway, cash danger, overdue AR, overspend, etc)
    - One central task board (collaboration_hub) for cross-page work, with filters
    - Tabs to collaborate on each main page (comments + tasks together)
    """
    top_header("Collaboration hub")

    # -------------------------------
    # Basic guards
    # -------------------------------
    if not selected_client_id:
        st.info("Select a business from the sidebar to see collaboration items.")
        return

    if not selected_month_start:
        st.info("Pick a focus month in the navbar to see alerts and tasks for that period.")
        return

# --- Board pack download block (top of the page) ---
    if selected_client_id:
        if st.button("Download board pack (PowerPoint)", key="btn_board_pack"):
            file_path = generate_board_pack_pptx(selected_client_id, selected_month_start)
            with open(file_path, "rb") as f:
                st.download_button(
                    label="Click to download board pack",
                    data=f,
                    file_name=os.path.basename(file_path),
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="dl_board_pack",
                )


    # -------------------------------
    # Keep automatic alerts in sync
    # -------------------------------
    try:
        ensure_runway_alert_for_month(selected_client_id, selected_month_start)
    except Exception as e:
        st.warning("Could not refresh runway alert for this month.")
        st.caption(str(e))

    try:
        ensure_cash_danger_alert_for_month(selected_client_id, selected_month_start)
    except Exception as e:
        st.warning("Could not refresh cash danger alert for this month.")
        st.caption(str(e))

    # Optionally keep AR + dept overspend alerts in sync too
    try:
        ensure_overdue_ar_alert(
            selected_client_id,
            as_of=date.today(),
            min_days_overdue=0,
        )
    except Exception:
        pass

    

    # -------------------------------
    # 🔥 Key alerts snapshot
    # -------------------------------
    st.subheader("🚨 Key alerts (snapshot)")

    alerts = fetch_alerts_for_client(
        selected_client_id,
        only_active=True,
        limit=100,
    )

    if not alerts:
        st.caption("No open alerts right now. 🎉")
    else:
        # Try to focus on this month; if nothing, fall back to all
        alerts_for_month = []
        for a in alerts:
            month_value = a.get("month_date")
            if month_value and selected_month_start:
                try:
                    mdt = pd.to_datetime(month_value).date().replace(day=1)
                    if mdt == selected_month_start:
                        alerts_for_month.append(a)
                except Exception:
                    alerts_for_month.append(a)
            else:
                alerts_for_month.append(a)

        if not alerts_for_month:
            alerts_for_month = alerts

        alerts_sorted = sort_alerts_by_severity(alerts_for_month)

        # Show only the top 4 in detail to keep this clean
        top_alerts = alerts_sorted[:4]

        col_a1, col_a2 = st.columns(2)
        alert_cols = [col_a1, col_a2]
        idx = 0

        for alert in top_alerts:
            sev = str(alert.get("severity", "medium")).lower()
            atype = alert.get("alert_type", "alert")
            msg = alert.get("message", "")
            created = alert.get("created_at", "")
            page_name = alert.get("page_name") or "general"

            label = f"[{atype}] {msg}"
            meta = f"Page: {page_name} · Created: {created}"

            if sev in ["critical", "high"]:
                box = alert_cols[idx % 2].error
            elif sev == "medium":
                box = alert_cols[idx % 2].warning
            else:
                box = alert_cols[idx % 2].info

            box(label)
            alert_cols[idx % 2].caption(meta)
            idx += 1

        # Small meta summary
        total_crit = sum(
            1
            for a in alerts_for_month
            if str(a.get("severity", "")).lower() in ["critical", "high"]
        )
        total_med = sum(
            1
            for a in alerts_for_month
            if str(a.get("severity", "")).lower() == "medium"
        )
        total_low = len(alerts_for_month) - total_crit - total_med

    st.caption(
            f"Summary this period: {total_crit} high/critical · "
            f"{total_med} medium · {total_low} low/other alerts."
        )


    st.markdown("---")
    st.subheader("➕ Create a task (Page → KPI → Issue)")

    page_keys = list(PAGE_KPI_MAP.keys())
    page_labels = [PAGE_KPI_MAP[k]["label"] for k in page_keys]

    sel_page_label = st.selectbox("Page", options=page_labels, key="task_new_page")
    sel_page_key = page_keys[page_labels.index(sel_page_label)]

    kpi_map = PAGE_KPI_MAP[sel_page_key]["kpis"]
    kpi_keys = list(kpi_map.keys())
    kpi_labels = [kpi_map[k] for k in kpi_keys]

    sel_kpi_label = st.selectbox("KPI / area", options=kpi_labels, key="task_new_kpi")
    sel_kpi_key = kpi_keys[kpi_labels.index(sel_kpi_label)]

    tcol1, tcol2 = st.columns(2)
    with tcol1:
        task_title = st.text_input("Task title", key="task_new_title")
    with tcol2:
        task_owner = st.text_input("Owner", value="", key="task_new_owner")

    task_desc = st.text_area("Describe the issue", key="task_new_desc")

    tcol3, tcol4 = st.columns(2)
    with tcol3:
        task_priority = st.selectbox("Priority", ["normal", "high", "critical", "low"], index=0, key="task_new_priority")
    with tcol4:
        task_due = st.date_input("Due date (optional)", value=None, key="task_new_due")

    if st.button("Create task", key="task_new_create_btn"):
        ok = save_task(
            client_id=selected_client_id,
            page_name=sel_page_key,
            title=task_title,
            month_date=selected_month_start,
            kpi_key=sel_kpi_key,
            description=task_desc,
            owner_name=task_owner,
            priority=task_priority,
        )
        if ok:
            st.success("Task created.")
            st.rerun()
        else:
            st.error("Could not create task. Check required fields and try again.")

        # -------------------------------
    # 🧩 Central workboard (Kanban-style)
    # -------------------------------
    st.subheader("✅ My workboard (cross-page tasks)")

    st.caption(
        "See tasks from all pages in one place. "
        "Filter by status, page or owner to focus your week."
    )

    # ----- Filters row -----
    col_f1, col_f2, col_f3 = st.columns([1.2, 1.2, 1.6])

    with col_f1:
        status_filter = st.multiselect(
            "Status",
            options=["open", "in_progress", "done", "cancelled"],
            default=["open", "in_progress"],  # hide done/cancelled by default
            key="collab_status_filter",
        )

    with col_f2:
        page_options = get_task_page_options(selected_client_id)
        page_filter = st.multiselect(
            "Page",
            options=page_options,
            default=[],
            help="Limit to tasks from specific pages (or leave blank for all).",
            key="collab_page_filter",
        )

    with col_f3:
        owner_options = get_task_owner_options(selected_client_id)
        owner_filter = st.multiselect(
            "Owner",
            options=owner_options,
            default=[],
            help="Focus on tasks assigned to specific people.",
            key="collab_owner_filter",
        )

    search_text = st.text_input(
        "Search in titles / descriptions",
        value="",
        key="collab_search_text",
    )

    # ----- Fetch tasks from Supabase -----
    df_tasks = fetch_tasks_for_collab_hub(
        client_id=selected_client_id,
        status_filter=status_filter or None,
        page_filter=page_filter or None,
        owner_filter=owner_filter or None,
        search_text=search_text,
    )

    if df_tasks.empty:
        st.info("No tasks match the current filters.")
    else:
        # Map page_name to friendly labels
        page_label_map = {
            "business_overview": "Business overview",
            "sales_deals": "Sales & deals",
            "team_spending": "Team & spending",
            "cash_bills": "Cash & bills",
            "alerts_todos": "Alerts & to-dos",
            "collaboration_hub": "Collaboration Hub",
        }
        df_tasks = df_tasks.copy()
        df_tasks["page_label"] = df_tasks["page_name"].map(page_label_map).fillna("Other")

        # Priority label
        def _prio_label(p: str) -> str:
            p = (p or "").lower()
            if p == "low":
                return "Low"
            if p == "normal":
                return "Normal"
            if p == "high":
                return "High"
            if p == "critical":
                return "Critical"
            return "Normal"

        df_tasks["priority_label"] = df_tasks["priority"].apply(_prio_label)

        # Kanban columns by status
        statuses_order = ["open", "in_progress", "done", "cancelled"]
        status_headers = {
            "open": "📝 To-do",
            "in_progress": "🚧 In progress",
            "done": "✅ Done",
            "cancelled": "🗑️ Cancelled",
        }

        cols = st.columns(4)

        for idx, status in enumerate(statuses_order):
            col = cols[idx]
            subset = df_tasks[df_tasks["status"] == status]

            with col:
                col.markdown(f"**{status_headers.get(status, status.title())}**")

                if subset.empty:
                    col.caption("No tasks.")
                    continue

                # Sort by due date (nulls last), then priority, then created_at
                subset = subset.sort_values(
                    by=["due_date", "priority", "created_at"],
                    ascending=[True, False, True],
                )

                for _, row in subset.iterrows():
                    title = row.get("title", "Untitled task")
                    page_label = row.get("page_label", "Unknown")
                    owner = row.get("owner_name") or "Unassigned"
                    due = row.get("due_date")
                    priority = row.get("priority_label", "Normal")

                    with st.container(border=True):
                        st.markdown(f"**{title}**")
                        task_id = row.get("id")
                        kpi_key = row.get("kpi_key") or ""
                        kpi_label = None
                        try:
                            # best-effort find label from map
                            pkey = row.get("page_name")
                            if pkey in PAGE_KPI_MAP:
                                kpi_label = PAGE_KPI_MAP[pkey]["kpis"].get(kpi_key)
                        except Exception:
                            pass

                        with st.container(border=True):
                            st.markdown(f"**{title}**")
                            if kpi_label:
                                st.caption(f"KPI: {kpi_label}")
                            elif kpi_key:
                                st.caption(f"KPI: {kpi_key}")

                            meta_parts = [page_label, owner, f"Priority: {priority}"]
                            if due:
                                meta_parts.append(f"Due: {due}")
                            st.caption(" · ".join(meta_parts))

                            # Quick actions
                            a1, a2, a3 = st.columns(3)
                            with a1:
                                if st.button("✅ Done", key=f"task_done_{task_id}"):
                                    update_task_status(task_id, "done")
                                    st.rerun()
                            with a2:
                                if st.button("🚧 In progress", key=f"task_prog_{task_id}"):
                                    update_task_status(task_id, "in_progress")
                                    st.rerun()
                            with a3:
                                if st.button("🗑 Cancel", key=f"task_cancel_{task_id}"):
                                    update_task_status(task_id, "cancelled")
                                    st.rerun()

                            # Replies thread
                            with st.expander("💬 Replies", expanded=False):
                                replies = fetch_task_replies(task_id)
                                if not replies:
                                    st.caption("No replies yet.")
                                else:
                                    for r in replies:
                                        who = r.get("author_name") or "Someone"
                                        ts = r.get("created_at") or ""
                                        msg = r.get("message") or ""
                                        st.markdown(f"**{who}** · {ts}")
                                        st.write(msg)
                                        st.markdown("---")

                                reply_author = st.text_input("Your name", key=f"reply_author_{task_id}")
                                reply_msg = st.text_area("Reply", key=f"reply_msg_{task_id}")

                                if st.button("Send reply", key=f"reply_send_{task_id}"):
                                    ok_r = add_task_reply(selected_client_id, task_id, reply_author, reply_msg)
                                    if ok_r:
                                        st.success("Reply added.")
                                        st.rerun()
                                    else:
                                        st.warning("Reply cannot be empty.")


                        # Later you can add buttons here (Mark done, etc.)

    st.markdown("---")

   
#...  For commentary....#
    render_final_commentary_box(selected_client_id, selected_month_start, df_tasks)

# --------------------------------------------------
# MAIN ENTRY POINT
# --------------------------------------------------




def run_app():


    page_names = [
        "Business at a Glance",
        "Investing & Financing Strategy",
        "Sales & Deals",
        "Team & Spending",
        "Cash & Bills",
        "Collaboration Hub",
        "Configuration",
        "Client Settings",
    ]

    page_choice = st.sidebar.radio("Go to page", page_names)

    if page_choice == "Business at a Glance":
        page_business_overview()
    elif page_choice == "Investing & Financing Strategy":
        page_investing_financing()
    elif page_choice == "Sales & Deals":
        page_sales_deals()
    elif page_choice == "Team & Spending":
        page_team_spending()
    elif page_choice == "Cash & Bills":
        page_cash_bills()    
    elif page_choice == "Collaboration Hub":
        page_alerts_todos()    
    elif page_choice == "Configuration":
        page_configuration()
    elif page_choice == "Client Settings":
        page_client_settings()


if __name__ == "__main__":
    run_app()






def supabase_select_with_debug(table: str, client_id: str, order_col: str | None = None, max_retries: int = 3):
    """
    Returns (df, meta) where meta contains debug info about failures.
    """
    meta = {
        "table": table,
        "client_id": client_id,
        "attempts": [],
        "ok": False,
        "row_count": 0,
        "error": None,
    }

    for attempt in range(1, max_retries + 1):
        t0 = time.time()
        try:
            q = supabase.table(table).select("*").eq("client_id", str(client_id))
            if order_col:
                q = q.order(order_col)

            res = q.execute()
            rows = res.data or []
            df = pd.DataFrame(rows)

            meta["ok"] = True
            meta["row_count"] = len(df)
            meta["attempts"].append({
                "attempt": attempt,
                "ms": int((time.time() - t0) * 1000),
                "status": "ok",
                "rows": len(df),
            })
            return df, meta

        except Exception as e:
            meta["error"] = str(e)
            meta["attempts"].append({
                "attempt": attempt,
                "ms": int((time.time() - t0) * 1000),
                "status": "fail",
                "error": str(e),
            })

            print(f"[WARN] fetch {table} attempt {attempt}/{max_retries} failed:", e)
            # small backoff
            time.sleep(0.4 * attempt)

    # all attempts failed
    return pd.DataFrame(), meta

# ---------- Board Pack Helper ----------


def _safe_float(val, default=0.0):
    try:
        if val is None:
            return default
        return float(val)
    except Exception:
        return default


def _safe_pct(numerator, denominator):
    try:
        if not denominator:
            return None
        return float(numerator) / float(denominator) * 100.0
    except Exception:
        return None


def _month_start(d: date) -> date:
    return pd.to_datetime(d).to_period("M").to_timestamp().date()





